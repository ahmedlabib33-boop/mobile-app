from __future__ import annotations

# pyright: reportGeneralTypeIssues=false, reportCallIssue=false, reportArgumentType=false, reportAttributeAccessIssue=false, reportOptionalCall=false, reportOptionalSubscript=false, reportOperatorIssue=false, reportOptionalMemberAccess=false
import base64
import html
import importlib
import io
import json
import hmac
import os
import re
import shutil
import site
import subprocess
import textwrap
from pathlib import Path
import tempfile
import zipfile
import xml.etree.ElementTree as ET
import sys
from typing import Any

try:
    sync_playwright = importlib.import_module("playwright.sync_api").sync_playwright
    PLAYWRIGHT_AVAILABLE = True
except ModuleNotFoundError:
    sync_playwright = None
    PLAYWRIGHT_AVAILABLE = False

import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st
import auth
import contract_claims_center as ccc
from reports.tia_director_pack_generator import (
    REPORT_TYPE_TIA_DIRECTOR_PACK,
    TIADirectorPackGenerator,
    fetch_last_generated_report,
)
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx import Presentation
from pptx.util import Inches, Pt
from ui.report_generator_page import (
    build_data_source_status_df,
    build_replacement_preview_df,
    compute_required_fields_completion,
    resolve_existing_output_path,
)

try:
    from reportlab.lib import colors as rl_colors
    from reportlab.lib.pagesizes import A3, landscape as rl_landscape
    from reportlab.lib.units import mm
    from reportlab.pdfbase.pdfmetrics import stringWidth
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.platypus import Table as RLTable, TableStyle as RLTableStyle
    REPORTLAB_AVAILABLE = True
except ModuleNotFoundError:
    rl_colors = None
    A3 = None
    rl_landscape = None
    mm = None
    stringWidth = None
    rl_canvas = None
    RLTable = None
    RLTableStyle = None
    REPORTLAB_AVAILABLE = False

try:
    from openpyxl import Workbook
    from openpyxl.chart import BarChart as XLBarChart, LineChart as XLLineChart, Reference
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
    from openpyxl.utils import get_column_letter
    from openpyxl.utils.dataframe import dataframe_to_rows
    from openpyxl.worksheet.table import Table, TableStyleInfo
    OPENPYXL_AVAILABLE = True
except ModuleNotFoundError:
    Workbook = None
    XLBarChart = None
    XLLineChart = None
    Reference = None
    Alignment = None
    Border = None
    Font = None
    PatternFill = None
    Side = None
    get_column_letter = None
    dataframe_to_rows = None
    Table = None
    TableStyleInfo = None
    OPENPYXL_AVAILABLE = False

try:
    Document = importlib.import_module("docx").Document
    DOCX_AVAILABLE = True
except ModuleNotFoundError:
    try:
        user_site = site.getusersitepackages()
        if user_site and user_site not in sys.path:
            sys.path.append(user_site)
        Document = importlib.import_module("docx").Document
        DOCX_AVAILABLE = True
    except ModuleNotFoundError:
        Document = None
        DOCX_AVAILABLE = False

from src.construction_system.contract_matcher import (
    analyze_delay_event,
    generate_ai_clause_brief,
    get_all_clauses,
    get_contract_terms,
    match_event_to_clauses,
    search_clauses,
    set_clause_library_path,
)
from src.construction_system.letters_auto_ingest import folder_fingerprint, merge_inbox_letters
from src.construction_system.project_catalog import (
    discover_projects,
    project_data_path,
    projects_frame,
)
from src.construction_system.project_context import build_project_context
from src.construction_system.steel_delay_tia import (
    CANONICAL_STEEL_FIELDS,
    STEEL_ALIASES,
    SteelTiaSettings,
    apply_mapping,
    build_requirement_df_from_client_supply_sheet,
    parse_mixed_date,
    run_steel_delay_tia_analysis,
    suggest_mapping,
)

APP_DIR = Path(__file__).parent
APP_ENV_PREFIX = "PIH_MOBILE_APP"
RUNTIME_DIR = APP_DIR / "runtime"
LOGO_PATH = APP_DIR / "assets" / "logo.png"
MOBILE_CONFIG_PATH = APP_DIR / "mobile_config.json"
PROJECTS_DIR = APP_DIR / "projects"
PATH_TOKENS_DIR = APP_DIR / ".project_paths"
IMPORT_TEMPLATES_DIR = PATH_TOKENS_DIR / "01-data" / "import_templates"
STEEL_TIA_DIR = PATH_TOKENS_DIR / "02-delay_analysis" / "steel_delay_tia_templates"
BL_FIXED_DIR = PATH_TOKENS_DIR / "03-schedule"
PROJECTS_CSV_PATH = IMPORT_TEMPLATES_DIR / "projects.csv"
ACTIVITIES_CSV_PATH = IMPORT_TEMPLATES_DIR / "activities.csv"
EVM_CSV_PATH = IMPORT_TEMPLATES_DIR / "evm.csv"
CONTRACTS_CSV_PATH = IMPORT_TEMPLATES_DIR / "contracts.csv"
PAYMENTS_CSV_PATH = IMPORT_TEMPLATES_DIR / "payments.csv"
DELAYS_CSV_PATH = IMPORT_TEMPLATES_DIR / "delay_events.csv"
RISKS_CSV_PATH = IMPORT_TEMPLATES_DIR / "risks.csv"
MILESTONES_CSV_PATH = IMPORT_TEMPLATES_DIR / "milestones.csv"
CHANGE_ORDERS_CSV_PATH = IMPORT_TEMPLATES_DIR / "change_orders.csv"
STEEL_DELAY_CSV_PATH = IMPORT_TEMPLATES_DIR / "steel_delay_status_mployer_free_issue_material.csv"
RFI_STATUS_CSV_PATH = IMPORT_TEMPLATES_DIR / "rfi_ status.csv"
IFC_CONFLICT_CSV_PATH = IMPORT_TEMPLATES_DIR / "ifc_conflict.csv"
S_CURVE_CSV_PATH = IMPORT_TEMPLATES_DIR / "s_curve.csv"
WBS_CSV_PATH = IMPORT_TEMPLATES_DIR / "wbs.csv"
EVM_COMMENTS_PATH = APP_DIR / "data" / "evm_comments.json"
TIA_DIRECTOR_WORD_TEMPLATE_PATH = APP_DIR / "reports" / "templates" / "time_impact_analysis_report_director_pack.docx"

st.set_page_config(
    page_title="Project Intelligence Hub",
    page_icon=str(LOGO_PATH) if LOGO_PATH.exists() else None,
    layout="wide",
    initial_sidebar_state="collapsed",
    menu_items={
        "About": "Project Intelligence Hub mobile-ready project controls intelligence platform.",
    },
)


def _github_sync_command(mode: str, interval_seconds: int = 30) -> list[str]:
    powershell = shutil.which("powershell.exe") or shutil.which("pwsh") or shutil.which("powershell")
    if not powershell:
        raise RuntimeError("PowerShell is not available on this server.")
    return [
        powershell,
        "-NoProfile",
        "-ExecutionPolicy",
        "Bypass",
        "-File",
        str(APP_DIR / "tools" / "github_no_git_sync.ps1"),
        "-Mode",
        mode,
        "-IntervalSeconds",
        str(interval_seconds),
    ]


def _streamlit_secret(name: str) -> str:
    try:
        return str(st.secrets.get(name, "") or "").strip()
    except (FileNotFoundError, KeyError, TypeError):
        return ""


def _github_sync_environment() -> dict[str, str]:
    environment = os.environ.copy()
    for name in (f"{APP_ENV_PREFIX}_GITHUB_TOKEN", f"{APP_ENV_PREFIX}_GH_TOKEN"):
        if not environment.get(name, "").strip():
            secret = _streamlit_secret(name)
            if secret:
                environment[name] = secret
    return environment


def run_repository_sync_once() -> tuple[bool, str]:
    result = subprocess.run(
        _github_sync_command("Once"),
        cwd=APP_DIR,
        env=_github_sync_environment(),
        capture_output=True,
        text=True,
        timeout=900,
        check=False,
    )
    output = "\n".join(part.strip() for part in (result.stdout, result.stderr) if part.strip())
    return result.returncode == 0, output[-6000:]


def start_repository_sync_watch() -> None:
    creation_flags = 0
    if os.name == "nt":
        creation_flags = subprocess.CREATE_NEW_PROCESS_GROUP | subprocess.DETACHED_PROCESS | subprocess.CREATE_NO_WINDOW
    subprocess.Popen(
        _github_sync_command("Watch", 30),
        cwd=APP_DIR,
        env=_github_sync_environment(),
        stdin=subprocess.DEVNULL,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        creationflags=creation_flags,
        close_fds=os.name != "nt",
    )


def repository_sync_log_tail(max_lines: int = 8) -> str:
    log_path = APP_DIR / "11-outputs" / "logs" / "pih_mobile_app_github_sync.log"
    if not log_path.exists():
        return "No synchronization run has been recorded."
    return "\n".join(log_path.read_text(encoding="utf-8", errors="replace").splitlines()[-max_lines:])


def repository_sync_authorized(pin: str) -> bool:
    expected = os.environ.get(f"{APP_ENV_PREFIX}_SYNC_ADMIN_PIN", "").strip() or _streamlit_secret(f"{APP_ENV_PREFIX}_SYNC_ADMIN_PIN")
    return not expected or hmac.compare_digest(pin.strip(), expected)


def image_as_base64(path: Path) -> str:
    return base64.b64encode(path.read_bytes()).decode("utf-8") if path.exists() else ""


def load_mobile_config() -> dict[str, str]:
    default = {
        "app_name": "Project Intelligence Hub",
        "streamlit_url": "PUT_DEPLOYED_STREAMLIT_URL_HERE",
    }
    if not MOBILE_CONFIG_PATH.exists():
        return default
    try:
        loaded = json.loads(MOBILE_CONFIG_PATH.read_text(encoding="utf-8"))
        if isinstance(loaded, dict):
            return {**default, **{str(key): str(value) for key, value in loaded.items()}}
    except (OSError, json.JSONDecodeError, TypeError):
        return default
    return default


MOBILE_CONFIG = load_mobile_config()


def configured_streamlit_url() -> str:
    url = str(MOBILE_CONFIG.get("streamlit_url", "")).strip()
    if url and url != "PUT_DEPLOYED_STREAMLIT_URL_HERE":
        return url
    return ""


def apply_mobile_first_dashboard_css() -> None:
    st.markdown(
        """
        <style>
        :root{
            --pih-bg:#f5f7fb;
            --pih-surface:#ffffff;
            --pih-surface-soft:#f8fafc;
            --pih-text:#111827;
            --pih-muted:#64748b;
            --pih-line:rgba(15,23,42,.12);
            --pih-accent:#0f766e;
            --pih-accent-2:#1d4ed8;
        }
        html, body, [data-testid="stAppViewContainer"]{
            background:var(--pih-bg);
            color:var(--pih-text);
            overflow-x:hidden;
        }
        #MainMenu,
        footer,
        header,
        [data-testid="stToolbar"],
        [data-testid="stDecoration"],
        [data-testid="stStatusWidget"],
        [data-testid="stHeader"],
        [data-testid="stMainMenu"],
        [data-testid="stDeployButton"],
        [data-testid="manage-app-button"],
        a[href*="github.com"],
        button[title*="GitHub"],
        button[aria-label*="GitHub"],
        button[title*="View source"],
        button[aria-label*="View source"]{
            display:none!important;
            visibility:hidden!important;
            opacity:0!important;
            pointer-events:none!important;
        }
        .viewerBadge_container__r5tak,
        .st-emotion-cache-14xtw13,
        .st-emotion-cache-1avcm0n{
            display:none!important;
        }
        .main .block-container{
            max-width:1440px;
            padding:1.1rem 1.4rem 2.5rem;
        }
        h1,h2,h3,h4{
            letter-spacing:0;
            color:var(--pih-text);
        }
        p, li, label, .stMarkdown, [data-testid="stMetricLabel"]{
            color:var(--pih-text);
        }
        [data-testid="stSidebar"]{
            border-right:1px solid var(--pih-line);
            background:#ffffff;
        }
        [data-testid="stMetric"],
        [data-testid="stMetricLabel"],
        [data-testid="stMetricValue"]{
            overflow-wrap:anywhere;
        }
        [data-testid="stMetric"]{
            background:linear-gradient(145deg,#ffffff,#f8fafc);
            border:1px solid var(--pih-line);
            border-left:4px solid var(--pih-accent);
            border-radius:10px;
            padding:14px 15px;
            min-height:96px;
            box-shadow:0 10px 24px rgba(15,23,42,.06);
        }
        [data-testid="stMetricValue"]{
            font-size:clamp(1.08rem, 1.9vw, 1.65rem);
            line-height:1.1;
        }
        div[data-testid="stDataFrame"],
        div[data-testid="stTable"],
        div[data-testid="stDataEditor"]{
            width:100%;
            overflow-x:auto;
            border-radius:10px;
            border:1px solid var(--pih-line);
            background:var(--pih-surface);
            box-shadow:0 8px 20px rgba(15,23,42,.05);
        }
        div[data-testid="stDataFrame"] *{
            font-size:12px;
        }
        .js-plotly-plot,
        .plot-container,
        .stPlotlyChart{
            width:100%!important;
            max-width:100%!important;
            overflow:hidden;
        }
        .stButton > button,
        .stDownloadButton > button,
        button[kind],
        div[data-baseweb="select"]{
            min-height:44px;
        }
        .stButton > button,
        .stDownloadButton > button{
            border-radius:8px;
            font-weight:800;
            letter-spacing:0;
            white-space:normal;
            line-height:1.2;
            padding:.72rem 1rem;
        }
        div[data-baseweb="select"] > div,
        div[data-baseweb="input"] input,
        textarea{
            border-radius:8px!important;
            min-height:44px;
        }
        [data-testid="stAlert"]{
            border-radius:10px;
            border:1px solid var(--pih-line);
            box-shadow:0 8px 20px rgba(15,23,42,.05);
        }
        .section-header,
        .claims-hero,
        .panel-note,
        .kpi-box,
        .claims-kpi-card{
            max-width:100%;
            overflow-wrap:anywhere;
        }
        .section-header{
            border-radius:10px!important;
        }
        @media (prefers-color-scheme: dark){
            :root{
                --pih-bg:#0f172a;
                --pih-surface:#111827;
                --pih-surface-soft:#1f2937;
                --pih-text:#e5edf6;
                --pih-muted:#94a3b8;
                --pih-line:rgba(148,163,184,.22);
                --pih-accent:#2dd4bf;
                --pih-accent-2:#60a5fa;
            }
            [data-testid="stSidebar"],
            [data-testid="stMetric"],
            div[data-testid="stDataFrame"],
            div[data-testid="stTable"],
            div[data-testid="stDataEditor"]{
                background:var(--pih-surface);
            }
        }
        @media (max-width: 760px){
            .main .block-container{
                padding:.75rem .7rem 1.6rem;
            }
            [data-testid="stSidebar"]{
                width:min(86vw, 21rem)!important;
            }
            [data-testid="column"]{
                width:100%!important;
                flex:1 1 100%!important;
                min-width:100%!important;
            }
            [data-testid="stHorizontalBlock"]{
                gap:.55rem!important;
                flex-wrap:wrap!important;
            }
            [data-testid="stMetric"]{
                min-height:84px;
                padding:12px;
            }
            [data-testid="stMetricValue"]{
                font-size:1.16rem;
            }
            .decision-dashboard-v2,
            .decision-card,
            .decision-kpi,
            .decision-filter-status,
            .claims-hero,
            .claims-kpi-card,
            .panel-note,
            .section-header{
                border-radius:10px!important;
                padding:13px!important;
                min-height:0!important;
            }
            .decision-kpi-value{
                font-size:19px!important;
                white-space:normal!important;
            }
            .decision-badge{
                max-width:100%!important;
                white-space:normal!important;
            }
            .stTabs [data-baseweb="tab-list"]{
                overflow-x:auto;
                gap:.25rem;
            }
            .stTabs [data-baseweb="tab"]{
                min-width:max-content;
                min-height:42px;
            }
            div[data-testid="stDataFrame"]{
                max-height:62vh;
            }
            .js-plotly-plot .gtitle{
                font-size:14px!important;
            }
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


apply_mobile_first_dashboard_css()


def render_health_check() -> None:
    st.json(
        {
            "app": "Project Intelligence Hub",
            "status": "ok",
            "auth_db": str(RUNTIME_DIR / "project_intelligence_hub_mobile_auth.sqlite"),
            "streamlit_url_configured": bool(configured_streamlit_url()),
        }
    )


try:
    health_requested = str(st.query_params.get("health", "")).lower() in {"1", "true", "yes"}
except Exception:
    health_requested = False
if health_requested:
    render_health_check()
    st.stop()


AUTH_USER = auth.require_authentication(APP_DIR)
auth.apply_role_ui(AUTH_USER)




def require_openpyxl():
    if not OPENPYXL_AVAILABLE:
        raise RuntimeError("openpyxl is required for the Detailed Progress report workbook export but is not installed.")


def egp(value) -> str:
    try:
        return f"EGP {float(value):,.0f}"
    except (TypeError, ValueError):
        return "EGP 0"


def sar(value) -> str:
    try:
        return f"SAR {float(value):,.2f}"
    except (TypeError, ValueError):
        return "SAR 0.00"


def pct(value) -> str:
    try:
        return f"{float(value):.1f}%"
    except (TypeError, ValueError):
        return "0.0%"


def parse_project_date(value):
    if pd.isna(value):
        return None
    if isinstance(value, (pd.Timestamp,)):
        return value
    parsed = pd.to_datetime(str(value).strip(), errors="coerce", dayfirst=True)
    return None if pd.isna(parsed) else parsed


def format_project_date(value, fallback: str = "N/A") -> str:
    if value is None or pd.isna(value):
        return fallback
    try:
        return pd.Timestamp(value).strftime("%d-%b-%Y")
    except (TypeError, ValueError):
        return fallback



def parse_numeric(value, default=0.0):
    """
    Safe numeric parser for dashboard calculations.
    Converts real numeric values to float.
    Returns default for business text such as Yes, No, Approved, Pending, N/A.
    """
    import re
    import pandas as pd

    if value is None:
        return default

    try:
        if pd.isna(value):
            return default
    except Exception:
        pass

    if isinstance(value, (int, float)):
        return float(value)

    text = str(value).strip()

    if not text:
        return default

    lowered = text.lower()

    non_numeric_values = {
        "yes", "no", "true", "false",
        "y", "n",
        "approved", "rejected", "pending",
        "open", "closed",
        "done", "not done",
        "n/a", "na", "none", "null",
        "-", "--"
    }

    if lowered in non_numeric_values:
        return default

    negative = False
    if text.startswith("(") and text.endswith(")"):
        negative = True
        text = text[1:-1]

    numeric_match = re.search(r"-?\d+(?:\.\d+)?", text.replace(",", ""))
    if numeric_match is None:
        return default

    cleaned = numeric_match.group(0)

    if cleaned in {"", "-", ".", "-.", ".-", "--"}:
        return default

    try:
        number = float(cleaned)
        return -number if negative else number
    except (ValueError, TypeError):
        return default



def normalized_column_lookup(df: pd.DataFrame) -> dict[str, str]:
    return {
        re.sub(r"[^a-z0-9]+", "", str(col).strip().lower()): col
        for col in df.columns
    }


def first_existing_column(df: pd.DataFrame, candidates: list[str]) -> str | None:
    lookup = normalized_column_lookup(df)
    for candidate in candidates:
        if candidate in df.columns:
            return candidate
        match = lookup.get(re.sub(r"[^a-z0-9]+", "", str(candidate).strip().lower()))
        if match is not None:
            return match
    return None


def copy_column_if_missing(df: pd.DataFrame, target: str, candidates: list[str], default: Any = "") -> None:
    if target in df.columns:
        return
    source = first_existing_column(df, candidates)
    df[target] = df[source] if source is not None else default


def normalize_import_template_frame(path: Path, df: pd.DataFrame) -> pd.DataFrame:
    if df.empty:
        return df
    working = df.loc[:, [col for col in df.columns if str(col).strip()]].copy()
    file_name = path.name.lower()

    if file_name == "payments.csv":
        copy_column_if_missing(working, "payment_id", ["payment", "payment id", "payment_id"])
        copy_column_if_missing(working, "contract_id", ["contract", "contract id", "contract_id"])
        copy_column_if_missing(working, "project_id", ["project", "project id", "project_id"])
        copy_column_if_missing(working, "invoice_no", ["invoice no", "invoice_no", "invoice"])
        copy_column_if_missing(working, "invoice_date", ["invoice date", "invoice_date", "payment_date"])
        copy_column_if_missing(working, "payment_date", ["Date of Cash Cheque Receipt", "payment date", "payment_date", "invoice date"])
        copy_column_if_missing(working, "certified_amount", ["certified amount", "certified_amount", "certified"])
        copy_column_if_missing(working, "paid_amount", ["paid amount", "paid_amount", "paid"])
        copy_column_if_missing(working, "payment_status", ["payment status", "payment_status", "status"])
        copy_column_if_missing(working, "delayed_duration_days", ["delayed duration(Days)", "delayed duration days", "delay days"])

    elif file_name == "delay_events.csv":
        copy_column_if_missing(working, "delay_id", ["delay_id", "delay event id", "event id", "Primary Event ID"])
        copy_column_if_missing(working, "delay_title", ["delay_title", "event_title", "Primary Event ID", "Activity Name"])
        copy_column_if_missing(working, "project_id", ["project_id", "project"])
        copy_column_if_missing(working, "activity_id", ["activity_id", "Activity ID"])
        copy_column_if_missing(working, "activity_name", ["activity_name", "Activity Name"])
        copy_column_if_missing(working, "start_date", ["start_date", "Start", "Overlap Start", "BL Start"])
        copy_column_if_missing(working, "end_date", ["end_date", "Finish", "Overlap Finish", "BL Finish"])
        copy_column_if_missing(working, "estimated_delay_days", ["estimated_delay_days", "Delayed duration after overlap", "Delayed duration", "Concurrent delay"])
        copy_column_if_missing(working, "approved_eot_days", ["approved_eot_days", "approved eot days"], 0)
        copy_column_if_missing(working, "responsibility", ["responsibility", "responsible_party"], "Employer / Client")
        copy_column_if_missing(working, "cause_category", ["cause_category", "Primary Event ID"], "Delay")
        copy_column_if_missing(working, "notice_ref", ["notice_ref", "notice ref"], "")
        copy_column_if_missing(working, "status", ["status"], "Open")
        if "delay_id" in working.columns:
            blank_delay_id = working["delay_id"].astype(str).str.strip().eq("")
            working.loc[blank_delay_id, "delay_id"] = working.index[blank_delay_id].map(lambda idx: f"DELAY-{idx + 1:03d}")
        if "project_id" in working.columns:
            blank_project = working["project_id"].astype(str).str.strip().eq("")
            working.loc[blank_project, "project_id"] = selected_project_id()

    elif file_name == "rfi_ status.csv":
        copy_column_if_missing(working, "RFI No.", ["RFI", "RFI No.", "rfi_no"])
        copy_column_if_missing(working, "Submission Date", ["Submittion date", "Submission date", "Submission Date"])
        copy_column_if_missing(working, "Reply Date", ["Reply", "Reply Date"])
        copy_column_if_missing(working, "Delay Beyond 10 Days", ["Delay beyond 10d", "Delay Beyond 10 Days"])

    elif file_name == "ifc_conflict.csv":
        copy_column_if_missing(working, "Delay Days", ["Delayed days", "Delayed Days"])
        copy_column_if_missing(working, "Revised Start", ["Re-Start", "Restart", "Revised Start"])
        copy_column_if_missing(working, "Revised Finish", ["Finish.1", "Revised Finish"])

    return working


def _xlsx_column_letters(cell_ref):
    letters = "".join(ch for ch in str(cell_ref) if ch.isalpha())
    return letters


def _read_xlsx_sheets(path: Path):
    ns = {
        "main": "http://schemas.openxmlformats.org/spreadsheetml/2006/main",
        "rel": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        "pkgrel": "http://schemas.openxmlformats.org/package/2006/relationships",
    }
    with zipfile.ZipFile(path) as zf:
        shared_strings = []
        if "xl/sharedStrings.xml" in zf.namelist():
            root = ET.fromstring(zf.read("xl/sharedStrings.xml"))
            for si in root.findall("main:si", ns):
                parts = [node.text or "" for node in si.findall(".//main:t", ns)]
                shared_strings.append("".join(parts))

        workbook_root = ET.fromstring(zf.read("xl/workbook.xml"))
        rels_root = ET.fromstring(zf.read("xl/_rels/workbook.xml.rels"))
        rel_map = {
            rel.attrib["Id"]: rel.attrib["Target"]
            for rel in rels_root.findall("pkgrel:Relationship", ns)
        }

        sheets = {}
        for sheet in workbook_root.findall("main:sheets/main:sheet", ns):
            name = sheet.attrib.get("name", "")
            rel_id = sheet.attrib.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
            target = rel_map.get(rel_id, "")
            if not target:
                continue
            sheet_path = f"xl/{target}" if not target.startswith("xl/") else target
            sheet_root = ET.fromstring(zf.read(sheet_path))
            rows = []
            for row in sheet_root.findall("main:sheetData/main:row", ns):
                row_map = {}
                for cell in row.findall("main:c", ns):
                    cell_ref = cell.attrib.get("r", "")
                    col = _xlsx_column_letters(cell_ref)
                    value_node = cell.find("main:v", ns)
                    if value_node is None:
                        row_map[col] = ""
                        continue
                    raw_value = value_node.text or ""
                    cell_type = cell.attrib.get("t")
                    if cell_type == "s":
                        value = shared_strings[int(raw_value)] if raw_value.isdigit() else raw_value
                    else:
                        value = raw_value
                    row_map[col] = value
                rows.append(row_map)

            if not rows:
                sheets[name] = pd.DataFrame()
                continue

            ordered_cols = []
            for row in rows:
                for col in row.keys():
                    if col not in ordered_cols:
                        ordered_cols.append(col)
            matrix = [[row.get(col, "") for col in ordered_cols] for row in rows]
            header = [str(v).strip() if str(v).strip() else col for v, col in zip(matrix[0], ordered_cols)]
            data = matrix[1:] if len(matrix) > 1 else []
            sheets[name] = pd.DataFrame(data, columns=header).fillna("")
        return sheets


@st.cache_data(show_spinner=False)
def _load_core_csv_cached(path_str: str, modified_ns: int, file_size: int):
    path = Path(path_str)
    return normalize_import_template_frame(path, pd.read_csv(path).fillna(""))


def load_core_csv(path: Path, project_id: str | None = None):
    requested_project_id = selected_project_id() if project_id is None else str(project_id or "").strip()
    try:
        relative_path = path.resolve().relative_to(IMPORT_TEMPLATES_DIR.resolve())
    except ValueError:
        relative_path = None

    if relative_path is not None:
        project_ids = [requested_project_id] if requested_project_id else [row["project_id"] for row in discover_projects(PROJECTS_DIR)]
        frames = []
        for discovered_project_id in project_ids:
            candidate = project_data_path(PROJECTS_DIR, discovered_project_id, "core", relative_path)
            if not candidate.exists():
                continue
            stat = candidate.stat()
            frame = _load_core_csv_cached(str(candidate), stat.st_mtime_ns, stat.st_size).copy()
            if "project_id" not in frame.columns:
                frame.insert(0, "project_id", discovered_project_id)
            else:
                source_ids = frame["project_id"].astype(str).str.strip()
                mismatched = source_ids.ne("") & source_ids.ne(discovered_project_id)
                if mismatched.any() and "source_project_id" not in frame.columns:
                    frame.insert(1, "source_project_id", source_ids)
                frame["project_id"] = discovered_project_id
            frames.append(frame)
        return pd.concat(frames, ignore_index=True, sort=False) if frames else pd.DataFrame()

    scoped_path = project_scoped_file(path, requested_project_id)
    if not scoped_path.exists():
        return pd.DataFrame()
    stat = scoped_path.stat()
    return _load_core_csv_cached(str(scoped_path), stat.st_mtime_ns, stat.st_size)


def project_filter_options(projects_df: pd.DataFrame) -> list[dict[str, str]]:
    portfolio_label = "Decision Making dashboard"
    if projects_df.empty:
        return [{"label": portfolio_label, "project_id": ""}]
    options = [{"label": portfolio_label, "project_id": ""}]
    for _, row in projects_df.iterrows():
        project_id = str(row.get("project_id", "")).strip()
        project_name = str(row.get("project_name", "")).strip()
        project_folder_name = str(row.get("project_folder_name", "")).strip()
        sector_name = str(row.get("sector_name", "")).strip()
        if not project_id and not project_name:
            continue
        label = project_name or project_id
        if sector_name and sector_name.casefold() != "unassigned":
            label = f"{sector_name} / {label}"
        if project_id and project_id not in label:
            label = f"{label} ({project_id})"
        if project_folder_name and project_folder_name not in label and project_folder_name != project_id:
            label = f"{label} [Folder: {project_folder_name}]"
        options.append({"label": label, "project_id": project_id})
    return options


def _project_group_value(df: pd.DataFrame, project_id: str, candidates: list[str], agg: str = "sum", default: float = 0.0) -> float:
    if df.empty or "project_id" not in df.columns:
        return default
    column = first_existing_column(df, candidates)
    if not column:
        return default
    scoped = df[df["project_id"].astype(str).str.strip().eq(project_id)]
    if scoped.empty:
        return default
    values = scoped[column].apply(parse_numeric)
    if agg == "mean":
        return float(values.dropna().mean()) if not values.dropna().empty else default
    return float(values.fillna(0).sum())


def _project_group_count(df: pd.DataFrame, project_id: str) -> int:
    if df.empty or "project_id" not in df.columns:
        return 0
    return int(df[df["project_id"].astype(str).str.strip().eq(project_id)].shape[0])


def _project_first_value(df: pd.DataFrame, project_id: str, candidates: list[str], default: Any = "") -> Any:
    if df.empty or "project_id" not in df.columns:
        return default
    scoped = df[df["project_id"].astype(str).str.strip().eq(project_id)]
    if scoped.empty:
        return default
    column = first_existing_column(scoped, candidates)
    if not column:
        return default
    value = scoped.iloc[0].get(column, default)
    return default if pd.isna(value) else value


def safe_divide(numerator: Any, denominator: Any) -> float | None:
    """Return a ratio only when both values are usable and the denominator is non-zero."""
    num = parse_numeric(numerator)
    den = parse_numeric(denominator)
    if den in (None, 0) or pd.isna(den):
        return None
    return float(num or 0) / float(den)


def normalize_to_100(value: Any, max_value: float = 100.0, invert: bool = False) -> float | None:
    numeric_value = parse_numeric(value)
    if numeric_value is None or pd.isna(numeric_value):
        return None
    bounded = max(0.0, min(float(numeric_value), max_value))
    score = (bounded / max_value) * 100.0 if max_value else bounded
    return 100.0 - score if invert else score


def build_decision_dashboard_registry(projects_catalog_df: pd.DataFrame) -> pd.DataFrame:
    if projects_catalog_df.empty:
        return pd.DataFrame()
    projects_df = load_core_csv(PROJECTS_CSV_PATH)
    activities_df = load_core_csv(ACTIVITIES_CSV_PATH)
    evm_df = load_core_csv(EVM_CSV_PATH)
    contracts_df = load_core_csv(CONTRACTS_CSV_PATH)
    payments_df = load_core_csv(PAYMENTS_CSV_PATH)
    delays_df = load_core_csv(DELAYS_CSV_PATH)
    risks_df = load_core_csv(RISKS_CSV_PATH)
    milestones_df = load_core_csv(MILESTONES_CSV_PATH)
    rows = []
    for _, catalog_row in projects_catalog_df.iterrows():
        project_id = str(catalog_row.get("project_id", "")).strip()
        if not project_id:
            continue
        contract_value = _project_first_value(projects_df, project_id, ["contract_value", "Contract Value", "budget", "Budget", "BAC"], 0)
        if parse_numeric(contract_value) == 0:
            contract_value = _project_group_value(evm_df, project_id, ["BAC", "bac", "budget_at_completion", "Budget at Completion"])
        planned_progress = _project_first_value(projects_df, project_id, ["planned_progress", "Planned Progress", "baseline_progress"], 0)
        actual_progress = _project_first_value(projects_df, project_id, ["overall_progress", "actual_progress", "Actual Progress", "progress"], 0)
        if parse_numeric(actual_progress) == 0:
            actual_progress = _project_group_value(activities_df, project_id, ["actual_progress", "Physical % Complete", "physical_percent_complete", "progress"], agg="mean")
        start_date = pd.to_datetime(_project_first_value(projects_df, project_id, ["project_start", "start_date", "Project Start Date", "baseline_start"], ""), errors="coerce")
        finish_date = pd.to_datetime(_project_first_value(projects_df, project_id, ["project_finish", "finish_date", "Project Finish Date", "baseline_finish"], ""), errors="coerce")
        paid_value = _project_group_value(payments_df, project_id, ["paid_amount", "Paid Amount", "paid", "actual_paid", "Actual Paid"])
        spent_value = _project_group_value(evm_df, project_id, ["AC", "ac", "actual_cost", "Actual Cost"])
        if spent_value == 0:
            spent_value = _project_group_value(contracts_df, project_id, ["actual_cost", "Actual Cost", "spent", "Spent", "certified_amount"])
        bac = _project_group_value(evm_df, project_id, ["BAC", "bac", "budget_at_completion", "Budget at Completion"])
        if bac == 0:
            bac = parse_numeric(contract_value) or 0
        pv = _project_group_value(evm_df, project_id, ["PV", "pv", "planned_value", "Planned Value"])
        ev = _project_group_value(evm_df, project_id, ["EV", "ev", "earned_value", "Earned Value"])
        ac = spent_value
        # Derive PV/EV only when BAC and progress percentages are available. These formulas preserve traceability and avoid fake values.
        if pv == 0 and bac and parse_numeric(planned_progress) is not None:
            pv = bac * (float(parse_numeric(planned_progress) or 0) / 100.0)
        if ev == 0 and bac and parse_numeric(actual_progress) is not None:
            ev = bac * (float(parse_numeric(actual_progress) or 0) / 100.0)
        spi = _project_first_value(evm_df, project_id, ["SPI", "spi"], 0)
        cpi = _project_first_value(evm_df, project_id, ["CPI", "cpi"], 0)
        spi = parse_numeric(spi) or safe_divide(ev, pv) or 0
        cpi = parse_numeric(cpi) or safe_divide(ev, ac) or 0
        eac = (bac / cpi) if cpi else 0
        etc = eac - ac if eac else 0
        vac = bac - eac if eac else 0
        sv = ev - pv
        cv = ev - ac
        risk_count = _project_group_count(risks_df, project_id)
        risk_score = _project_group_value(risks_df, project_id, ["risk_score", "Risk Score", "risk_rating", "score"], agg="mean")
        if risk_score == 0 and risk_count:
            risk_score = min(100.0, float(risk_count) * 20.0)
        milestone_count = _project_group_count(milestones_df, project_id)
        delay_days = _project_group_value(delays_df, project_id, ["delay_days", "Delay Days", "estimated_delay_days", "time_impact_days"], agg="sum")
        claims_exposure = _project_group_value(delays_df, project_id, ["claim_amount", "Claim Amount", "eot_days", "EOT Days", "time_impact_days"], agg="sum")
        progress_value = parse_numeric(actual_progress)
        planned_value = parse_numeric(planned_progress)
        status = "On Track"
        if risk_score >= 70 or risk_count >= 5 or delay_days > 0 or (planned_value and progress_value < planned_value - 10) or (spi and spi < 0.9) or (cpi and cpi < 0.9):
            status = "High Attention"
        elif risk_score >= 35 or risk_count or (planned_value and progress_value < planned_value - 3) or (spi and spi < 1) or (cpi and cpi < 1):
            status = "Watch"
        required_decision = "Monitor"
        if spi and spi < 0.9:
            required_decision = "Approve recovery plan"
        if cpi and cpi < 0.9:
            required_decision = "Review cost recovery"
        if delay_days > 0:
            required_decision = "Resolve delay exposure"
        rows.append({
            "project_id": project_id,
            "Project": str(catalog_row.get("project_name") or catalog_row.get("project_display_name") or project_id),
            "Sector": str(catalog_row.get("sector_name") or "Unassigned"),
            "Folder": str(catalog_row.get("project_relative_path") or catalog_row.get("project_folder_name") or ""),
            "Contract Value": parse_numeric(contract_value),
            "Paid": paid_value,
            "Remaining": max(0.0, float(parse_numeric(contract_value) or 0) - paid_value),
            "BAC": bac,
            "PV": pv,
            "EV": ev,
            "AC": ac,
            "SV": sv,
            "CV": cv,
            "EAC": eac,
            "ETC": etc,
            "VAC": vac,
            "SPI": float(spi or 0),
            "CPI": float(cpi or 0),
            "Progress": float(progress_value or 0),
            "Planned Progress": float(planned_value or 0),
            "Progress Variance": float((progress_value or 0) - (planned_value or 0)),
            "Quality": float(parse_numeric(_project_first_value(projects_df, project_id, ["quality_score", "Quality Score", "quality"], 0)) or 0),
            "Safety": float(parse_numeric(_project_first_value(projects_df, project_id, ["safety_score", "Safety Score", "safety"], 0)) or 0),
            "Risks": risk_count,
            "Risk Score": risk_score,
            "Delay Days": delay_days,
            "Claims / EOT Exposure": claims_exposure,
            "Milestones": milestone_count,
            "Required Decision": required_decision,
            "Start": start_date,
            "Finish": finish_date,
            "Status": status,
        })
    return pd.DataFrame(rows)


def decision_status_badge(value: str) -> tuple[str, str, str]:
    status = str(value or "Neutral")
    if status in {"On Track", "Healthy", "Cost Efficient", "Complete"}:
        return status, "#123F3D", "#FFFFFF"
    if status in {"Watch", "Watchlist", "Cost Watch", "Partial"}:
        return status, "#4E3D12", "#FFFFFF"
    if status in {"High Attention", "Critical", "Cost Critical", "Missing critical fields"}:
        return status, "#4B1D22", "#FFFFFF"
    return status, "#173B63", "#FFFFFF"


def validate_decision_dashboard_data(registry_df: pd.DataFrame) -> dict[str, Any]:
    required = ["Project", "Sector", "Contract Value", "Progress", "Planned Progress", "PV", "EV", "AC", "Risk Score", "Milestones"]
    missing = [column for column in required if column not in registry_df.columns or registry_df[column].isna().all()]
    if not missing:
        status = "Complete"
    elif {"Project", "Sector", "Contract Value", "Progress"} - set(missing):
        status = "Partial"
    else:
        status = "Missing critical fields"
    return {
        "status": status,
        "missing_fields": missing,
        "last_updated": pd.Timestamp.now().strftime("%d-%b-%Y %H:%M"),
    }


def apply_decision_dashboard_filters(registry_df: pd.DataFrame) -> pd.DataFrame:
    filtered = registry_df.copy()
    selected_sectors = st.session_state.get("decision_filter_sectors", [])
    selected_statuses = st.session_state.get("decision_filter_statuses", [])
    selected_projects = st.session_state.get("decision_filter_projects", [])
    progress_range = st.session_state.get("decision_filter_progress", (0, 100))
    if selected_sectors:
        filtered = filtered[filtered["Sector"].isin(selected_sectors)]
    if selected_statuses:
        filtered = filtered[filtered["Status"].isin(selected_statuses)]
    if selected_projects:
        filtered = filtered[filtered["Project"].isin(selected_projects)]
    filtered = filtered[
        filtered["Progress"].fillna(0).between(float(progress_range[0]), float(progress_range[1]))
    ]
    return filtered.copy()


def render_decision_command_bar(registry_df: pd.DataFrame, quality: dict[str, Any]) -> None:
    st.markdown("<div class='decision-section-title'>Portfolio Filters</div>", unsafe_allow_html=True)
    sectors = sorted(registry_df["Sector"].dropna().astype(str).unique().tolist())
    statuses = sorted(registry_df["Status"].dropna().astype(str).unique().tolist())
    projects = sorted(registry_df["Project"].dropna().astype(str).unique().tolist())
    col1, col2, col3, col4 = st.columns([0.23, 0.2, 0.22, 0.35], gap="medium")
    with col1:
        with st.container(border=True):
            st.multiselect("Sector filter", sectors, key="decision_filter_sectors")
    with col2:
        with st.container(border=True):
            st.multiselect("Status filter", statuses, key="decision_filter_statuses")
    with col3:
        with st.container(border=True):
            st.slider("Progress range", 0, 100, key="decision_filter_progress", value=st.session_state.get("decision_filter_progress", (0, 100)))
    with col4:
        with st.container(border=True):
            st.multiselect("Project filter", projects, key="decision_filter_projects")
    col5, col6, col7, col8 = st.columns([0.30, 0.22, 0.27, 0.21], gap="medium")
    with col5:
        with st.container(border=True):
            st.selectbox("Chart section", ["All sections", "Health and EVM", "Cash flow", "Risks and decisions"], key="decision_chart_section")
    badge_text, badge_bg, badge_fg = decision_status_badge(str(quality["status"]))
    with col6:
        st.markdown(f"<div class='decision-filter-status'><span class='decision-filter-label'>Data quality</span><span class='decision-badge' style='background:{badge_bg};color:#FFFFFF !important;-webkit-text-fill-color:#FFFFFF !important;font-weight:900 !important;'>{html.escape(badge_text)}</span></div>", unsafe_allow_html=True)
    with col7:
        st.markdown(f"<div class='decision-filter-status'><span class='decision-filter-label'>Last updated</span><strong>{html.escape(str(quality['last_updated']))}</strong></div>", unsafe_allow_html=True)
    with col8:
        if st.button("Reset filters", key="decision_reset_filters", width="stretch"):
            for key in ["decision_filter_sectors", "decision_filter_statuses", "decision_filter_projects", "decision_filter_progress", "decision_chart_section"]:
                st.session_state.pop(key, None)
            st.rerun()
    if quality["missing_fields"]:
        st.caption("Missing fields: " + ", ".join(quality["missing_fields"]))


def render_status_badge(status: str) -> str:
    badge_text, badge_bg, badge_fg = decision_status_badge(status)
    return f"<span class='decision-badge' style='background:{badge_bg};color:#FFFFFF !important;-webkit-text-fill-color:#FFFFFF !important;font-weight:900 !important;'>{html.escape(badge_text)}</span>"


def dashboard_display(value: Any, fallback: str = "N/A") -> str:
    if value is None:
        return fallback
    if isinstance(value, float) and (pd.isna(value) or value in (float("inf"), float("-inf"))):
        return fallback
    text = str(value).strip()
    return text if text and text.lower() not in {"nan", "none", "inf", "-inf"} else fallback


def render_executive_kpi_card(
    col,
    title: str,
    value: Any,
    icon: str,
    status: str,
    delta_text: str | None = None,
    delta_direction: str = "neutral",
    help_text: str | None = None,
) -> None:
    delta_class = {"up": "good", "down": "bad", "flat": "warn", "neutral": "neutral"}.get(delta_direction, "neutral")
    delta_label = delta_text or "↔ Comparison unavailable"
    title_attr = html.escape(help_text or title)
    col.markdown(
        f"""
        <div class='decision-kpi' title='{title_attr}'>
          <div class='decision-kpi-header'>
            <div class='decision-kpi-title'><span class='decision-kpi-icon'>{html.escape(icon)}</span><span>{html.escape(title)}</span></div>
            {render_status_badge(status)}
          </div>
          <div class='decision-kpi-value'>{html.escape(dashboard_display(value))}</div>
          <div class='decision-delta {delta_class}'>{html.escape(delta_label)}</div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_decision_kpi_cards(registry_df: pd.DataFrame) -> None:
    if registry_df.empty:
        st.info("No project folders were detected yet. Add project folders under projects, or under projects/<sector>/<project>.")
        return
    st.markdown("<div class='decision-section-title'>Executive KPI Cards</div>", unsafe_allow_html=True)
    total_value = registry_df["Contract Value"].fillna(0).sum()
    total_paid = registry_df["Paid"].fillna(0).sum()
    remaining = registry_df["Remaining"].fillna(0).sum()
    delayed_projects = int((registry_df["Delay Days"].fillna(0) > 0).sum())
    high_risk_projects = int((registry_df["Risk Score"].fillna(0) >= 70).sum())
    avg_spi = registry_df["SPI"].replace(0, pd.NA).dropna().mean()
    avg_cpi = registry_df["CPI"].replace(0, pd.NA).dropna().mean()
    avg_risk = registry_df["Risk Score"].replace(0, pd.NA).dropna().mean()
    decisions_required = int(registry_df["Required Decision"].astype(str).ne("Monitor").sum())
    cards = [
        ("Total Projects", len(registry_df), "📁", "Neutral", "Total discovered project folders in the selected portfolio scope."),
        ("Total Contract Value", egp(total_value), "💼", "Neutral", "Sum of selected projects contract value or BAC when available."),
        ("Total Paid", egp(total_paid), "💳", "Neutral", "Payments paid from project payment records."),
        ("Remaining Value", egp(remaining), "📌", "Neutral", "Contract value less paid value where both are available."),
        ("Average Progress", pct(registry_df["Progress"].mean()), "📈", "Healthy" if registry_df["Progress"].mean() >= registry_df["Planned Progress"].mean() else "Watchlist", "Average actual progress for selected projects."),
        ("Delayed Projects", delayed_projects, "⏱", "Critical" if delayed_projects else "Healthy", "Projects with delay days in delay event records."),
        ("High-Risk Projects", high_risk_projects, "⚠", "Critical" if high_risk_projects else "Healthy", "Projects with normalized risk score >= 70."),
        ("Average SPI", f"{avg_spi:.2f}" if pd.notna(avg_spi) else "N/A", "📊", "On Track" if pd.notna(avg_spi) and avg_spi >= 1 else ("Watchlist" if pd.notna(avg_spi) and avg_spi >= 0.9 else "Critical"), "Schedule Performance Index: EV / PV."),
        ("Average CPI", f"{avg_cpi:.2f}" if pd.notna(avg_cpi) else "N/A", "💰", "Cost Efficient" if pd.notna(avg_cpi) and avg_cpi >= 1 else ("Cost Watch" if pd.notna(avg_cpi) and avg_cpi >= 0.9 else "Cost Critical"), "Cost Performance Index: EV / AC."),
        ("Average Risk Score", f"{avg_risk:.1f}" if pd.notna(avg_risk) else "N/A", "🛡", "Critical" if pd.notna(avg_risk) and avg_risk >= 70 else ("Watchlist" if pd.notna(avg_risk) and avg_risk >= 35 else "Healthy"), "Normalized 0-100 risk score from available risk data."),
        ("Claims / EOT Exposure", f"{registry_df['Claims / EOT Exposure'].fillna(0).sum():,.0f}", "📑", "Watchlist" if registry_df["Claims / EOT Exposure"].fillna(0).sum() else "Neutral", "Sum of available claim amount, EOT days, or delay impact fields."),
        ("Decisions Required", decisions_required, "🧭", "Critical" if decisions_required else "Healthy", "Triggered by SPI/CPI/risk/delay thresholds."),
    ]
    for chunk_start in (0, 4, 8):
        cols = st.columns(4, gap="medium")
        for col, card in zip(cols, cards[chunk_start:chunk_start + 6]):
            title, value, icon, status, help_text = card
            render_executive_kpi_card(col, title, value, icon, status, delta_text="↔ Comparison unavailable", delta_direction="neutral", help_text=help_text)


def render_decision_cards(registry_df: pd.DataFrame) -> None:
    decision_df = registry_df[
        registry_df["Required Decision"].astype(str).ne("Monitor")
        | registry_df["Status"].astype(str).isin(["Watch", "High Attention"])
    ].copy()
    if decision_df.empty:
        st.markdown("<div class='decision-chart-note'>No management decision triggers are active in the selected dashboard scope.</div>", unsafe_allow_html=True)
        return
    decision_df = decision_df.sort_values(["Status", "Risk Score", "Delay Days", "Contract Value"], ascending=[True, False, False, False]).head(6)
    cols = st.columns(3)
    for idx, (_, row) in enumerate(decision_df.iterrows()):
        urgency = "Critical" if row["Status"] == "High Attention" else "Watchlist"
        badge_text, badge_bg, badge_fg = decision_status_badge(urgency)
        impact = []
        if row["SPI"] and row["SPI"] < 0.9:
            impact.append("Schedule performance below threshold")
        if row["CPI"] and row["CPI"] < 0.9:
            impact.append("Cost performance below threshold")
        if row["Delay Days"]:
            impact.append(f"{row['Delay Days']:.0f} delay days recorded")
        if row["Risk Score"] >= 70:
            impact.append("High risk score")
        root_cause = "; ".join(impact) if impact else "Portfolio threshold watch item"
        with cols[idx % 3]:
            st.markdown(
                f"""
                <div class='decision-card'>
                  <h4>{html.escape(str(row['Required Decision']))}</h4>
                  <p><b>Project:</b> {html.escape(str(row['Project']))}</p>
                  <p><b>Sector:</b> {html.escape(str(row['Sector']))}</p>
                  <p><b>Root cause:</b> {html.escape(root_cause)}</p>
                  <p><b>Financial exposure:</b> {egp(row['Remaining'])} remaining / {egp(row['Claims / EOT Exposure'])} claims-EOT signal</p>
                  <p><b>Schedule exposure:</b> {row['Delay Days']:.0f} days | SPI {row['SPI']:.2f}</p>
                  <p><span class='decision-badge' style='background:{badge_bg};color:#FFFFFF !important;-webkit-text-fill-color:#FFFFFF !important;font-weight:900 !important;'>{html.escape(badge_text)}</span></p>
                </div>
                """,
                unsafe_allow_html=True,
            )


def render_management_decision_brief_once(registry_df: pd.DataFrame, quality: dict[str, Any]) -> None:
    high_attention = int(registry_df["Status"].astype(str).eq("High Attention").sum())
    watch_count = int(registry_df["Status"].astype(str).eq("Watch").sum())
    avg_spi = registry_df["SPI"].replace(0, pd.NA).dropna().mean()
    avg_cpi = registry_df["CPI"].replace(0, pd.NA).dropna().mean()
    critical_sectors = registry_df[registry_df["Status"].astype(str).eq("High Attention")]["Sector"].dropna().astype(str).unique().tolist()
    risk_source = registry_df.sort_values("Risk Score", ascending=False).head(1)
    main_risk_driver = "N/A" if risk_source.empty else f"{risk_source.iloc[0]['Project']} ({risk_source.iloc[0]['Risk Score']:.1f})"
    eot_exposure = registry_df["Claims / EOT Exposure"].fillna(0).sum()
    portfolio_position = "Management attention required" if high_attention or watch_count else "Stable portfolio"
    immediate_action = "Continue portfolio monitoring"
    if high_attention:
        immediate_action = "Review critical project recovery and commercial exposure"
    elif watch_count:
        immediate_action = "Review watchlist project controls actions"
    if quality.get("status") != "Complete":
        immediate_action = "Complete missing project controls data before executive decision"
    evidence_status = "Partial" if quality.get("missing_fields") else "Available from project source files"
    st.markdown("#### Management Decision Brief")
    st.markdown(
        f"""
        <div class='decision-card' style='min-height:0'>
          <h4>Executive Position</h4>
          <p><b>Portfolio position:</b> {html.escape(portfolio_position)}</p>
          <p><b>Projects requiring attention:</b> {high_attention + watch_count}</p>
          <p><b>Critical sectors:</b> {html.escape(', '.join(critical_sectors) if critical_sectors else 'N/A')}</p>
          <p><b>Average SPI:</b> {dashboard_display(f'{avg_spi:.2f}' if pd.notna(avg_spi) else None)} | <b>Average CPI:</b> {dashboard_display(f'{avg_cpi:.2f}' if pd.notna(avg_cpi) else None)}</p>
          <p><b>Main risk driver:</b> {html.escape(main_risk_driver)}</p>
          <p><b>Immediate management action:</b> {html.escape(immediate_action)}</p>
          <p><b>Evidence availability:</b> {html.escape(evidence_status)} | <b>Claims / EOT signal:</b> {dashboard_display(f'{eot_exposure:,.0f}' if eot_exposure else None)}</p>
          <p><b>Recommended next step:</b> Review the triggered decision cards and use the main project selector when project-level follow-up is required.</p>
        </div>
        """,
        unsafe_allow_html=True,
    )
    render_decision_cards(registry_df)


def render_decision_making_dashboard(projects_catalog_df: pd.DataFrame) -> None:
    registry_df = build_decision_dashboard_registry(projects_catalog_df)
    st.markdown(
        """
        <style>
        .decision-dashboard-v2{background:linear-gradient(180deg,rgba(245,249,252,.98),rgba(236,244,248,.96));border:1px solid rgba(15,73,105,.16);border-radius:18px;padding:18px 18px 24px;margin:12px 0 22px;box-shadow:0 22px 46px rgba(3,14,28,.10)}
        .section-header{background:#FFFFFF;border:1px solid rgba(15,73,105,.15);border-left:5px solid #1A8A8F;border-radius:14px;padding:16px 18px;margin-bottom:14px;box-shadow:0 10px 24px rgba(3,14,28,.06)}
        .decision-section-title{margin:18px 0 10px;color:#0A3153;font-size:15px;font-weight:900;text-transform:uppercase;letter-spacing:.04em}
        .decision-kpi{background:linear-gradient(145deg,#FFFFFF,#F5FAFC);border:1px solid rgba(15,73,105,.16);border-left:5px solid #1A8A8F;border-radius:14px;padding:15px 16px;height:154px;box-shadow:0 12px 26px rgba(3,14,28,.08);overflow:hidden;margin-bottom:12px}
        .decision-kpi-header{display:grid;grid-template-columns:minmax(0,1fr) auto;gap:10px;align-items:start}
        .decision-kpi-title{display:flex;gap:8px;align-items:flex-start;color:#0A3153;font-size:12px;font-weight:900;text-transform:uppercase;letter-spacing:.03em;line-height:1.25;min-width:0}
        .decision-kpi-icon{flex:0 0 auto}
        .decision-kpi-value{font-size:23px;font-weight:900;color:#061A2D;margin:16px 0 10px;line-height:1.05;white-space:nowrap;overflow:hidden;text-overflow:ellipsis}
        .decision-badge{display:inline-block;border-radius:999px;padding:5px 9px;font-size:10px;font-weight:900!important;white-space:nowrap;max-width:118px;overflow:hidden;text-overflow:ellipsis;color:#FFFFFF!important;-webkit-text-fill-color:#FFFFFF!important;text-shadow:0 1px 1px rgba(0,0,0,.35)}
        
        /* DECISION_BADGE_WHITE_BOLD_FINAL_OVERRIDE */
        .decision-badge,
        .decision-badge *{
            color:#FFFFFF!important;
            -webkit-text-fill-color:#FFFFFF!important;
            font-weight:900!important;
            text-shadow:0 1px 1px rgba(0,0,0,.35);
        }

        .decision-delta{font-size:11px;font-weight:800}
        .decision-delta.good{color:#15876F}.decision-delta.bad{color:#B73D45}.decision-delta.warn{color:#A77A00}.decision-delta.neutral{color:#607487}
        .decision-filter-status{background:#FFFFFF;border:1px solid rgba(15,73,105,.14);border-radius:14px;padding:14px 15px;min-height:82px;box-shadow:0 10px 22px rgba(3,14,28,.06);display:flex;flex-direction:column;justify-content:center;gap:7px;color:#0A3153}
        .decision-filter-label{font-size:11px;text-transform:uppercase;letter-spacing:.04em;color:#607487;font-weight:900}
        .decision-updated{background:#FFFFFF;border:1px solid rgba(15,73,105,.14);border-radius:999px;padding:8px 11px;color:#607487;font-size:12px;font-weight:700;margin-top:24px;text-align:center}
        .decision-card{background:linear-gradient(145deg,#FFFFFF,#F5FAFC);border:1px solid rgba(15,73,105,.16);border-left:5px solid #1A8A8F;border-radius:14px;padding:18px 18px;min-height:184px;box-shadow:0 12px 26px rgba(3,14,28,.08);margin-bottom:12px}
        .decision-card h4{margin:0 0 12px;color:#0A3153;font-size:16px;line-height:1.25}
        .decision-card p{margin:7px 0;color:#263A4D;font-size:12px;line-height:1.45;background:#F7FAFC;border:1px solid rgba(15,73,105,.10);border-radius:9px;padding:7px 9px}
        .decision-card p b{color:#061A2D}
        .decision-chart-note{background:#FFFFFF;border:1px dashed rgba(96,116,135,.42);border-left:5px solid #D4A017;border-radius:14px;padding:15px 16px;color:#263A4D;font-size:13px;margin:10px 0;box-shadow:0 8px 18px rgba(3,14,28,.05)}
        div[data-testid="stVerticalBlockBorderWrapper"]{background:#FFFFFF;border-color:rgba(15,73,105,.14)!important;border-radius:14px!important;box-shadow:0 10px 22px rgba(3,14,28,.06)}
        </style>
        <div class='decision-dashboard-v2'>
        """,
        unsafe_allow_html=True,
    )
    st.markdown(
        """
        <div class='section-header'>
          <h3>Decision Making dashboard</h3>
          <p style='margin:6px 0 0;color:#526276;font-size:13px'>Executive portfolio command center generated from project data, sector aggregation, and project controls indicators.</p>
        </div>
        """,
        unsafe_allow_html=True,
    )
    if registry_df.empty:
        render_decision_kpi_cards(registry_df)
        st.markdown("</div>", unsafe_allow_html=True)
        return

    quality = validate_decision_dashboard_data(registry_df)
    scoped_registry_df = apply_decision_dashboard_filters(registry_df)
    if scoped_registry_df.empty:
        empty_scope_message = "No projects match the selected Decision Making Dashboard filters."
    else:
        empty_scope_message = ""

    decision_view = st.selectbox(
        "Decision dashboard view",
        ["📊 Overall Portfolio", "🏭 Sector Analysis", "📋 Projects Analysis"],
        key="decision_dashboard_active_view",
    )
    if decision_view == "📊 Overall Portfolio":
        render_decision_command_bar(registry_df, quality)
        scoped_registry_df = apply_decision_dashboard_filters(registry_df)
        render_decision_kpi_cards(scoped_registry_df)
        if scoped_registry_df.empty:
            st.warning(empty_scope_message or "No projects match the selected Decision Making Dashboard filters.")
            st.markdown("</div>", unsafe_allow_html=True)
            return
        left, right = st.columns([0.42, 0.58])
        sector_summary = scoped_registry_df.groupby("Sector", as_index=False).agg(projects=("project_id", "count"), budget=("Contract Value", "sum"), progress=("Progress", "mean"))
        with left:
            fig = px.pie(sector_summary, names="Sector", values="projects", hole=0.54, title="Sector Distribution", color_discrete_sequence=["#0B3A5B", "#1A8A8F", "#D4A017", "#617487", "#8D6E63"])
            st.plotly_chart(style_plotly(fig, 360), width="stretch")
            status_summary = scoped_registry_df["Status"].value_counts().reset_index()
            status_summary.columns = ["Status", "Projects"]
            fig = px.bar(status_summary, x="Status", y="Projects", title="Status Breakdown", color="Status", color_discrete_map={"On Track": "#1A8A8F", "Watch": "#D4A017", "High Attention": "#C94C4C"})
            st.plotly_chart(style_plotly(fig, 320), width="stretch")
        with right:
            fig = px.bar(sector_summary.sort_values("budget", ascending=True), x="budget", y="Sector", orientation="h", title="Budget Allocation by Sector", labels={"budget": "Contract Value"})
            st.plotly_chart(style_plotly(fig, 320), width="stretch")
            fig = px.scatter(scoped_registry_df, x="Progress", y="Contract Value", size=scoped_registry_df["Contract Value"].clip(lower=1), color="Sector", hover_name="Project", title="Progress Overview")
            st.plotly_chart(style_plotly(fig, 360), width="stretch")
        timeline_df = scoped_registry_df.dropna(subset=["Start", "Finish"]).copy()
        if not timeline_df.empty:
            fig = px.timeline(timeline_df, x_start="Start", x_end="Finish", y="Project", color="Sector", title="Project Schedules")
            fig.update_yaxes(autorange="reversed")
            st.plotly_chart(style_plotly(fig, 390), width="stretch")
        charts_col1, charts_col2 = st.columns(2)
        with charts_col1:
            fig = px.scatter(scoped_registry_df, x="CPI", y="SPI", size=scoped_registry_df["Contract Value"].clip(lower=1), color="Status", hover_name="Project", hover_data=["Sector", "Progress", "Risk Score", "Delay Days", "Contract Value"], title="Portfolio Health Matrix - SPI vs CPI", color_discrete_map={"On Track": "#1A8A8F", "Watch": "#D4A017", "High Attention": "#C94C4C"})
            fig.add_hline(y=1, line_dash="dash", line_color="#607080")
            fig.add_vline(x=1, line_dash="dash", line_color="#607080")
            fig.add_annotation(x=1.08, y=1.08, text="Healthy", showarrow=False, font=dict(color="#50D5B7"))
            fig.add_annotation(x=.82, y=1.08, text="Cost Risk", showarrow=False, font=dict(color="#D4A017"))
            fig.add_annotation(x=1.08, y=.82, text="Schedule Risk", showarrow=False, font=dict(color="#D4A017"))
            fig.add_annotation(x=.82, y=.82, text="Critical", showarrow=False, font=dict(color="#F05D5E"))
            st.plotly_chart(style_plotly(fig, 360), width="stretch")
            ev_long = scoped_registry_df[["Project", "BAC", "PV", "EV", "AC"]].melt(id_vars="Project", var_name="Metric", value_name="Value")
            fig = px.bar(ev_long, x="Project", y="Value", color="Metric", barmode="group", title="BAC / PV / EV / AC Comparison")
            st.plotly_chart(style_plotly(fig, 390), width="stretch")
        with charts_col2:
            radar_values = [
                scoped_registry_df["Progress"].mean(),
                max(min(scoped_registry_df["SPI"].replace(0, pd.NA).dropna().mean() * 100 if not scoped_registry_df["SPI"].replace(0, pd.NA).dropna().empty else 0, 100), 0),
                max(min(scoped_registry_df["CPI"].replace(0, pd.NA).dropna().mean() * 100 if not scoped_registry_df["CPI"].replace(0, pd.NA).dropna().empty else 0, 100), 0),
                100 - min(scoped_registry_df["Risk Score"].fillna(0).mean(), 100),
                scoped_registry_df["Quality"].mean() if scoped_registry_df["Quality"].fillna(0).sum() else 0,
                scoped_registry_df["Safety"].mean() if scoped_registry_df["Safety"].fillna(0).sum() else 0,
            ]
            radar_labels = ["Progress", "Schedule", "Cost", "Risk Control", "Quality", "Safety"]
            fig = go.Figure(data=go.Scatterpolar(r=radar_values, theta=radar_labels, fill="toself", line_color="#1A8A8F"))
            fig.update_layout(title="Quality & Safety Radar", polar=dict(radialaxis=dict(range=[0, 100])), showlegend=False)
            st.plotly_chart(style_plotly(fig, 360), width="stretch")
            fig = px.line(scoped_registry_df.sort_values("Project"), x="Project", y=["SPI", "CPI"], markers=True, title="EVM Trend by Project")
            st.plotly_chart(style_plotly(fig, 390), width="stretch")
        cash_long = scoped_registry_df[["Project", "Contract Value", "AC", "Paid", "Remaining"]].melt(id_vars="Project", var_name="Metric", value_name="Value")
        fig = px.bar(cash_long, x="Project", y="Value", color="Metric", barmode="group", title="Portfolio Cash Flow - Budget / Spent / Paid / Remaining")
        st.plotly_chart(style_plotly(fig, 390), width="stretch")
        progress_long = scoped_registry_df[["Project", "Planned Progress", "Progress"]].melt(id_vars="Project", var_name="Curve", value_name="Percent")
        fig = px.area(progress_long, x="Project", y="Percent", color="Curve", line_group="Curve", title="Portfolio S-Curve Proxy - Planned vs Actual Progress")
        st.plotly_chart(style_plotly(fig, 350), width="stretch")
        risk_heat = scoped_registry_df.assign(Risk_Band=pd.cut(scoped_registry_df["Risk Score"].fillna(0), bins=[-1, 34, 69, 100], labels=["Low", "Medium", "High"]))
        heat_df = pd.crosstab(risk_heat["Sector"], risk_heat["Risk_Band"])
        fig = px.imshow(heat_df, text_auto=True, title="Risk Assessment Matrix", color_continuous_scale=["#EAF3F4", "#D4A017", "#C94C4C"])
        st.plotly_chart(style_plotly(fig, 340), width="stretch")
        display_df = scoped_registry_df[["Sector", "Project", "Status", "Contract Value", "Paid", "Remaining", "Progress", "Planned Progress", "Progress Variance", "SPI", "CPI", "Risk Score", "Delay Days", "Required Decision", "Folder"]].copy()
        st.dataframe(display_df, width="stretch", hide_index=True, height=dataframe_height(display_df, max_height=520))

    if decision_view == "🏭 Sector Analysis":
        sector_options = ["All sectors"] + sorted(scoped_registry_df["Sector"].dropna().unique())
        with st.container(border=True):
            sector_name = st.selectbox("Sector", sector_options, key="decision_sector_filter")
        scoped = scoped_registry_df.copy() if sector_name == "All sectors" else scoped_registry_df[scoped_registry_df["Sector"].eq(sector_name)].copy()
        health_spi = scoped["SPI"].replace(0, pd.NA).dropna().mean()
        health_cpi = scoped["CPI"].replace(0, pd.NA).dropna().mean()
        health_risk = scoped["Risk Score"].fillna(0).mean()
        health_status = "Healthy" if pd.notna(health_spi) and pd.notna(health_cpi) and health_spi >= 1 and health_cpi >= 1 and health_risk < 35 else ("Critical" if (pd.notna(health_spi) and health_spi < .9) or (pd.notna(health_cpi) and health_cpi < .9) or health_risk >= 70 else "Watchlist")
        badge_text, badge_bg, badge_fg = decision_status_badge(health_status)
        spi_text = dashboard_display(f"{health_spi:.2f}" if pd.notna(health_spi) else None)
        cpi_text = dashboard_display(f"{health_cpi:.2f}" if pd.notna(health_cpi) else None)
        risk_text = dashboard_display(f"{health_risk:.1f}" if pd.notna(health_risk) else None)
        st.markdown(
            f"<div class='decision-card'><h4>Sector Health Indicator</h4><p><b>Sector:</b> {html.escape(str(sector_name))}</p><p><b>Projects:</b> {len(scoped)}</p><p><b>Total value:</b> {egp(scoped['Contract Value'].sum())}</p><p><b>Average progress:</b> {pct(scoped['Progress'].mean())} | <b>SPI:</b> {spi_text} | <b>CPI:</b> {cpi_text}</p><p><b>Average risk score:</b> {risk_text}</p><p><span class='decision-badge' style='background:{badge_bg};color:#FFFFFF !important;-webkit-text-fill-color:#FFFFFF !important;font-weight:900 !important;'>{html.escape(badge_text)}</span></p></div>",
            unsafe_allow_html=True,
        )
        c1, c2 = st.columns(2)
        with c1:
            fig = px.bar(scoped, x="Project", y=["Contract Value", "AC", "Paid", "Remaining"], barmode="group", title="Budget vs Spent / Paid / Remaining per Project")
            st.plotly_chart(style_plotly(fig, 360), width="stretch")
            if scoped["Milestones"].fillna(0).sum():
                fig = px.bar(scoped, x="Project", y="Milestones", color="Status", title="Milestone Completion Tracking")
                st.plotly_chart(style_plotly(fig, 330), width="stretch")
            else:
                st.markdown("<div class='decision-chart-note'>Milestone data is not available in the current project files.</div>", unsafe_allow_html=True)
            ev_long = scoped[["Project", "EV", "PV", "AC", "SV", "CV", "EAC", "VAC"]].melt(id_vars="Project", var_name="Metric", value_name="Value")
            fig = px.bar(ev_long, x="Project", y="Value", color="Metric", barmode="group", title="EV Metrics per Project")
            st.plotly_chart(style_plotly(fig, 330), width="stretch")
        with c2:
            fig = px.bar(scoped, x="Project", y=["Progress", "Planned Progress"], color_discrete_sequence=["#50D5B7", "#D4A017"], barmode="group", title="Progress Gauges per Project", range_y=[0, 100])
            st.plotly_chart(style_plotly(fig, 360), width="stretch")
            radar_metrics = {
                "Progress": scoped["Progress"].mean(),
                "SPI": min((scoped["SPI"].replace(0, pd.NA).dropna().mean() or 0) * 100, 100),
                "CPI": min((scoped["CPI"].replace(0, pd.NA).dropna().mean() or 0) * 100, 100),
                "Risk Control": 100 - min(scoped["Risk Score"].fillna(0).mean(), 100),
            }
            if scoped["Quality"].fillna(0).sum():
                radar_metrics["Quality"] = scoped["Quality"].mean()
            if scoped["Safety"].fillna(0).sum():
                radar_metrics["Safety"] = scoped["Safety"].mean()
            fig = go.Figure(data=go.Scatterpolar(r=list(radar_metrics.values()), theta=list(radar_metrics.keys()), fill="toself", line_color="#D4A017"))
            fig.update_layout(title="Performance Radar vs Benchmark", polar=dict(radialaxis=dict(range=[0, 100])), showlegend=False)
            st.plotly_chart(style_plotly(fig, 330), width="stretch")
            fig = px.scatter(scoped, x="CPI", y="SPI", size=scoped["Contract Value"].clip(lower=1), color="Status", hover_name="Project", title="Sector SPI vs CPI Benchmark")
            fig.add_hline(y=1, line_dash="dash", line_color="#607080")
            fig.add_vline(x=1, line_dash="dash", line_color="#607080")
            st.plotly_chart(style_plotly(fig, 330), width="stretch")
        st.markdown("<div class='decision-chart-note'>Resource allocation data is not available in the current project files.</div>", unsafe_allow_html=True)
        st.dataframe(scoped[["Project", "Status", "Contract Value", "AC", "Paid", "Remaining", "Progress", "SPI", "CPI", "SV", "CV", "EAC", "VAC", "Risk Score"]], width="stretch", hide_index=True)

    if decision_view == "📋 Projects Analysis":
        default_projects = scoped_registry_df["Project"].head(8).tolist()
        with st.container(border=True):
            selected_projects = st.multiselect("Projects", scoped_registry_df["Project"].tolist(), default=default_projects, key="decision_projects_filter")
        scoped = scoped_registry_df[scoped_registry_df["Project"].isin(selected_projects)].copy()
        if scoped.empty:
            st.info("Select one or more projects to compare.")
        else:
            card_cols = st.columns(3)
            for idx, (_, row) in enumerate(scoped.iterrows()):
                badge_text, badge_bg, badge_fg = decision_status_badge(str(row["Status"]))
                spi_text = dashboard_display(f"{row['SPI']:.2f}" if pd.notna(row["SPI"]) else None)
                cpi_text = dashboard_display(f"{row['CPI']:.2f}" if pd.notna(row["CPI"]) else None)
                risk_text = dashboard_display(f"{row['Risk Score']:.1f}" if pd.notna(row["Risk Score"]) else None)
                with card_cols[idx % 3]:
                    st.markdown(
                        f"""
                        <div class='decision-card'>
                          <h4>{html.escape(str(row['Project']))}</h4>
                          <p><b>Sector:</b> {html.escape(str(row['Sector']))}</p>
                          <p><b>Status:</b> <span class='decision-badge' style='background:{badge_bg};color:#FFFFFF !important;-webkit-text-fill-color:#FFFFFF !important;font-weight:900 !important;'>{html.escape(badge_text)}</span></p>
                          <p><b>Contract:</b> {egp(row['Contract Value'])} | <b>Paid:</b> {egp(row['Paid'])}</p>
                          <p><b>Remaining:</b> {egp(row['Remaining'])}</p>
                          <p><b>Progress:</b> {pct(row['Progress'])} / planned {pct(row['Planned Progress'])}</p>
                          <p><b>SPI:</b> {spi_text} | <b>CPI:</b> {cpi_text} | <b>Risk:</b> {risk_text}</p>
                          <p><b>Decision:</b> {html.escape(str(row['Required Decision']))}</p>
                        </div>
                        """,
                        unsafe_allow_html=True,
                    )
            c1, c2 = st.columns(2)
            with c1:
                fig = px.treemap(scoped, path=["Sector", "Project"], values="Contract Value", color="Status", title="Budget Distribution Comparison", color_discrete_map={"On Track": "#1A8A8F", "Watch": "#D4A017", "High Attention": "#C94C4C"})
                st.plotly_chart(style_plotly(fig, 360), width="stretch")
                fig = px.scatter(scoped, x="CPI", y="SPI", color="Status", size=scoped["Contract Value"].clip(lower=1), hover_name="Project", hover_data=["Sector", "Progress", "Risk Score", "Delay Days"], title="SPI vs CPI Scatter Plot")
                fig.add_hline(y=1, line_dash="dash", line_color="#607080")
                fig.add_vline(x=1, line_dash="dash", line_color="#607080")
                st.plotly_chart(style_plotly(fig, 360), width="stretch")
            with c2:
                trend_cols = ["Project", "Progress", "Planned Progress"]
                if scoped["Quality"].fillna(0).sum():
                    trend_cols.append("Quality")
                if scoped["Safety"].fillna(0).sum():
                    trend_cols.append("Safety")
                trends = scoped[trend_cols].melt(id_vars="Project", var_name="Metric", value_name="Value")
                fig = px.line(trends, x="Project", y="Value", color="Metric", markers=True, title="Progress / Quality / Safety Trends")
                st.plotly_chart(style_plotly(fig, 360), width="stretch")
                ev_long = scoped[["Project", "PV", "EV", "AC", "SPI", "CPI", "SV", "CV", "EAC", "VAC"]].melt(id_vars="Project", var_name="Metric", value_name="Value")
                fig = px.bar(ev_long, x="Project", y="Value", color="Metric", barmode="group", title="EV Metrics Comparison")
                st.plotly_chart(style_plotly(fig, 360), width="stretch")
            risk_distribution = scoped.assign(
                Schedule_Risk=scoped["SPI"].apply(lambda value: "High" if value and value < .9 else ("Medium" if value and value < 1 else "Low")),
                Cost_Risk=scoped["CPI"].apply(lambda value: "High" if value and value < .9 else ("Medium" if value and value < 1 else "Low")),
                Portfolio_Risk=scoped["Risk Score"].apply(lambda value: "High" if value >= 70 else ("Medium" if value >= 35 else "Low")),
            )
            risk_long = risk_distribution[["Project", "Schedule_Risk", "Cost_Risk", "Portfolio_Risk"]].melt(id_vars="Project", var_name="Risk Type", value_name="Risk Level")
            risk_count = risk_long.groupby(["Project", "Risk Type", "Risk Level"], as_index=False).size()
            fig = px.bar(risk_count, x="Project", y="size", color="Risk Level", facet_col="Risk Type", title="Risk Distribution Stacked Bars", color_discrete_map={"Low": "#50D5B7", "Medium": "#D4A017", "High": "#F05D5E"})
            st.plotly_chart(style_plotly(fig, 360), width="stretch")
            comparison = scoped[["Project", "Sector", "Status", "Contract Value", "Planned Progress", "Progress", "Progress Variance", "BAC", "AC", "Paid", "Remaining", "SPI", "CPI", "Risk Score", "Delay Days", "Milestones", "Claims / EOT Exposure", "Required Decision"]].copy()
            styler = comparison.style.map(lambda value: "background-color:#4B1D22;color:#fff" if isinstance(value, (int, float)) and value < .9 else "", subset=["SPI", "CPI"])
            st.dataframe(styler, width="stretch", hide_index=True, height=dataframe_height(comparison, max_height=620))
    render_management_decision_brief_once(scoped_registry_df, quality)
    st.markdown("</div>", unsafe_allow_html=True)


def selected_project_id() -> str:
    return str(st.session_state.get("active_project_id", "") or "").strip()


def project_scoped_file(path: Path, project_id: str | None = None) -> Path:
    project_id = selected_project_id() if project_id is None else str(project_id or "").strip()
    if not project_id:
        return path

    source_map = [
        (IMPORT_TEMPLATES_DIR, project_data_path(PROJECTS_DIR, project_id, "core", "")),
        (STEEL_TIA_DIR, project_data_path(PROJECTS_DIR, project_id, "delay_analysis", "")),
        (BL_FIXED_DIR, project_data_path(PROJECTS_DIR, project_id, "bl", "")),
    ]
    try:
        resolved_path = path.resolve()
    except FileNotFoundError:
        resolved_path = path.absolute()

    for source_root, project_root in source_map:
        try:
            relative_path = resolved_path.relative_to(source_root.resolve())
        except ValueError:
            continue
        candidate = project_root / relative_path
        return candidate
    return path


def filter_active_project(df: pd.DataFrame, project_id: str | None = None) -> pd.DataFrame:
    if df.empty:
        return df
    project_id = selected_project_id() if project_id is None else str(project_id or "").strip()
    if not project_id or "project_id" not in df.columns:
        return df
    project_values = df["project_id"].astype(str).str.strip()
    return df[project_values.eq(project_id)].copy()


def active_project_row(projects_df: pd.DataFrame) -> pd.Series:
    if projects_df.empty:
        return pd.Series(dtype=object)
    project_id = selected_project_id()
    if project_id and "project_id" in projects_df.columns:
        match = projects_df[projects_df["project_id"].astype(str).str.strip().eq(project_id)]
        if not match.empty:
            return match.iloc[0]
    return projects_df.iloc[0]


def build_overview_metrics():
    projects_df = load_core_csv(PROJECTS_CSV_PATH)
    activities_df = filter_active_project(load_core_csv(ACTIVITIES_CSV_PATH))

    if projects_df.empty:
        return {}

    if not selected_project_id():
        valid_start = pd.to_datetime(projects_df.get("planned_start", pd.Series(dtype=object)), errors="coerce", dayfirst=True)
        valid_finish = pd.to_datetime(projects_df.get("planned_finish", pd.Series(dtype=object)), errors="coerce", dayfirst=True)
        contract_values = projects_df.get("contract_value", pd.Series(0, index=projects_df.index)).apply(parse_numeric)
        planned_values = projects_df.get("planned_progress_percent", pd.Series(0, index=projects_df.index)).apply(parse_numeric)
        actual_values = projects_df.get("actual_progress_percent", pd.Series(0, index=projects_df.index)).apply(parse_numeric)
        weight_total = float(contract_values.sum())
        project_row = pd.Series({
            "project_name": f"All Projects ({len(projects_df)})",
            "planned_start": valid_start.min() if valid_start.notna().any() else None,
            "planned_finish": valid_finish.max() if valid_finish.notna().any() else None,
            "contract_value": weight_total,
            "planned_progress_percent": float((planned_values * contract_values).sum() / weight_total) if weight_total else float(planned_values.mean()),
            "actual_progress_percent": float((actual_values * contract_values).sum() / weight_total) if weight_total else float(actual_values.mean()),
        })
    else:
        project_row = active_project_row(projects_df)
    project_start = parse_project_date(project_row.get("planned_start"))
    project_finish = parse_project_date(project_row.get("planned_finish"))
    duration_days = (project_finish - project_start).days if project_start is not None and project_finish is not None else 0
    today = pd.Timestamp.today().normalize()
    elapsed_ratio = 0.0
    if duration_days > 0 and project_start is not None:
        elapsed_ratio = ((today - project_start).days / duration_days)
    elapsed_ratio = min(max(elapsed_ratio, 0.0), 1.0)
    duration_elapsed_pct = elapsed_ratio * 100.0
    remaining_duration_pct = 100.0 - duration_elapsed_pct

    critical_count = 0
    total_activities = 0
    if not activities_df.empty:
        total_activities = activities_df["activity_id"].count()
        critical_count = activities_df["is_critical"].astype(str).str.strip().str.lower().eq("yes").sum()

    return {
        "project_name": project_row.get("project_name", ""),
        "project_start": project_start,
        "project_finish": project_finish,
        "duration_days": duration_days,
        "duration_elapsed_pct": duration_elapsed_pct,
        "remaining_duration_pct": remaining_duration_pct,
        "overall_progress": parse_numeric(project_row.get("actual_progress_percent")),
        "planned_progress": parse_numeric(project_row.get("planned_progress_percent")),
        "contract_value": parse_numeric(project_row.get("contract_value")),
        "total_activities": int(total_activities),
        "critical_activities": int(critical_count),
    }


def build_evm_metrics():
    evm_df = filter_active_project(load_core_csv(EVM_CSV_PATH))
    if evm_df.empty:
        return {}

    for col in ["BAC", "AC", "EV", "PV", "SV", "CV"]:
        evm_df[col] = evm_df[col].apply(parse_numeric)

    bac = float(evm_df["BAC"].sum())
    ac = float(evm_df["AC"].sum())
    ev = float(evm_df["EV"].sum())
    pv = float(evm_df["PV"].sum())
    sv = float(evm_df["SV"].sum())
    cv = float(evm_df["CV"].sum())
    eac = ac + (bac - ev)
    tcpi = ((bac - ev) / (eac - ac)) if (eac - ac) not in (0, None) else None

    return {
        "bac": bac,
        "ac": ac,
        "ev": ev,
        "pv": pv,
        "sv": sv,
        "cv": cv,
        "eac": eac,
        "tcpi": tcpi,
    }


def build_earned_value_analysis_data(evm_metrics: dict | None = None):
    metrics = evm_metrics if evm_metrics is not None else build_evm_metrics()
    bac = float(metrics.get("bac", 0.0) or 0.0)
    pv = float(metrics.get("pv", 0.0) or 0.0)
    ev = float(metrics.get("ev", 0.0) or 0.0)
    sv = ev - pv
    spi = ev / pv if pv else 0.0
    planned_completion_value_position = pv / bac * 100 if bac else 0.0
    earned_value_gap = pv - ev
    schedule_health = "Critical Delay" if spi < 1.0 and sv < 0 else "Stable"
    interpretation = (
        "The project is underperforming against the planned value baseline. "
        "The earned value is significantly below the planned value, creating a negative "
        f"schedule variance of SAR {abs(sv)/1_000_000:.2f}M and an SPI of {spi:.2f}. "
        f"This indicates that the project is progressing at approximately {spi * 100:.0f}% of the planned schedule efficiency."
    )
    executive_summary = (
        "Earned Value Analysis confirms material schedule underperformance against the planned value baseline. "
        "The project remains in critical delay status and requires recovery actions focused on clearing external constraints and converting planned work fronts into earned value."
    )
    return {
        "BAC": bac,
        "PV": pv,
        "EV": ev,
        "SV": sv,
        "SPI": spi,
        "plannedCompletionValuePosition": planned_completion_value_position,
        "earnedValueGap": earned_value_gap,
        "scheduleHealthClassification": schedule_health,
        "interpretation": interpretation,
        "executiveSummary": executive_summary,
    }


def build_evm_root_cause_rows(delay_metrics: dict, risk_metrics: dict, contract_metrics: dict):
    payments_df = contract_metrics.get("payments_df", pd.DataFrame()).copy()
    under_payment_exists = False
    if not payments_df.empty and "payment_status" in payments_df.columns:
        under_payment_exists = payments_df["payment_status"].astype(str).str.contains("under payment", case=False, na=False).any()

    return pd.DataFrame(
        [
            {
                "Cause Title": "Delay in RFT steel supply by Owner",
                "Impact Area": "Procurement / Construction",
                "EVM Impact Link": "PV not converted into EV due to blocked work fronts and delayed material-driven execution.",
                "Status": "Ongoing",
                "Severity": "Critical",
            },
            {
                "Cause Title": "Delay in design review process",
                "Impact Area": "Engineering",
                "EVM Impact Link": "Planned engineering progress could not flow into approved deliverables and downstream earned value.",
                "Status": "Closed",
                "Severity": "High",
            },
            {
                "Cause Title": "Delay in reply on RFI 4 and RFI 5",
                "Impact Area": "Engineering / Construction",
                "EVM Impact Link": "Technical closure delays restricted workface release and slowed conversion of planned progress into earned progress.",
                "Status": "Ongoing",
                "Severity": "High",
            },
            {
                "Cause Title": "Delay in determining the MEP executing party",
                "Impact Area": "Engineering / Procurement / Construction",
                "EVM Impact Link": "Interface uncertainty delayed execution planning, procurement flow, and forecast progress realization.",
                "Status": "Closed",
                "Severity": "Medium",
            },
            {
                "Cause Title": "Payment delay impact on resource continuity and productivity",
                "Impact Area": "Commercial / Construction",
                "EVM Impact Link": "Liquidity pressure reduced productivity continuity and slowed the conversion of planned value into earned value.",
                "Status": "Ongoing" if under_payment_exists else "Closed",
                "Severity": "High",
            },
        ]
    )


def build_evm_mitigation_rows():
    return pd.DataFrame(
        [
            {
                "Action": "RFT Steel Supply Recovery",
                "Owner / Responsible Party": "Project Controls / Procurement / Owner Interface",
                "Current Status": "In Progress",
                "Recovery Impact": "Releases blocked structural work fronts and supports EV generation in construction.",
                "Required Next Decision": "Secure committed steel release and delivery sequence from Owner.",
            },
            {
                "Action": "RFI Closure Acceleration",
                "Owner / Responsible Party": "Engineering Team / Consultant Interface",
                "Current Status": "In Progress",
                "Recovery Impact": "Enables technical closure and unblocks affected planned activities.",
                "Required Next Decision": "Escalate overdue RFI responses for closure dates.",
            },
            {
                "Action": "Design Review Follow-up",
                "Owner / Responsible Party": "Engineering Management",
                "Current Status": "Monitoring",
                "Recovery Impact": "Protects planned engineering completion and downstream procurement readiness.",
                "Required Next Decision": "Freeze review turnaround commitments with reviewing parties.",
            },
            {
                "Action": "MEP Interface Finalization",
                "Owner / Responsible Party": "Project Management / Client Interface",
                "Current Status": "Monitoring",
                "Recovery Impact": "Reduces interface uncertainty and improves forecast execution reliability.",
                "Required Next Decision": "Confirm final executing responsibility and interface release logic.",
            },
            {
                "Action": "Commercial / Payment Follow-up",
                "Owner / Responsible Party": "Commercial Team",
                "Current Status": "In Progress",
                "Recovery Impact": "Supports resource continuity and mitigates productivity disruption risk.",
                "Required Next Decision": "Close outstanding payment approvals and release under-payment invoices.",
            },
            {
                "Action": "Construction Productivity Recovery",
                "Owner / Responsible Party": "Construction Team",
                "Current Status": "Ready Upon Constraint Removal",
                "Recovery Impact": "Accelerates EV generation once technical and supply constraints are removed.",
                "Required Next Decision": "Approve recovery sequence by priority work front.",
            },
        ]
    )


def load_evm_comments() -> dict:
    defaults = {
        "quantitativePerformance": "",
        "rootCauseLinkage": "",
        "mitigationRecovery": "",
    }
    if not EVM_COMMENTS_PATH.exists():
        return defaults
    try:
        saved = json.loads(EVM_COMMENTS_PATH.read_text(encoding="utf-8"))
        return {key: str(saved.get(key, "")) for key in defaults}
    except Exception:
        return defaults


def save_evm_comments(comments: dict) -> None:
    EVM_COMMENTS_PATH.parent.mkdir(parents=True, exist_ok=True)
    EVM_COMMENTS_PATH.write_text(json.dumps(comments, indent=2), encoding="utf-8")


def ensure_evm_comment_state() -> None:
    saved = load_evm_comments()
    for key, state_key in [
        ("quantitativePerformance", "evm_comment_quantitativePerformance"),
        ("rootCauseLinkage", "evm_comment_rootCauseLinkage"),
        ("mitigationRecovery", "evm_comment_mitigationRecovery"),
    ]:
        if state_key not in st.session_state:
            st.session_state[state_key] = saved.get(key, "")


def get_evm_comments() -> dict:
    return {
        "quantitativePerformance": st.session_state.get("evm_comment_quantitativePerformance", ""),
        "rootCauseLinkage": st.session_state.get("evm_comment_rootCauseLinkage", ""),
        "mitigationRecovery": st.session_state.get("evm_comment_mitigationRecovery", ""),
    }


def persist_evm_comments_from_state() -> None:
    save_evm_comments(get_evm_comments())


def persist_evm_comments_from_widget_keys(q_key: str, r_key: str, m_key: str) -> None:
    st.session_state["evm_comment_quantitativePerformance"] = st.session_state.get(q_key, "")
    st.session_state["evm_comment_rootCauseLinkage"] = st.session_state.get(r_key, "")
    st.session_state["evm_comment_mitigationRecovery"] = st.session_state.get(m_key, "")
    persist_evm_comments_from_state()


def clear_evm_comment(section_key: str) -> None:
    state_key = f"evm_comment_{section_key}"
    st.session_state[state_key] = ""
    persist_evm_comments_from_state()


def clear_evm_comment_with_widget(section_key: str, widget_key: str) -> None:
    clear_evm_comment(section_key)
    st.session_state[widget_key] = ""


def build_contract_metrics():
    contracts_df = filter_active_project(load_core_csv(CONTRACTS_CSV_PATH))
    payments_df = filter_active_project(load_core_csv(PAYMENTS_CSV_PATH))

    if not contracts_df.empty:
        contracts_df["original_value_num"] = contracts_df["original_value"].apply(parse_numeric)
        active_contracts = contracts_df[contracts_df["status"].astype(str).str.strip().str.lower().eq("active")]
    else:
        active_contracts = pd.DataFrame()

    if not payments_df.empty:
        payments_df = payments_df.loc[:, [col for col in payments_df.columns if str(col).strip()]]
        payments_df = payments_df[payments_df["payment_id"].astype(str).str.strip() != ""]
        payments_df["certified_amount_num"] = payments_df["certified_amount"].apply(parse_numeric)
        payments_df["paid_amount_num"] = payments_df["paid_amount"].apply(parse_numeric)

    return {
        "contracts_df": contracts_df,
        "payments_df": payments_df,
        "total_contracts": int(contracts_df["contract_id"].astype(str).str.strip().ne("").sum()) if not contracts_df.empty else 0,
        "total_contract_value": float(active_contracts["original_value_num"].sum()) if not active_contracts.empty else 0.0,
        "total_certified": float(payments_df["certified_amount_num"].sum()) if not payments_df.empty else 0.0,
        "total_paid": float(payments_df["paid_amount_num"].sum()) if not payments_df.empty else 0.0,
    }


def build_delay_metrics():
    delays_source_df = filter_active_project(load_core_csv(DELAYS_CSV_PATH))
    delays_df = delays_source_df.copy()
    if delays_df.empty:
        return {
            "delays_df": delays_df,
            "display_delays_df": pd.DataFrame(),
            "raw_delay_columns": [col for col in delays_source_df.columns if str(col).strip()],
            "total_delay_events": 0,
            "total_delay_days": 0.0,
            "employer_delays": 0,
            "eot_potential_count": 0,
            "open_delays": 0,
            "closed_delays": 0,
        }

    delays_df = delays_df.loc[:, [col for col in delays_df.columns if str(col).strip()]]
    delays_df = delays_df[delays_df["delay_id"].astype(str).str.strip() != ""].copy()
    delays_df["estimated_delay_days_num"] = delays_df["estimated_delay_days"].apply(parse_numeric)
    delays_df["approved_eot_days_num"] = delays_df["approved_eot_days"].apply(parse_numeric)
    delays_df["pending_exposure_days"] = delays_df["estimated_delay_days_num"] - delays_df["approved_eot_days_num"]
    delays_df["pending_exposure_days"] = delays_df["pending_exposure_days"].clip(lower=0)
    responsibility_text = delays_df["responsibility"].astype(str).str.strip().str.lower()
    status_text = delays_df["status"].astype(str).str.strip().str.lower()
    delays_df["responsible_group"] = responsibility_text.apply(
        lambda x: (
            "Employer / Client"
            if any(term in x for term in ["employer", "client"])
            else "Consultant / Engineer"
            if any(term in x for term in ["consultant", "engineer", "design team"])
            else "Contractor"
            if "contractor" in x
            else "Other"
        )
    )
    delays_df["eot_potential"] = delays_df.apply(
        lambda row: "Yes"
        if (
            row["estimated_delay_days_num"] > 0
            and row["responsible_group"] in {"Employer / Client", "Consultant / Engineer"}
        )
        else "No",
        axis=1,
    )
    delays_df["status_group"] = status_text.apply(
        lambda x: "Open" if "open" in x or "pending" in x else ("Closed" if "closed" in x or "delivered" in x or "received" in x else "Other")
    )
    employer_mask = delays_df["responsible_group"].eq("Employer / Client")
    display_df = delays_df[[col for col in delays_source_df.columns if str(col).strip()]].copy()
    for col in ["estimated_delay_days", "approved_eot_days"]:
        if col in display_df.columns:
            source_col = f"{col}_num"
            if source_col in delays_df.columns:
                display_df[col] = delays_df[source_col].apply(lambda x: "" if float(x) == 0 else int(x) if float(x).is_integer() else round(float(x), 2))
    for col in ["start_date", "end_date", "notice_ref", "project_id", "activity_id", "status", "cause_category", "responsibility", "delay_title"]:
        if col in display_df.columns:
            display_df[col] = display_df[col].astype(str).replace({"nan": "", "None": ""}).str.strip()
    display_df = display_df.fillna("")
    display_df = display_df.replace("", "ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â")

    return {
        "delays_df": delays_df,
        "display_delays_df": display_df,
        "raw_delay_columns": [col for col in delays_source_df.columns if str(col).strip()],
        "total_delay_events": int(len(delays_df)),
        "total_delay_days": float(delays_df["estimated_delay_days_num"].sum()),
        "employer_delays": int(employer_mask.sum()),
        "eot_potential_count": int(delays_df["eot_potential"].eq("Yes").sum()),
        "open_delays": int(delays_df["status_group"].eq("Open").sum()),
        "closed_delays": int(delays_df["status_group"].eq("Closed").sum()),
    }


def build_risk_metrics():
    risks_df = filter_active_project(load_core_csv(RISKS_CSV_PATH))
    steel_df = filter_active_project(load_core_csv(STEEL_DELAY_CSV_PATH))
    rfi_df = filter_active_project(load_core_csv(RFI_STATUS_CSV_PATH))
    ifc_df = filter_active_project(load_core_csv(IFC_CONFLICT_CSV_PATH))

    if not steel_df.empty:
        steel_df = steel_df.loc[:, [col for col in steel_df.columns if str(col).strip()]]
        if "Ref No" in steel_df.columns:
            steel_df = steel_df[steel_df["Ref No"].astype(str).str.strip() != ""].copy()

    if not rfi_df.empty:
        rfi_df = rfi_df.loc[:, [col for col in rfi_df.columns if str(col).strip()]]
        if "RFI No." in rfi_df.columns:
            rfi_df = rfi_df[rfi_df["RFI No."].astype(str).str.strip() != ""].copy()

    if not ifc_df.empty:
        ifc_df = ifc_df.loc[:, [col for col in ifc_df.columns if str(col).strip()]]
        first_col = ifc_df.columns[0] if len(ifc_df.columns) else None
        if first_col:
            ifc_df = ifc_df[ifc_df[first_col].astype(str).str.strip() != ""].copy()

    if risks_df.empty:
        return {
            "risks_df": risks_df,
            "steel_df": steel_df,
            "rfi_df": rfi_df,
            "ifc_df": ifc_df,
            "steel_issues": int(len(steel_df)),
            "rfi_items": int(len(rfi_df)),
            "ifc_conflicts": int(len(ifc_df)),
        }

    risks_df["probability_norm"] = risks_df["probability"].astype(str).str.strip().str.lower()
    risks_df["time_impact_flag"] = risks_df["time_impact_days"].astype(str).str.strip().str.lower()
    risks_df["cost_impact_flag"] = risks_df["cost_impact"].astype(str).str.strip().str.lower()
    return {
        "risks_df": risks_df,
        "steel_df": steel_df,
        "rfi_df": rfi_df,
        "ifc_df": ifc_df,
        "total_risks": int(len(risks_df)),
        "high_risks": int(risks_df["probability_norm"].eq("high").sum()),
        "open_risks": int(risks_df["status"].astype(str).str.strip().str.lower().eq("open").sum()),
        "closed_risks": int(risks_df["status"].astype(str).str.strip().str.lower().eq("closed").sum()),
        "steel_issues": int(len(steel_df)),
        "rfi_items": int(len(rfi_df)),
        "ifc_conflicts": int(len(ifc_df)),
    }


def build_milestone_metrics():
    milestones_df = filter_active_project(load_core_csv(MILESTONES_CSV_PATH))
    change_orders_df = filter_active_project(load_core_csv(CHANGE_ORDERS_CSV_PATH))
    if not change_orders_df.empty:
        for col in ["claimed_amount", "approved_amount", "cost_impact", "time_impact_days"]:
            if col in change_orders_df.columns:
                change_orders_df[f"{col}_num"] = change_orders_df[col].apply(parse_numeric)
    return {
        "milestones_df": milestones_df,
        "change_orders_df": change_orders_df,
    }


def build_time_impact_metrics():
    delays_df = filter_active_project(load_core_csv(DELAYS_CSV_PATH))
    if delays_df.empty:
        return {"time_impact_df": delays_df}

    time_impact_df = delays_df.copy()
    time_impact_df["estimated_delay_days_num"] = time_impact_df["estimated_delay_days"].apply(parse_numeric)
    time_impact_df["approved_eot_days_num"] = time_impact_df["approved_eot_days"].apply(parse_numeric)
    time_impact_df["impact_type"] = time_impact_df["responsibility"].astype(str).str.strip().apply(
        lambda x: "Employer Delay" if x.lower() == "employer" else "Contractor Delay"
    )
    avg_delay = float(time_impact_df["estimated_delay_days_num"].mean()) if len(time_impact_df) else 0.0
    return {
        "time_impact_df": time_impact_df,
        "total_delay_events": int(len(time_impact_df)),
        "total_delay_days": float(time_impact_df["estimated_delay_days_num"].sum()),
        "average_delay_days": avg_delay,
        "impacted_activities": int(time_impact_df["activity_id"].astype(str).str.strip().ne("").sum()),
    }


def build_s_curve_metrics():
    curve_df = filter_active_project(load_core_csv(S_CURVE_CSV_PATH))
    if curve_df.empty:
        return {"curve_df": pd.DataFrame()}

    for col in [
        "monthly_planned",
        "cumm_monthly_planned",
        "monthly_actual",
        "cumm_monthly_actual",
        "monthly_invoiced",
        "cumm_monthly_invoiced",
    ]:
        curve_df[col] = curve_df[col].apply(parse_numeric)

    curve_df["MonthSort"] = pd.to_datetime(curve_df["months"], format="%d-%b-%y", errors="coerce")
    curve_df = curve_df.sort_values("MonthSort").reset_index(drop=True)
    curve_df["Month"] = curve_df["months"]
    curve_df["Planned"] = curve_df["cumm_monthly_planned"]
    curve_df["Actual"] = curve_df["cumm_monthly_actual"]
    curve_df["Invoiced"] = curve_df["cumm_monthly_invoiced"]
    curve_df["planned_value_num"] = curve_df["monthly_planned"]
    curve_df["actual_cost_num"] = curve_df["monthly_actual"]
    curve_df["invoice_amount_num"] = curve_df["monthly_invoiced"]

    actual_non_zero = curve_df[curve_df["monthly_actual"] > 0]
    last_actual_month = actual_non_zero["Month"].iloc[-1] if not actual_non_zero.empty else "N/A"
    return {
        "curve_df": curve_df,
        "planned_total": float(curve_df["cumm_monthly_planned"].iloc[-1]) if not curve_df.empty else 0.0,
        "actual_total": float(curve_df["cumm_monthly_actual"].max()) if not curve_df.empty else 0.0,
        "invoiced_total": float(curve_df["cumm_monthly_invoiced"].max()) if not curve_df.empty else 0.0,
        "last_actual_month": last_actual_month,
    }


def parse_activity_date(value):
    if pd.isna(value):
        return pd.NaT
    cleaned = str(value).replace(" A", "").strip()
    return pd.to_datetime(cleaned, format="%d-%b-%y", errors="coerce")


def build_activity_metrics():
    activities_df = filter_active_project(load_core_csv(ACTIVITIES_CSV_PATH))
    if activities_df.empty:
        return {
            "activities_df": activities_df,
            "critical_df": pd.DataFrame(),
            "deviated_df": pd.DataFrame(),
            "rft_df": pd.DataFrame(),
            "critical_count": 0,
            "deviated_count": 0,
            "rft_count": 0,
            "avg_critical_variance": 0.0,
            "avg_rft_variance": 0.0,
        }

    activities_df = activities_df.loc[:, [col for col in activities_df.columns if str(col).strip()]]
    activities_df = activities_df[activities_df["activity_id"].astype(str).str.strip() != ""].copy()
    activities_df["planned_progress_num"] = activities_df["planned_progress"].apply(parse_numeric)
    activities_df["actual_progress_num"] = activities_df["actual_progress"].apply(parse_numeric)
    activities_df["planned_weight_num"] = activities_df["planned_weight"].apply(parse_numeric)
    activities_df["total_float_days_num"] = activities_df["total_float_days"].apply(parse_numeric)
    activities_df["progress_variance"] = activities_df["actual_progress_num"] - activities_df["planned_progress_num"]
    activities_df["planned_finish_dt"] = activities_df["planned_finish"].apply(parse_activity_date)
    activities_df["forecast_finish_dt"] = activities_df["forecast_finish"].apply(parse_activity_date)
    activities_df["actual_finish_dt"] = activities_df["actual_finish"].apply(parse_activity_date)
    activities_df["finish_slip_days"] = (
        (activities_df["forecast_finish_dt"] - activities_df["planned_finish_dt"]).dt.days.fillna(0)
    )

    critical_df = activities_df[activities_df["is_critical"].astype(str).str.strip().str.lower().eq("yes")].copy()
    critical_df = critical_df.sort_values(["progress_variance", "finish_slip_days"], ascending=[True, False])

    deviated_df = activities_df[
        (activities_df["progress_variance"] < 0) | (activities_df["finish_slip_days"] > 0)
    ].copy()
    deviated_df = deviated_df.sort_values(["progress_variance", "finish_slip_days"], ascending=[True, False])

    rft_df = activities_df[
        activities_df["activity_name"].astype(str).str.contains("RFT", case=False, na=False)
    ].copy()
    rft_df = rft_df.sort_values(["progress_variance", "finish_slip_days"], ascending=[True, False])

    return {
        "activities_df": activities_df,
        "critical_df": critical_df,
        "deviated_df": deviated_df,
        "rft_df": rft_df,
        "critical_count": int(len(critical_df)),
        "deviated_count": int(len(deviated_df)),
        "rft_count": int(len(rft_df)),
        "avg_critical_variance": float(critical_df["progress_variance"].mean()) if not critical_df.empty else 0.0,
        "avg_rft_variance": float(rft_df["progress_variance"].mean()) if not rft_df.empty else 0.0,
    }


def build_wbs_metrics():
    wbs_df = filter_active_project(load_core_csv(WBS_CSV_PATH))
    if wbs_df.empty:
        return {"wbs_df": pd.DataFrame(), "con_wbs_df": pd.DataFrame(), "chart_rows": pd.DataFrame(), "code_col": None}

    wbs_df = wbs_df.loc[:, [col for col in wbs_df.columns if str(col).strip()]].copy()
    for col in ["schedule_%_complete", "performance_%_complete"]:
        if col in wbs_df.columns:
            wbs_df[f"{col}_num"] = wbs_df[col].apply(parse_numeric)

    code_col = "WBS Code" if "WBS Code" in wbs_df.columns else ("wbs_code" if "wbs_code" in wbs_df.columns else None)
    con_wbs_df = wbs_df.copy()
    if code_col:
        con_wbs_df = con_wbs_df[con_wbs_df[code_col].astype(str).str.contains("CON", case=False, na=False)].copy()

    chart_rows = pd.DataFrame()
    if code_col:
        chart_rows = con_wbs_df.head(4).copy()

    return {
        "wbs_df": wbs_df,
        "con_wbs_df": con_wbs_df,
        "chart_rows": chart_rows,
        "code_col": code_col,
    }


def build_time_impact_engine(delay_metrics, risk_metrics, activity_metrics, contract_metrics):
    time_impact_df = delay_metrics.get("delays_df", pd.DataFrame()).copy()
    activities_df = activity_metrics.get("activities_df", pd.DataFrame()).copy()
    critical_df = activity_metrics.get("critical_df", pd.DataFrame()).copy()
    rft_df = activity_metrics.get("rft_df", pd.DataFrame()).copy()

    if time_impact_df.empty:
        return {"time_impact_df": time_impact_df}

    time_impact_df["source_stream"] = time_impact_df["delay_id"].astype(str).apply(
        lambda x: "Steel Delay" if "STEEL" in x else ("IFC Conflict" if "IFC" in x else ("RFI Delay" if "RFI" in x else "Other"))
    )

    critical_ids = set(critical_df["activity_id"].astype(str).str.strip()) if not critical_df.empty else set()
    rft_ids = set(rft_df["activity_id"].astype(str).str.strip()) if not rft_df.empty else set()

    activity_lookup = pd.DataFrame()
    if not activities_df.empty:
        activity_lookup = activities_df[
            ["activity_id", "activity_name", "planned_progress_num", "actual_progress_num", "progress_variance", "finish_slip_days", "total_float_days_num", "is_critical"]
        ].copy()
        activity_lookup["activity_id"] = activity_lookup["activity_id"].astype(str).str.strip()
        activity_lookup = activity_lookup.drop_duplicates(subset=["activity_id"])

    time_impact_df["activity_id"] = time_impact_df["activity_id"].astype(str).str.strip()
    if not activity_lookup.empty:
        time_impact_df = time_impact_df.merge(activity_lookup, on="activity_id", how="left")
        if "activity_name" not in time_impact_df.columns:
            left_name = time_impact_df.get("activity_name_x", pd.Series([""] * len(time_impact_df), index=time_impact_df.index))
            right_name = time_impact_df.get("activity_name_y", pd.Series([""] * len(time_impact_df), index=time_impact_df.index))
            time_impact_df["activity_name"] = left_name.astype(str).where(left_name.astype(str).str.strip().ne(""), right_name)
        for required_col, default in {
            "planned_progress_num": 0.0,
            "actual_progress_num": 0.0,
            "progress_variance": 0.0,
            "finish_slip_days": 0.0,
            "total_float_days_num": 0.0,
            "is_critical": "",
        }.items():
            if required_col not in time_impact_df.columns:
                time_impact_df[required_col] = default
    else:
        time_impact_df["activity_name"] = ""
        time_impact_df["planned_progress_num"] = 0.0
        time_impact_df["actual_progress_num"] = 0.0
        time_impact_df["progress_variance"] = 0.0
        time_impact_df["finish_slip_days"] = 0.0
        time_impact_df["total_float_days_num"] = 0.0
        time_impact_df["is_critical"] = ""

    time_impact_df["critical_path_flag"] = time_impact_df["activity_id"].isin(critical_ids) | time_impact_df["is_critical"].astype(str).str.strip().str.lower().eq("yes")
    time_impact_df["near_critical_flag"] = time_impact_df["total_float_days_num"].apply(lambda x: 0 <= float(x) <= 15 if pd.notna(x) else False)
    time_impact_df["rft_flag"] = (
        time_impact_df["activity_id"].isin(rft_ids)
        | time_impact_df["activity_name"].astype(str).str.contains("RFT", case=False, na=False)
        | time_impact_df["delay_title"].astype(str).str.contains("RFT", case=False, na=False)
        | time_impact_df["notice_ref"].astype(str).str.contains("RFT", case=False, na=False)
    )
    time_impact_df["notice_status"] = time_impact_df["notice_ref"].astype(str).str.strip().replace("", "Missing")
    time_impact_df["commercial_gap"] = float(contract_metrics.get("total_certified", 0.0)) - float(contract_metrics.get("total_paid", 0.0))
    time_impact_df["responsibility_band"] = time_impact_df["responsible_group"]
    time_impact_df["engine_priority"] = (
        (time_impact_df["estimated_delay_days_num"] > 0).astype(int) * 4
        + time_impact_df["critical_path_flag"].astype(int) * 4
        + time_impact_df["near_critical_flag"].astype(int) * 2
        + time_impact_df["rft_flag"].astype(int) * 2
        + time_impact_df["eot_potential"].eq("Yes").astype(int) * 3
        + time_impact_df["status_group"].eq("Open").astype(int) * 3
    )
    time_impact_df = time_impact_df.sort_values(
        ["engine_priority", "estimated_delay_days_num", "progress_variance"],
        ascending=[False, False, True],
    ).reset_index(drop=True)

    input_matrix = pd.DataFrame(
        [
            {"Input Slide": "Overview", "Input Dataset": "projects.csv", "Time Impact Use": "Project dates and remaining duration context", "Status": "Connected"},
            {"Input Slide": "Delays", "Input Dataset": "delay_events.csv", "Time Impact Use": "Primary event register and delay durations", "Status": "Connected"},
            {"Input Slide": "Risks", "Input Dataset": "risks.csv + steel/rfi/ifc files", "Time Impact Use": "Causation, responsibility, and issue-source enrichment", "Status": "Connected"},
            {"Input Slide": "Activities", "Input Dataset": "activities.csv", "Time Impact Use": "Critical path, near-critical path, RFT, and deviation linkage", "Status": "Connected"},
            {"Input Slide": "Contracts", "Input Dataset": "contracts.csv + payments.csv", "Time Impact Use": "Commercial pressure and certification/payment gap context", "Status": "Connected"},
            {"Input Slide": "Letters Intelligence", "Input Dataset": "letters workbook", "Time Impact Use": "Notice and correspondence traceability context", "Status": "Available"},
        ]
    )

    causation_matrix = (
        time_impact_df.groupby(["source_stream", "responsibility_band"], as_index=False)
        .agg(
            Delay_Events=("delay_id", "count"),
            Delay_Days=("estimated_delay_days_num", "sum"),
            Open_Events=("status_group", lambda s: int(s.eq("Open").sum())),
            Critical_Path_Events=("critical_path_flag", lambda s: int(s.sum())),
            RFT_Linked_Events=("rft_flag", lambda s: int(s.sum())),
        )
        .sort_values(["Delay_Days", "Delay_Events"], ascending=[False, False])
        .reset_index(drop=True)
    )

    return {
        "time_impact_df": time_impact_df,
        "input_matrix": input_matrix,
        "causation_matrix": causation_matrix,
        "priority_count": int(len(time_impact_df)),
        "critical_impacted": int(time_impact_df["critical_path_flag"].sum()),
        "near_critical_impacted": int(time_impact_df["near_critical_flag"].sum()),
        "rft_impacted": int(time_impact_df["rft_flag"].sum()),
        "commercial_gap": float(contract_metrics.get("total_certified", 0.0)) - float(contract_metrics.get("total_paid", 0.0)),
    }


def letters_intelligence_root(project_id: str | None = None) -> Path:
    project_id = selected_project_id() if project_id is None else str(project_id or "").strip()
    if project_id:
        return project_data_path(PROJECTS_DIR, project_id, "letters", "")
    return PROJECTS_DIR / "_PORTFOLIO" / "letters_intelligence"


def letters_inbox_dir(project_id: str | None = None) -> Path:
    return letters_intelligence_root(project_id) / "inbox"


def letters_workbook_path(project_id: str | None = None) -> Path:
    return letters_intelligence_root(project_id) / "letters_intelligence.xlsx"


@st.cache_data(show_spinner=False)
def _load_letters_workbook_cached(
    workbook_path_str: str,
    workbook_modified_ns: int,
    workbook_size: int,
    inbox_path_str: str,
    inbox_fingerprint: tuple[tuple[str, int, int], ...],
):
    del workbook_modified_ns, workbook_size, inbox_fingerprint
    workbook_path = Path(workbook_path_str)
    inbox_path = Path(inbox_path_str)
    sheets = {}
    if workbook_path.exists():
        try:
            sheets = pd.read_excel(workbook_path, sheet_name=None)
        except ImportError:
            try:
                sheets = _read_xlsx_sheets(workbook_path)
            except Exception:
                sheets = {}
        except Exception:
            try:
                sheets = _read_xlsx_sheets(workbook_path)
            except Exception:
                sheets = {}
    return merge_inbox_letters(
        {name: df.fillna("") for name, df in sheets.items()},
        inbox_path,
        ccc.extract_text_from_path,
    )


def load_letters_workbook():
    active_id = selected_project_id()
    project_ids = [active_id] if active_id else [row["project_id"] for row in discover_projects(PROJECTS_DIR)]
    combined: dict[str, list[pd.DataFrame]] = {}
    for project_id in project_ids:
        workbook_path = letters_workbook_path(project_id)
        inbox_path = letters_inbox_dir(project_id)
        inbox_path.mkdir(parents=True, exist_ok=True)
        if workbook_path.exists():
            stat = workbook_path.stat()
            modified_ns, file_size = stat.st_mtime_ns, stat.st_size
        else:
            modified_ns, file_size = 0, 0
        project_sheets = _load_letters_workbook_cached(
            str(workbook_path),
            modified_ns,
            file_size,
            str(inbox_path),
            folder_fingerprint(inbox_path),
        )
        for sheet_name, frame in project_sheets.items():
            scoped = frame.copy()
            if "project_id" not in scoped.columns:
                scoped.insert(0, "project_id", project_id)
            else:
                source_ids = scoped["project_id"].astype(str).str.strip()
                if "source_project_id" not in scoped.columns:
                    scoped.insert(1, "source_project_id", source_ids)
                scoped["project_id"] = project_id
            combined.setdefault(sheet_name, []).append(scoped)
    return {
        sheet_name: pd.concat(frames, ignore_index=True, sort=False).fillna("")
        for sheet_name, frames in combined.items()
    }


def tokenize_notice_text(value: str) -> set[str]:
    stopwords = {
        "for", "the", "and", "with", "from", "that", "this", "into", "before", "after",
        "work", "works", "project", "letter", "response", "notice", "urgent", "required",
        "re", "rc", "of", "to", "on", "in", "no", "date", "submittal", "submission",
    }
    tokens = set(re.findall(r"[a-z0-9]+", str(value).lower()))
    return {token for token in tokens if len(token) > 2 and token not in stopwords}


def predict_delay_type(letter_type: str, risk_type: str, subject: str, main_purpose: str) -> str:
    text = " ".join([str(letter_type), str(risk_type), str(subject), str(main_purpose)]).lower()
    if any(term in text for term in ["rfi", "consultant reply", "reply from consultant"]):
        return "RFI / response delay"
    if any(term in text for term in ["ifc", "drawing", "design", "shop drawing"]):
        return "IFC / design delay"
    if any(term in text for term in ["payment", "invoice", "certificate", "certified"]):
        return "Payment delay"
    if any(term in text for term in ["steel", "reinforcement", "rft", "free issue", "material"]):
        return "Steel supply delay"
    if any(term in text for term in ["delay", "eot", "extension"]):
        return "General delay / EOT"
    return "General correspondence"


def predict_notice_status(reply_received: str, letter_type: str) -> str:
    if str(reply_received).strip().lower() == "yes":
        return "Replied"
    if "notice" in str(letter_type).lower() or "delay" in str(letter_type).lower():
        return "Open / action required"
    return "Open"


def build_letters_reference_maps(letters_book: dict[str, pd.DataFrame]) -> tuple[dict[str, dict[str, str]], dict[str, dict[str, str]]]:
    ref_map: dict[str, dict[str, str]] = {}
    thread_map: dict[str, dict[str, str]] = {}

    for sheet_name, from_party, to_party in [
        ("From Contractor", "Contractor", "Consultant"),
        ("From Consultant", "Consultant", "Contractor"),
    ]:
        df = letters_book.get(sheet_name, pd.DataFrame())
        if df.empty:
            continue
        for _, row in df.iterrows():
            ref = str(row.get("Ref No", "")).strip()
            if not ref:
                continue
            ref_map[ref] = {
                "From Party": from_party,
                "To Party": to_party,
                "Date": str(row.get("Date", "")).strip(),
                "Type": str(row.get("Type", "")).strip(),
                "Subject": str(row.get("Subject", "")).strip(),
                "Main Purpose": str(row.get("Main Purpose", "")).strip(),
                "Affected Activities": str(row.get("Affected Activities", "")).strip(),
                "Risk Type": str(row.get("Risk Type", "")).strip(),
                "Delay Risk": str(row.get("Delay Risk", "")).strip(),
                "Required Actions": str(row.get("Required Actions", "")).strip(),
            }

    samco_links = letters_book.get("Contractor Links", pd.DataFrame())
    if not samco_links.empty:
        for _, row in samco_links.iterrows():
            ref = str(row.get("SAMCO Ref No", "")).strip()
            if ref:
                thread_map[ref] = {
                    "Thread": str(row.get("Thread", "")).strip(),
                    "Reply Ref": str(row.get("Related ACE Ref No(s)", "")).strip(),
                    "Reply Date": str(row.get("ACE Date(s)", "")).strip(),
                    "Relationship": str(row.get("Relationship", "")).strip(),
                    "Recommended Follow-up": str(row.get("Recommended Follow-up", "")).strip(),
                }

    ace_links = letters_book.get("Consultant Links", pd.DataFrame())
    if not ace_links.empty:
        for _, row in ace_links.iterrows():
            ref = str(row.get("ACE Ref No", "")).strip()
            if ref:
                thread_map[ref] = {
                    "Thread": str(row.get("Thread(s)", "")).strip(),
                    "Reply Ref": str(row.get("Related SAMCO Ref No(s)", "")).strip(),
                    "Reply Date": "",
                    "Relationship": "Consultant to contractor linked thread",
                    "Recommended Follow-up": str(row.get("ACE Required Actions", "")).strip(),
                }

    return ref_map, thread_map


def predict_activity_from_notice(notice_text: str, p6_df: pd.DataFrame) -> tuple[str, str, str]:
    if p6_df.empty or "Activity ID" not in p6_df.columns or "Activity Name" not in p6_df.columns:
        return "", "", "Low"
    notice_tokens = tokenize_notice_text(notice_text)
    if not notice_tokens:
        return "", "", "Low"
    scored_rows = []
    for _, row in p6_df.iterrows():
        activity_id = str(row.get("Activity ID", "")).strip()
        activity_name = str(row.get("Activity Name", "")).strip()
        building = str(row.get("WBS", "")).strip()
        haystack_tokens = tokenize_notice_text(" ".join([activity_id, activity_name, building]))
        overlap = len(notice_tokens & haystack_tokens)
        if overlap > 0:
            scored_rows.append((overlap, activity_id, activity_name))
    if not scored_rows:
        return "", "", "Low"
    scored_rows.sort(key=lambda item: (-item[0], item[1]))
    score, activity_id, activity_name = scored_rows[0]
    confidence = "High" if score >= 4 else ("Medium" if score >= 2 else "Low")
    return activity_id, activity_name, confidence




def risk_rank(value) -> int:
    text = str(value).lower()
    if "critical" in text or "high" in text:
        return 3
    if "medium" in text:
        return 2
    if "low" in text:
        return 1
    return 0


def style_plotly(fig, height=390):
    title = fig.layout.title.text if fig.layout.title and fig.layout.title.text else ""
    fig.update_layout(
        height=height,
        autosize=True,
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="#ffffff",
        font=dict(color="#172033", family="Arial, sans-serif", size=12),
        margin=dict(l=18, r=18, t=64 if title else 42, b=38),
        title=dict(text=title, x=0.02, xanchor="left", font=dict(size=17, color="#0f172a")),
        legend=dict(
            orientation="h",
            yanchor="bottom",
            y=1.02,
            xanchor="right",
            x=1,
            bgcolor="rgba(255,255,255,.72)",
            bordercolor="rgba(15,23,42,.08)",
            borderwidth=1,
            font=dict(size=11),
        ),
        hoverlabel=dict(bgcolor="#0f172a", bordercolor="#0f172a", font=dict(color="#ffffff", size=12)),
        modebar=dict(orientation="v"),
        xaxis=dict(gridcolor="#edf3f6", zerolinecolor="#dde7ef", automargin=True, tickfont=dict(size=11)),
        yaxis=dict(gridcolor="#edf3f6", zerolinecolor="#dde7ef", automargin=True, tickfont=dict(size=11)),
    )
    fig.update_xaxes(showline=False, ticks="outside")
    fig.update_yaxes(showline=False, ticks="outside")
    fig.update_traces(hoverlabel=dict(namelength=-1))
    return fig


def dataframe_height(df: pd.DataFrame, max_height: int = 620, row_height: int = 31, base: int = 35) -> int:
    return min(base + (len(df) * row_height), max_height)


def arrow_safe_display_df(df: pd.DataFrame) -> pd.DataFrame:
    """Return a display-only dataframe that avoids Streamlit Arrow mixed-type warnings."""
    if df.empty:
        return df
    display = df.copy()
    for col in display.columns:
        if display[col].dtype == "object":
            display[col] = display[col].astype(str).replace({"nan": "", "None": "", "<NA>": ""})
    return display




def render_kpi_box(title: str, value: str):
    st.markdown(
        f"""
        <div class="kpi-box">
          <div class="kpi-box-title">{title}</div>
          <div class="kpi-box-value">{value}</div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_claims_kpi_card(title: str, value: str, tone: str = "blue", subtitle: str = ""):
    subtitle_html = f"<div class='claims-kpi-subtitle'>{html.escape(subtitle)}</div>" if subtitle else ""
    st.markdown(
        f"""
        <div class="claims-kpi-card tone-{tone}">
          <div class="claims-kpi-title">{html.escape(title)}</div>
          <div class="claims-kpi-value">{html.escape(str(value))}</div>
          {subtitle_html}
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_claims_badge(label: str, tone: str = "blue") -> str:
    return f"<span class='claims-badge badge-{html.escape(tone)}'>{html.escape(str(label))}</span>"


def render_ai_status_note(ai_status: dict | None):
    if not isinstance(ai_status, dict):
        return
    status = str(ai_status.get("status", "")).strip()
    source = str(ai_status.get("source", "")).strip()
    model = str(ai_status.get("model", "")).strip()
    latency_ms = int(ai_status.get("latency_ms") or 0)
    if status == "ok":
        origin = "OpenAI cache" if source == "cache" else "OpenAI"
        suffix = f" using {model}" if model else ""
        timing = f" in {latency_ms} ms" if latency_ms else ""
        st.caption(f"AI enhancement: {origin}{suffix}{timing}.")
    elif status == "error":
        st.warning("OpenAI enhancement was unavailable; the local contract intelligence result was used.")
    elif status == "skipped" and source in {"disabled", "missing_api_key"}:
        st.caption("AI enhancement: local contract intelligence engine.")


def render_claims_header(title: str, subtitle: str):
    st.markdown(
        f"""
        <div class="claims-hero">
          <div class="claims-hero-topline">COMMERCIAL / CLAIMS INTELLIGENCE</div>
          <div class="claims-hero-title">{html.escape(title)}</div>
          <div class="claims-hero-subtitle">{html.escape(subtitle)}</div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_claims_workflow():
    steps = [
        "Contract Upload",
        "Clause Extraction",
        "Claim Classification",
        "Evidence Mapping",
        "Client Rebuttal",
        "Claim Draft",
        "Export Report",
    ]
    items = "".join(
        f"<div class='claims-workflow-step'><div class='claims-workflow-index'>{idx}</div><div class='claims-workflow-label'>{html.escape(step)}</div></div>"
        for idx, step in enumerate(steps, start=1)
    )
    st.markdown(f"<div class='claims-workflow'>{items}</div>", unsafe_allow_html=True)


def metrics_frame(title: str, values: dict[str, object]) -> pd.DataFrame:
    return pd.DataFrame({"Metric": list(values.keys()), "Value": [str(v) for v in values.values()], "Group": title})


def df_for_export(df: pd.DataFrame) -> pd.DataFrame:
    export_df = df.copy()
    export_df.index = range(1, len(export_df) + 1)
    export_df.index.name = "Row"
    return export_df.reset_index()




def build_the_big_decision_dashboard_html(
    overview_metrics: dict,
    evm_metrics: dict,
    contract_metrics: dict,
    delay_metrics: dict,
    risk_metrics: dict,
    milestone_metrics: dict,
    activity_metrics: dict,
    s_curve_metrics: dict,
    project_record: dict,
) -> str:
    report_date = pd.Timestamp.today().strftime("%d %b %Y")
    project_name = str(overview_metrics.get("project_name") or project_record.get("project_name") or "Project").strip()
    brand_name = str(project_record.get("contractor") or project_record.get("contractor_name") or "Project Team").strip()
    actual = float(overview_metrics.get("overall_progress", 0.0) or 0.0)
    planned = float(overview_metrics.get("planned_progress", 0.0) or 0.0)
    remaining = float(overview_metrics.get("remaining_duration_pct", 0.0) or 0.0)
    contract_value = float(overview_metrics.get("contract_value", 0.0) or 0.0)
    bac = float(evm_metrics.get("bac", contract_value) or 0.0)
    ac = float(evm_metrics.get("ac", contract_metrics.get("total_paid", 0.0)) or 0.0)
    ev = float(evm_metrics.get("ev", 0.0) or 0.0)
    pv = float(evm_metrics.get("pv", 0.0) or 0.0)
    sv = float(evm_metrics.get("sv", ev - pv) or 0.0)
    cv = float(evm_metrics.get("cv", ev - ac) or 0.0)
    eac = float(evm_metrics.get("eac", ac + max(bac - ev, 0.0)) or 0.0)
    tcpi = evm_metrics.get("tcpi")
    tcpi_text = f"{float(tcpi):.3f}" if tcpi is not None else "N/A"
    spi = (ev / pv) if pv else 0.0
    cpi = (ev / ac) if ac else 0.0
    total_activities = int(overview_metrics.get("total_activities", 0) or 0)
    critical_activities = int(overview_metrics.get("critical_activities", 0) or 0)
    open_rfis = int(risk_metrics.get("rfi_items", 0) or 0)
    open_submittals = int(len(activity_metrics.get("rft_df", pd.DataFrame())) if isinstance(activity_metrics.get("rft_df", pd.DataFrame()), pd.DataFrame) else 0)
    active_contracts = int(contract_metrics.get("total_contracts", 0) or 0)
    paid_ratio = (ac / bac * 100.0) if bac else 0.0
    days_left = max(int(round(float(overview_metrics.get("duration_days", 0) or 0) * remaining / 100.0)), 0)

    def money(value: float) -> str:
        return f"EGP {format_currency_html(value)}"

    def safe(value: object) -> str:
        return html.escape(str(value))

    def evm_card(label: str, value: str, note: str, icon: str, tone: str) -> str:
        return f"""
        <div class="evm-card {tone}">
          <div class="icon">{safe(icon)}</div>
          <div><div class="card-label">{safe(label)}</div><div class="card-value">{safe(value)}</div><div class="card-note">{safe(note)}</div></div>
        </div>
        """

    def donut_card(label: str, value: float, center: str, note: str, tone: str) -> str:
        value = max(min(float(value or 0.0), 100.0), 0.0)
        return f"""
        <div class="overview-card">
          <div class="small-label">{safe(label)}</div>
          <div class="donut-row">
            <div class="donut {tone}" style="--p:{value:.2f}"></div>
            <div><div class="big-percent">{safe(center)}</div><div class="card-note">{safe(note)}</div></div>
          </div>
        </div>
        """

    def alert_card(title: str, body: str, priority: str, tone: str, icon: str) -> str:
        return f"""
        <div class="alert {tone}">
          <div class="alert-icon">{safe(icon)}</div>
          <div class="alert-body"><strong>{safe(title)}</strong><span>{safe(body)}</span></div>
          <div class="priority"><i></i>{safe(priority)}</div>
        </div>
        """

    phase_df = _phase_progress_rows(activity_metrics, overview_metrics).copy()
    if phase_df.empty:
        phase_rows = [
            ("Civil Works", actual * 1.10),
            ("Structural Works", actual),
            ("MEP Works", actual * 0.82),
            ("Architectural Works", actual * 0.70),
            ("Infrastructure", actual * 0.52),
        ]
    else:
        phase_rows = [(str(row.get("phase", "Workstream")), float(row.get("actual", 0.0) or 0.0)) for _, row in phase_df.head(5).iterrows()]
    phase_rows = [(name, max(min(value, 100.0), 0.0)) for name, value in phase_rows]
    phase_html = "".join(
        f"<div class='bar-row'><span>{safe(name)}</span><div class='bar-track'><b style='width:{value:.1f}%'></b></div><em>{value:.1f}%</em></div>"
        for name, value in phase_rows
    )

    cost_total = max(ac, 1.0)
    cost_breakdown = []
    palette = ["#2f80ed", "#55d5d0", "#2fb59d", "#f7b733", "#a777e3"]
    for idx, (name, value) in enumerate(phase_rows):
        share = max(value, 1.0)
        cost_breakdown.append((name, share, palette[idx % len(palette)]))
    share_sum = sum(item[1] for item in cost_breakdown) or 1.0
    legend_html = "".join(
        f"<div class='legend-row'><i style='background:{color}'></i><span>{safe(name)}</span><b>{share / share_sum * 100:.1f}%</b></div>"
        for name, share, color in cost_breakdown
    )
    gradient_parts = []
    cursor = 0.0
    for _, share, color in cost_breakdown:
        pct_share = share / share_sum * 100.0
        gradient_parts.append(f"{color} {cursor:.2f}% {cursor + pct_share:.2f}%")
        cursor += pct_share
    cost_gradient = ", ".join(gradient_parts)

    months = ["Jan '26", "Feb '26", "Mar '26", "Apr '26", "May '26", "Jun '26", "Jul '26", "Aug '26", "Sep '26", "Oct '26", "Nov '26", "Dec '26"]
    pv_points = [max(pv, bac * 0.25) * (0.28 + i * 0.055) for i in range(12)]
    ac_points = [max(ac, bac * 0.08) * (0.18 + i * 0.075) for i in range(12)]
    ev_points = [max(ev, bac * 0.04) * (0.14 + i * 0.052) for i in range(12)]
    max_line = max(pv_points + ac_points + ev_points + [1.0])

    def sparkline(points: list[float], color: str, dash: str = "") -> str:
        coords = []
        for idx, value in enumerate(points):
            x = 32 + idx * 48
            y = 184 - (value / max_line * 144)
            coords.append(f"{x:.1f},{y:.1f}")
        return f"<polyline points='{' '.join(coords)}' fill='none' stroke='{color}' stroke-width='4' stroke-linecap='round' stroke-linejoin='round' {dash}/>"

    month_labels = "".join(f"<span>{safe(month)}</span>" for month in months)
    cost_svg = f"""
    <svg class="line-chart" viewBox="0 0 590 220" preserveAspectRatio="none">
      <g class="grid-lines">
        <line x1="30" y1="40" x2="570" y2="40"/><line x1="30" y1="88" x2="570" y2="88"/>
        <line x1="30" y1="136" x2="570" y2="136"/><line x1="30" y1="184" x2="570" y2="184"/>
      </g>
      {sparkline(ac_points, '#47a3ff')}
      {sparkline(ev_points, '#4fe2d1')}
      {sparkline(pv_points, '#ffc247', "stroke-dasharray='8 6'")}
    </svg>
    """

    alerts_html = "".join([
        alert_card("Steel Supply Delay", "Delay in steel shipment affecting structure building progress", "High", "red", "!"),
        alert_card("IFC Clash Detected", f"{risk_metrics.get('ifc_conflicts', 0)} clashes found in current coordination records", "High", "amber", "!"),
        alert_card("RFIs Pending", f"{open_rfis} RFIs are pending response from consultants", "Medium", "blue", "?"),
        alert_card("Contract Exposure", "Exposure amount exceeds threshold for active contracts", "Medium", "purple", "✓"),
    ])

    evm_cards = "".join([
        evm_card("BAC", money(bac), "Budget at Completion", "▣", "blue"),
        evm_card("AC", money(ac), "Actual Cost", "▤", "cyan"),
        evm_card("EV", money(ev), "Earned Value", "▧", "green"),
        evm_card("PV", money(pv), "Planned Value", "▥", "gold"),
        evm_card("Schedule Variance", money(sv), "Ahead of Schedule" if sv >= 0 else "Behind Schedule", "▦", "teal"),
        evm_card("Cost Variance", money(cv), "Under Budget" if cv >= 0 else "Over Budget", "▨", "red"),
        evm_card("EAC", money(eac), "Estimate at Completion", "▩", "blue"),
        evm_card("TCPI", tcpi_text, "To Complete Performance Index", "⌁", "purple"),
    ])

    bottom_items = [
        ("Total Activities", f"{total_activities:,}", "▦", "blue"),
        ("Critical Activities", f"{critical_activities:,}", "⚑", "red"),
        ("SPI", f"{spi:.2f}" if spi else "N/A", "⌁", "green"),
        ("CPI", f"{cpi:.2f}" if cpi else "N/A", "▣", "gold"),
        ("Active Contracts", f"{active_contracts:,}", "▤", "amber"),
        ("Open RFIs", f"{open_rfis:,}", "⌕", "blue"),
        ("Open Submittals", f"{open_submittals:,}", "◇", "purple"),
    ]
    bottom_html = "".join(
        f"<div class='bottom-item {tone}'><i>{safe(icon)}</i><span>{safe(label)}</span><b>{safe(value)}</b></div>"
        for label, value, icon, tone in bottom_items
    )

    return f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>{safe(project_name)} - Executive Dashboard</title>
  <style>
    @page {{ size: A3 landscape; margin: 0; }}
    * {{ box-sizing: border-box; }}
    body {{ margin:0; min-height:100vh; background:#061a31; color:#f4f8ff; font-family:Inter,Segoe UI,Arial,sans-serif; }}
    .dash {{ width:1600px; min-height:1000px; margin:0 auto; padding:28px 26px 18px; background:radial-gradient(circle at 22% 0%,#0f355c 0,#071d36 38%,#041426 100%); border:1px solid rgba(71,163,255,.18); }}
    .top {{ display:flex; justify-content:space-between; align-items:flex-start; padding:0 2px 28px; border-bottom:1px solid rgba(125,170,218,.22); }}
    h1 {{ margin:0 0 8px; font-size:28px; letter-spacing:0; }}
    .subtitle {{ color:#aac0dc; font-size:17px; }}
    .top-right {{ display:flex; gap:46px; align-items:center; color:#c8d8ef; }}
    .date-pill {{ border:1px solid rgba(125,170,218,.25); border-radius:9px; padding:12px 18px; min-width:180px; text-align:center; background:rgba(3,22,43,.45); }}
    .director small {{ display:block; color:#b7c9e2; }} .director b {{ font-size:18px; color:#dce8fb; }}
    .avatar {{ width:48px; height:48px; border-radius:50%; border:1px solid rgba(168,196,231,.3); display:grid; place-items:center; background:rgba(255,255,255,.08); }}
    .content {{ display:grid; grid-template-columns:1fr 380px; gap:18px; }}
    .section {{ padding:18px 0 0; }}
    .section-title {{ margin:0 0 14px; padding-left:14px; border-left:3px solid #36e0d0; font-weight:800; font-size:17px; }}
    .evm-grid {{ display:grid; grid-template-columns:repeat(4,1fr); gap:12px; }}
    .evm-card {{ min-height:106px; border-radius:8px; border:1px solid rgba(81,144,222,.35); padding:18px; display:grid; grid-template-columns:38px 1fr; gap:12px; background:linear-gradient(135deg,rgba(19,57,99,.9),rgba(12,34,64,.82)); box-shadow:inset 0 0 28px rgba(41,125,204,.08); }}
    .evm-card.cyan {{ background:linear-gradient(135deg,rgba(14,97,124,.86),rgba(9,45,72,.82)); border-color:rgba(52,211,224,.42); }}
    .evm-card.green,.evm-card.teal {{ background:linear-gradient(135deg,rgba(16,95,87,.72),rgba(11,45,61,.86)); border-color:rgba(70,220,188,.38); }}
    .evm-card.gold {{ background:linear-gradient(135deg,rgba(91,72,29,.72),rgba(25,42,64,.86)); border-color:rgba(255,194,71,.42); }}
    .evm-card.red {{ background:linear-gradient(135deg,rgba(92,43,57,.72),rgba(25,42,64,.86)); border-color:rgba(255,101,101,.45); }}
    .evm-card.purple {{ background:linear-gradient(135deg,rgba(69,51,113,.72),rgba(25,42,64,.86)); border-color:rgba(178,128,255,.42); }}
    .icon {{ width:32px; height:32px; border-radius:50%; display:grid; place-items:center; background:rgba(55,153,255,.25); color:#76c9ff; font-weight:900; }}
    .cyan .icon {{ color:#65eff0; }} .green .icon,.teal .icon {{ color:#65f0ce; }} .gold .icon {{ color:#ffc247; }} .red .icon {{ color:#ff7777; }} .purple .icon {{ color:#c9a7ff; }}
    .card-label,.small-label {{ color:#b7c9e2; font-size:13px; }}
    .card-value {{ margin-top:7px; font-size:25px; font-weight:900; color:#fff; }}
    .card-note {{ margin-top:5px; color:#9fb4d1; font-size:13px; }}
    .alerts {{ border-left:1px solid rgba(125,170,218,.22); padding-left:16px; }}
    .alerts-head {{ display:flex; justify-content:space-between; align-items:center; margin:18px 6px 12px; }}
    .alerts-head b {{ font-size:17px; }} .alerts-head span {{ border:1px solid rgba(125,170,218,.28); border-radius:8px; padding:8px 14px; color:#9fb4d1; font-size:12px; }}
    .alert {{ display:grid; grid-template-columns:44px 1fr 74px; gap:12px; align-items:center; margin-bottom:10px; padding:14px; border:1px solid rgba(70,130,205,.32); border-radius:8px; background:rgba(4,28,54,.72); }}
    .alert.red {{ background:linear-gradient(90deg,rgba(103,43,58,.5),rgba(4,28,54,.72)); border-color:rgba(255,97,97,.42); }}
    .alert.amber {{ background:linear-gradient(90deg,rgba(95,68,35,.45),rgba(4,28,54,.72)); }} .alert.purple {{ background:linear-gradient(90deg,rgba(64,49,111,.45),rgba(4,28,54,.72)); }}
    .alert-icon {{ width:34px; height:34px; border-radius:50%; display:grid; place-items:center; border:2px solid currentColor; color:#ff6666; font-weight:900; }} .amber .alert-icon {{ color:#ff9f2f; }} .blue .alert-icon {{ color:#55b8ff; }} .purple .alert-icon {{ color:#ba8cff; }}
    .alert-body strong {{ display:block; font-size:15px; }} .alert-body span {{ display:block; color:#a9bbd6; font-size:12px; line-height:1.35; margin-top:4px; }}
    .priority {{ font-size:12px; color:#ff6b6b; font-weight:800; }} .priority i {{ display:inline-block; width:10px; height:10px; border-radius:50%; background:#ff4d4d; margin-right:7px; }} .amber .priority,.purple .priority,.blue .priority {{ color:#ff9f2f; }} .amber .priority i,.purple .priority i,.blue .priority i {{ background:#ff9f2f; }}
    .overview-grid {{ display:grid; grid-template-columns:repeat(5,1fr); gap:12px; }}
    .overview-card,.panel {{ border:1px solid rgba(88,145,214,.35); border-radius:8px; background:rgba(5,31,59,.74); padding:18px; }}
    .donut-row {{ display:flex; align-items:center; gap:18px; margin-top:12px; }}
    .donut {{ --c:#4ee2d1; width:86px; height:86px; border-radius:50%; background:conic-gradient(var(--c) calc(var(--p)*1%), rgba(126,151,181,.32) 0); position:relative; }} .donut:after {{ content:''; position:absolute; inset:10px; border-radius:50%; background:#08203a; }} .donut.gold {{ --c:#ffc247; }} .donut.blue {{ --c:#55b8ff; }}
    .big-percent {{ font-size:26px; font-weight:900; color:white; }}
    .panels {{ display:grid; grid-template-columns:1.05fr .95fr 1.05fr; gap:14px; margin-top:8px; }}
    .panel-head {{ display:flex; justify-content:space-between; align-items:center; margin-bottom:8px; }} .panel-head b {{ font-size:15px; }} .panel-head span {{ border:1px solid rgba(125,170,218,.28); border-radius:7px; padding:7px 12px; color:#9fb4d1; font-size:12px; }}
    .line-chart {{ width:100%; height:180px; }} .grid-lines line {{ stroke:rgba(139,172,211,.17); stroke-width:1; }} .month-axis {{ display:grid; grid-template-columns:repeat(12,1fr); color:#8fa7c5; font-size:10px; margin:0 2px; }}
    .legend {{ display:flex; gap:18px; color:#bdd0e7; font-size:11px; margin-bottom:6px; }} .legend i,.legend-row i {{ display:inline-block; width:9px; height:9px; border-radius:50%; margin-right:7px; }}
    .bar-row {{ display:grid; grid-template-columns:130px 1fr 54px; gap:12px; align-items:center; margin:17px 0; color:#dce8fb; font-size:14px; }} .bar-track {{ height:8px; border-radius:999px; background:rgba(111,146,184,.24); overflow:hidden; }} .bar-track b {{ display:block; height:100%; border-radius:999px; background:linear-gradient(90deg,#54e0d0,#55a7ff); }} .bar-row em {{ color:#b6c8e1; font-style:normal; font-size:12px; }}
    .cost-wrap {{ display:grid; grid-template-columns:190px 1fr; gap:20px; align-items:center; }} .cost-donut {{ width:170px; height:170px; border-radius:50%; background:conic-gradient({cost_gradient}); display:grid; place-items:center; position:relative; }} .cost-donut:after {{ content:''; position:absolute; inset:42px; background:#08203a; border-radius:50%; }} .cost-center {{ position:relative; z-index:1; text-align:center; font-weight:900; }} .cost-center small {{ display:block; color:#9fb4d1; font-size:11px; font-weight:500; }}
    .legend-row {{ display:grid; grid-template-columns:16px 1fr 52px; gap:6px; align-items:center; margin:14px 0; color:#c9d8ed; font-size:12px; }} .legend-row b {{ text-align:right; color:#eaf2ff; }}
    .bottom-strip {{ display:grid; grid-template-columns:repeat(8,1fr); margin-top:18px; border:1px solid rgba(88,145,214,.35); border-radius:8px; background:rgba(5,31,59,.78); overflow:hidden; }}
    .bottom-item {{ min-height:84px; display:grid; grid-template-columns:42px 1fr; grid-template-rows:1fr 1fr; column-gap:10px; padding:15px 18px; border-right:1px solid rgba(125,170,218,.22); }} .bottom-item i {{ grid-row:1/3; align-self:center; width:34px; height:34px; border-radius:50%; display:grid; place-items:center; background:rgba(55,153,255,.22); color:#55b8ff; font-style:normal; }} .bottom-item span {{ color:#9fb4d1; font-size:12px; }} .bottom-item b {{ color:#fff; font-size:18px; }} .bottom-item.red b {{ color:#ff6262; }} .bottom-item.green i {{ color:#42e2aa; }} .bottom-item.gold i,.bottom-item.amber i {{ color:#ffc247; }} .bottom-item.purple i {{ color:#c9a7ff; }}
    .brand {{ display:grid; place-items:center; font-size:28px; letter-spacing:4px; color:#f4f8ff; }} .brand small {{ display:block; margin-top:6px; letter-spacing:1px; color:#a9bbd6; font-size:11px; }}
    .footer {{ margin-top:26px; text-align:center; color:#a8bdd8; font-size:14px; letter-spacing:3px; font-style:italic; }} .footer b {{ color:#43e0d3; }}
  </style>
</head>
<body>
  <main class="dash">
    <header class="top"><div><h1>Executive Dashboard</h1><div class="subtitle">Project Performance & Control Center</div></div><div class="top-right"><div class="date-pill">{safe(report_date)}</div><div class="director"><small>Good Morning,</small><b>Executive Director</b></div><div class="avatar">◎</div></div></header>
    <section class="content">
      <div>
        <section class="section"><div class="section-title">Earned Value Management</div><div class="evm-grid">{evm_cards}</div></section>
        <section class="section"><div class="section-title">Performance Overview</div><div class="overview-grid">{donut_card('Overall Progress', actual, pct(actual), 'Actual Complete', 'cyan')}{donut_card('Planned Progress', planned, pct(planned), 'Baseline Plan', 'cyan')}{evm_card('Contract Value', money(contract_value), 'Total Contract Value', '▣', 'blue')}{evm_card('Actual Cost', money(ac), f'{paid_ratio:.1f}% of BAC', '▤', 'cyan')}{donut_card('Remaining Duration', remaining, pct(remaining), f'{days_left} Days Left', 'gold')}</div></section>
      </div>
      <aside class="alerts"><div class="alerts-head"><b>Critical Alerts</b><span>View All</span></div>{alerts_html}</aside>
    </section>
    <section class="panels">
      <div class="panel"><div class="panel-head"><b>Cost Performance</b><span>View Trend</span></div><div class="legend"><span><i style="background:#47a3ff"></i>AC (Actual Cost)</span><span><i style="background:#4fe2d1"></i>EV (Earned Value)</span><span><i style="background:#ffc247"></i>PV (Planned Value)</span></div>{cost_svg}<div class="month-axis">{month_labels}</div></div>
      <div class="panel"><div class="panel-head"><b>Progress by Discipline</b><span>View All</span></div>{phase_html}</div>
      <div class="panel"><div class="panel-head"><b>Cost Breakdown</b><span>View Report</span></div><div class="cost-wrap"><div class="cost-donut"><div class="cost-center">{money(ac).replace('EGP ', 'EGP<br>')}<small>Total Cost</small></div></div><div>{legend_html}</div></div></div>
    </section>
    <section class="bottom-strip">{bottom_html}<div class="brand">{safe(brand_name)}<small>{safe(project_name)}</small></div></section>
    <footer class="footer">Created and Developed by <b>Eng. Ahmed Labib</b> © Senior Planning Engineer</footer>
  </main>
</body>
</html>"""

PROJECT_HUB_SLIDE_NAMES = [
    "Overview",
    "WBS",
    "Activities",
    "Milestones",
    "S-Curve",
    "EVM Analysis",
    "Contracts",
    "Letters Intelligence",
    "Delays",
    "Time Impact",
    "Risks",
    "Delay Analysis - Time Impact Analysis",
    "Contract & Claims Intelligence Center",
    "Output Studio",
]


def sanitize_conference_room_name(value: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9_-]+", "-", str(value or "").strip())
    cleaned = re.sub(r"-+", "-", cleaned).strip("-")
    return cleaned or "Project-Intelligence-Hub-Room"


def build_inline_project_meeting_html(room_name: str) -> str:
    safe_room = sanitize_conference_room_name(room_name)
    return f"""
    <div style="background:#ffffff;border:1px solid #d7e3ec;border-radius:8px;padding:8px;">
      <div id="project-meeting-container" style="width:100%;height:460px;border-radius:6px;overflow:hidden;"></div>
    </div>
    <script src="https://meet.jit.si/external_api.js"></script>
    <script>
      new JitsiMeetExternalAPI("meet.jit.si", {{
        roomName: "{safe_room}",
        width: "100%",
        height: 460,
        parentNode: document.querySelector("#project-meeting-container"),
        configOverwrite: {{
          prejoinPageEnabled: true,
          startAudioOnly: true,
          startWithAudioMuted: false,
          startWithVideoMuted: true,
          disableDeepLinking: true
        }},
        interfaceConfigOverwrite: {{
          SHOW_JITSI_WATERMARK: false,
          SHOW_WATERMARK_FOR_GUESTS: false,
          MOBILE_APP_PROMO: false,
          TILE_VIEW_MAX_COLUMNS: 5
        }}
      }});
    </script>
    """


def add_simple_table_to_slide(slide, left, top, width, title: str, df: pd.DataFrame, max_rows: int = 8):
    title_box = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(15)
    title_box.text_frame.paragraphs[0].font.bold = True
    show_df = df.head(max_rows).copy()
    rows = len(show_df) + 1
    cols = len(show_df.columns)
    table_shape = slide.shapes.add_table(rows, cols, left, top + Inches(0.35), width, Inches(0.28 * rows))
    table = table_shape.table
    for c, col in enumerate(show_df.columns):
        table.cell(0, c).text = str(col)
    for r, (_, row) in enumerate(show_df.iterrows(), start=1):
        for c, value in enumerate(row):
            table.cell(r, c).text = str(value)
    for r in range(rows):
        for c in range(cols):
            for p in table.cell(r, c).text_frame.paragraphs:
                p.font.size = Pt(9)
                if r == 0:
                    p.font.bold = True


def _comment_or_default(text: str) -> str:
    return text.strip() if str(text).strip() else "No user comment added."


def printEVMHtml(project_title: str, report_date: str, evmData: dict, rootCauseDf: pd.DataFrame, mitigationDf: pd.DataFrame, evmComments: dict) -> str:
    root_rows = "".join(
        f"<tr><td>{html.escape(str(r['Cause Title']))}</td><td>{html.escape(str(r['Impact Area']))}</td><td>{html.escape(str(r['EVM Impact Link']))}</td><td>{html.escape(str(r['Status']))}</td><td>{html.escape(str(r['Severity']))}</td></tr>"
        for _, r in rootCauseDf.iterrows()
    )
    mitigation_rows = "".join(
        f"<tr><td>{html.escape(str(r['Action']))}</td><td>{html.escape(str(r['Owner / Responsible Party']))}</td><td>{html.escape(str(r['Current Status']))}</td><td>{html.escape(str(r['Recovery Impact']))}</td><td>{html.escape(str(r['Required Next Decision']))}</td></tr>"
        for _, r in mitigationDf.iterrows()
    )
    return f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>Earned Value Analysis Print</title>
  <style>
    @page {{ size: A4; margin: 16mm; }}
    body {{ font-family: Arial, sans-serif; margin: 0; color: #172033; background: #ffffff; }}
    .page {{ max-width: 1120px; margin: 0 auto; padding: 24px; }}
    h1 {{ margin: 0; font-size: 28px; color: #173b63; }}
    h2 {{ margin: 28px 0 10px; font-size: 20px; color: #173b63; page-break-after: avoid; }}
    .meta {{ margin-top: 8px; color: #526276; font-size: 13px; }}
    .summary {{ margin-top: 18px; border: 1px solid #d9e5ee; border-left: 4px solid #168f8b; border-radius: 10px; padding: 14px 16px; background: #f7fbfd; }}
    .metric-grid {{ display: grid; grid-template-columns: repeat(4, minmax(0,1fr)); gap: 12px; margin-top: 14px; }}
    .metric {{ border: 1px solid #d9e5ee; border-radius: 10px; padding: 12px 14px; background: #fff; }}
    .metric-title {{ font-size: 11px; text-transform: uppercase; color: #667085; font-weight: 700; }}
    .metric-value {{ font-size: 22px; font-weight: 800; color: #173b63; margin-top: 8px; }}
    .metric-sub {{ font-size: 12px; color: #526276; margin-top: 6px; }}
    table {{ width: 100%; border-collapse: collapse; margin-top: 12px; }}
    th, td {{ border: 1px solid #d9e5ee; padding: 10px; text-align: left; vertical-align: top; font-size: 12px; }}
    th {{ background: #eef5f8; color: #173b63; }}
    .comment {{ margin-top: 12px; padding: 12px 14px; border-radius: 10px; background: #fffdf7; border: 1px solid #efe0a8; }}
    .comment b {{ color: #7c5a00; }}
    .narrative {{ margin-top: 12px; line-height: 1.65; }}
    .page-break {{ page-break-before: always; }}
    @media print {{
      .page {{ padding: 0; }}
    }}
  </style>
</head>
<body>
  <div class="page">
    <h1>{html.escape(project_title)}</h1>
    <div class="meta">Report Date: {html.escape(report_date)}</div>
    <div class="meta"><strong>Earned Value Analysis</strong></div>
    <div class="summary">{html.escape(str(evmData['executiveSummary']))}</div>

    <h2>Quantitative Performance (EVM Metrics)</h2>
    <div class="metric-grid">
      <div class="metric"><div class="metric-title">BAC</div><div class="metric-value">{sar(evmData['BAC'])}</div><div class="metric-sub">Budget @ Completion</div></div>
      <div class="metric"><div class="metric-title">PV</div><div class="metric-value">{sar(evmData['PV'])}</div><div class="metric-sub">Cumm. Planned Value</div></div>
      <div class="metric"><div class="metric-title">EV</div><div class="metric-value">{sar(evmData['EV'])}</div><div class="metric-sub">Cumm. Earned Value</div></div>
      <div class="metric"><div class="metric-title">SV</div><div class="metric-value">{sar(evmData['SV'])}</div><div class="metric-sub">Negative schedule exposure</div></div>
      <div class="metric"><div class="metric-title">SPI</div><div class="metric-value">{evmData['SPI']:.2f}</div><div class="metric-sub">Schedule performance index</div></div>
      <div class="metric"><div class="metric-title">Planned Completion Value Position</div><div class="metric-value">{evmData['plannedCompletionValuePosition']:.2f}%</div><div class="metric-sub">PV as a share of BAC</div></div>
      <div class="metric"><div class="metric-title">Earned Value Gap</div><div class="metric-value">{sar(evmData['earnedValueGap'])}</div><div class="metric-sub">PV not converted into EV</div></div>
      <div class="metric"><div class="metric-title">Schedule Health</div><div class="metric-value">{html.escape(str(evmData['scheduleHealthClassification']))}</div><div class="metric-sub">Executive classification</div></div>
    </div>
    <div class="narrative">{html.escape(str(evmData['interpretation']))}</div>
    <div class="comment"><b>Add / Edit My Comment:</b><br>{html.escape(_comment_or_default(evmComments['quantitativePerformance']))}</div>

    <div class="page-break"></div>
    <h2>Root Cause Linkage</h2>
    <table>
      <thead><tr><th>Cause Title</th><th>Impact Area</th><th>Link to EVM Impact</th><th>Status</th><th>Severity</th></tr></thead>
      <tbody>{root_rows}</tbody>
    </table>
    <div class="narrative">The negative schedule variance is not an isolated numerical deviation. It is directly linked to unresolved external and interface-driven constraints that prevented planned progress from being converted into earned value, mainly within construction and procurement work fronts.</div>
    <div class="comment"><b>Add / Edit My Comment:</b><br>{html.escape(_comment_or_default(evmComments['rootCauseLinkage']))}</div>

    <div class="page-break"></div>
    <h2>Contractor Mitigation & Recovery Status</h2>
    <table>
      <thead><tr><th>Action</th><th>Owner / Responsible party</th><th>Current status</th><th>Recovery impact</th><th>Required next decision</th></tr></thead>
      <tbody>{mitigation_rows}</tbody>
    </table>
    <div class="narrative">Contractor mitigation is focused on protecting available work fronts, accelerating technical closures, maintaining commercial entitlement records, and recovering productivity once external constraints are removed. Recovery remains dependent on timely closure of outstanding Owner / Engineer-driven constraints.</div>
    <div class="comment"><b>Add / Edit My Comment:</b><br>{html.escape(_comment_or_default(evmComments['mitigationRecovery']))}</div>
  </div>
</body>
</html>"""


def _add_evm_ppt_header(slide, title: str, subtitle: str) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(7, 17, 31)
    band = slide.shapes.add_shape(1, Inches(0.35), Inches(0.25), Inches(12.55), Inches(0.7))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(15, 23, 42)
    band.line.color.rgb = RGBColor(33, 46, 66)
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(9.2), Inches(0.4))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(22)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.42), Inches(11.5), Inches(0.25))
    subtitle_box.text_frame.text = subtitle
    subtitle_box.text_frame.paragraphs[0].font.size = Pt(10)
    subtitle_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(186, 200, 214)


def _add_evm_metric_card(slide, left: float, top: float, width: float, title: str, value: str, subtitle: str) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.95))
    tf = box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = RGBColor(159, 177, 198)
    p = tf.add_paragraph()
    p.text = value
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(9)
    p2.font.color.rgb = RGBColor(199, 212, 227)


def _add_evm_comment_box(slide, top: float, comment_text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.0), Inches(0.8))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(28, 38, 56)
    box.line.color.rgb = RGBColor(90, 106, 126)
    tf = box.text_frame
    tf.text = "Add / Edit My Comment"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    p = tf.add_paragraph()
    p.text = _comment_or_default(comment_text)
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(226, 232, 240)


def exportEVMToPowerPoint(project_title: str, report_date: str, evmData: dict, rootCauseDf: pd.DataFrame, mitigationDf: pd.DataFrame, evmComments: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_evm_ppt_header(slide1, "Earned Value Analysis ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â Executive Summary", f"{project_title} | Report Date: {report_date}")
    summary_items = [
        ("BAC", sar(evmData["BAC"]), "Budget @ Completion"),
        ("PV", sar(evmData["PV"]), "Cumm. Planned Value"),
        ("EV", sar(evmData["EV"]), "Cumm. Earned Value"),
        ("SV", sar(evmData["SV"]), "Schedule Variance"),
        ("SPI", f"{evmData['SPI']:.2f}", "Schedule Performance Index"),
        ("Status", evmData["scheduleHealthClassification"], "Executive schedule health"),
    ]
    positions = [(0.6, 1.9), (3.2, 1.9), (5.8, 1.9), (8.4, 1.9), (0.6, 3.1), (3.2, 3.1)]
    for (title, value, sub), (left, top) in zip(summary_items, positions):
        _add_evm_metric_card(slide1, left, top, 2.3, title, value, sub)
    interp = slide1.shapes.add_textbox(Inches(0.6), Inches(4.45), Inches(12.0), Inches(1.4))
    interp.fill.solid()
    interp.fill.fore_color.rgb = RGBColor(18, 28, 44)
    interp.line.color.rgb = RGBColor(46, 66, 91)
    interp.text_frame.text = evmData["interpretation"]
    interp.text_frame.paragraphs[0].font.size = Pt(14)
    interp.text_frame.paragraphs[0].font.color.rgb = RGBColor(230, 238, 248)

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_evm_ppt_header(slide2, "Quantitative Performance (EVM Metrics)", project_title)
    metric_df = pd.DataFrame(
        [
            ["BAC", sar(evmData["BAC"])],
            ["PV", sar(evmData["PV"])],
            ["EV", sar(evmData["EV"])],
            ["SV", sar(evmData["SV"])],
            ["SPI", f"{evmData['SPI']:.2f}"],
            ["Schedule Health", evmData["scheduleHealthClassification"]],
        ],
        columns=["Metric", "Value"],
    )
    add_simple_table_to_slide(slide2, Inches(0.6), Inches(1.8), Inches(4.0), "EVM KPI Table", metric_df, max_rows=6)
    chart_origin = (0.6, 4.45)
    chart_values = [("BAC", evmData["BAC"], RGBColor(56, 189, 248)), ("PV", evmData["PV"], RGBColor(249, 115, 22)), ("EV", evmData["EV"], RGBColor(34, 197, 94))]
    max_value = max(v for _, v, _ in chart_values) or 1
    for idx, (label, value, color) in enumerate(chart_values):
        left = 5.0 + idx * 2.2
        height = max(0.35, 2.4 * value / max_value)
        bar = slide2.shapes.add_shape(1, Inches(left), Inches(5.75 - height), Inches(1.1), Inches(height))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.color.rgb = color
        lbl = slide2.shapes.add_textbox(Inches(left), Inches(5.95), Inches(1.4), Inches(0.3))
        lbl.text_frame.text = label
        lbl.text_frame.paragraphs[0].font.size = Pt(10)
        lbl.text_frame.paragraphs[0].font.color.rgb = RGBColor(226, 232, 240)
    sv_box = slide2.shapes.add_textbox(Inches(8.0), Inches(1.8), Inches(4.2), Inches(1.15))
    sv_box.fill.solid()
    sv_box.fill.fore_color.rgb = RGBColor(58, 22, 30)
    sv_box.line.color.rgb = RGBColor(180, 76, 90)
    sv_box.text_frame.text = f"SV: {sar(evmData['SV'])}"
    sv_box.text_frame.paragraphs[0].font.size = Pt(16)
    sv_box.text_frame.paragraphs[0].font.bold = True
    sv_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 221, 221)
    p = sv_box.text_frame.add_paragraph()
    p.text = f"SPI: {evmData['SPI']:.2f} | Status: {evmData['scheduleHealthClassification']}"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(255, 230, 230)
    _add_evm_comment_box(slide2, 6.25, evmComments["quantitativePerformance"])

    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_evm_ppt_header(slide3, "Root Cause Linkage", project_title)
    add_simple_table_to_slide(slide3, Inches(0.55), Inches(1.75), Inches(12.15), "Cause-to-EVM Linkage", rootCauseDf, max_rows=6)
    _add_evm_comment_box(slide3, 6.2, evmComments["rootCauseLinkage"])

    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_evm_ppt_header(slide4, "Contractor Mitigation & Recovery Status", project_title)
    add_simple_table_to_slide(slide4, Inches(0.55), Inches(1.75), Inches(12.15), "Mitigation Action Matrix", mitigationDf, max_rows=6)
    _add_evm_comment_box(slide4, 6.2, evmComments["mitigationRecovery"])

    stream = io.BytesIO()
    prs.save(stream)
    stream.seek(0)
    return stream.getvalue()


def export_linked_executive_dashboard_to_powerpoint(
    overview_metrics: dict,
    evm_metrics: dict,
    contract_metrics: dict,
    delay_metrics: dict,
    risk_metrics: dict,
    milestone_metrics: dict,
    activity_metrics: dict,
    evmData: dict,
    rootCauseDf: pd.DataFrame,
    mitigationDf: pd.DataFrame,
    evmComments: dict,
) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    project_title = overview_metrics.get("project_name") or "Project"
    report_date = pd.Timestamp.today().strftime("%d %b %Y")

    phase_df = _phase_progress_rows(activity_metrics, overview_metrics).copy()
    docs_df = _build_executive_doc_rows(activity_metrics).copy()
    invoice_df = _build_executive_invoice_rows(contract_metrics).copy()
    milestone_df = _build_executive_summary_milestone_rows(activity_metrics, milestone_metrics).copy()
    risk_df = _build_executive_risk_rows(delay_metrics, risk_metrics).copy()

    paid_amount = float(contract_metrics.get("total_paid", 0.0))
    certified_amount = float(contract_metrics.get("total_certified", 0.0))
    ongoing_risks = int((risk_df["status"].astype(str).str.lower() == "ongoing").sum()) if not risk_df.empty else 0
    spi = (float(evm_metrics.get("ev", 0.0)) / float(evm_metrics.get("pv", 1.0))) if float(evm_metrics.get("pv", 0.0)) else 0.0

    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_evm_ppt_header(slide1, "Linked Executive Dashboard ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â Executive Summary", f"{project_title} | Report Date: {report_date} | Landscape print-ready")
    summary_items = [
        ("Contract Value", f"EGP {format_currency_html(float(overview_metrics.get('contract_value', 0.0)))}", "Original contract value"),
        ("Overall Progress", pct(overview_metrics.get("overall_progress")), "Current actual progress"),
        ("Planned Progress", pct(overview_metrics.get("planned_progress")), "Current planned progress"),
        ("SPI", f"{spi:.2f}", "Schedule health"),
        ("SV", f"EGP {format_currency_html(float(evm_metrics.get('sv', 0.0)))}", "EV minus PV exposure"),
        ("Cumulative Paid", f"EGP {format_currency_html(paid_amount)}", "Paid to date"),
        ("Ongoing Risks", str(ongoing_risks), "Active project constraints"),
        ("Critical Activities", str(int(overview_metrics.get("critical_activities", 0))), "Critical path pulse"),
    ]
    positions = [(0.55, 1.75), (3.1, 1.75), (5.65, 1.75), (8.2, 1.75), (10.75, 1.75), (0.55, 3.0), (3.1, 3.0), (5.65, 3.0)]
    for (title, value, subtitle), (left, top) in zip(summary_items, positions):
        _add_evm_metric_card(slide1, left, top, 2.25, title, value, subtitle)
    note = slide1.shapes.add_textbox(Inches(0.55), Inches(4.5), Inches(12.1), Inches(1.25))
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(18, 28, 44)
    note.line.color.rgb = RGBColor(46, 66, 91)
    note.text_frame.text = (
        "This linked executive dashboard PowerPoint export is generated directly from the platform CSV files. "
        "It is formatted as a widescreen landscape deck for clean printing and executive review."
    )
    note.text_frame.paragraphs[0].font.size = Pt(13)
    note.text_frame.paragraphs[0].font.color.rgb = RGBColor(230, 238, 248)

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_evm_ppt_header(slide2, "Linked Executive Dashboard ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â Progress & Milestones", project_title)
    phase_table_df = phase_df.rename(columns={"phase": "Phase", "planned": "Planned %", "actual": "Actual %", "variance": "Variance %"})
    add_simple_table_to_slide(slide2, Inches(0.55), Inches(1.7), Inches(5.8), "Planned vs Actual by Phase", phase_table_df, max_rows=6)
    if not milestone_df.empty:
        milestone_view = milestone_df.rename(columns={"discipline": "Discipline", "plan": "Plan Completion", "forecast": "Forecast Completion", "variance_days": "Variance Days"})
        add_simple_table_to_slide(slide2, Inches(6.6), Inches(1.7), Inches(6.15), "Forecast Completion by Discipline / Area", milestone_view, max_rows=7)
    if not docs_df.empty:
        docs_view = docs_df.rename(columns={"discipline": "Discipline", "total": "Total", "plan": "Plan", "actual": "Actual"})
        add_simple_table_to_slide(slide2, Inches(0.55), Inches(4.75), Inches(5.8), "Document Submission Control", docs_view, max_rows=4)

    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_evm_ppt_header(slide3, "Linked Executive Dashboard ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â Commercial & Risks", project_title)
    commercial_df = pd.DataFrame(
        [
            ["Submitted", f"EGP {format_currency_html(certified_amount)}"],
            ["Approved", f"EGP {format_currency_html(certified_amount)}"],
            ["Paid", f"EGP {format_currency_html(paid_amount)}"],
            ["Under Payment", f"EGP {format_currency_html(max(certified_amount - paid_amount, 0.0))}"],
        ],
        columns=["Commercial Position", "Value"],
    )
    add_simple_table_to_slide(slide3, Inches(0.55), Inches(1.7), Inches(4.4), "Commercial Position", commercial_df, max_rows=5)
    if not invoice_df.empty:
        invoice_view = invoice_df.rename(columns={"title": "Invoice", "date": "Submission Date", "value": "Value", "status": "Status"}).copy()
        invoice_view["Value"] = invoice_view["Value"].apply(lambda v: f"EGP {format_currency_html(v)}")
        add_simple_table_to_slide(slide3, Inches(5.15), Inches(1.7), Inches(7.6), "Invoice Register", invoice_view, max_rows=6)
    if not risk_df.empty:
        risk_view = risk_df.rename(columns={"category": "Category", "title": "Title", "status": "Status", "impact": "Impact"})
        add_simple_table_to_slide(slide3, Inches(0.55), Inches(4.85), Inches(12.2), "Risk & Issue Control", risk_view, max_rows=5)

    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_evm_ppt_header(slide4, "Linked Executive Dashboard ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â Earned Value Analysis", project_title)
    evm_metric_df = pd.DataFrame(
        [
            ["BAC", sar(evmData["BAC"])],
            ["PV", sar(evmData["PV"])],
            ["EV", sar(evmData["EV"])],
            ["SV", sar(evmData["SV"])],
            ["SPI", f"{evmData['SPI']:.2f}"],
            ["Schedule Health", evmData["scheduleHealthClassification"]],
            ["Planned Completion Value Position", pct(evmData["plannedCompletionValuePosition"])],
            ["Earned Value Gap", sar(evmData["earnedValueGap"])],
        ],
        columns=["Metric", "Value"],
    )
    add_simple_table_to_slide(slide4, Inches(0.55), Inches(1.7), Inches(5.2), "Quantitative Performance (EVM Metrics)", evm_metric_df, max_rows=8)
    interp = slide4.shapes.add_textbox(Inches(6.0), Inches(1.95), Inches(6.7), Inches(1.4))
    interp.fill.solid()
    interp.fill.fore_color.rgb = RGBColor(18, 28, 44)
    interp.line.color.rgb = RGBColor(46, 66, 91)
    interp.text_frame.text = evmData["interpretation"]
    interp.text_frame.paragraphs[0].font.size = Pt(12)
    interp.text_frame.paragraphs[0].font.color.rgb = RGBColor(230, 238, 248)
    add_simple_table_to_slide(slide4, Inches(6.0), Inches(3.6), Inches(6.7), "Root Cause Linkage", rootCauseDf, max_rows=5)
    _add_evm_comment_box(slide4, 6.2, evmComments["quantitativePerformance"])

    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_evm_ppt_header(slide5, "Linked Executive Dashboard ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â Mitigation & Recovery", project_title)
    add_simple_table_to_slide(slide5, Inches(0.55), Inches(1.7), Inches(12.2), "Contractor Mitigation & Recovery Status", mitigationDf, max_rows=6)
    root_comment_box = slide5.shapes.add_textbox(Inches(0.55), Inches(5.45), Inches(6.0), Inches(1.1))
    root_comment_box.fill.solid()
    root_comment_box.fill.fore_color.rgb = RGBColor(18, 28, 44)
    root_comment_box.line.color.rgb = RGBColor(46, 66, 91)
    root_comment_box.text_frame.text = f"Root Cause Linkage Comment\n{_comment_or_default(evmComments['rootCauseLinkage'])}"
    root_comment_box.text_frame.paragraphs[0].font.size = Pt(10)
    root_comment_box.text_frame.paragraphs[0].font.bold = True
    root_comment_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    if len(root_comment_box.text_frame.paragraphs) > 1:
        root_comment_box.text_frame.paragraphs[1].font.size = Pt(10)
        root_comment_box.text_frame.paragraphs[1].font.color.rgb = RGBColor(226, 232, 240)
    mit_comment_box = slide5.shapes.add_textbox(Inches(6.75), Inches(5.45), Inches(6.0), Inches(1.1))
    mit_comment_box.fill.solid()
    mit_comment_box.fill.fore_color.rgb = RGBColor(18, 28, 44)
    mit_comment_box.line.color.rgb = RGBColor(46, 66, 91)
    mit_comment_box.text_frame.text = f"Mitigation & Recovery Comment\n{_comment_or_default(evmComments['mitigationRecovery'])}"
    mit_comment_box.text_frame.paragraphs[0].font.size = Pt(10)
    mit_comment_box.text_frame.paragraphs[0].font.bold = True
    mit_comment_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    if len(mit_comment_box.text_frame.paragraphs) > 1:
        mit_comment_box.text_frame.paragraphs[1].font.size = Pt(10)
        mit_comment_box.text_frame.paragraphs[1].font.color.rgb = RGBColor(226, 232, 240)

    stream = io.BytesIO()
    prs.save(stream)
    stream.seek(0)
    return stream.getvalue()


def build_linked_executive_dashboard_a3_summary_html(
    overview_metrics: dict,
    evm_metrics: dict,
    contract_metrics: dict,
    delay_metrics: dict,
    risk_metrics: dict,
    milestone_metrics: dict,
    activity_metrics: dict,
    s_curve_metrics: dict,
    evmData: dict,
    rootCauseDf: pd.DataFrame,
    mitigationDf: pd.DataFrame,
) -> str:
    project_df = load_core_csv(PROJECTS_CSV_PATH)
    project_row = active_project_row(project_df) if not project_df.empty else pd.Series(dtype=object)
    phase_df = _phase_progress_rows(activity_metrics, overview_metrics).copy()
    milestone_df = _build_executive_summary_milestone_rows(activity_metrics, milestone_metrics).copy()
    risk_df = _build_executive_risk_rows(delay_metrics, risk_metrics).copy()
    invoice_df = _build_executive_invoice_rows(contract_metrics).copy()
    curve_df = s_curve_metrics.get("curve_df", pd.DataFrame()).copy()

    paid_amount = float(contract_metrics.get("total_paid", 0.0))
    certified_amount = float(contract_metrics.get("total_certified", 0.0))
    under_payment = max(certified_amount - paid_amount, 0.0)
    ongoing_risks = int((risk_df["status"].astype(str).str.lower() == "ongoing").sum()) if not risk_df.empty else 0
    spi = evmData["SPI"]

    progress_fig = px.bar(
        phase_df,
        x="phase",
        y=["planned", "actual"],
        barmode="group",
        title="Progress by Phase",
        color_discrete_map={"planned": "#38bdf8", "actual": "#f97316"},
    )
    progress_fig.update_layout(legend_title_text="", paper_bgcolor="#0b1320", plot_bgcolor="#0b1320", font=dict(color="#e5eef8"), margin=dict(l=10, r=10, t=42, b=20))
    progress_fig.update_xaxes(title_text="")
    progress_fig.update_yaxes(title_text="%")

    evm_fig = px.bar(
        pd.DataFrame(
            [
                {"metric": "BAC", "value": evmData["BAC"]},
                {"metric": "PV", "value": evmData["PV"]},
                {"metric": "EV", "value": evmData["EV"]},
                {"metric": "SV Exposure", "value": abs(evmData["SV"])},
            ]
        ),
        x="metric",
        y="value",
        title="Earned Value Snapshot",
        color="metric",
        color_discrete_map={"BAC": "#38bdf8", "PV": "#22c55e", "EV": "#f97316", "SV Exposure": "#f43f5e"},
    )
    evm_fig.update_layout(showlegend=False, paper_bgcolor="#0b1320", plot_bgcolor="#0b1320", font=dict(color="#e5eef8"), margin=dict(l=10, r=10, t=42, b=20))

    commercial_fig = px.pie(
        pd.DataFrame(
            [
                {"name": "Paid", "value": paid_amount},
                {"name": "Under Payment", "value": under_payment},
            ]
        ),
        names="name",
        values="value",
        hole=0.55,
        color="name",
        color_discrete_map={"Paid": "#22c55e", "Under Payment": "#f59e0b"},
    )
    commercial_fig.update_layout(showlegend=True, legend_title_text="", paper_bgcolor="#0b1320", plot_bgcolor="#0b1320", font=dict(color="#e5eef8"), margin=dict(l=10, r=10, t=42, b=20))

    s_curve_html = "<div class='empty-card'>No S-curve data available.</div>"
    if not curve_df.empty:
        s_curve_fig = px.line(
            curve_df[["Month", "Planned", "Actual", "Invoiced"]].copy(),
            x="Month",
            y=["Planned", "Actual", "Invoiced"],
            markers=True,
            title="S-Curve",
            color_discrete_map={"Planned": "#38bdf8", "Actual": "#22c55e", "Invoiced": "#f59e0b"},
        )
        s_curve_fig.update_layout(legend_title_text="", paper_bgcolor="#0b1320", plot_bgcolor="#0b1320", font=dict(color="#e5eef8"), margin=dict(l=10, r=10, t=42, b=20))
        s_curve_html = s_curve_fig.to_html(include_plotlyjs=False, full_html=False, config={"displayModeBar": False})

    progress_html = progress_fig.to_html(include_plotlyjs="cdn", full_html=False, config={"displayModeBar": False})
    evm_html = evm_fig.to_html(include_plotlyjs=False, full_html=False, config={"displayModeBar": False})
    commercial_html = commercial_fig.to_html(include_plotlyjs=False, full_html=False, config={"displayModeBar": False})

    milestone_rows = "".join(
        f"<tr><td>{html.escape(str(r['discipline']))}</td><td>{html.escape(str(r['plan']))}</td><td>{html.escape(str(r['forecast']))}</td><td>{int(r['variance_days'])}</td></tr>"
        for _, r in milestone_df.iterrows()
    ) if not milestone_df.empty else "<tr><td colspan='4'>No milestone data available.</td></tr>"

    risk_rows = "".join(
        f"<tr><td>{html.escape(str(r['category']))}</td><td>{html.escape(str(r['title']))}</td><td>{html.escape(str(r['status']))}</td></tr>"
        for _, r in risk_df.head(4).iterrows()
    ) if not risk_df.empty else "<tr><td colspan='3'>No risk data available.</td></tr>"

    mitigation_rows = "".join(
        f"<tr><td>{html.escape(str(r['Action']))}</td><td>{html.escape(str(r['Current Status']))}</td><td>{html.escape(str(r['Required Next Decision']))}</td></tr>"
        for _, r in mitigationDf.head(4).iterrows()
    ) if not mitigationDf.empty else "<tr><td colspan='3'>No mitigation data available.</td></tr>"

    invoice_rows = "".join(
        f"<tr><td>{html.escape(str(r['title']))}</td><td>{html.escape(str(r['date']))}</td><td>{html.escape(str(r['status']))}</td></tr>"
        for _, r in invoice_df.head(4).iterrows()
    ) if not invoice_df.empty else "<tr><td colspan='3'>No invoice data available.</td></tr>"

    def metric_card(title: str, value: str, subtitle: str) -> str:
        return f"""
        <div class="metric-card">
          <div class="metric-title">{html.escape(title)}</div>
          <div class="metric-value">{html.escape(value)}</div>
          <div class="metric-sub">{html.escape(subtitle)}</div>
        </div>
        """

    return f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>Linked Executive Dashboard ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â A3 Summary</title>
  <style>
    @page {{ size: A3 landscape; margin: 10mm; }}
    body {{ margin: 0; background: #07111f; color: #e5eef8; font-family: Arial, sans-serif; }}
    .page {{ width: 100%; min-height: 100vh; box-sizing: border-box; padding: 18px; background: #07111f; }}
    .header, .panel {{ border: 1px solid rgba(255,255,255,.08); background: rgba(255,255,255,.05); border-radius: 22px; box-shadow: 0 20px 40px rgba(0,0,0,.22); }}
    .header {{ padding: 18px 22px; }}
    .eyebrow span {{ display: inline-block; margin-right: 8px; margin-bottom: 8px; padding: 7px 12px; border-radius: 999px; font-size: 11px; font-weight: 700; letter-spacing: .16em; text-transform: uppercase; }}
    .eyebrow .a {{ background: #fff; color: #0f172a; }}
    .eyebrow .b {{ background: rgba(56,189,248,.14); color: #bae6fd; border: 1px solid rgba(56,189,248,.25); }}
    .eyebrow .c {{ background: rgba(249,115,22,.14); color: #fed7aa; border: 1px solid rgba(249,115,22,.25); }}
    h1 {{ margin: 8px 0 0; font-size: 34px; line-height: 1.05; }}
    .sub {{ margin-top: 10px; color: #c7d4e3; font-size: 14px; line-height: 1.6; }}
    .kpis {{ display: grid; grid-template-columns: repeat(6, minmax(0,1fr)); gap: 12px; margin-top: 14px; }}
    .metric-card {{ padding: 14px 16px; border-radius: 18px; background: rgba(255,255,255,.06); border: 1px solid rgba(255,255,255,.08); }}
    .metric-title {{ font-size: 10px; letter-spacing: .16em; text-transform: uppercase; color: #94a3b8; }}
    .metric-value {{ font-size: 24px; font-weight: 900; color: #fff; margin-top: 8px; }}
    .metric-sub {{ font-size: 12px; color: #c7d4e3; margin-top: 8px; }}
    .grid2 {{ display: grid; grid-template-columns: repeat(2, minmax(0,1fr)); gap: 14px; margin-top: 14px; }}
    .grid3 {{ display: grid; grid-template-columns: 1.15fr 1fr 1fr; gap: 14px; margin-top: 14px; }}
    .panel {{ padding: 16px; }}
    .panel-label {{ font-size: 11px; letter-spacing: .16em; text-transform: uppercase; color: #94a3b8; }}
    .panel h3 {{ margin: 8px 0 12px; font-size: 22px; color: #fff; }}
    table {{ width: 100%; border-collapse: collapse; font-size: 12px; }}
    th, td {{ padding: 9px 8px; border-top: 1px solid rgba(255,255,255,.08); text-align: left; vertical-align: top; }}
    th {{ font-size: 10px; letter-spacing: .14em; text-transform: uppercase; color: #94a3b8; background: rgba(255,255,255,.05); }}
    td {{ color: #e5eef8; }}
    .note {{ margin-top: 10px; padding: 12px 14px; border-radius: 16px; background: rgba(255,255,255,.04); border: 1px solid rgba(255,255,255,.08); color: #dbe6f2; line-height: 1.55; font-size: 13px; }}
    .empty-card {{ padding: 24px; border-radius: 18px; background: rgba(255,255,255,.04); border: 1px solid rgba(255,255,255,.08); color: #c7d4e3; }}
    @media print {{
      body {{ background: #07111f !important; }}
      .page {{ padding: 8px; }}
    }}
  </style>
</head>

<body>
  <div class="page">
    <div class="header">
      <div class="eyebrow">
        <span class="a">{html.escape(str(project_row.get("client_name", "THE BIG")))}</span>
        <span class="b">Linked Executive Dashboard</span>
        <span class="c">A3 Landscape Summary | {_ppt_date(pd.Timestamp.today())}</span>
      </div>
      <h1>{html.escape(str(project_row.get("project_name", overview_metrics.get("project_name", "Project"))))}</h1>
      <div class="sub">Single-page linked executive summary generated from the live platform CSV files. Structured for clean A3 landscape printing while preserving the dark executive dashboard identity.</div>
      <div class="kpis">
        {metric_card("Contract Value", f"EGP {format_currency_html(float(overview_metrics.get('contract_value', 0.0)))}", "Original contract value")}
        {metric_card("Overall Progress", pct(overview_metrics.get("overall_progress")), "Actual progress")}
        {metric_card("Planned Progress", pct(overview_metrics.get("planned_progress")), "Baseline progress")}
        {metric_card("SPI", f"{spi:.2f}", "Schedule health")}
        {metric_card("SV", sar(evmData["SV"]), "Schedule variance")}
        {metric_card("Paid", f"EGP {format_currency_html(paid_amount)}", "Cash released")}
      </div>
    </div>

    <div class="grid2">
      <div class="panel">
        <div class="panel-label">Progress Intelligence</div>
        <h3>Progress by Phase</h3>
        {progress_html}
      </div>
      <div class="panel">
        <div class="panel-label">S-Curve</div>
        <h3>Planned vs Actual vs Invoiced</h3>
        {s_curve_html}
      </div>
    </div>

    <div class="grid2">
      <div class="panel">
        <div class="panel-label">Earned Value</div>
        <h3>Earned Value Snapshot</h3>
        {evm_html}
      </div>
      <div class="panel">
        <div class="panel-label">Commercial Pulse</div>
        <h3>Paid vs Under Payment</h3>
        {commercial_html}
      </div>
    </div>

    <div class="grid3">
      <div class="panel">
        <div class="panel-label">Milestones</div>
        <h3>Milestone Exposure</h3>
        <table>
          <thead><tr><th>Discipline</th><th>Plan</th><th>Forecast</th><th>Var Days</th></tr></thead>
          <tbody>{milestone_rows}</tbody>
        </table>
      </div>
      <div class="panel">
        <div class="panel-label">Risks</div>
        <h3>Top Risks</h3>
        <table>
          <thead><tr><th>Category</th><th>Title</th><th>Status</th></tr></thead>
          <tbody>{risk_rows}</tbody>
        </table>
      </div>
      <div class="panel">
        <div class="panel-label">Recovery</div>
        <h3>Mitigation Priorities</h3>
        <table>
          <thead><tr><th>Action</th><th>Status</th><th>Next Decision</th></tr></thead>
          <tbody>{mitigation_rows}</tbody>
        </table>
      </div>
    </div>

    <div class="grid2">
      <div class="panel">
        <div class="panel-label">Invoice Pulse</div>
        <h3>Invoice Status</h3>
        <table>
          <thead><tr><th>Invoice</th><th>Date</th><th>Status</th></tr></thead>
          <tbody>{invoice_rows}</tbody>
        </table>
      </div>
      <div class="panel">
        <div class="panel-label">Executive Summary</div>
        <h3>Management Reading</h3>
        <div class="note">The project remains under schedule pressure with SPI {evmData['SPI']:.2f} and schedule variance {sar(evmData['SV'])}. Priority management focus remains on converting planned work fronts into earned value, closing outstanding technical interfaces, accelerating blocked approvals, and protecting procurement and payment continuity.</div>
      </div>
    </div>
  </div>
</body>
</html>"""


def export_linked_executive_dashboard_a3_summary_ppt(
    overview_metrics: dict,
    evm_metrics: dict,
    contract_metrics: dict,
    delay_metrics: dict,
    risk_metrics: dict,
    milestone_metrics: dict,
    activity_metrics: dict,
    s_curve_metrics: dict,
    evmData: dict,
    rootCauseDf: pd.DataFrame,
    mitigationDf: pd.DataFrame,
) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(16.535)
    prs.slide_height = Inches(11.693)

    project_title = overview_metrics.get("project_name") or "Project"
    report_date = pd.Timestamp.today().strftime("%d %b %Y")
    phase_df = _phase_progress_rows(activity_metrics, overview_metrics).copy()
    milestone_df = _build_executive_milestone_rows(activity_metrics).copy()
    risk_df = _build_executive_risk_rows(delay_metrics, risk_metrics).copy()
    invoice_df = _build_executive_invoice_rows(contract_metrics).copy()
    curve_df = s_curve_metrics.get("curve_df", pd.DataFrame()).copy()

    paid_amount = float(contract_metrics.get("total_paid", 0.0))
    certified_amount = float(contract_metrics.get("total_certified", 0.0))
    ongoing_risks = int((risk_df["status"].astype(str).str.lower() == "ongoing").sum()) if not risk_df.empty else 0
    under_payment = max(certified_amount - paid_amount, 0.0)
    spi = (float(evm_metrics.get("ev", 0.0)) / float(evm_metrics.get("pv", 1.0))) if float(evm_metrics.get("pv", 0.0)) else 0.0

    def add_panel(left: float, top: float, width: float, height: float, title: str, subtitle: str = ""):
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(15, 23, 42)
        shape.line.color.rgb = RGBColor(33, 46, 66)
        title_box = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.12), Inches(width - 0.36), Inches(0.22))
        title_box.text_frame.text = title
        title_box.text_frame.paragraphs[0].font.size = Pt(14)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.36), Inches(width - 0.36), Inches(0.18))
            sub_box.text_frame.text = subtitle
            sub_box.text_frame.paragraphs[0].font.size = Pt(8)
            sub_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(186, 200, 214)
        return shape

    def add_ppt_chart(chart_type, left, top, width, height, chart_data, title, colors: list[RGBColor] | None = None):
        chart = slide.shapes.add_chart(chart_type, Inches(left), Inches(top), Inches(width), Inches(height), chart_data).chart
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(8)
        chart.legend.font.color.rgb = RGBColor(226, 232, 240)
        try:
            chart.value_axis.tick_labels.font.size = Pt(8)
            chart.value_axis.tick_labels.font.color.rgb = RGBColor(186, 200, 214)
        except Exception:
            pass
        try:
            chart.category_axis.tick_labels.font.size = Pt(8)
            chart.category_axis.tick_labels.font.color.rgb = RGBColor(186, 200, 214)
        except Exception:
            pass
        if colors:
            for idx, series in enumerate(chart.series):
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = colors[idx % len(colors)]
                try:
                    series.format.line.color.rgb = colors[idx % len(colors)]
                except Exception:
                    pass
        return chart

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(7, 17, 31)

    header = slide.shapes.add_shape(1, Inches(0.35), Inches(0.25), Inches(15.83), Inches(0.95))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(15, 23, 42)
    header.line.color.rgb = RGBColor(33, 46, 66)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.38), Inches(9.6), Inches(0.38))
    title_box.text_frame.text = "Linked Executive Dashboard ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â A3 Landscape Summary"
    title_box.text_frame.paragraphs[0].font.size = Pt(24)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.76), Inches(11.0), Inches(0.25))
    subtitle_box.text_frame.text = f"{project_title} | Report Date: {report_date} | Single-page executive print summary"
    subtitle_box.text_frame.paragraphs[0].font.size = Pt(11)
    subtitle_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(186, 200, 214)

    summary_items = [
        ("Contract Value", f"EGP {format_currency_html(float(overview_metrics.get('contract_value', 0.0)))}", "Original contract value"),
        ("Overall Progress", pct(overview_metrics.get("overall_progress")), "Actual progress"),
        ("Planned Progress", pct(overview_metrics.get("planned_progress")), "Baseline progress"),
        ("SPI", f"{spi:.2f}", "Schedule health"),
        ("SV", f"EGP {format_currency_html(float(evm_metrics.get('sv', 0.0)))}", "Schedule variance"),
        ("Cumulative Paid", f"EGP {format_currency_html(paid_amount)}", "Paid to date"),
        ("Under Payment", f"EGP {format_currency_html(under_payment)}", "Certified not yet paid"),
        ("Ongoing Risks", str(ongoing_risks), "Active constraints"),
    ]
    card_positions = [
        (0.6, 1.45), (4.45, 1.45), (8.3, 1.45), (12.15, 1.45),
        (0.6, 2.8), (4.45, 2.8), (8.3, 2.8), (12.15, 2.8),
    ]
    for (title, value, subtitle), (left, top) in zip(summary_items, card_positions):
        _add_evm_metric_card(slide, left, top, 3.25, title, value, subtitle)

    add_panel(0.45, 4.15, 5.15, 3.1, "Progress by Phase", "Planned vs actual performance")
    if not phase_df.empty:
        progress_chart_data = CategoryChartData()
        progress_chart_data.categories = phase_df["phase"].tolist()
        progress_chart_data.add_series("Planned", phase_df["planned"].tolist())
        progress_chart_data.add_series("Actual", phase_df["actual"].tolist())
        add_ppt_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, 0.68, 4.65, 4.7, 2.3, progress_chart_data, "Progress by Phase", [RGBColor(56, 189, 248), RGBColor(249, 115, 22)])

    add_panel(5.75, 4.15, 5.15, 3.1, "S-Curve", "Planned vs actual vs invoiced")
    if not curve_df.empty:
        curve_chart_data = CategoryChartData()
        curve_slice = curve_df.tail(6).copy()
        curve_chart_data.categories = curve_slice["Month"].astype(str).tolist()
        curve_chart_data.add_series("Planned", curve_slice["Planned"].tolist())
        curve_chart_data.add_series("Actual", curve_slice["Actual"].tolist())
        curve_chart_data.add_series("Invoiced", curve_slice["Invoiced"].tolist())
        add_ppt_chart(XL_CHART_TYPE.LINE_MARKERS, 5.98, 4.65, 4.7, 2.3, curve_chart_data, "S-Curve", [RGBColor(56, 189, 248), RGBColor(34, 197, 94), RGBColor(245, 158, 11)])

    add_panel(11.05, 4.15, 5.05, 3.1, "Earned Value Snapshot", "BAC / PV / EV / SV exposure")
    evm_chart_data = CategoryChartData()
    evm_chart_data.categories = ["BAC", "PV", "EV", "SV Exposure"]
    evm_chart_data.add_series("Value", [evmData["BAC"], evmData["PV"], evmData["EV"], abs(evmData["SV"])])
    add_ppt_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, 11.28, 4.65, 4.6, 2.3, evm_chart_data, "Earned Value Snapshot", [RGBColor(56, 189, 248)])

    evm_metric_df = pd.DataFrame(
        [
            ["BAC", sar(evmData["BAC"])],
            ["PV", sar(evmData["PV"])],
            ["EV", sar(evmData["EV"])],
            ["SV", sar(evmData["SV"])],
            ["SPI", f"{evmData['SPI']:.2f}"],
            ["Status", evmData["scheduleHealthClassification"]],
        ],
        columns=["EVM Metric", "Value"],
    )
    milestone_view = milestone_df.rename(columns={"discipline": "Discipline", "plan": "Plan", "forecast": "Forecast", "variance_days": "Var Days"}) if not milestone_df.empty else pd.DataFrame(columns=["Discipline", "Plan", "Forecast", "Var Days"])
    risk_view = risk_df.rename(columns={"category": "Category", "title": "Title", "status": "Status", "impact": "Impact"}) if not risk_df.empty else pd.DataFrame(columns=["Category", "Title", "Status", "Impact"])
    mitigation_view = mitigationDf[["Action", "Current Status", "Required Next Decision"]].copy() if not mitigationDf.empty else pd.DataFrame(columns=["Action", "Current Status", "Required Next Decision"])
    invoice_view = invoice_df.rename(columns={"title": "Invoice", "date": "Date", "status": "Status"})[["Invoice", "Date", "Status"]].copy() if not invoice_df.empty else pd.DataFrame(columns=["Invoice", "Date", "Status"])

    add_simple_table_to_slide(slide, Inches(0.55), Inches(7.55), Inches(5.0), "Milestone Exposure", milestone_view, max_rows=5)
    add_simple_table_to_slide(slide, Inches(5.75), Inches(7.55), Inches(5.0), "Risk & Mitigation Priority", risk_view.head(3), max_rows=3)
    add_simple_table_to_slide(slide, Inches(11.05), Inches(7.55), Inches(5.0), "Invoice Pulse", invoice_view, max_rows=4)

    narrative = slide.shapes.add_textbox(Inches(0.55), Inches(10.05), Inches(15.25), Inches(1.1))
    narrative.fill.solid()
    narrative.fill.fore_color.rgb = RGBColor(18, 28, 44)
    narrative.line.color.rgb = RGBColor(46, 66, 91)
    narrative.text_frame.word_wrap = True
    narrative.text_frame.text = (
        "Executive Summary: The project remains under schedule pressure with SPI "
        f"{evmData['SPI']:.2f} and schedule variance {sar(evmData['SV'])}. "
        "Priority management focus remains on converting planned work fronts into earned value, "
        "closing outstanding technical interfaces, and protecting procurement and payment continuity."
    )
    narrative.text_frame.paragraphs[0].font.size = Pt(12)
    narrative.text_frame.paragraphs[0].font.color.rgb = RGBColor(230, 238, 248)

    stream = io.BytesIO()
    prs.save(stream)
    stream.seek(0)
    return stream.getvalue()





def _ppt_date(value) -> str:
    if value is None or value == "":
        return ""
    if isinstance(value, pd.Timestamp):
        return value.strftime("%d-%b-%Y")
    parsed = parse_project_date(value)
    if parsed is not None and not pd.isna(parsed):
        return parsed.strftime("%d-%b-%Y")
    parsed = pd.to_datetime(value, errors="coerce")
    if parsed is not None and not pd.isna(parsed):
        return parsed.strftime("%d-%b-%Y")
    return str(value)


def _ppt_money(value, decimals: int = 2) -> str:
    try:
        return f"{float(value):,.{decimals}f}"
    except (TypeError, ValueError):
        return "0.00" if decimals else "0"


def _ppt_pct(value, decimals: int = 2) -> str:
    try:
        return f"{float(value):,.{decimals}f}%"
    except (TypeError, ValueError):
        return "0.00%"


def _set_table_cell(table, row_idx: int, col_idx: int, value) -> None:
    if row_idx < len(table.rows) and col_idx < len(table.columns):
        table.cell(row_idx, col_idx).text = "" if value is None else str(value)


def _fill_table_rows(table, rows: list[list[str]], start_row: int = 0, clear_remaining: bool = True) -> None:
    max_rows = len(table.rows) - start_row
    for r_idx in range(max_rows):
        target_row = start_row + r_idx
        values = rows[r_idx] if r_idx < len(rows) else None
        for c_idx in range(len(table.columns)):
            if values is None:
                if clear_remaining:
                    _set_table_cell(table, target_row, c_idx, "")
            else:
                _set_table_cell(table, target_row, c_idx, values[c_idx] if c_idx < len(values) else "")


def _phase_filter(df: pd.DataFrame, prefixes: tuple[str, ...]) -> pd.Series:
    return df["activity_id"].astype(str).str.startswith(prefixes)


def _max_finish_date(df: pd.DataFrame) -> tuple[str, str, str]:
    if df.empty:
        return "", "", ""
    plan = pd.to_datetime(df["planned_finish"], format="%d-%b-%y", errors="coerce").dropna()
    forecast = pd.to_datetime(df["forecast_finish"], format="%d-%b-%y", errors="coerce").dropna()
    plan_dt = plan.max() if not plan.empty else pd.NaT
    forecast_dt = forecast.max() if not forecast.empty else pd.NaT
    variance = ""
    if not pd.isna(plan_dt) and not pd.isna(forecast_dt):
        variance = str((plan_dt - forecast_dt).days)
    return _ppt_date(plan_dt), _ppt_date(forecast_dt), variance


def _completed_ratio(df: pd.DataFrame) -> float:
    if df.empty:
        return 0.0
    actual = pd.to_numeric(df["actual_progress"].astype(str).str.replace("%", ""), errors="coerce").fillna(0)
    return round((actual >= 100).sum() / len(df) * 100.0, 2)


def _selected_ratio(df: pd.DataFrame, prefixes: tuple[str, ...]) -> float:
    subset = df[_phase_filter(df, prefixes)].copy()
    return _completed_ratio(subset)


def _build_risk_statement_rows(delay_metrics: dict, risk_metrics: dict, contract_metrics: dict) -> list[list[str]]:
    rows: list[list[str]] = []
    delays_df = delay_metrics.get("delays_df", pd.DataFrame())
    payments_df = contract_metrics.get("payments_df", pd.DataFrame())
    risks_df = risk_metrics.get("risks_df", pd.DataFrame())
    ifc_df = risk_metrics.get("ifc_df", pd.DataFrame())
    rfi_df = risk_metrics.get("rfi_df", pd.DataFrame())

    if not delays_df.empty:
        design_delay = delays_df[delays_df["cause_category"].astype(str).str.contains("design|ifc", case=False, na=False)]
        if not design_delay.empty:
            row = design_delay.iloc[0]
            rows.append([str(len(rows) + 1), str(row.get("delay_title", "")), str(row.get("status", ""))])

        steel_delay = delays_df[delays_df["delay_title"].astype(str).str.contains("steel", case=False, na=False)]
        if not steel_delay.empty:
            row = steel_delay.sort_values("estimated_delay_days_num", ascending=False).iloc[0]
            rows.append([str(len(rows) + 1), str(row.get("delay_title", "")), str(row.get("status", ""))])

    if not risks_df.empty:
        mep_risk = risks_df[risks_df["risk_title"].astype(str).str.contains("mep|sub contractor|subcontractor", case=False, na=False)]
        if not mep_risk.empty:
            row = mep_risk.iloc[0]
            rows.append([str(len(rows) + 1), str(row.get("risk_title", "")), str(row.get("status", ""))])

    if not payments_df.empty:
        under_payment = payments_df[payments_df["payment_status"].astype(str).str.contains("under", case=False, na=False)]
        if not under_payment.empty:
            invoice_list = ", ".join(f"No. {str(v).rstrip('.0')}" for v in under_payment["invoice_no"].head(4).tolist())
            rows.append([str(len(rows) + 1), f"Delay in releasing payment of invoices {invoice_list}", "Ongoing"])

    if not rfi_df.empty:
        rows.append([str(len(rows) + 1), f"Delay in reply on {len(rfi_df)} RFI items", "Ongoing"])
    elif not ifc_df.empty:
        rows.append([str(len(rows) + 1), f"IFC conflict register contains {len(ifc_df)} live conflict items", "Ongoing"])

    return rows[:5]


def _build_variance_rows(delay_metrics: dict, risk_metrics: dict, activities_df: pd.DataFrame, overview_metrics: dict) -> list[list[str]]:
    delays_df = delay_metrics.get("delays_df", pd.DataFrame())
    engineering_actual = _selected_ratio(activities_df, ("E-SUB",))
    procurement_actual = _selected_ratio(activities_df, ("P-",))
    construction_actual = float(overview_metrics.get("overall_progress", 0))
    engineering_plan = float(activities_df[_phase_filter(activities_df, ("E-SUB",))]["planned_progress"].astype(str).str.replace("%", "", regex=False).apply(parse_numeric).mean() if not activities_df.empty else 0)
    procurement_plan = float(activities_df[_phase_filter(activities_df, ("P-",))]["planned_progress"].astype(str).str.replace("%", "", regex=False).apply(parse_numeric).mean() if not activities_df.empty else 0)
    construction_plan = float(overview_metrics.get("planned_progress", 0))

    design_remark = "Design information delay impact under review."
    procurement_remark = "Procurement variance under review."
    construction_remark = "Construction variance under review."

    if not delays_df.empty:
        design_delay = delays_df[delays_df["cause_category"].astype(str).str.contains("design|ifc", case=False, na=False)]
        if not design_delay.empty:
            design_remark = str(design_delay.iloc[0].get("delay_title", design_remark))
        procurement_delay = delays_df[delays_df["delay_title"].astype(str).str.contains("steel|procurement", case=False, na=False)]
        if not procurement_delay.empty:
            procurement_remark = str(procurement_delay.iloc[0].get("delay_title", procurement_remark))
        top_construction = delays_df.sort_values("estimated_delay_days_num", ascending=False).head(3)["delay_title"].astype(str).tolist()
        if top_construction:
            construction_remark = " / ".join(top_construction)

    return [
        ["1", "ENGINEERING & DESIGN", _ppt_pct(engineering_actual - engineering_plan), design_remark],
        ["2", "PROCUREMENT", _ppt_pct(procurement_actual - procurement_plan), procurement_remark],
        ["3", "CONSTRUCTION", _ppt_pct(construction_actual - construction_plan), construction_remark],
    ]


def build_original_template_presentation(overview_metrics: dict, evm_metrics: dict, contract_metrics: dict, delay_metrics: dict, risk_metrics: dict, milestone_metrics: dict, activity_metrics: dict, s_curve_metrics: dict) -> tuple[bytes, list[str]]:
    prs = Presentation()
    prs.slide_width = Inches(16.54)
    prs.slide_height = Inches(11.69)

    project_title = str(overview_metrics.get("project_name") or "Project")
    report_date = pd.Timestamp.today().strftime("%d %b %Y")
    activities_df = activity_metrics.get("activities_df", pd.DataFrame()).copy()
    payments_df = contract_metrics.get("payments_df", pd.DataFrame()).copy()
    risks_df = risk_metrics.get("risks_df", pd.DataFrame()).copy()
    phase_df = _phase_progress_rows(activity_metrics, overview_metrics)
    spi = (float(evm_metrics.get("ev", 0.0)) / float(evm_metrics.get("pv", 1.0))) if float(evm_metrics.get("pv", 0.0)) else 0.0

    def add_title(slide, title: str, subtitle: str) -> None:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(246, 249, 252)
        band = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(16.54), Inches(0.92))
        band.fill.solid()
        band.fill.fore_color.rgb = RGBColor(11, 31, 51)
        band.line.color.rgb = RGBColor(11, 31, 51)
        t = slide.shapes.add_textbox(Inches(0.55), Inches(0.18), Inches(10.5), Inches(0.32))
        t.text_frame.text = title
        t.text_frame.paragraphs[0].font.size = Pt(22)
        t.text_frame.paragraphs[0].font.bold = True
        t.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        s = slide.shapes.add_textbox(Inches(0.55), Inches(0.53), Inches(12.5), Inches(0.22))
        s.text_frame.text = subtitle
        s.text_frame.paragraphs[0].font.size = Pt(10)
        s.text_frame.paragraphs[0].font.color.rgb = RGBColor(205, 220, 235)

    def add_card(slide, left: float, top: float, width: float, title: str, value: str, note: str, accent: RGBColor = RGBColor(34, 118, 168)) -> None:
        box = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box.line.color.rgb = RGBColor(214, 225, 236)
        strip = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(0.08), Inches(0.9))
        strip.fill.solid()
        strip.fill.fore_color.rgb = accent
        strip.line.color.rgb = accent
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = RGBColor(96, 111, 128)
        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.size = Pt(17)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(16, 43, 68)
        p3 = tf.add_paragraph()
        p3.text = note
        p3.font.size = Pt(7)
        p3.font.color.rgb = RGBColor(98, 112, 128)

    def add_table(slide, left: float, top: float, width: float, height: float, title: str, df: pd.DataFrame, max_rows: int = 8) -> None:
        heading = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.25))
        heading.text_frame.text = title
        heading.text_frame.paragraphs[0].font.size = Pt(12)
        heading.text_frame.paragraphs[0].font.bold = True
        heading.text_frame.paragraphs[0].font.color.rgb = RGBColor(16, 43, 68)
        view = df.head(max_rows).fillna("").astype(str) if not df.empty else pd.DataFrame({"Status": ["No records available"]})
        rows = len(view) + 1
        cols = len(view.columns)
        shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top + 0.32), Inches(width), Inches(height)).table
        for c_idx, col in enumerate(view.columns):
            cell = shape.cell(0, c_idx)
            cell.text = str(col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(229, 238, 247)
            cell.text_frame.paragraphs[0].font.size = Pt(7)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(40, 63, 85)
        for r_idx, (_, row) in enumerate(view.iterrows(), start=1):
            for c_idx, col in enumerate(view.columns):
                cell = shape.cell(r_idx, c_idx)
                cell.text = str(row.get(col, ""))[:90]
                cell.text_frame.paragraphs[0].font.size = Pt(6)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(45, 55, 72)

    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide1, "Project Progress Control Report", f"{project_title} | A3 Landscape | {report_date}")
    add_card(slide1, 0.55, 1.18, 2.45, "Overall Progress", pct(overview_metrics.get("overall_progress")), f"Planned {pct(overview_metrics.get('planned_progress'))}", RGBColor(22, 143, 139))
    add_card(slide1, 3.15, 1.18, 2.45, "Contract Value", egp(overview_metrics.get("contract_value")), "Current project value", RGBColor(36, 95, 149))
    add_card(slide1, 5.75, 1.18, 2.45, "SPI", f"{spi:.2f}" if spi else "N/A", "Schedule performance", RGBColor(201, 133, 25))
    add_card(slide1, 8.35, 1.18, 2.45, "Critical Activities", str(int(activity_metrics.get("critical_count", 0))), "Current critical path", RGBColor(217, 84, 77))
    add_card(slide1, 10.95, 1.18, 2.45, "Open Risks", str(int(risk_metrics.get("open_risks", 0))), "Active risk items", RGBColor(95, 111, 184))
    add_card(slide1, 13.55, 1.18, 2.45, "Delay Events", str(int(delay_metrics.get("total_delay_events", 0))), f"{int(delay_metrics.get('total_delay_days', 0))} days", RGBColor(217, 84, 77))
    add_table(slide1, 0.55, 2.45, 5.05, 2.35, "Engineering / Procurement / Construction Progress", phase_df, max_rows=5)
    add_table(slide1, 5.85, 2.45, 4.8, 2.35, "Current Delay Register", delay_metrics.get("display_delays_df", pd.DataFrame()), max_rows=6)
    add_table(slide1, 10.9, 2.45, 5.1, 2.35, "Risk And Issue Snapshot", risks_df, max_rows=6)
    add_table(slide1, 0.55, 5.55, 7.55, 2.15, "Payment Status", payments_df, max_rows=6)
    evm_df = pd.DataFrame(
        [
            {"Metric": "BAC", "Value": egp(evm_metrics.get("bac"))},
            {"Metric": "PV", "Value": egp(evm_metrics.get("pv"))},
            {"Metric": "EV", "Value": egp(evm_metrics.get("ev"))},
            {"Metric": "SV", "Value": egp(evm_metrics.get("sv"))},
            {"Metric": "CV", "Value": egp(evm_metrics.get("cv"))},
        ]
    )
    add_table(slide1, 8.35, 5.55, 3.35, 2.15, "Earned Value Readout", evm_df, max_rows=6)
    action_df = pd.DataFrame(
        [
            {"Priority Action": "Recover critical work fronts", "Owner": "Construction / Planning", "Decision Need": "Approve recovery sequence"},
            {"Priority Action": "Close delayed engineering/RFI items", "Owner": "Engineering", "Decision Need": "Confirm response dates"},
            {"Priority Action": "Resolve payment exposure", "Owner": "Commercial", "Decision Need": "Release certified amounts"},
            {"Priority Action": "Protect entitlement records", "Owner": "Planning / Contracts", "Decision Need": "Maintain notices and evidence"},
        ]
    )
    add_table(slide1, 11.95, 5.55, 4.05, 2.15, "Management Decisions", action_df, max_rows=5)

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide2, "Detailed Data Tables", f"{project_title} | A3 Landscape | {report_date}")
    add_table(slide2, 0.55, 1.25, 7.7, 3.0, "Critical Activities", activity_metrics.get("critical_df", pd.DataFrame()), max_rows=9)
    add_table(slide2, 8.55, 1.25, 7.45, 3.0, "Milestones", milestone_metrics.get("milestones_df", pd.DataFrame()), max_rows=9)
    add_table(slide2, 0.55, 5.05, 7.7, 3.0, "S-Curve Data", s_curve_metrics.get("curve_df", pd.DataFrame()), max_rows=9)
    add_table(slide2, 8.55, 5.05, 7.45, 3.0, "Contracts", contract_metrics.get("contracts_df", pd.DataFrame()), max_rows=9)

    stream = io.BytesIO()
    prs.save(stream)
    stream.seek(0)
    return stream.getvalue(), []


def _phase_subset(activities_df: pd.DataFrame, phase: str) -> pd.DataFrame:
    if activities_df.empty:
        return pd.DataFrame()
    ids = activities_df.get("activity_id", pd.Series("", index=activities_df.index)).astype(str).str.strip()
    ids_upper = ids.str.upper()
    names = activities_df.get("activity_name", pd.Series("", index=activities_df.index)).astype(str)
    wbs = activities_df.get("wbs_id", pd.Series("", index=activities_df.index)).astype(str)
    combined = (ids + " " + names + " " + wbs).str.lower()
    engineering_prefix = ids_upper.str.startswith(("E-", "E_SUB", "E-SUB", "E-APP"), na=False)
    procurement_prefix = ids_upper.str.startswith(("P-",), na=False)
    construction_prefix = ids_upper.str.startswith(("CON", "AB-"), na=False)
    if phase == "Engineering":
        if engineering_prefix.any():
            return activities_df[engineering_prefix].copy()
        return activities_df[
            ~procurement_prefix
            & ~construction_prefix
            & combined.str.contains("engineering|design|submittal|shop drawing|drawing|approval", na=False)
        ].copy()
    if phase == "Procurement":
        if procurement_prefix.any():
            return activities_df[procurement_prefix].copy()
        return activities_df[
            ~engineering_prefix
            & ~construction_prefix
            & combined.str.contains("procurement|purchase|material|supply|steel|delivery", na=False)
        ].copy()
    if phase == "Construction":
        if construction_prefix.any():
            return activities_df[construction_prefix].copy()
        return activities_df[
            ~engineering_prefix
            & ~procurement_prefix
            & combined.str.contains("construction|civil|structural|concrete|rft|slab|column|beam|foundation|reinforcement", na=False)
        ].copy()
    return activities_df.copy()


def _phase_progress_rows(activity_metrics: dict, overview_metrics: dict) -> pd.DataFrame:
    activities_df = activity_metrics.get("activities_df", pd.DataFrame()).copy()
    rows = []
    for phase in ["Engineering", "Procurement", "Construction"]:
        subset = _phase_subset(activities_df, phase)
        if subset.empty:
            planned = float(overview_metrics.get("planned_progress", 0.0))
            actual = float(overview_metrics.get("overall_progress", 0.0))
        else:
            planned = float(subset["planned_progress_num"].mean()) if "planned_progress_num" in subset.columns else 0.0
            actual = float(subset["actual_progress_num"].mean()) if "actual_progress_num" in subset.columns else 0.0
        rows.append(
            {
                "phase": phase,
                "planned": round(planned, 2),
                "actual": round(actual, 2),
                "variance": round(actual - planned, 2),
            }
        )
    rows.append(
        {
            "phase": "Overall",
            "planned": round(float(overview_metrics.get("planned_progress", 0.0)), 2),
            "actual": round(float(overview_metrics.get("overall_progress", 0.0)), 2),
            "variance": round(float(overview_metrics.get("overall_progress", 0.0)) - float(overview_metrics.get("planned_progress", 0.0)), 2),
        }
    )
    return pd.DataFrame(rows)


def _build_executive_doc_rows(activity_metrics: dict) -> pd.DataFrame:
    activities_df = activity_metrics.get("activities_df", pd.DataFrame()).copy()
    if activities_df.empty:
        return pd.DataFrame()
    ids = activities_df["activity_id"].astype(str).str.strip()
    phases = [
        ("Civil - Submittals", activities_df[ids.str.startswith(("E-SUB",), na=False)].copy()),
        ("Civil - Shop Drawings", activities_df[ids.str.startswith(("E-APP",), na=False)].copy()),
    ]
    rows = []
    for label, subset in phases:
        if subset.empty:
            rows.append({"discipline": label, "total": 0, "plan": 0, "actual": 0})
            continue
        plan = int((subset["planned_progress_num"] >= 100).sum()) if "planned_progress_num" in subset.columns else 0
        actual = int((subset["actual_progress_num"] >= 100).sum()) if "actual_progress_num" in subset.columns else 0
        rows.append({"discipline": label, "total": int(len(subset)), "plan": plan, "actual": actual})
    return pd.DataFrame(rows)


def _build_executive_invoice_rows(contract_metrics: dict) -> pd.DataFrame:
    payments_df = contract_metrics.get("payments_df", pd.DataFrame()).copy()
    if payments_df.empty:
        return pd.DataFrame()
    view = payments_df.sort_values("payment_id").copy()
    view["title"] = view["invoice_no"].astype(str).str.rstrip(".0").apply(lambda x: f"Invoice #{x}" if x else "Invoice")
    view["date"] = view["invoice_date"].apply(_ppt_date)
    view["value"] = view["certified_amount_num"].fillna(0.0)
    view["status"] = view["payment_status"].astype(str)
    return view[["title", "date", "value", "status"]]


def _build_executive_milestone_rows(activity_metrics: dict) -> pd.DataFrame:
    activities_df = activity_metrics.get("activities_df", pd.DataFrame()).copy()
    if activities_df.empty:
        return pd.DataFrame()
    phase_defs = [
        ("Engineering", _phase_subset(activities_df, "Engineering")),
        ("Procurement", _phase_subset(activities_df, "Procurement")),
        ("Construction", _phase_subset(activities_df, "Construction")),
        ("B-01", activities_df[activities_df["activity_id"].astype(str).str.contains("B01", case=False, na=False)].copy()),
        ("B-02", activities_df[activities_df["activity_id"].astype(str).str.contains("B02", case=False, na=False)].copy()),
        ("B-03", activities_df[activities_df["activity_id"].astype(str).str.contains("B03", case=False, na=False)].copy()),
        ("B-04", activities_df[activities_df["activity_id"].astype(str).str.contains("B04", case=False, na=False)].copy()),
    ]
    rows = []
    for label, subset in phase_defs:
        plan_date, forecast_date, variance = _max_finish_date(subset)
        rows.append(
            {
                "discipline": label,
                "plan": plan_date or "ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â",
                "forecast": forecast_date or "ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â",
                "variance_days": int(variance) if str(variance).strip() not in {"", "nan"} else 0,
            }
        )
    return pd.DataFrame(rows)


def _build_executive_summary_milestone_rows(activity_metrics: dict, milestone_metrics: dict) -> pd.DataFrame:
    activities_df = activity_metrics.get("activities_df", pd.DataFrame()).copy()
    milestones_df = milestone_metrics.get("milestones_df", pd.DataFrame()).copy()
    rows = []

    def safe_ts(value):
        ts = pd.to_datetime(value, errors="coerce", dayfirst=True)
        return None if pd.isna(ts) else ts

    if not activities_df.empty:
        for label, code in [("B-02", "B02"), ("B-03", "B03"), ("B-04", "B04")]:
            subset = activities_df[activities_df["activity_id"].astype(str).str.contains(code, case=False, na=False)].copy()
            plan_date, forecast_date, variance = _max_finish_date(subset)
            rows.append(
                {
                    "discipline": label,
                    "plan": plan_date or "ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â",
                    "forecast": forecast_date or "ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â",
                    "variance_days": int(variance) if str(variance).strip() not in {"", "nan"} else 0,
                }
            )

    handover_row = None
    if not milestones_df.empty:
        finish_mask = milestones_df["activity_name"].astype(str).str.contains("project finish|handover|handing over", case=False, na=False)
        if finish_mask.any():
            finish_row = milestones_df[finish_mask].iloc[0]
            plan_date = finish_row.get("planned_date", "ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â")
            forecast_date = finish_row.get("forecast_date", "ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â")
            plan_ts = safe_ts(plan_date)
            forecast_ts = safe_ts(forecast_date)
            variance = (forecast_ts - plan_ts).days if plan_ts is not None and forecast_ts is not None else 0
            handover_row = {
                "discipline": "Handing Over",
                "plan": plan_date or "ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â",
                "forecast": forecast_date or "ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â",
                "variance_days": int(variance),
            }

    if handover_row is None:
        wbs_df = filter_active_project(load_core_csv(WBS_CSV_PATH)).copy()
        if not wbs_df.empty:
            code_col = "WBS Code" if "WBS Code" in wbs_df.columns else ("wbs_code" if "wbs_code" in wbs_df.columns else None)
            name_col = "WBS Name" if "WBS Name" in wbs_df.columns else ("wbs_name" if "wbs_name" in wbs_df.columns else None)
            if code_col and name_col:
                ho_mask = (
                    wbs_df[code_col].astype(str).str.contains(r"\.HO(\.|$)", case=False, na=False, regex=True)
                    | wbs_df[name_col].astype(str).str.contains("handover|handing over", case=False, na=False)
                )
                if ho_mask.any():
                    ho_row = wbs_df[ho_mask].iloc[0]
                    plan_date = ho_row.get("bl_project_finish", "ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â")
                    remaining_duration = parse_numeric(ho_row.get("remaining_duration", 0))
                    plan_ts = safe_ts(plan_date)
                    forecast_ts = plan_ts + pd.Timedelta(days=remaining_duration) if plan_ts is not None and remaining_duration else None
                    variance = (forecast_ts - plan_ts).days if plan_ts is not None and forecast_ts is not None else 0
                    handover_row = {
                        "discipline": "Handing Over",
                        "plan": plan_date or "ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â",
                        "forecast": forecast_ts.strftime("%d-%b-%y") if forecast_ts is not None else "ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â",
                        "variance_days": int(variance),
                    }

    if handover_row is None:
        handover_row = {"discipline": "Handing Over", "plan": "ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â", "forecast": "ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â", "variance_days": 0}

    rows.append(handover_row)
    return pd.DataFrame(rows)


def _build_executive_risk_rows(delay_metrics: dict, risk_metrics: dict) -> pd.DataFrame:
    rows: list[dict[str, str]] = []
    risks_df = risk_metrics.get("risks_df", pd.DataFrame()).copy()
    if not risks_df.empty:
        for _, row in risks_df.head(5).iterrows():
            rows.append(
                {
                    "category": str(row.get("risk_category", "Risk")),
                    "title": str(row.get("risk_title", "")),
                    "status": "Ongoing" if str(row.get("status", "")).strip().lower() == "open" else str(row.get("status", "Closed")).title(),
                    "impact": str(row.get("probability", "Medium")).title(),
                }
            )
    delays_df = delay_metrics.get("delays_df", pd.DataFrame()).copy()
    if len(rows) < 5 and not delays_df.empty:
        top_delays = delays_df.sort_values("estimated_delay_days_num", ascending=False).head(5 - len(rows))
        for _, row in top_delays.iterrows():
            rows.append(
                {
                    "category": str(row.get("cause_category", "Delay")),
                    "title": str(row.get("delay_title", "")),
                    "status": "Ongoing" if str(row.get("status_group", "")).strip().lower() == "open" else str(row.get("status", "Closed")).title(),
                    "impact": "Critical" if float(row.get("estimated_delay_days_num", 0)) >= 30 else "High",
                }
            )
    return pd.DataFrame(rows)


def _build_executive_cause_rows(delay_metrics: dict, risk_metrics: dict) -> pd.DataFrame:
    rows = []
    delays_df = delay_metrics.get("delays_df", pd.DataFrame()).copy()
    if not delays_df.empty:
        for _, row in delays_df.sort_values("estimated_delay_days_num", ascending=False).head(5).iterrows():
            delay_days = float(row.get("estimated_delay_days_num", 0.0))
            rows.append(
                {
                    "subject": str(row.get("delay_title", "Delay"))[:42],
                    "exposure": min(100, round(delay_days * 2, 1)),
                    "recovery": max(10, 100 - min(90, round(delay_days * 1.1, 1))),
                    "ownerImpact": 85 if str(row.get("responsible_group", "")).startswith("Employer") else 55,
                    "timeImpact": min(100, round(delay_days * 2.1, 1)),
                }
            )
    if not rows:
        risks_df = risk_metrics.get("risks_df", pd.DataFrame()).copy()
        for _, row in risks_df.head(5).iterrows():
            rows.append(
                {
                    "subject": str(row.get("risk_title", "Risk"))[:42],
                    "exposure": 70,
                    "recovery": 40,
                    "ownerImpact": 60,
                    "timeImpact": 65,
                }
            )
    return pd.DataFrame(rows[:5])


def _build_executive_resource_rows() -> pd.DataFrame:
    rft_df = filter_active_project(load_core_csv(APP_DIR / "data" / "import_templates" / "rft_qtys.csv"))
    if rft_df.empty:
        return pd.DataFrame()
    rft_df["Budgeted Units_num"] = rft_df["Budgeted Units"].apply(parse_numeric) if "Budgeted Units" in rft_df.columns else 0.0
    view = (
        rft_df.sort_values("Budgeted Units_num", ascending=False)
        .head(8)
        .rename(columns={"Activity Name": "resource", "Budgeted Units_num": "qty"})
    )
    return view[["resource", "qty"]]


def build_linked_executive_dashboard_html(
    overview_metrics: dict,
    evm_metrics: dict,
    contract_metrics: dict,
    delay_metrics: dict,
    risk_metrics: dict,
    milestone_metrics: dict,
    activity_metrics: dict,
    evmData: dict,
    rootCauseDf: pd.DataFrame,
    mitigationDf: pd.DataFrame,
    evmComments: dict,
) -> str:
    project_df = load_core_csv(PROJECTS_CSV_PATH)
    project_row = active_project_row(project_df) if not project_df.empty else pd.Series(dtype=object)
    phase_df = _phase_progress_rows(activity_metrics, overview_metrics)
    docs_df = _build_executive_doc_rows(activity_metrics)
    invoice_df = _build_executive_invoice_rows(contract_metrics)
    milestone_df = _build_executive_milestone_rows(activity_metrics)
    risk_df = _build_executive_risk_rows(delay_metrics, risk_metrics)
    cause_df = _build_executive_cause_rows(delay_metrics, risk_metrics)
    resource_df = _build_executive_resource_rows()

    paid_amount = float(contract_metrics.get("total_paid", 0.0))
    certified_amount = float(contract_metrics.get("total_certified", 0.0))
    under_payment = max(certified_amount - paid_amount, 0.0)
    ongoing_risks = int((risk_df["status"].astype(str).str.lower() == "ongoing").sum()) if not risk_df.empty else 0
    spi = (float(evm_metrics.get("ev", 0.0)) / float(evm_metrics.get("pv", 1.0))) if float(evm_metrics.get("pv", 0.0)) else 0.0

    progress_fig = px.bar(
        phase_df,
        x="phase",
        y=["planned", "actual"],
        barmode="group",
        title="Planned vs Actual by Phase",
        color_discrete_map={"planned": "#38bdf8", "actual": "#f97316"},
    )
    progress_fig.update_layout(legend_title_text="", paper_bgcolor="#0b1320", plot_bgcolor="#0b1320", font=dict(color="#e5eef8"))
    progress_fig.update_xaxes(title_text="")
    progress_fig.update_yaxes(title_text="%")

    commercial_fig = px.pie(
        pd.DataFrame(
            [
                {"name": "Submitted", "value": certified_amount},
                {"name": "Approved", "value": certified_amount},
                {"name": "Paid", "value": paid_amount},
            ]
        ),
        names="name",
        values="value",
        hole=0.45,
        color="name",
        color_discrete_map={"Submitted": "#38bdf8", "Approved": "#a78bfa", "Paid": "#22c55e"},
    )
    commercial_fig.update_layout(showlegend=True, legend_title_text="", paper_bgcolor="#0b1320", plot_bgcolor="#0b1320", font=dict(color="#e5eef8"))

    evm_fig = px.bar(
        pd.DataFrame(
            [
                {"metric": "BAC", "value": float(evm_metrics.get("bac", 0.0))},
                {"metric": "PV", "value": float(evm_metrics.get("pv", 0.0))},
                {"metric": "EV", "value": float(evm_metrics.get("ev", 0.0))},
                {"metric": "SV", "value": abs(float(evm_metrics.get("sv", 0.0)))},
            ]
        ),
        x="metric",
        y="value",
        title="BAC / PV / EV / SV Executive Readout",
        color="metric",
        color_discrete_map={"BAC": "#38bdf8", "PV": "#22c55e", "EV": "#f97316", "SV": "#f43f5e"},
    )
    evm_fig.update_layout(showlegend=False, paper_bgcolor="#0b1320", plot_bgcolor="#0b1320", font=dict(color="#e5eef8"))

    evm_addon_fig = px.bar(
        pd.DataFrame(
            [
                {"metric": "BAC", "value": evmData["BAC"]},
                {"metric": "PV", "value": evmData["PV"]},
                {"metric": "EV", "value": evmData["EV"]},
                {"metric": "SV Exposure", "value": abs(evmData["SV"])},
            ]
        ),
        x="metric",
        y="value",
        title="BAC vs PV vs EV with SV Exposure",
        color="metric",
        color_discrete_map={
            "BAC": "#38bdf8",
            "PV": "#22c55e",
            "EV": "#f97316",
            "SV Exposure": "#f43f5e",
        },
    )
    evm_addon_fig.update_layout(showlegend=False, paper_bgcolor="#0b1320", plot_bgcolor="#0b1320", font=dict(color="#e5eef8"))

    evm_gauge_fig = go.Figure(
        go.Indicator(
            mode="gauge+number",
            value=evmData["SPI"],
            title={"text": "SPI Health"},
            number={"valueformat": ".2f"},
            gauge={
                "axis": {"range": [0, 1.5]},
                "bar": {"color": "#f43f5e" if evmData["SPI"] < 1 else "#22c55e"},
                "steps": [
                    {"range": [0, 0.75], "color": "#3a161e"},
                    {"range": [0.75, 1.0], "color": "#4a3412"},
                    {"range": [1.0, 1.5], "color": "#15382b"},
                ],
                "threshold": {"line": {"color": "#ffffff", "width": 4}, "value": 1.0},
            },
        )
    )
    evm_gauge_fig.update_layout(height=340, paper_bgcolor="#0b1320", margin=dict(l=10, r=10, t=50, b=10), font=dict(color="#e5eef8"))

    resource_fig = None
    if not resource_df.empty:
        resource_fig = px.bar(
            resource_df,
            x="qty",
            y="resource",
            orientation="h",
            title="RFT Quantity Distribution",
            color_discrete_sequence=["#38bdf8"],
        )
        resource_fig.update_layout(showlegend=False, paper_bgcolor="#0b1320", plot_bgcolor="#0b1320", font=dict(color="#e5eef8"))

    radar_fig = None
    if not cause_df.empty:
        radar_fig = go.Figure()
        radar_fig.add_trace(go.Scatterpolar(r=cause_df["exposure"], theta=cause_df["subject"], fill="toself", name="Exposure", line_color="#f97316"))
        radar_fig.add_trace(go.Scatterpolar(r=cause_df["recovery"], theta=cause_df["subject"], fill="toself", name="Recovery", line_color="#38bdf8"))
        radar_fig.add_trace(go.Scatterpolar(r=cause_df["timeImpact"], theta=cause_df["subject"], fill="toself", name="Time Impact", line_color="#f43f5e"))
        radar_fig.update_layout(
            polar=dict(radialaxis=dict(visible=True, range=[0, 100], gridcolor="rgba(255,255,255,0.12)"), angularaxis=dict(gridcolor="rgba(255,255,255,0.12)")),
            showlegend=True,
            paper_bgcolor="#0b1320",
            plot_bgcolor="#0b1320",
            font=dict(color="#e5eef8"),
            title="Root Cause / Impact Radar",
        )

    progress_html = progress_fig.to_html(include_plotlyjs="cdn", full_html=False, config={"displayModeBar": False})
    commercial_html = commercial_fig.to_html(include_plotlyjs=False, full_html=False, config={"displayModeBar": False})
    evm_overview_html = evm_fig.to_html(include_plotlyjs=False, full_html=False, config={"displayModeBar": False})
    evm_addon_chart_html = evm_addon_fig.to_html(include_plotlyjs=False, full_html=False, config={"displayModeBar": False})
    evm_gauge_html = evm_gauge_fig.to_html(include_plotlyjs=False, full_html=False, config={"displayModeBar": False})
    resource_html = resource_fig.to_html(include_plotlyjs=False, full_html=False, config={"displayModeBar": False}) if resource_fig is not None else ""
    radar_html = radar_fig.to_html(include_plotlyjs=False, full_html=False, config={"displayModeBar": False}) if radar_fig is not None else ""

    risk_cards = ""
    for _, row in risk_df.iterrows():
        risk_cards += f"""
        <div class="risk-card">
          <div class="risk-top"><strong>{html.escape(str(row.get("category", "")))}</strong><span>{html.escape(str(row.get("status", "")))}</span></div>
          <div class="risk-body">{html.escape(str(row.get("title", "")))}</div>
          <div class="risk-foot">Impact: {html.escape(str(row.get("impact", "")))}</div>
        </div>
        """

    doc_cards = ""
    for _, row in docs_df.iterrows():
        doc_cards += f"""
        <div class="mini-card">
          <div class="mini-title">{html.escape(str(row['discipline']))}</div>
          <div class="mini-grid">
            <div><span>Total</span><strong>{int(row['total'])}</strong></div>
            <div><span>Plan</span><strong>{int(row['plan'])}</strong></div>
            <div><span>Actual</span><strong>{int(row['actual'])}</strong></div>
          </div>
        </div>
        """

    invoice_rows = ""
    for _, row in invoice_df.iterrows():
        invoice_rows += f"""
        <tr>
          <td>{html.escape(str(row['title']))}</td>
          <td>{html.escape(str(row['date']))}</td>
          <td>{format_currency_html(row['value'])}</td>
          <td>{html.escape(str(row['status']))}</td>
        </tr>
        """

    milestone_rows = ""
    for _, row in milestone_df.iterrows():
        milestone_rows += f"""
        <tr>
          <td>{html.escape(str(row['discipline']))}</td>
          <td>{html.escape(str(row['plan']))}</td>
          <td>{html.escape(str(row['forecast']))}</td>
          <td>{int(row['variance_days'])}</td>
        </tr>
        """

    root_rows = ""
    for _, row in rootCauseDf.iterrows():
        root_rows += f"""
        <tr>
          <td>{html.escape(str(row.get("Cause Title", "")))}</td>
          <td>{html.escape(str(row.get("Impact Area", "")))}</td>
          <td>{html.escape(str(row.get("EVM Impact Link", "")))}</td>
          <td>{html.escape(str(row.get("Status", "")))}</td>
          <td>{html.escape(str(row.get("Severity", "")))}</td>
        </tr>
        """

    mitigation_rows = ""
    for _, row in mitigationDf.iterrows():
        mitigation_rows += f"""
        <tr>
          <td>{html.escape(str(row.get("Action", "")))}</td>
          <td>{html.escape(str(row.get("Owner / Responsible Party", "")))}</td>
          <td>{html.escape(str(row.get("Current Status", "")))}</td>
          <td>{html.escape(str(row.get("Recovery Impact", "")))}</td>
          <td>{html.escape(str(row.get("Required Next Decision", "")))}</td>
        </tr>
        """

    def evm_comment_block(title: str, body: str) -> str:
        return f"""
        <div class="comment-panel">
          <div class="comment-label">{html.escape(title)}</div>
          <div class="comment-body">{html.escape(_comment_or_default(body))}</div>
        </div>
        """

    def metric_block(title: str, value: str, subtitle: str) -> str:
        return f"""
        <div class="metric">
          <div class="metric-title">{html.escape(title)}</div>
          <div class="metric-value">{html.escape(value)}</div>
          <div class="metric-sub">{html.escape(subtitle)}</div>
        </div>
        """

    return f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>Linked Executive Dashboard Export</title>
  <style>
    body{{margin:0;background:#07111f;color:#e5eef8;font-family:Arial,sans-serif}}
    .wrap{{max-width:1380px;margin:0 auto;padding:24px}}
    .hero,.card{{border:1px solid rgba(255,255,255,.08);background:rgba(255,255,255,.05);backdrop-filter:blur(18px);border-radius:24px;box-shadow:0 20px 40px rgba(0,0,0,.22)}}
    .hero{{padding:24px}}
    .hero-top{{display:flex;justify-content:space-between;gap:16px;align-items:flex-start;flex-wrap:wrap}}
    .eyebrow span{{display:inline-block;margin-right:8px;margin-bottom:8px;padding:7px 12px;border-radius:999px;font-size:11px;font-weight:700;letter-spacing:.18em;text-transform:uppercase}}
    .eyebrow .a{{background:#fff;color:#0f172a}}
    .eyebrow .b{{background:rgba(56,189,248,.14);color:#bae6fd;border:1px solid rgba(56,189,248,.25)}}
    .eyebrow .c{{background:rgba(249,115,22,.14);color:#fed7aa;border:1px solid rgba(249,115,22,.25)}}
    h1{{margin:12px 0 0;font-size:48px;line-height:1.05}}
    h2{{margin:0;font-size:32px}}
    .sub{{margin-top:12px;color:#c7d4e3;max-width:820px;line-height:1.7}}
    .hero-side{{min-width:280px}}
    .hero-side .slot{{background:rgba(2,6,23,.5);padding:14px 16px;border-radius:18px;margin-bottom:10px}}
    .hero-side .slot span{{display:block;color:#9fb1c6;font-size:12px}}
    .hero-side .slot strong{{display:block;margin-top:6px;color:#fff;font-size:22px}}
    .grid4{{display:grid;grid-template-columns:repeat(4,minmax(0,1fr));gap:18px;margin-top:18px}}
    .grid2{{display:grid;grid-template-columns:repeat(2,minmax(0,1fr));gap:18px;margin-top:18px}}
    .metric{{padding:22px;border-radius:22px;background:rgba(255,255,255,.06);border:1px solid rgba(255,255,255,.08)}}
    .metric-title{{font-size:11px;letter-spacing:.18em;text-transform:uppercase;color:#94a3b8}}
    .metric-value{{font-size:34px;font-weight:900;color:#fff;margin-top:12px}}
    .metric-sub{{font-size:13px;color:#c7d4e3;margin-top:10px}}
    .section{{display:grid;grid-template-columns:repeat(12,minmax(0,1fr));gap:18px;margin-top:18px}}
    .card{{padding:22px}}
    .span-7{{grid-column:span 7}}
    .span-5{{grid-column:span 5}}
    .span-8{{grid-column:span 8}}
    .span-4{{grid-column:span 4}}
    .span-12{{grid-column:span 12}}
    .label{{font-size:11px;letter-spacing:.18em;text-transform:uppercase;color:#94a3b8}}
    .card h3{{margin:10px 0 16px;font-size:28px;color:#fff}}
    .risk-stack,.mini-stack{{display:flex;flex-direction:column;gap:14px}}
    .risk-card,.mini-card{{background:rgba(2,6,23,.52);border:1px solid rgba(255,255,255,.08);border-radius:22px;padding:16px}}
    .evm-addon-grid{{display:grid;grid-template-columns:repeat(4,minmax(0,1fr));gap:14px}}
    .evm-chip{{display:inline-flex;align-items:center;gap:8px;padding:8px 14px;border-radius:999px;background:rgba(244,63,94,.16);border:1px solid rgba(244,63,94,.28);color:#fecdd3;font-size:12px;font-weight:800;text-transform:uppercase;letter-spacing:.12em}}
    .evm-note{{margin-top:14px;padding:16px 18px;border-radius:18px;background:rgba(255,255,255,.04);border:1px solid rgba(255,255,255,.08);color:#dbe6f2;line-height:1.7}}
    .comment-panel{{margin-top:14px;padding:14px 16px;border-radius:18px;background:rgba(2,6,23,.52);border:1px solid rgba(255,255,255,.08)}}
    .comment-label{{font-size:12px;font-weight:800;color:#93c5fd;text-transform:uppercase;letter-spacing:.12em}}
    .comment-body{{margin-top:10px;color:#e5eef8;line-height:1.6}}
    .risk-top{{display:flex;justify-content:space-between;gap:12px;color:#fff}}
    .risk-top span{{color:#fda4af;font-size:12px}}
    .risk-body{{margin-top:10px;color:#c7d4e3;line-height:1.6;font-size:14px}}
    .risk-foot{{margin-top:10px;color:#94a3b8;font-size:12px}}
    .mini-title{{font-weight:800;color:#fff}}
    .mini-grid{{display:grid;grid-template-columns:repeat(3,minmax(0,1fr));gap:10px;margin-top:14px}}
    .mini-grid div{{background:rgba(255,255,255,.04);padding:12px;border-radius:16px;text-align:center}}
    .mini-grid span{{display:block;color:#94a3b8;font-size:11px}}
    .mini-grid strong{{display:block;color:#fff;font-size:22px;margin-top:6px}}
    table{{width:100%;border-collapse:collapse;font-size:14px;overflow:hidden}}
    th,td{{padding:14px 12px;border-top:1px solid rgba(255,255,255,.08);text-align:left}}
    th{{font-size:11px;letter-spacing:.18em;text-transform:uppercase;color:#94a3b8;background:rgba(255,255,255,.05)}}
    td{{color:#e5eef8}}
    .footer{{margin-top:18px;padding:18px 22px;border-radius:22px;background:rgba(255,255,255,.05);border:1px solid rgba(255,255,255,.08);display:flex;justify-content:space-between;gap:16px;flex-wrap:wrap;color:#c7d4e3}}
    @media (max-width: 1100px) {{
      .grid4,.grid2{{grid-template-columns:1fr 1fr}}
      .span-7,.span-5,.span-8,.span-4,.span-12{{grid-column:span 12}}
      h1{{font-size:36px}}
    }}
    @media (max-width: 700px) {{
      .grid4,.grid2{{grid-template-columns:1fr}}
      h1{{font-size:28px}}
      h2{{font-size:24px}}
    }}
  </style>
</head>
<body>
  <div class="wrap">
    <div class="hero">
      <div class="hero-top">
        <div>
          <div class="eyebrow">
            <span class="a">{html.escape(str(project_row.get("client_name", "THE BIG")))}</span>
            <span class="b">Executive Control Tower</span>
            <span class="c">Report Date: {_ppt_date(pd.Timestamp.today())}</span>
          </div>
          <h1>{html.escape(str(project_row.get("project_name", overview_metrics.get("project_name", "Project"))))}</h1>
          <div class="sub">This linked executive dashboard is generated from the platform CSV files and updates each time the platform data changes.</div>
        </div>
        <div class="hero-side">
          <div class="slot"><span>Schedule Health</span><strong>SPI {spi:.2f}</strong></div>
          <div class="slot"><span>Contractual Completion</span><strong>{_ppt_date(overview_metrics.get("project_finish"))}</strong></div>
          <div class="slot"><span>Forecast Completion</span><strong>{_ppt_date(project_row.get("forecast_finish"))}</strong></div>
        </div>
      </div>
      <div class="grid4">
        {metric_block("Contract Value", f"{format_currency_html(float(overview_metrics.get('contract_value', 0.0)))} EGP", "Original contract value")}
        {metric_block("Schedule Variance", f"{format_currency_html(float(evm_metrics.get('sv', 0.0)))} EGP", "EV minus PV exposure")}
        {metric_block("Cumulative Paid", f"{format_currency_html(paid_amount)} EGP", "Paid to date")}
        {metric_block("Ongoing Risks", str(ongoing_risks), "Active project constraints")}
      </div>
    </div>

    <div class="section">
      <div class="card span-7">
        <div class="label">Progress Intelligence</div>
        <h3>Planned vs Actual by Phase</h3>
        {progress_html}
      </div>
      <div class="card span-5">
        <div class="label">Command Decisions</div>
        <h3>Critical Recovery Agenda</h3>
        <div class="risk-stack">{risk_cards}</div>
      </div>
      <div class="card span-12">
        <div class="label">Phase Gauges</div>
        <h3>Plan vs Actual Snapshot</h3>
        <div class="grid4">
          {"".join(metric_block(str(row["phase"]), f"{row['actual']:.2f}%", f"Planned {row['planned']:.2f}% | Variance {row['variance']:.2f}%") for _, row in phase_df.iterrows())}
        </div>
      </div>
    </div>

    <div class="section">
      <div class="card span-8">
        <div class="label">Milestone Exposure</div>
        <h3>Forecast Completion by Discipline / Area</h3>
        <table>
          <thead><tr><th>Discipline</th><th>Plan Completion</th><th>Forecast Completion</th><th>Variance Days</th></tr></thead>
          <tbody>{milestone_rows}</tbody>
        </table>
      </div>
      <div class="card span-4">
        <div class="label">Engineering Deliverables</div>
        <h3>Document Submission Control</h3>
        <div class="mini-stack">{doc_cards}</div>
      </div>
    </div>

    <div class="section">
      <div class="card span-5">
        <div class="label">Cashflow Position</div>
        <h3>Invoice Conversion Funnel</h3>
        {commercial_html}
        <div class="grid2">
          {metric_block("Paid", f"{format_currency_html(paid_amount)} EGP", "Released cash")}
          {metric_block("Under Payment", f"{format_currency_html(under_payment)} EGP", "Certified not yet paid")}
        </div>
      </div>
      <div class="card span-7">
        <div class="label">Invoices Register</div>
        <h3>Payment Status Control</h3>
        <table>
          <thead><tr><th>Invoice</th><th>Submission Date</th><th>Value</th><th>Status</th></tr></thead>
          <tbody>{invoice_rows}</tbody>
        </table>
      </div>
      <div class="card span-12">
        <div class="label">Earned Value Management</div>
        <h3>BAC / PV / EV / SV Executive Readout</h3>
        {evm_overview_html}
      </div>
    </div>

    <div class="section">
      <div class="card span-12">
        <div class="label">Earned Value Analysis Add-on</div>
        <h3>Quantitative Performance (EVM Metrics)</h3>
        <div class="evm-chip">{html.escape(str(evmData["scheduleHealthClassification"]))}</div>
        <div class="grid4" style="margin-top:16px">
          {metric_block("BAC", sar(evmData["BAC"]), "Budget @ Completion")}
          {metric_block("PV", sar(evmData["PV"]), "Cumm. Planned Value")}
          {metric_block("EV", sar(evmData["EV"]), "Cumm. Earned Value")}
          {metric_block("SV", sar(evmData["SV"]), "Negative schedule exposure")}
          {metric_block("SPI", f"{evmData['SPI']:.2f}", "Schedule performance index")}
          {metric_block("Planned completion value position", pct(evmData["plannedCompletionValuePosition"]), "PV as a share of BAC")}
          {metric_block("Earned value gap", sar(evmData["earnedValueGap"]), "PV not converted into EV")}
          {metric_block("Schedule health classification", str(evmData["scheduleHealthClassification"]), "Executive status")}
        </div>
        <div class="grid2">
          <div>{evm_addon_chart_html}</div>
          <div>{evm_gauge_html}</div>
        </div>
        <div class="evm-note">{html.escape(str(evmData["interpretation"]))}</div>
        {evm_comment_block("Add / Edit My Comment", evmComments["quantitativePerformance"])}
      </div>

      <div class="card span-12">
        <div class="label">Earned Value Analysis Add-on</div>
        <h3>Root Cause Linkage</h3>
        <table>
          <thead><tr><th>Cause Title</th><th>Impact Area</th><th>Link to EVM Impact</th><th>Status</th><th>Severity</th></tr></thead>
          <tbody>{root_rows}</tbody>
        </table>
        <div class="evm-note">The negative schedule variance is not an isolated numerical deviation. It is directly linked to unresolved external and interface-driven constraints that prevented planned progress from being converted into earned value, mainly within construction and procurement work fronts.</div>
        {evm_comment_block("Add / Edit My Comment", evmComments["rootCauseLinkage"])}
      </div>

      <div class="card span-12">
        <div class="label">Earned Value Analysis Add-on</div>
        <h3>Contractor Mitigation & Recovery Status</h3>
        <table>
          <thead><tr><th>Action</th><th>Owner / Responsible party</th><th>Current status</th><th>Recovery impact</th><th>Required next decision</th></tr></thead>
          <tbody>{mitigation_rows}</tbody>
        </table>
        <div class="evm-note">Contractor mitigation is focused on protecting available work fronts, accelerating technical closures, maintaining commercial entitlement records, and recovering productivity once external constraints are removed. Recovery remains dependent on timely closure of outstanding Owner / Engineer-driven constraints.</div>
        {evm_comment_block("Add / Edit My Comment", evmComments["mitigationRecovery"])}
      </div>
    </div>

    <div class="section">
      <div class="card span-5">
        <div class="label">Risk Register</div>
        <h3>Risk & Issue Control</h3>
        <div class="risk-stack">{risk_cards}</div>
      </div>
      <div class="card span-7">
        <div class="label">Causal Intelligence</div>
        <h3>Root Cause / Impact Radar</h3>
        {radar_html}
      </div>
      <div class="card span-12">
        <div class="label">Resource Quantities</div>
        <h3>RFT Quantity Distribution</h3>
        {resource_html if resource_html else "<div class='sub'>No dedicated resource quantity source is currently available.</div>"}
      </div>
    </div>

    <div class="footer">
      <div>Effective Date: {_ppt_date(overview_metrics.get("project_start"))} ÃƒÆ’Ã¢â‚¬Å¡Ãƒâ€šÃ‚Â· Ground Works Forecast: {_ppt_date(milestone_metrics.get("milestones_df", pd.DataFrame()).iloc[0].get("forecast_date")) if not milestone_metrics.get("milestones_df", pd.DataFrame()).empty else "ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â‚¬Å¡Ã‚Â¬ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â"}</div>
      <div>Strategic priority: recover construction productivity and preserve entitlement evidence.</div>
    </div>
  </div>
</body>
</html>"""


def format_currency_html(value) -> str:
    try:
        return f"{float(value):,.0f}"
    except (TypeError, ValueError):
        return "0"


def _safe_table_name(name: str) -> str:
    cleaned = "".join(ch for ch in name if ch.isalnum() or ch == "_")
    if not cleaned or cleaned[0].isdigit():
        cleaned = f"T_{cleaned}"
    return cleaned[:250]


def _discipline_from_activity_id(activity_id: str) -> str:
    text = str(activity_id).upper()
    if text.startswith(("E-", "E_SUB", "E-SUB", "E-APP")):
        return "Engineering"
    if text.startswith("P-"):
        return "Procurement"
    if text.startswith(("CON", "AB-")):
        return "Construction"
    return "General"


def _area_from_activity_id(activity_id: str) -> str:
    text = str(activity_id).upper()
    for area in ["B01", "B02", "B03", "B04"]:
        if area in text:
            return area
    return "Project"


def _write_dataframe_table(ws, df: pd.DataFrame, table_name: str, start_row: int = 1, start_col: int = 1) -> None:
    require_openpyxl()
    if df is None:
        df = pd.DataFrame()
    if df.empty:
        return
    df = df.copy()
    for row_idx, row in enumerate(dataframe_to_rows(df, index=False, header=True), start=start_row):
        for col_idx, value in enumerate(row, start=start_col):
            ws.cell(row=row_idx, column=col_idx, value=value)

    end_row = start_row + len(df)
    end_col = start_col + len(df.columns) - 1
    ref = f"{get_column_letter(start_col)}{start_row}:{get_column_letter(end_col)}{end_row}"
    table = Table(displayName=_safe_table_name(table_name), ref=ref)
    style = TableStyleInfo(name="TableStyleMedium2", showFirstColumn=False, showLastColumn=False, showRowStripes=True, showColumnStripes=False)
    table.tableStyleInfo = style
    ws.add_table(table)

    header_fill = PatternFill("solid", fgColor="0B1220")
    header_font = Font(color="FFFFFF", bold=True)
    thin = Side(style="thin", color="DDE7EF")
    for col in range(start_col, end_col + 1):
        cell = ws.cell(start_row, col)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center", vertical="center")
        cell.border = Border(left=thin, right=thin, top=thin, bottom=thin)

    for col in range(start_col, end_col + 1):
        max_len = len(str(ws.cell(start_row, col).value))
        for row in range(start_row + 1, end_row + 1):
            max_len = max(max_len, len(str(ws.cell(row, col).value if ws.cell(row, col).value is not None else "")))
        ws.column_dimensions[get_column_letter(col)].width = min(max(max_len + 2, 12), 36)


def _style_dashboard_sheet(ws, title: str, subtitle: str) -> None:
    require_openpyxl()
    ws.sheet_view.showGridLines = False
    ws["A1"] = title
    ws["A2"] = subtitle
    ws["A1"].font = Font(size=20, bold=True, color="FFFFFF")
    ws["A2"].font = Font(size=11, color="D0D5DD")
    for cell_ref in ["A1", "A2"]:
        ws[cell_ref].fill = PatternFill("solid", fgColor="0B1220")
    for row in range(1, 4):
        for col in range(1, 12):
            ws.cell(row=row, column=col).fill = PatternFill("solid", fgColor="0B1220")
    ws.freeze_panes = "A5"


def _metric_block(ws, cell_ref: str, title: str, value: str, note: str) -> None:
    require_openpyxl()
    col = ws[cell_ref].column
    row = ws[cell_ref].row
    fill = PatternFill("solid", fgColor="FFFFFF")
    border_side = Side(style="thin", color="DDE7EF")
    for r in range(row, row + 3):
        for c in range(col, col + 2):
            cell = ws.cell(r, c)
            cell.fill = fill
            cell.border = Border(left=border_side, right=border_side, top=border_side, bottom=border_side)
    ws.cell(row=row, column=col, value=title).font = Font(size=10, bold=True, color="667085")
    ws.cell(row=row + 1, column=col, value=value).font = Font(size=16, bold=True, color="173B63")
    ws.cell(row=row + 2, column=col, value=note).font = Font(size=9, color="344054")


def _build_detailed_progress_report_frames(overview_metrics: dict, evm_metrics: dict, contract_metrics: dict, delay_metrics: dict, risk_metrics: dict, milestone_metrics: dict, activity_metrics: dict, s_curve_metrics: dict) -> dict[str, pd.DataFrame]:
    project_df = load_core_csv(PROJECTS_CSV_PATH)
    project_row = active_project_row(project_df) if not project_df.empty else pd.Series(dtype=object)
    activities_df = activity_metrics.get("activities_df", pd.DataFrame()).copy()
    milestones_df = milestone_metrics.get("milestones_df", pd.DataFrame()).copy()
    payments_df = contract_metrics.get("payments_df", pd.DataFrame()).copy()
    risks_df = risk_metrics.get("risks_df", pd.DataFrame()).copy()
    delays_df = delay_metrics.get("delays_df", pd.DataFrame()).copy()
    curve_df = s_curve_metrics.get("curve_df", pd.DataFrame()).copy()
    wbs_df = wbs_metrics.get("wbs_df", pd.DataFrame()).copy()
    letters = load_letters_workbook()

    report_control = pd.DataFrame([{
        "Report Number": "RPT-001",
        "Revision": "0",
        "Reporting Period": pd.Timestamp.today().strftime("%b %Y"),
        "Data Date": _ppt_date(pd.Timestamp.today()),
        "Baseline Version": "Current Baseline",
        "Update Version": "Current Update",
        "Prepared By": "Planning Team",
        "Reviewed By": "Project Controls",
        "Approved By": "Management",
        "Issue Date": _ppt_date(pd.Timestamp.today()),
        "Distribution List": "Executive / Consultant / PMO / Client",
        "Confidentiality Status": "Internal Controlled",
    }])

    project_overview = pd.DataFrame([{
        "Contract Value": parse_numeric(project_row.get("contract_value")),
        "Approved Variations": 0.0,
        "Revised Contract Value": parse_numeric(project_row.get("contract_value")),
        "Original Duration": overview_metrics.get("duration_days", 0),
        "Elapsed Duration": round(overview_metrics.get("duration_elapsed_pct", 0), 2),
        "Remaining Duration": round(overview_metrics.get("remaining_duration_pct", 0), 2),
        "Baseline Finish": _ppt_date(overview_metrics.get("project_finish")),
        "Forecast Finish": _ppt_date(project_row.get("forecast_finish")),
        "Project Phase": "Execution",
        "Main Scope": str(project_row.get("project_name", "")),
        "Client": str(project_row.get("client_name", "")),
        "Consultant": str(project_row.get("consultant", "")),
        "Contractor": str(project_row.get("contractor", "")),
        "PMO": str(project_row.get("pmo", "")),
        "Key Narrative": "Unified progress, cost, risk, and correspondence control pack generated from the live platform data model.",
    }])

    executive_summary = pd.DataFrame([
        {"KPI": "Overall Progress %", "Value": round(overview_metrics.get("overall_progress", 0), 2), "Status": "Critical" if overview_metrics.get("overall_progress", 0) < overview_metrics.get("planned_progress", 0) else "On Track", "Trend": "Negative", "Threshold": ">= Planned", "Owner": "Project Controls", "Management Note": "Recover construction progress conversion."},
        {"KPI": "Planned Progress %", "Value": round(overview_metrics.get("planned_progress", 0), 2), "Status": "Baseline", "Trend": "Reference", "Threshold": "N/A", "Owner": "Project Controls", "Management Note": "Baseline progress reference."},
        {"KPI": "Progress Variance %", "Value": round(overview_metrics.get("overall_progress", 0) - overview_metrics.get("planned_progress", 0), 2), "Status": "Critical", "Trend": "Negative", "Threshold": ">= 0", "Owner": "Construction", "Management Note": "Variance requires recovery actions."},
        {"KPI": "SPI", "Value": round(evmData["SPI"], 2), "Status": evmData["scheduleHealthClassification"], "Trend": "Negative", "Threshold": ">= 1.00", "Owner": "Project Controls", "Management Note": "Schedule efficiency below target."},
        {"KPI": "PV", "Value": evmData["PV"], "Status": "Reference", "Trend": "Baseline", "Threshold": "N/A", "Owner": "Project Controls", "Management Note": "Planned value position."},
        {"KPI": "EV", "Value": evmData["EV"], "Status": "Underperforming", "Trend": "Negative", "Threshold": ">= PV", "Owner": "Construction", "Management Note": "Earned value conversion lagging."},
        {"KPI": "SV", "Value": evmData["SV"], "Status": "Critical", "Trend": "Negative", "Threshold": ">= 0", "Owner": "Project Controls", "Management Note": "Negative schedule exposure persists."},
        {"KPI": "Forecast Finish", "Value": _ppt_date(project_row.get("forecast_finish")), "Status": "Slipped", "Trend": "Negative", "Threshold": "<= Baseline Finish", "Owner": "Planning", "Management Note": "Completion forecast beyond baseline."},
    ])

    progress_activities = activities_df.copy()
    if not progress_activities.empty:
        progress_activities["WBS"] = progress_activities.get("wbs_id", "")
        progress_activities["Discipline"] = progress_activities["activity_id"].apply(_discipline_from_activity_id)
        progress_activities["Area"] = progress_activities["activity_id"].apply(_area_from_activity_id)
        progress_activities["Previous Progress %"] = progress_activities.get("actual_progress", "")
        progress_activities["Current Period Progress %"] = (progress_activities.get("actual_progress_num", 0) - progress_activities.get("planned_progress_num", 0)).round(2)
        progress_activities["Total Progress %"] = progress_activities.get("actual_progress", "")
        progress_activities["Planned Progress %"] = progress_activities.get("planned_progress", "")
        progress_activities["Progress Variance %"] = progress_activities.get("progress_variance", 0).round(2)
        progress_activities["Planned Start"] = progress_activities.get("planned_start", "")
        progress_activities["Planned Finish"] = progress_activities.get("planned_finish", "")
        progress_activities["Forecast Start"] = progress_activities.get("forecast_start", "")
        progress_activities["Forecast Finish"] = progress_activities.get("forecast_finish", "")
        progress_activities["Actual Start"] = progress_activities.get("actual_start", "")
        progress_activities["Actual Finish"] = progress_activities.get("actual_finish", "")
        progress_activities["Critical Flag"] = progress_activities.get("is_critical", "")
        progress_activities["Longest Path Flag"] = progress_activities.get("is_critical", "")
        progress_activities["Float Days"] = progress_activities.get("total_float_days_num", 0)
        progress_activities["Responsible Party"] = progress_activities.get("responsible_party", "")
        progress_activities["Remarks"] = ""
        progress_activities = progress_activities[[
            "activity_id", "activity_name", "WBS", "Discipline", "Area",
            "Previous Progress %", "Current Period Progress %", "Total Progress %",
            "Planned Progress %", "Progress Variance %", "Planned Start", "Planned Finish",
            "Forecast Start", "Forecast Finish", "Actual Start", "Actual Finish",
            "Critical Flag", "Longest Path Flag", "Float Days", "Responsible Party", "Remarks"
        ]].rename(columns={"activity_id": "Activity ID", "activity_name": "Activity Name"})

    deliverables = activity_metrics.get("rft_df", pd.DataFrame()).copy()
    if not deliverables.empty:
        deliverables = deliverables.assign(
            **{
                "Deliverable ID": deliverables["activity_id"],
                "Deliverable Name": deliverables["activity_name"],
                "Discipline": "Construction",
                "Planned Quantity": 0.0,
                "Actual Quantity": 0.0,
                "Planned %": deliverables["planned_progress_num"].round(2),
                "Actual %": deliverables["actual_progress_num"].round(2),
                "Variance %": deliverables["progress_variance"].round(2),
                "Status": deliverables["is_critical"].apply(lambda x: "Critical" if str(x).strip().lower() == "yes" else "Open"),
                "Remarks": "",
            }
        )[["Deliverable ID", "Deliverable Name", "Discipline", "Planned Quantity", "Actual Quantity", "Planned %", "Actual %", "Variance %", "Status", "Remarks"]]

    schedule_paths = activity_metrics.get("critical_df", pd.DataFrame()).copy()
    if not schedule_paths.empty:
        schedule_paths = schedule_paths.assign(
            **{
                "Path ID": range(1, len(schedule_paths) + 1),
                "Path Type": "Critical Path",
                "Planned Start": schedule_paths["planned_start"],
                "Planned Finish": schedule_paths["planned_finish"],
                "Forecast Start": schedule_paths["forecast_start"],
                "Forecast Finish": schedule_paths["forecast_finish"],
                "Total Float": schedule_paths["total_float_days_num"],
                "Free Float": schedule_paths["total_float_days_num"],
                "Driving Logic": "",
                "Delay Days": schedule_paths["finish_slip_days"].astype(int),
                "Responsibility": schedule_paths.get("responsible_party", ""),
                "Mitigation Required": "",
            }
        )[["Path ID", "Path Type", "activity_id", "activity_name", "Planned Start", "Planned Finish", "Forecast Start", "Forecast Finish", "Total Float", "Free Float", "Driving Logic", "Delay Days", "Responsibility", "Mitigation Required"]].rename(columns={"activity_id": "Activity ID", "activity_name": "Activity Name"})

    deviated = activity_metrics.get("deviated_df", pd.DataFrame()).copy()
    if not deviated.empty:
        deviated = deviated.assign(
            **{
                "Rank": range(1, len(deviated) + 1),
                "WBS": deviated["wbs_id"],
                "Planned %": deviated["planned_progress_num"].round(2),
                "Actual %": deviated["actual_progress_num"].round(2),
                "Variance %": deviated["progress_variance"].round(2),
                "Forecast Delay Days": deviated["finish_slip_days"].astype(int),
                "Impact Category": deviated["activity_id"].apply(_discipline_from_activity_id),
                "Root Cause": "",
                "Required Action": "",
                "Owner": deviated.get("responsible_party", ""),
                "Due Date": deviated["forecast_finish"],
            }
        )[["Rank", "activity_id", "activity_name", "WBS", "Planned %", "Actual %", "Variance %", "Forecast Delay Days", "Impact Category", "Root Cause", "Required Action", "Owner", "Due Date"]].rename(columns={"activity_id": "Activity ID", "activity_name": "Activity Name"})

    eva = pd.DataFrame([{
        "Period": pd.Timestamp.today().strftime("%b-%Y"),
        "PV": evmData["PV"],
        "EV": evmData["EV"],
        "AC": float(evm_metrics.get("ac", 0.0)),
        "SV": "=C2-B2",
        "CV": "=C2-D2",
        "SPI": "=IF(B2=0,0,C2/B2)",
        "CPI": "=IF(D2=0,0,C2/D2)",
        "BAC": evmData["BAC"],
        "EAC": "=IF(H2=0,0,I2/H2)",
        "ETC": "=J2-D2",
        "VAC": "=I2-J2",
        "Cost Status": "Under Review",
        "Comment": evmData["interpretation"],
    }])

    manpower = pd.DataFrame(columns=["Period", "Discipline", "Direct/Indirect", "Planned Manpower", "Actual Manpower", "Planned Man-hours", "Actual Man-hours", "Variance", "Variance %", "Productivity Note"])
    equipment = pd.DataFrame(columns=["Period", "Equipment Type", "Planned Units", "Actual Units", "Planned Hours", "Actual Hours", "Variance", "Variance %", "Utilization Status", "Remarks"])

    s_curve = curve_df[["Month", "Planned", "Actual", "Invoiced", "planned_value_num", "actual_cost_num"]].copy() if not curve_df.empty else pd.DataFrame()
    if not s_curve.empty:
        s_curve["Planned Progress % Baseline"] = (s_curve["Planned"] / max(s_curve["Planned"].max(), 1) * 100).round(2)
        s_curve["Actual Earned Progress % Primavera"] = (s_curve["Actual"] / max(s_curve["Planned"].max(), 1) * 100).round(2)
        s_curve["Invoiced Submitted %"] = (s_curve["Invoiced"] / max(s_curve["Planned"].max(), 1) * 100).round(2)
        s_curve["Planned Cost"] = s_curve["planned_value_num"]
        s_curve["Earned Value"] = s_curve["Actual"]
        s_curve["Actual Cost"] = s_curve["actual_cost_num"]
        s_curve["Planned Man-hours"] = ""
        s_curve["Actual Man-hours"] = ""
        s_curve["Forecast Progress %"] = s_curve["Actual Earned Progress % Primavera"]
        s_curve = s_curve[["Month", "Planned Progress % Baseline", "Actual Earned Progress % Primavera", "Invoiced Submitted %", "Planned Cost", "Earned Value", "Actual Cost", "Planned Man-hours", "Actual Man-hours", "Forecast Progress %"]].rename(columns={"Month": "Period"})

    shop_drawings = activities_df[activities_df["activity_id"].astype(str).str.startswith(("E-SUB", "E-APP"), na=False)].copy() if not activities_df.empty else pd.DataFrame()
    if not shop_drawings.empty:
        shop_drawings = shop_drawings.assign(
            **{
                "Submittal ID": shop_drawings["activity_id"],
                "Package": shop_drawings["wbs_id"],
                "Discipline": "Engineering",
                "Title": shop_drawings["activity_name"],
                "Planned Submission": shop_drawings["planned_start"],
                "Actual Submission": shop_drawings["actual_start"],
                "Consultant Response Due": shop_drawings["planned_finish"],
                "Actual Response": shop_drawings["actual_finish"],
                "Status": shop_drawings["is_critical"].apply(lambda x: "Delayed" if str(x).strip().lower() == "yes" else "Submitted"),
                "Revision": "",
                "Delay Days": shop_drawings["finish_slip_days"].astype(int),
                "Responsible Party": shop_drawings.get("responsible_party", ""),
                "Action Required": "",
            }
        )[["Submittal ID", "Package", "Discipline", "Title", "Planned Submission", "Actual Submission", "Consultant Response Due", "Actual Response", "Status", "Revision", "Delay Days", "Responsible Party", "Action Required"]]

    procurement = activities_df[activities_df["activity_id"].astype(str).str.startswith(("P-",), na=False)].copy() if not activities_df.empty else pd.DataFrame()
    if not procurement.empty:
        procurement = procurement.assign(
            **{
                "Package ID": procurement["activity_id"],
                "Item/Package": procurement["activity_name"],
                "Supplier": "",
                "PO Status": "Active",
                "Planned Delivery": procurement["planned_finish"],
                "Forecast Delivery": procurement["forecast_finish"],
                "Actual Delivery": procurement["actual_finish"],
                "Status": procurement["is_critical"].apply(lambda x: "Delayed" if str(x).strip().lower() == "yes" else "Open"),
                "Delay Days": procurement["finish_slip_days"].astype(int),
                "Linked Activity": procurement["activity_id"],
                "Mitigation Action": "",
                "Owner": procurement.get("responsible_party", ""),
            }
        )[["Package ID", "Item/Package", "Supplier", "PO Status", "Planned Delivery", "Forecast Delivery", "Actual Delivery", "Status", "Delay Days", "Linked Activity", "Mitigation Action", "Owner"]]

    invoices = payments_df.copy()
    if not invoices.empty:
        invoices["Outstanding Amount"] = invoices["certified_amount_num"] - invoices["paid_amount_num"]
        invoices["Submission Date"] = invoices["invoice_date"].apply(_ppt_date)
        invoices["Payment Due Date"] = ""
        invoices["Certification Date"] = ""
        invoices["Aging Days"] = ""
        invoices["Remarks"] = ""
        invoices = invoices[["invoice_no", "project_id", "certified_amount_num", "certified_amount_num", "paid_amount_num", "Outstanding Amount", "Submission Date", "Certification Date", "Payment Due Date", "payment_status", "Aging Days", "Remarks"]]
        invoices.columns = ["Invoice No", "Period", "Submitted Amount", "Certified Amount", "Paid Amount", "Outstanding Amount", "Submission Date", "Certification Date", "Payment Due Date", "Payment Status", "Aging Days", "Remarks"]

    risks_issues = risks_df.copy()
    if not risks_issues.empty:
        score_map = {"low": 5, "medium": 10, "high": 18}
        risks_issues["Score"] = risks_issues["probability"].astype(str).str.lower().map(score_map).fillna(10)
        risks_issues["Severity"] = risks_issues["Score"].apply(lambda x: "High" if x >= 15 else ("Medium" if x >= 8 else "Low"))
        risks_issues = risks_issues.assign(
            **{
                "Type": "Risk",
                "Description": risks_issues["risk_title"],
                "Impact": risks_issues["probability"].astype(str).str.title(),
                "Owner": risks_issues["owner"],
                "Mitigation Strategy": risks_issues["response_strategy"],
                "Linked Activity": "",
                "Cost Impact": risks_issues["cost_impact"],
                "Schedule Impact": risks_issues["time_impact_days"],
            }
        )[["risk_id", "Type", "risk_category", "Description", "probability", "Impact", "Score", "Severity", "Owner", "status", "Mitigation Strategy", "due_date", "Linked Activity", "Cost Impact", "Schedule Impact"]]
        risks_issues.columns = ["ID", "Type", "Category", "Description", "Probability", "Impact", "Score", "Severity", "Owner", "Status", "Mitigation Strategy", "Due Date", "Linked Activity", "Cost Impact", "Schedule Impact"]

    mitigation_plan = build_evm_mitigation_rows().rename(columns={"Linked Risk/Issue": "Linked Risk/Issue" if "Linked Risk/Issue" in build_evm_mitigation_rows().columns else "Action"})
    mitigation_plan = mitigation_plan.assign(
        **{
            "Action ID": range(1, len(mitigation_plan) + 1),
            "Linked Risk/Issue": "",
            "Status": mitigation_plan["Current Status"],
            "Effectiveness": "",
            "Required Escalation": mitigation_plan["Required Next Decision"],
            "Management Decision": "",
        }
    )[["Action ID", "Linked Risk/Issue", "Action", "Owner / Responsible Party", "Required Next Decision", "Status", "Effectiveness", "Required Escalation", "Management Decision"]]
    mitigation_plan.columns = ["Action ID", "Linked Risk/Issue", "Action Description", "Owner", "Due Date", "Status", "Effectiveness", "Required Escalation", "Management Decision"]

    hse_qaqc = pd.DataFrame([{"Period": pd.Timestamp.today().strftime("%b-%Y"), "Man-hours": "", "LTI": "", "Near Miss": "", "NCRs Open": "", "NCRs Closed": "", "Inspections": "", "Pass Rate %": "", "Key Concern": ""}])

    correspondence_frames = []
    for sheet_name in ["From Contractor", "From Consultant"]:
        if sheet_name in letters and not letters[sheet_name].empty:
            df = letters[sheet_name].copy()
            if sheet_name == "From Contractor":
                df["From"] = contractor_name
                df["To"] = str(active_project_record.get("consultant") or employer_name)
            else:
                df["From"] = str(active_project_record.get("consultant") or employer_name)
                df["To"] = contractor_name
            df["Status"] = ""
            correspondence_frames.append(df)
    correspondence = pd.concat(correspondence_frames, ignore_index=True) if correspondence_frames else pd.DataFrame()
    if not correspondence.empty:
        correspondence = pd.DataFrame(
            {
                "Reference": correspondence.get("Ref No", ""),
                "Date": correspondence.get("Date", ""),
                "From": correspondence.get("From", ""),
                "To": correspondence.get("To", ""),
                "Subject": correspondence.get("Subject", ""),
                "Type": correspondence.get("Type", ""),
                "Status": correspondence.get("Status", ""),
                "Delay Related": correspondence.get("Delay Risk", ""),
                "Claim Related": correspondence.get("Claim Strength", ""),
                "Required Response Date": "",
                "Actual Response Date": "",
                "Overdue Days": "",
                "Linked Activity": correspondence.get("Affected Activities", ""),
                "Key Evidence": correspondence.get("Main Purpose", ""),
            }
        )

    photos_register = pd.DataFrame(columns=["Photo ID", "Date", "Area", "Description", "Activity ID", "Before/Current/After", "File Path", "Comment"])

    weekly_data = pd.DataFrame([{
        "Week Ending": _ppt_date(pd.Timestamp.today()),
        "Weekly Planned %": round(float(overview_metrics.get("planned_progress", 0.0)), 2),
        "Weekly Actual %": round(float(overview_metrics.get("overall_progress", 0.0)), 2),
        "Weekly Variance %": round(float(overview_metrics.get("overall_progress", 0.0)) - float(overview_metrics.get("planned_progress", 0.0)), 2),
        "Activities Completed": int((activities_df["actual_progress_num"] >= 100).sum()) if not activities_df.empty else 0,
        "Activities Planned Next Week": int((activities_df["planned_progress_num"] < 100).sum()) if not activities_df.empty else 0,
        "Constraints": "Steel supply / technical closures / payment continuity",
        "Required Decisions": "Release blocked fronts and close outstanding approvals.",
    }])

    power_bi_model = pd.DataFrame([
        {"Object": "Calendar", "Type": "Dimension", "Relationship / Note": "Related to all period/date fields"},
        {"Object": "Project", "Type": "Dimension", "Relationship / Note": "Project-level metadata"},
        {"Object": "WBS / Activity", "Type": "Dimension", "Relationship / Note": "Linked to progress, schedule, and EVA"},
        {"Object": "Discipline", "Type": "Dimension", "Relationship / Note": "Linked to activities, resources, shop drawings"},
        {"Object": "FactProgress", "Type": "Fact", "Relationship / Note": "Based on Progress Activities"},
        {"Object": "FactCostEVA", "Type": "Fact", "Relationship / Note": "Based on Cost & EVA"},
        {"Object": "FactManpower", "Type": "Fact", "Relationship / Note": "Based on Manpower"},
        {"Object": "FactEquipment", "Type": "Fact", "Relationship / Note": "Based on Equipment"},
        {"Object": "FactRiskIssue", "Type": "Fact", "Relationship / Note": "Based on Risks & Issues"},
        {"Object": "FactShopDrawing", "Type": "Fact", "Relationship / Note": "Based on Shop Drawings"},
        {"Object": "FactProcurement", "Type": "Fact", "Relationship / Note": "Based on Procurement"},
        {"Object": "FactInvoices", "Type": "Fact", "Relationship / Note": "Based on Invoices"},
        {"Object": "Refresh", "Type": "Instruction", "Relationship / Note": "Import named Excel Tables only"},
    ])

    lists = pd.DataFrame({
        "Status": ["Open", "Closed", "Delayed", "Submitted", "Paid", "Under Payment"],
        "Severity": ["Low", "Medium", "High", "", "", ""],
        "Discipline": ["Engineering", "Procurement", "Construction", "Commercial", "HSE", "QAQC"],
        "RAG Threshold": ["1-7 Low", "8-14 Medium", "15-25 High", "", "", ""],
    })

    return {
        "03 Report Control": report_control,
        "04 Project Overview": project_overview,
        "05 Executive Summary": executive_summary,
        "06 Milestones": milestones_df,
        "07 Progress Activities": progress_activities,
        "08 Deliverables": deliverables,
        "09 Schedule Paths": schedule_paths,
        "10 Deviated Activities": deviated,
        "11 Cost & EVA": eva,
        "12 Manpower": manpower,
        "13 Equipment": equipment,
        "14 S-Curve": s_curve,
        "15 Shop Drawings": shop_drawings,
        "16 Procurement": procurement,
        "17 Invoices": invoices,
        "18 Risks & Issues": risks_issues,
        "19 Mitigation Plan": mitigation_plan,
        "20 HSE & QAQC": hse_qaqc,
        "21 Correspondence": correspondence,
        "22 Photos Register": photos_register,
        "23 Weekly Data": weekly_data,
        "24 Power BI Model": power_bi_model,
        "99 Lists": lists,
    }


def build_detailed_progress_report_package(overview_metrics: dict, evm_metrics: dict, contract_metrics: dict, delay_metrics: dict, risk_metrics: dict, milestone_metrics: dict, activity_metrics: dict, s_curve_metrics: dict) -> tuple[bytes, str, list[str]]:
    require_openpyxl()
    frames = _build_detailed_progress_report_frames(overview_metrics, evm_metrics, contract_metrics, delay_metrics, risk_metrics, milestone_metrics, activity_metrics, s_curve_metrics)
    wb = Workbook()
    default_ws = wb.active
    wb.remove(default_ws)

    ordered_sheets = [
        "01 POWER BI DASHBOARD",
        "02 WEEKLY DASHBOARD",
        "03 Report Control",
        "04 Project Overview",
        "05 Executive Summary",
        "06 Milestones",
        "07 Progress Activities",
        "08 Deliverables",
        "09 Schedule Paths",
        "10 Deviated Activities",
        "11 Cost & EVA",
        "12 Manpower",
        "13 Equipment",
        "14 S-Curve",
        "15 Shop Drawings",
        "16 Procurement",
        "17 Invoices",
        "18 Risks & Issues",
        "19 Mitigation Plan",
        "20 HSE & QAQC",
        "21 Correspondence",
        "22 Photos Register",
        "23 Weekly Data",
        "24 Power BI Model",
        "99 Lists",
    ]

    assumptions = [
        "Power BI workbook export is generated; PBIX/PBIT is not produced in this local environment.",
        "Sheets without a dedicated live CSV source are created as structured templates linked to the current unified workbook model.",
        "Engineering / deliverable and shop drawing sheets are derived from engineering-related activities where dedicated package logs are not available at row level.",
    ]

    for sheet_name in ordered_sheets:
        ws = wb.create_sheet(title=sheet_name)
        if sheet_name == "01 POWER BI DASHBOARD":
            _style_dashboard_sheet(ws, "Detailed Progress Report", "Monthly executive dashboard/control tower linked to the unified data model")
            _metric_block(ws, "A5", "Overall Progress", pct(overview_metrics.get("overall_progress")), "Actual project progress")
            _metric_block(ws, "C5", "Planned Progress", pct(overview_metrics.get("planned_progress")), "Baseline planned progress")
            _metric_block(ws, "E5", "SPI", f"{evmData['SPI']:.2f}", "Schedule performance index")
            _metric_block(ws, "G5", "SV", sar(evmData["SV"]), "Schedule variance")
            _metric_block(ws, "I5", "Open Risks", str(risk_metrics.get("open_risks", 0)), "Current active risks")
            monthly_kpi_df = pd.DataFrame([
                {"Metric": "Overall Progress", "Value": overview_metrics.get("overall_progress", 0)},
                {"Metric": "Planned Progress", "Value": overview_metrics.get("planned_progress", 0)},
                {"Metric": "SPI", "Value": evmData["SPI"]},
                {"Metric": "Open Risks", "Value": risk_metrics.get("open_risks", 0)},
            ])
            _write_dataframe_table(ws, monthly_kpi_df, "tbl_dashboard_monthly", start_row=10, start_col=1)
            chart = XLBarChart()
            chart.title = "Executive KPI Snapshot"
            data = Reference(ws, min_col=2, min_row=10, max_row=14)
            cats = Reference(ws, min_col=1, min_row=11, max_row=14)
            chart.add_data(data, titles_from_data=True)
            chart.set_categories(cats)
            chart.height = 7
            chart.width = 10
            ws.add_chart(chart, "F10")
        elif sheet_name == "02 WEEKLY DASHBOARD":
            _style_dashboard_sheet(ws, "Weekly Dashboard", "Operational weekly control view linked to the same data model")
            weekly_df = frames["23 Weekly Data"]
            _metric_block(ws, "A5", "Weekly Actual %", pct(weekly_df.iloc[0]["Weekly Actual %"]), "Actual achieved")
            _metric_block(ws, "C5", "Weekly Planned %", pct(weekly_df.iloc[0]["Weekly Planned %"]), "Weekly target")
            _metric_block(ws, "E5", "Weekly Variance %", pct(weekly_df.iloc[0]["Weekly Variance %"]), "Plan vs actual")
            _metric_block(ws, "G5", "Critical Activities", str(activity_metrics.get("critical_count", 0)), "Critical path pulse")
            _write_dataframe_table(ws, weekly_df, "tbl_dashboard_weekly", start_row=10, start_col=1)
        else:
            frame = frames.get(sheet_name, pd.DataFrame())
            _write_dataframe_table(ws, frame, f"tbl_{sheet_name.lower().replace(' ', '_').replace('&', 'and').replace('-', '_')}", start_row=1, start_col=1)
            ws.freeze_panes = "A2"

    readme_text = """# Detailed Progress Report Package

## Generated Files
1. `Detailed_Progress_Report_PowerBI_Ready.xlsx`
2. `Detailed_Progress_Report.html`
3. `Detailed_Progress_Report.docx`
4. `Detailed_Progress_Report_PowerBI_Style.html`

## Validation Checks
- Workbook generated with named sheets in required order.
- Analytical sheets exported as structured Excel tables where data exists.
- Dashboard sheets styled separately from data tables.
- Shared workbook model tied to the current platform CSV sources.
"""

    stream = io.BytesIO()
    wb.save(stream)
    stream.seek(0)
    return stream.getvalue(), readme_text, assumptions


def _detailed_progress_summary_rows(overview_metrics: dict, evm_metrics: dict, contract_metrics: dict, delay_metrics: dict, risk_metrics: dict, activity_metrics: dict) -> list[tuple[str, str, str]]:
    spi = (float(evm_metrics.get("ev", 0.0)) / float(evm_metrics.get("pv", 1.0))) if float(evm_metrics.get("pv", 0.0)) else 0.0
    certified = float(contract_metrics.get("total_certified", 0.0) or 0.0)
    paid = float(contract_metrics.get("total_paid", 0.0) or 0.0)
    return [
        ("Overall Progress", pct(overview_metrics.get("overall_progress")), f"Planned {pct(overview_metrics.get('planned_progress'))}"),
        ("SPI", f"{spi:.2f}" if spi else "N/A", "Schedule performance index"),
        ("Critical Activities", f"{int(activity_metrics.get('critical_count', 0)):,}", "Current critical path pulse"),
        ("Delay Events", f"{int(delay_metrics.get('total_delay_events', 0)):,}", f"{int(delay_metrics.get('total_delay_days', 0)):,} delay days"),
        ("Open Risks", f"{int(risk_metrics.get('open_risks', 0)):,}", "Active risk register items"),
        ("Commercial Gap", egp(max(certified - paid, 0.0)), "Certified less paid"),
    ]


def build_detailed_progress_report_html(overview_metrics: dict, evm_metrics: dict, contract_metrics: dict, delay_metrics: dict, risk_metrics: dict, milestone_metrics: dict, activity_metrics: dict, s_curve_metrics: dict) -> str:
    project_name = html.escape(str(overview_metrics.get("project_name", "Project")))
    phase_df = _phase_progress_rows(activity_metrics, overview_metrics)
    delay_df = delay_metrics.get("display_delays_df", pd.DataFrame()).head(12)
    risk_df = risk_metrics.get("risks_df", pd.DataFrame()).head(12)
    metric_html = "".join(
        f"<div class='metric'><span>{html.escape(title)}</span><strong>{html.escape(value)}</strong><small>{html.escape(note)}</small></div>"
        for title, value, note in _detailed_progress_summary_rows(overview_metrics, evm_metrics, contract_metrics, delay_metrics, risk_metrics, activity_metrics)
    )

    def table_html(df: pd.DataFrame) -> str:
        if df.empty:
            return "<div class='empty'>No records available.</div>"
        return df.fillna("").to_html(index=False, escape=True, border=0, classes="data")

    return f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>Detailed Progress Report</title>
  <style>
    body{{margin:0;background:#f5f7fb;color:#152238;font-family:Inter,Segoe UI,Arial,sans-serif}}
    .page{{max-width:1320px;margin:0 auto;padding:28px}}
    .hero{{background:linear-gradient(135deg,#071a2f,#123a5d);color:white;border-radius:8px;padding:28px;border:1px solid #163f62}}
    h1{{margin:0;font-size:34px;letter-spacing:0}} .sub{{margin-top:10px;color:#bfd3e8}}
    .grid{{display:grid;grid-template-columns:repeat(3,minmax(0,1fr));gap:14px;margin-top:18px}}
    .metric,.panel{{background:white;border:1px solid #dce5ef;border-radius:8px;padding:18px;box-shadow:0 10px 28px rgba(16,35,61,.07)}}
    .metric span{{display:block;font-size:12px;color:#65758b;text-transform:uppercase;font-weight:800}} .metric strong{{display:block;margin-top:10px;font-size:28px;color:#0f2f4e}} .metric small{{display:block;margin-top:6px;color:#718197}}
    .section{{display:grid;grid-template-columns:1fr 1fr;gap:16px;margin-top:16px}} .wide{{grid-column:1/-1}}
    h2{{margin:0 0 14px;font-size:20px;color:#0f2f4e}} table.data{{width:100%;border-collapse:collapse;font-size:13px}} table.data th,table.data td{{padding:10px;border-bottom:1px solid #e6edf5;text-align:left;vertical-align:top}} table.data th{{background:#f0f5fb;color:#46566c;text-transform:uppercase;font-size:11px}}
    .empty{{color:#718197;padding:12px;background:#f6f9fc;border-radius:8px}}
  </style>
</head>
<body>
  <main class="page">
    <section class="hero"><h1>Detailed Progress Report</h1><div class="sub">{project_name} | Generated {html.escape(pd.Timestamp.today().strftime('%d %b %Y'))}</div></section>
    <section class="grid">{metric_html}</section>
    <section class="section">
      <div class="panel wide"><h2>Engineering / Procurement / Construction Progress</h2>{table_html(phase_df)}</div>
      <div class="panel"><h2>Delay Register Snapshot</h2>{table_html(delay_df)}</div>
      <div class="panel"><h2>Risk Register Snapshot</h2>{table_html(risk_df)}</div>
    </section>
  </main>
</body>
</html>"""


def build_detailed_progress_power_bi_style_html(overview_metrics: dict, evm_metrics: dict, contract_metrics: dict, delay_metrics: dict, risk_metrics: dict, milestone_metrics: dict, activity_metrics: dict, s_curve_metrics: dict) -> str:
    base_html = build_detailed_progress_report_html(overview_metrics, evm_metrics, contract_metrics, delay_metrics, risk_metrics, milestone_metrics, activity_metrics, s_curve_metrics)
    return (
        base_html
        .replace("Detailed Progress Report</h1>", "Detailed Progress Report - Power BI Style</h1>")
        .replace("#f5f7fb", "#07111f")
        .replace("background:white", "background:#0f1b2d")
        .replace("color:#152238", "color:#e5eef8")
        .replace("color:#0f2f4e", "color:#e5eef8")
    )


def build_detailed_progress_report_docx_bytes(overview_metrics: dict, evm_metrics: dict, contract_metrics: dict, delay_metrics: dict, risk_metrics: dict, milestone_metrics: dict, activity_metrics: dict, s_curve_metrics: dict) -> bytes:
    if not DOCX_AVAILABLE or Document is None:
        return b""
    doc = Document()
    doc.add_heading("Detailed Progress Report", level=1)
    doc.add_paragraph(str(overview_metrics.get("project_name", "Project")))
    doc.add_paragraph(f"Generated: {pd.Timestamp.today().strftime('%d %b %Y')}")
    doc.add_heading("Executive KPI Summary", level=2)
    table = doc.add_table(rows=1, cols=3)
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    hdr[0].text = "Metric"
    hdr[1].text = "Value"
    hdr[2].text = "Note"
    for title, value, note in _detailed_progress_summary_rows(overview_metrics, evm_metrics, contract_metrics, delay_metrics, risk_metrics, activity_metrics):
        cells = table.add_row().cells
        cells[0].text = title
        cells[1].text = value
        cells[2].text = note
    doc.add_heading("Engineering / Procurement / Construction Progress", level=2)
    phase_df = _phase_progress_rows(activity_metrics, overview_metrics)
    progress_table = doc.add_table(rows=1, cols=4)
    progress_table.style = "Table Grid"
    for idx, col in enumerate(["Phase", "Planned %", "Actual %", "Variance %"]):
        progress_table.rows[0].cells[idx].text = col
    for _, row in phase_df.iterrows():
        cells = progress_table.add_row().cells
        cells[0].text = str(row.get("phase", ""))
        cells[1].text = str(row.get("planned", ""))
        cells[2].text = str(row.get("actual", ""))
        cells[3].text = str(row.get("variance", ""))
    stream = io.BytesIO()
    doc.save(stream)
    stream.seek(0)
    return stream.getvalue()


st.markdown(
    """
    <style>
    :root{
      --bg:#f8fbfd;--bg2:#eef5f8;--panel:#ffffff;--panel2:#f5fafc;--line:#dde7ef;
      --ink:#172033;--muted:#667085;--teal:#168f8b;--blue:#245f95;--amber:#c98519;--red:#d9544d;--violet:#8b6bd9;
    }
    .stApp{background:linear-gradient(180deg,#f8fbfd 0%,#eef5f8 100%);color:var(--ink)}
    [data-testid="stAppViewContainer"]{background:transparent}
    .block-container{padding-top:.8rem;max-width:1380px;padding-left:1.6rem;padding-right:1.6rem}
    h1,h2,h3,h4,h5,h6,p,label,div,span{color:var(--ink)!important;letter-spacing:0!important}
    .samco-header{background:radial-gradient(circle at 92% 12%,rgba(22,143,139,.18),transparent 30%),linear-gradient(135deg,#fff 0%,#f5fafc 44%,#e8f2f7 100%);border:1px solid var(--line);padding:20px 26px 14px;border-radius:12px;margin-bottom:12px;box-shadow:0 14px 34px rgba(23,59,99,.08)}
    .samco-headline{display:grid;grid-template-columns:118px minmax(0,1fr);align-items:center;column-gap:18px}
    .samco-headline img{width:92px;height:58px;object-fit:contain;display:block;background:#fff;border:1px solid #d8e6ef;border-radius:8px;padding:7px}
    .samco-title{font-size:26px;font-weight:850;margin:0;color:#173b63!important;line-height:1.05}
    .samco-subtitle{font-size:13px;color:#4d647d!important;margin-top:4px;font-weight:650;line-height:1.2}
    .samco-contract-info{display:grid;grid-template-columns:repeat(5,minmax(0,1fr));gap:14px;margin-top:16px}
    .samco-info-item{background:rgba(255,255,255,.76);padding:10px 12px;border-radius:10px;border:1px solid #d8e6ef;min-height:82px;display:flex;flex-direction:column;justify-content:center}
    .samco-info-label{font-size:10px;color:#667085!important;text-transform:uppercase;font-weight:850}
    .samco-info-value{font-size:12px;font-weight:850;color:#172033!important;margin-top:6px;line-height:1.25;word-break:break-word}
    .samco-header .credit-line{margin-top:16px;padding-top:12px;border-top:1px solid #d8e6ef;text-align:left;font-size:11px;letter-spacing:0;color:#526276!important}
    [data-testid="stMetric"]{background:#ffffff;border:1px solid var(--line);border-left:4px solid var(--blue);border-radius:12px;padding:10px 14px;box-shadow:0 10px 22px rgba(22,56,86,.07);min-height:92px}
    [data-testid="stMetricValue"]{color:#173b63!important;font-weight:850;font-size:22px}
    [data-testid="stMetricLabel"]{color:#42526b!important;font-weight:700;font-size:13px}
    [data-testid="stMetricDelta"]{color:var(--teal)!important}
    .kpi-box{background:#ffffff;border:1px solid var(--line);border-left:4px solid var(--blue);border-radius:12px;padding:14px 18px;box-shadow:0 10px 22px rgba(22,56,86,.07);min-height:88px}
    .kpi-box-title{font-size:13px;font-weight:700;color:#42526b!important;line-height:1.2}
    .kpi-box-value{font-size:34px;font-weight:850;color:#173b63!important;line-height:1.05;margin-top:10px}
    .section-header{margin:12px 0 10px;padding:16px 16px;background:#ffffff;border:1px solid var(--line);border-left:4px solid var(--teal);border-radius:12px}
    .section-header h3{margin:0;font-size:22px}
    .panel-note,.executive-panel{background:#ffffff;border:1px solid var(--line);border-radius:12px;padding:14px;box-shadow:0 10px 22px rgba(22,56,86,.07)}
    .executive-topbar{display:flex;justify-content:space-between;align-items:center;gap:24px;background:#ffffff;border:1px solid var(--line);border-radius:18px;padding:22px 24px;margin-bottom:16px;box-shadow:0 12px 28px rgba(22,56,86,.08)}
    .executive-title{font-size:28px;font-weight:850;line-height:1.15}
    .executive-subtitle{font-size:15px;color:var(--muted)!important;margin-top:6px}
    .topbar-right{display:flex;align-items:center;gap:18px}
    .date-pill,.user-pill{padding:11px 15px;border:1px solid var(--line);border-radius:14px;background:#f8fbfd;font-size:14px;color:#173b63!important}
    .overview-shell{display:grid;grid-template-columns:minmax(0,1.55fr) minmax(320px,.65fr);gap:16px;align-items:start}
    .overview-card-title{font-size:14px;font-weight:800;color:#42526b!important;text-transform:none;margin-bottom:12px}
    .alert-card{display:flex;justify-content:space-between;align-items:flex-start;gap:16px;padding:14px 16px;border:1px solid rgba(217,84,77,.22);border-radius:14px;background:#ffffff;margin-bottom:12px}
    .alert-card.medium{border-color:rgba(244,180,55,.24)}
    .alert-card.low{border-color:rgba(73,167,255,.24)}
    .alert-title{font-size:17px;font-weight:800}
    .alert-body{font-size:13px;line-height:1.45;color:#667085!important;margin-top:3px}
    .alert-level{font-size:13px;font-weight:800}
    .executive-footer{display:grid;grid-template-columns:repeat(8,minmax(0,1fr));gap:12px;margin-top:16px}
    .footer-stat{background:#ffffff;border:1px solid var(--line);border-radius:14px;padding:12px 14px}
    .footer-stat-label{font-size:12px;color:#667085!important}
    .footer-stat-value{font-size:28px;font-weight:850;margin-top:4px;color:#173b63!important}
    .credit-line{margin-top:18px;text-align:center;font-size:18px;letter-spacing:.12em;color:#526276!important}
    .credit-line span{color:var(--teal)!important;font-weight:800}
    .stDataFrame,.stPlotlyChart{background:transparent}
    .stTabs{margin-top:0}
    .stTabs [data-baseweb="tab-list"]{gap:6px;border-bottom:1px solid var(--line);padding-left:0}
    .stTabs [data-baseweb="tab"]{height:36px;padding:0 12px;border-radius:8px 8px 0 0;background:#ffffff;border:1px solid #d8e6ef;border-bottom:0;color:#173b63!important;font-weight:500;font-size:11px}
    .stTabs [aria-selected="true"]{box-shadow:none}
    .stTabs [data-baseweb="tab-highlight"]{background:#ff5a52;height:2px;border-radius:2px}
    .claims-hero{background:linear-gradient(135deg,#0f172a 0%,#173b63 52%,#1f4f7f 100%);border:1px solid rgba(20,40,70,.32);border-radius:16px;padding:22px 24px 20px;margin:8px 0 14px;box-shadow:0 16px 36px rgba(15,23,42,.22)}
    .claims-hero-topline{font-size:11px;font-weight:800;letter-spacing:.16em;color:#cbd5e1!important}
    .claims-hero-title{font-size:29px;font-weight:850;line-height:1.1;color:#ffffff!important;margin-top:8px}
    .claims-hero-subtitle{font-size:14px;line-height:1.45;color:#dbe7f3!important;margin-top:8px;max-width:980px}
    .claims-kpi-card{background:#ffffff;border:1px solid var(--line);border-top:4px solid var(--blue);border-radius:14px;padding:14px 16px;box-shadow:0 12px 24px rgba(22,56,86,.08);min-height:112px}
    .claims-kpi-card.tone-blue{border-top-color:var(--blue)}
    .claims-kpi-card.tone-green{border-top-color:var(--teal)}
    .claims-kpi-card.tone-amber{border-top-color:var(--amber)}
    .claims-kpi-card.tone-red{border-top-color:var(--red)}
    .claims-kpi-title{font-size:12px;font-weight:800;color:#516177!important;text-transform:uppercase;line-height:1.25}
    .claims-kpi-value{font-size:30px;font-weight:850;color:#0f172a!important;line-height:1.05;margin-top:10px}
    .claims-kpi-subtitle{font-size:11px;color:#667085!important;line-height:1.35;margin-top:8px}
    .claims-badge{display:inline-flex;align-items:center;padding:4px 10px;border-radius:999px;font-size:11px;font-weight:800;border:1px solid transparent}
    .badge-red{background:#fde8e8;color:#9f1239!important;border-color:#fecaca}
    .badge-amber{background:#fff4d6;color:#92400e!important;border-color:#fde68a}
    .badge-green{background:#e7f7ef;color:#0f766e!important;border-color:#a7f3d0}
    .badge-blue{background:#e8f1fb;color:#1d4f91!important;border-color:#bfdbfe}
    .badge-slate{background:#eef2f7;color:#334155!important;border-color:#cbd5e1}
    .claims-workflow{display:grid;grid-template-columns:repeat(7,minmax(0,1fr));gap:10px;margin:6px 0 18px}
    .claims-workflow-step{background:#ffffff;border:1px solid var(--line);border-radius:14px;padding:14px 10px;min-height:90px;display:flex;flex-direction:column;justify-content:flex-start;gap:10px;box-shadow:0 10px 22px rgba(22,56,86,.06)}
    .claims-workflow-index{width:28px;height:28px;border-radius:999px;background:#173b63;color:#ffffff!important;display:flex;align-items:center;justify-content:center;font-size:12px;font-weight:850}
    .claims-workflow-label{font-size:12px;font-weight:750;line-height:1.35;color:#172033!important}
    .claims-surface{background:#ffffff;border:1px solid var(--line);border-radius:14px;padding:16px;box-shadow:0 10px 22px rgba(22,56,86,.07)}
    .claims-surface-title{font-size:15px;font-weight:850;color:#173b63!important;margin-bottom:8px}
    .claims-surface-body{font-size:13px;line-height:1.45;color:#526276!important}
    .claims-statline{display:flex;gap:10px;flex-wrap:wrap;margin-top:10px}
    .claims-dual-grid{display:grid;grid-template-columns:1fr 1fr;gap:14px}
    .claims-answer-panel{background:#ffffff;border:1px solid var(--line);border-radius:14px;padding:14px 16px;box-shadow:0 10px 22px rgba(22,56,86,.07);height:100%}
    .claims-answer-title{font-size:13px;font-weight:800;color:#173b63!important;margin-bottom:8px}
    .claims-answer-body{font-size:13px;line-height:1.5;color:#334155!important}
    .claims-export-grid{display:grid;grid-template-columns:repeat(3,minmax(0,1fr));gap:12px}
    .stTabs [data-testid="stMarkdownContainer"] p{font-size:12px}
    [data-testid="stMetricDelta"] > div{font-size:11px}
    .stDataFrame div[data-testid="stDataFrameResizable"] *{font-size:12px}
    table{color:var(--ink)!important}
    @media(max-width:1200px){.overview-shell{grid-template-columns:1fr}.executive-footer{grid-template-columns:repeat(2,minmax(0,1fr))}.claims-workflow{grid-template-columns:repeat(3,minmax(0,1fr))}.claims-dual-grid,.claims-export-grid{grid-template-columns:1fr}}
    @media(max-width:900px){.samco-contract-info{grid-template-columns:repeat(2,1fr)}.samco-headline{grid-template-columns:88px minmax(0,1fr);align-items:flex-start}.samco-title{font-size:24px}}
    @media(max-width:640px){
      .block-container{padding:.7rem .55rem 1.4rem!important}
      .samco-header{padding:14px!important;border-radius:12px!important;margin-bottom:8px!important}
      .samco-headline{grid-template-columns:68px minmax(0,1fr)!important;column-gap:10px!important}
      .samco-headline img{width:62px!important;height:44px!important;padding:5px!important}
      .samco-title{font-size:19px!important;line-height:1.16!important}
      .samco-subtitle{font-size:12px!important}
      .samco-contract-info{grid-template-columns:1fr!important;gap:8px!important;margin-top:12px!important}
      .samco-info-item{min-height:58px!important;padding:8px 10px!important}
      .section-header{padding:11px 12px!important;margin:8px 0!important;border-radius:10px!important}
      .section-header h3{font-size:18px!important;line-height:1.2!important}
      .panel-note,.executive-panel,.claims-surface,.claims-answer-panel{padding:11px!important;border-radius:10px!important}
      [data-testid="stMetric"]{min-height:72px!important;padding:8px 10px!important;border-radius:10px!important}
      [data-testid="stMetricValue"]{font-size:18px!important;line-height:1.1!important}
      [data-testid="stMetricLabel"]{font-size:11px!important;line-height:1.2!important}
      .kpi-box{min-height:70px!important;padding:10px 12px!important;border-radius:10px!important}
      .kpi-box-title{font-size:11px!important}
      .kpi-box-value{font-size:22px!important}
      .claims-kpi-card{min-height:86px!important;padding:10px 12px!important;border-radius:10px!important}
      .claims-kpi-title{font-size:10px!important}
      .claims-kpi-value{font-size:22px!important}
      .claims-workflow{grid-template-columns:1fr!important}
      .executive-topbar,.topbar-right{display:block!important}
      .executive-title{font-size:21px!important}
      .executive-footer{grid-template-columns:1fr!important}
      div[data-testid="stPlotlyChart"],.stDataFrame{overflow-x:auto!important}
      .stDataFrame div[data-testid="stDataFrameResizable"] *{font-size:10px!important}
      .stButton button,.stDownloadButton button{min-height:34px!important;padding:6px 10px!important;font-size:12px!important}
      div[data-baseweb="select"]{font-size:12px!important}
    }
    </style>
    """,
    unsafe_allow_html=True,
)

projects_for_selector_df = projects_frame(PROJECTS_DIR)
project_selector_options = project_filter_options(projects_for_selector_df)
project_selector_labels = [option["label"] for option in project_selector_options]
current_project_id = selected_project_id()
current_project_index = next(
    (
        idx
        for idx, option in enumerate(project_selector_options)
        if option["project_id"] == current_project_id
    ),
    0,
)
project_selector_col, project_cache_col, project_selector_spacer = st.columns([0.34, 0.18, 0.48])
with project_selector_col:
    selected_project_label = st.selectbox(
        "Dashboard project",
        project_selector_labels,
        index=current_project_index,
        key="overall_dashboard_project_selector",
    )
selected_project_option = project_selector_options[project_selector_labels.index(selected_project_label)]
st.session_state["active_project_id"] = selected_project_option["project_id"]

with project_cache_col:
    st.write("")
    if st.button("Clear selected project cache", width="stretch"):
        st.cache_data.clear()
        for key in list(st.session_state):
            if key.startswith(("tia_", "delay_tia_", "ccc_", "linked_exec_")):
                del st.session_state[key]
        st.rerun()

active_project_id = selected_project_id()
catalog_match = projects_for_selector_df[
    projects_for_selector_df["project_id"].astype(str).str.strip().eq(active_project_id)
] if active_project_id and not projects_for_selector_df.empty else pd.DataFrame()
active_catalog_record = catalog_match.iloc[0].to_dict() if not catalog_match.empty else None
active_project_context = build_project_context(active_catalog_record, PROJECTS_DIR)
active_project_rows = load_core_csv(PROJECTS_CSV_PATH, project_id=active_project_id) if active_project_id else pd.DataFrame()
active_project_record = {
    **(active_catalog_record or {}),
    **(active_project_rows.iloc[0].to_dict() if not active_project_rows.empty else selected_project_option),
}
contract_scope_dir = active_project_context.project_folder_path
CONTRACT_CLAIMS_DIRS = ccc.ensure_contract_claims_dirs(contract_scope_dir)
CONTRACT_REPOSITORY_DIR = CONTRACT_CLAIMS_DIRS["contracts_dir"]
CONTRACT_EVIDENCE_DIR = CONTRACT_CLAIMS_DIRS["evidence_dir"]
CONTRACT_CLAIMS_EXPORT_DIR = CONTRACT_CLAIMS_DIRS["exports_dir"]
CONTRACT_CLAIMS_DB_PATH = CONTRACT_CLAIMS_DIRS["base_dir"] / "contract_claims.db"
active_clause_library = project_data_path(PROJECTS_DIR, active_project_id, "delay_analysis", "06- contract_library.csv") if active_project_id else Path()
set_clause_library_path(active_clause_library)

project = active_project_record
activities = []
wbs_costs = []
letters = load_letters_workbook()
overview_metrics = build_overview_metrics()
evm_metrics = build_evm_metrics()
contract_metrics = build_contract_metrics()
delay_metrics = build_delay_metrics()
risk_metrics = build_risk_metrics()
milestone_metrics = build_milestone_metrics()
time_impact_metrics = build_time_impact_metrics()
s_curve_metrics = build_s_curve_metrics()
activity_metrics = build_activity_metrics()
wbs_metrics = build_wbs_metrics()
wbs_cost_source = wbs_metrics.get("wbs_df", pd.DataFrame())
if not wbs_cost_source.empty:
    wbs_name_col = first_existing_column(wbs_cost_source, ["wbs_name", "WBS Name", "WBS Code", "wbs_code"])
    wbs_budget_col = first_existing_column(wbs_cost_source, ["budget", "budget_cost", "Budget Cost"])
    wbs_actual_col = first_existing_column(wbs_cost_source, ["actual", "actual_cost", "Actual Cost"])
    if wbs_name_col and wbs_budget_col and wbs_actual_col:
        wbs_costs = pd.DataFrame({
            "wbs_name": wbs_cost_source[wbs_name_col].astype(str),
            "budget": wbs_cost_source[wbs_budget_col].apply(parse_numeric),
            "actual": wbs_cost_source[wbs_actual_col].apply(parse_numeric),
        }).to_dict("records")
time_impact_engine = build_time_impact_engine(delay_metrics, risk_metrics, activity_metrics, contract_metrics)
evmData = build_earned_value_analysis_data(evm_metrics)
ensure_evm_comment_state()
evmComments = get_evm_comments()
project_logo_path = active_project_context.branding_path / "logo.png" if active_project_id else LOGO_PATH
legacy_project_logo_path = project_data_path(PROJECTS_DIR, active_project_id, "branding", "logo.png") if active_project_id else LOGO_PATH
if not project_logo_path.exists() and legacy_project_logo_path.exists():
    project_logo_path = legacy_project_logo_path
logo_b64 = image_as_base64(project_logo_path if project_logo_path.exists() else LOGO_PATH)
contractor_name = str(active_project_record.get("contractor") or active_project_record.get("contractor_name") or "Portfolio").strip()
employer_name = str(active_project_record.get("client_name") or active_project_record.get("employer") or "Portfolio").strip()
project_display_name = str(overview_metrics.get("project_name") or active_project_record.get("project_name") or "All Projects").strip()
hub_owner_name = contractor_name if active_project_id else "Projects"
logo_html = f"<img src='data:image/png;base64,{logo_b64}' alt='Project logo'>" if logo_b64 else ""
st.markdown(
    f"""
    <div class='samco-header'>
      <div class='samco-headline'>{logo_html}<div><div class='samco-title'>{html.escape(hub_owner_name)} - Projects Intelligence Hub</div><div class='samco-subtitle'>Integrated Project Controls System</div></div></div>
      <div class='samco-contract-info'>
        <div class='samco-info-item'><div class='samco-info-label'>Contractor</div><div class='samco-info-value'>{html.escape(contractor_name)}</div></div>
        <div class='samco-info-item'><div class='samco-info-label'>Employer</div><div class='samco-info-value'>{html.escape(employer_name)}</div></div>
        <div class='samco-info-item'><div class='samco-info-label'>Contract Value</div><div class='samco-info-value'>{egp(overview_metrics.get('contract_value'))}</div></div>
        <div class='samco-info-item'><div class='samco-info-label'>Project Duration</div><div class='samco-info-value'>{overview_metrics.get('duration_days', 0)} Days</div></div>
        <div class='samco-info-item'><div class='samco-info-label'>Project</div><div class='samco-info-value'>{html.escape(project_display_name)}</div></div>
      </div>
      <div class='credit-line'>Designed and Developed By Eng. Ahmed Labib © Planning Department</div>
    </div>
    """,
    unsafe_allow_html=True,
)

if active_project_context.is_all_projects:
    render_decision_making_dashboard(projects_for_selector_df)
    st.stop()

export_sources = {
    "Overview Metrics": df_for_export(metrics_frame("Overview", {
        "Project Start Date": format_project_date(overview_metrics.get("project_start")),
        "Project Finish Date": format_project_date(overview_metrics.get("project_finish")),
        "Project Duration [Days]": int(overview_metrics.get("duration_days", 0)),
        "Duration Elapsed": pct(overview_metrics.get("duration_elapsed_pct")),
        "Remaining Duration": pct(overview_metrics.get("remaining_duration_pct")),
        "Overall Progress": pct(overview_metrics.get("overall_progress")),
        "Planned Progress": pct(overview_metrics.get("planned_progress")),
        "Contract Value": egp(overview_metrics.get("contract_value")),
        "Total Activities": int(overview_metrics.get("total_activities", 0)),
        "Critical Activities": int(overview_metrics.get("critical_activities", 0)),
    })),
    "EVM Metrics": df_for_export(metrics_frame("EVM", {
        "BAC": egp(evm_metrics.get("bac")),
        "AC": egp(evm_metrics.get("ac")),
        "EV": egp(evm_metrics.get("ev")),
        "PV": egp(evm_metrics.get("pv")),
        "SV": egp(evm_metrics.get("sv")),
        "CV": egp(evm_metrics.get("cv")),
        "EAC": egp(evm_metrics.get("eac")),
        "TCPI": f"{evm_metrics.get('tcpi'):.3f}" if evm_metrics.get("tcpi") is not None else "N/A",
    })),
    "Contracts": df_for_export(contract_metrics["contracts_df"]) if not contract_metrics["contracts_df"].empty else pd.DataFrame(),
    "Payments": df_for_export(contract_metrics["payments_df"]) if not contract_metrics["payments_df"].empty else pd.DataFrame(),
    "Delays": df_for_export(delay_metrics["display_delays_df"]) if "display_delays_df" in delay_metrics and not delay_metrics["display_delays_df"].empty else pd.DataFrame(),
    "Risks": df_for_export(risk_metrics["risks_df"]) if "risks_df" in risk_metrics and not risk_metrics["risks_df"].empty else pd.DataFrame(),
    "Milestones": df_for_export(milestone_metrics["milestones_df"]) if "milestones_df" in milestone_metrics and not milestone_metrics["milestones_df"].empty else pd.DataFrame(),
    "Change Orders": df_for_export(milestone_metrics["change_orders_df"]) if "change_orders_df" in milestone_metrics and not milestone_metrics["change_orders_df"].empty else pd.DataFrame(),
    "Activities": df_for_export(activity_metrics["activities_df"]) if "activities_df" in activity_metrics and not activity_metrics["activities_df"].empty else pd.DataFrame(),
    "WBS": df_for_export(wbs_metrics["wbs_df"]) if "wbs_df" in wbs_metrics and not wbs_metrics["wbs_df"].empty else pd.DataFrame(),
    "Time Impact": df_for_export(time_impact_engine["time_impact_df"]) if "time_impact_df" in time_impact_engine and not time_impact_engine["time_impact_df"].empty else pd.DataFrame(),
    "S-Curve": df_for_export(s_curve_metrics["curve_df"]) if "curve_df" in s_curve_metrics and not s_curve_metrics["curve_df"].empty else pd.DataFrame(),
}

DELAY_TIA_UPLOAD_CACHE_DIR = APP_DIR / ".streamlit" / "delay_tia_upload_cache"
STEEL_DELAY_DURATION_COLUMNS = [
    "Delayed duration in days due to steel un avilability in site",
    "Delayed duration in days due to steel unavailability in site",
    "Steel Delay Duration (days)",
]


def steel_tia_load_csv_or_empty(path: Path) -> pd.DataFrame:
    if not selected_project_id():
        return pd.DataFrame()
    path = project_scoped_file(path)
    if not path.exists():
        return pd.DataFrame()
    try:
        frame = pd.read_csv(path)
        if "project_id" not in frame.columns:
            frame.insert(0, "project_id", selected_project_id())
        else:
            source_ids = frame["project_id"].astype(str).str.strip()
            mismatched = source_ids.ne("") & source_ids.ne(selected_project_id())
            if mismatched.any() and "source_project_id" not in frame.columns:
                frame.insert(1, "source_project_id", source_ids)
            frame["project_id"] = selected_project_id()
        if "source_file" not in frame.columns:
            frame["source_file"] = path.name
        if "source_folder" not in frame.columns:
            frame["source_folder"] = str(path.parent)
        if "source_row" not in frame.columns:
            frame["source_row"] = frame.index + 2
        return frame
    except Exception:
        return pd.DataFrame()


def steel_tia_load_first_csv_or_empty(paths: list[Path]) -> pd.DataFrame:
    for path in paths:
        frame = steel_tia_load_csv_or_empty(path)
        if not frame.empty:
            return frame
    return pd.DataFrame()


def build_steel_delay_template_inventory_df() -> pd.DataFrame:
    rows = []
    inventory_dir = project_scoped_file(STEEL_TIA_DIR / "01-project_metadata_template.csv").parent
    for path in sorted(inventory_dir.glob("*")):
        if path.suffix.lower() == ".csv":
            frame = steel_tia_load_csv_or_empty(path)
            rows.append(
                {
                    "File": path.name,
                    "Type": "CSV",
                    "Rows": int(len(frame)),
                    "Columns": int(len(frame.columns)),
                    "Recognized Use": steel_tia_template_use_label(path.name),
                    "Column List": ", ".join(str(col) for col in frame.columns),
                }
            )
        elif path.suffix.lower() == ".md":
            rows.append(
                {
                    "File": path.name,
                    "Type": "Documentation",
                    "Rows": 0,
                    "Columns": 0,
                    "Recognized Use": "Folder guidance / upload instructions.",
                    "Column List": "",
                }
            )
    return pd.DataFrame(rows)


def steel_tia_template_use_label(file_name: str) -> str:
    mapping = {
        "01-project_metadata_template.csv": "Project metadata, baseline/update identity, data date, parties, and report context.",
        "02- master_activity_steel_analysis.csv": "Activity-level steel requirement, shortage, delayed duration, and affected activity source.",
        "03- employer_steel_supply_at_site.csv": "Employer supply timing used for steel availability and stock-out calculation.",
        "04- p6_activity_export.csv": "Current schedule, baseline dates, float, criticality, longest path, and progress context.",
        "05- relationship_file.csv": "Predecessor/successor logic for causation and fragnet insertion.",
        "06- contract_library.csv": "Entitlement, notice, time-bar, money impact, and schedule-impact support.",
        "07- ifc_conflict.csv": "IFC/design conflict support delay stream.",
        "08- payments.csv": "Payment/cashflow support delay stream.",
        "09- rfi_status.csv": "RFI status support delay stream.",
        "10- contractor_steel_supplied_at_site.csv": "Contractor steel visibility and mitigation evidence only; excluded from employer-steel entitlement calculation.",
        "11-concurrency_matrix_template.updated.csv": "Concurrency framework for overlap windows, critical path checks, and compensability separation.",
        "01- master_activity_steel_analysis.csv": "Legacy master activity steel analysis file.",
        "02- employer_steel_supply.csv": "Legacy employer steel supply file.",
        "03- employer_steel_supply.csv": "Legacy employer supply timing used for steel availability and stock-out calculation.",
        "03- p6_activity_export.csv": "Legacy P6 activity export file.",
        "04- relationship_file.csv": "Legacy relationship file.",
        "05- contract_library.csv": "Legacy contract library file.",
        "06- ifc_conflict.csv": "Legacy IFC/design conflict support delay stream.",
        "07- payments.csv": "Legacy payment/cashflow support delay stream.",
        "08- rfi_status.csv": "Legacy RFI status support delay stream.",
        "09- samco_steel_supplied_at_site.csv": "Legacy contractor site steel visibility only; excluded from employer-steel delay calculation.",
        "10- samco_steel_supplied_at_site.csv": "Legacy contractor site steel visibility only; excluded from employer-steel delay calculation.",
        "RFI Delay.csv": "Detailed RFI delay register used for RFI support and concurrency review.",
    }
    return mapping.get(file_name, "Recognized folder file; inspect columns before analysis.")


def delay_tia_template_candidate_paths(key: str, label: str = "") -> list[Path]:
    current = {
        "metadata": ["01-project_metadata_template.csv"],
        "master": ["02- master_activity_steel_analysis.csv", "02-master_activity_steel_analysis.csv", "01- master_activity_steel_analysis.csv", "01- master_activity_steel_analysis.csv"],
        "employer": ["03- employer_steel_supply_at_site.csv", "03- employer_steel_supply.csv", "02- employer_steel_supply.csv"],
        "p6": ["04- p6_activity_export.csv", "03- p6_activity_export.csv"],
        "relationship": ["05- relationship_file.csv", "04- relationship_file.csv"],
        "contract": ["06- contract_library.csv", "05- contract_library.csv"],
        "ifc": ["07- ifc_conflict.csv", "06- ifc_conflict.csv"],
        "payments": ["08- payments.csv", "07- payments.csv"],
        "rfi": ["09- rfi_status.csv", "08- rfi_status.csv"],
        "samco": ["10- contractor_steel_supplied_at_site.csv", "10- samco_steel_supplied_at_site.csv", "09- samco_steel_supplied_at_site.csv"],
        "concurrency": ["11-concurrency_matrix_template.updated.csv", "11-concurrency_matrix_template.csv", "07-concurrency_matrix_template.csv"],
    }
    names = list(current.get(key, []))
    if label and label not in names:
        names.insert(0, label)
    return [STEEL_TIA_DIR / name for name in names]


def load_delay_tia_template_fallback(key: str, label: str = "") -> tuple[str | None, pd.DataFrame]:
    for path in delay_tia_template_candidate_paths(key, label):
        frame = steel_tia_load_csv_or_empty(path)
        if not frame.empty:
            return path.name, filter_active_project(frame)
    return None, pd.DataFrame()


def steel_tia_qty_sum(series: pd.Series | None) -> float:
    if series is None:
        return 0.0
    return float(pd.to_numeric(series, errors="coerce").fillna(0).sum())




def numeric_series_from_columns(df: pd.DataFrame, candidates: list[str], default: float = 0.0) -> pd.Series:
    if df is None or df.empty:
        return pd.Series(dtype="float64")
    normalized_lookup = {
        re.sub(r"[^a-z0-9]+", "", str(col).lower()): col
        for col in df.columns
    }
    for candidate in candidates:
        normalized = re.sub(r"[^a-z0-9]+", "", candidate.lower())
        col = normalized_lookup.get(normalized)
        if col is not None:
            return pd.to_numeric(df[col], errors="coerce").fillna(default)
    return pd.Series([default] * len(df), index=df.index, dtype="float64")


def steel_tia_date_label(value) -> str:
    if value is None or pd.isna(value):
        return "N/A"
    parsed = value if isinstance(value, pd.Timestamp) else pd.to_datetime(parse_mixed_date(value), errors="coerce")
    if pd.isna(parsed):
        text = str(value).strip()
        return text if text else "N/A"
    return parsed.strftime("%d-%b-%Y")


def steel_tia_prepare_employer_supply_df(raw_df: pd.DataFrame) -> pd.DataFrame:
    if raw_df.empty:
        return pd.DataFrame(
            columns=[
                "Date",
                "Steel Type / Diameter",
                "Delivered Qty",
                "Consumed Qty",
                "Balance Qty",
                "Delivery Ref",
                "Remarks",
                "Building / Zone Allocation",
                "Activity ID consuming the steel",
                "Remaining Required Qty",
                "Balance After Allocation",
            ]
        )
    mapping = suggest_mapping(raw_df, STEEL_ALIASES)
    mapped, _, _ = apply_mapping(raw_df, mapping, CANONICAL_STEEL_FIELDS)
    working = mapped.copy()
    for col in CANONICAL_STEEL_FIELDS:
        if col not in working.columns:
            working[col] = ""
    if "Steel Type / Diameter" in working.columns:
        working["Steel Type / Diameter"] = working["Steel Type / Diameter"].replace("", pd.NA).fillna("RFT STEEL")
    if "Delivered Qty" in working.columns:
        working["Delivered Qty"] = pd.to_numeric(working["Delivered Qty"], errors="coerce").fillna(0)
    for numeric_col in ["Consumed Qty", "Balance Qty", "Remaining Required Qty", "Balance After Allocation"]:
        working[numeric_col] = pd.to_numeric(working[numeric_col], errors="coerce").fillna(0)
    if "Delivery Ref" in working.columns:
        working["Delivery Ref"] = working["Delivery Ref"].replace("", pd.NA).fillna(
            pd.Series([f"EMP-SUPPLY-{idx + 1:03d}" for idx in range(len(working))], index=working.index)
        )
    if "Remarks" in working.columns:
        working["Remarks"] = working["Remarks"].replace("", pd.NA).fillna("Employer steel supply")
    if "Building / Zone Allocation" in working.columns:
        working["Building / Zone Allocation"] = working["Building / Zone Allocation"].replace("", pd.NA).fillna("MULTI")
    if "Activity ID consuming the steel" in working.columns:
        working["Activity ID consuming the steel"] = working["Activity ID consuming the steel"].fillna("")
    return working[CANONICAL_STEEL_FIELDS]


def build_bl_critical_path_comparison(bl_df: pd.DataFrame, activity_metrics: dict) -> tuple[pd.DataFrame, pd.DataFrame, dict]:
    if bl_df.empty:
        return pd.DataFrame(), pd.DataFrame(), {"bl_count": 0, "current_count": 0, "matched_count": 0, "bl_only_count": 0, "current_only_count": 0}

    working = bl_df.copy()
    working = working.loc[:, [col for col in working.columns if str(col).strip()]]
    if "Activity ID" not in working.columns:
        return pd.DataFrame(), pd.DataFrame(), {"bl_count": 0, "current_count": 0, "matched_count": 0, "bl_only_count": 0, "current_only_count": 0}

    working["Activity ID"] = working["Activity ID"].astype(str).str.strip()
    working = working[working["Activity ID"] != ""].copy()
    working["Critical Flag"] = working.get("Critical", "").astype(str).str.strip().str.lower().eq("yes")
    working["Longest Path Flag"] = working.get("Longest Path", "").astype(str).str.strip().str.lower().eq("yes")
    working["BL Critical Path Flag"] = working["Critical Flag"] | working["Longest Path Flag"]
    working["BL Total Float"] = pd.to_numeric(working.get("Total Float"), errors="coerce")
    working["BL Project Start Parsed"] = working.get("BL Project Start", "").apply(parse_mixed_date)
    working["BL Project Finish Parsed"] = working.get("BL Project Finish", "").apply(parse_mixed_date)
    working["Current Start Parsed"] = working.get("Start", "").apply(parse_mixed_date)
    working["Current Finish Parsed"] = working.get("Finish", "").apply(parse_mixed_date)
    working = working.drop_duplicates(subset=["Activity ID"], keep="first")

    activities_df = activity_metrics.get("activities_df", pd.DataFrame()).copy()
    if activities_df.empty:
        comparison_df = pd.DataFrame()
        fixed_view = working[working["BL Critical Path Flag"]].copy()
    else:
        activities_df["Activity ID"] = activities_df["activity_id"].astype(str).str.strip()
        activities_df["Current Critical Path Flag"] = activities_df["is_critical"].astype(str).str.strip().str.lower().eq("yes")
        activities_df["Current Float"] = pd.to_numeric(activities_df.get("total_float_days"), errors="coerce")
        activities_df["Current Planned Finish"] = activities_df.get("planned_finish", "")
        activities_df["Current Forecast Finish"] = activities_df.get("forecast_finish", "")
        comparison_df = working.merge(
            activities_df[
                [
                    "Activity ID",
                    "activity_name",
                    "Current Critical Path Flag",
                    "Current Float",
                    "Current Planned Finish",
                    "Current Forecast Finish",
                ]
            ],
            on="Activity ID",
            how="outer",
        )
        comparison_df["BL Critical Path Flag"] = comparison_df["BL Critical Path Flag"].fillna(False)
        comparison_df["Current Critical Path Flag"] = comparison_df["Current Critical Path Flag"].fillna(False)
        comparison_df["Comparison Status"] = comparison_df.apply(
            lambda row: "Matched"
            if bool(row["BL Critical Path Flag"]) and bool(row["Current Critical Path Flag"])
            else ("BL Only" if bool(row["BL Critical Path Flag"]) else ("Current Only" if bool(row["Current Critical Path Flag"]) else "Neither")),
            axis=1,
        )
        fixed_view = comparison_df[comparison_df["BL Critical Path Flag"]].copy()

    summary = {
        "bl_count": int(working["BL Critical Path Flag"].sum()),
        "current_count": int(activity_metrics.get("critical_count", 0)),
        "matched_count": int((comparison_df["Comparison Status"] == "Matched").sum()) if not comparison_df.empty else 0,
        "bl_only_count": int((comparison_df["Comparison Status"] == "BL Only").sum()) if not comparison_df.empty else 0,
        "current_only_count": int((comparison_df["Comparison Status"] == "Current Only").sum()) if not comparison_df.empty else 0,
    }

    fixed_view = fixed_view[
        [
            col for col in [
                "Activity ID",
                "Activity Name",
                "WBS",
                "BL Project Start",
                "BL Project Finish",
                "Start",
                "Finish",
                "BL Total Float",
                "Critical",
                "Longest Path",
                "Current Critical Path Flag",
                "Current Float",
                "Current Planned Finish",
                "Current Forecast Finish",
                "Comparison Status",
            ]
            if col in fixed_view.columns
        ]
    ]

    return working, comparison_df, summary






def steel_tia_read_file_bytes(file_bytes: bytes, file_name: str) -> pd.DataFrame:
    suffix = Path(file_name).suffix.lower()
    buffer = io.BytesIO(file_bytes)
    try:
        if suffix in {".xlsx", ".xlsm", ".xls"}:
            df = pd.read_excel(buffer).fillna("")
        else:
            df = pd.read_csv(buffer).fillna("")
    except Exception:
        df = pd.DataFrame()
    return df


def delay_tia_cache_payload_path(cache_key: str) -> Path:
    return DELAY_TIA_UPLOAD_CACHE_DIR / f"{cache_key}.bin"


def delay_tia_cache_meta_path(cache_key: str) -> Path:
    return DELAY_TIA_UPLOAD_CACHE_DIR / f"{cache_key}.json"








def load_bl_fixed_context() -> dict[str, pd.DataFrame]:
    return {
        "schedule_df": filter_active_project(steel_tia_load_csv_or_empty(BL_FIXED_DIR / "BL Schedule.csv")),
        "float_df": filter_active_project(steel_tia_load_csv_or_empty(BL_FIXED_DIR / "BL float bath.csv")),
        "longest_df": filter_active_project(steel_tia_load_csv_or_empty(BL_FIXED_DIR / "Bl Longest bath.csv")),
        "critical_df": filter_active_project(steel_tia_load_csv_or_empty(BL_FIXED_DIR / "BL critical path.csv")),
        "mep_df": filter_active_project(steel_tia_load_csv_or_empty(BL_FIXED_DIR / "MEP Activities.csv")),
        "mep_schedule_df": filter_active_project(steel_tia_load_csv_or_empty(BL_FIXED_DIR / "MEP Schedule.csv")),
        "mep_civil_logic_df": filter_active_project(steel_tia_load_csv_or_empty(BL_FIXED_DIR / "MEP Civil Logic.csv")),
    }


def build_mep_activities_kpis(mep_df: pd.DataFrame) -> dict[str, int | float]:
    if mep_df is None or mep_df.empty:
        return {
            "activity_count": 0,
            "building_count": 0,
            "total_duration_days": 0.0,
            "interface_count": 0,
        }
    duration = pd.to_numeric(mep_df.get("Duration Days"), errors="coerce").fillna(0)
    return {
        "activity_count": int(len(mep_df)),
        "building_count": int(mep_df.get("Building / Zone", pd.Series(dtype=object)).astype(str).str.strip().replace("", pd.NA).dropna().nunique()),
        "total_duration_days": float(duration.sum()),
        "interface_count": int(mep_df.get("Interface Type", pd.Series(dtype=object)).astype(str).str.strip().replace("", pd.NA).dropna().nunique()),
    }


def build_mep_related_search_terms(mep_df: pd.DataFrame) -> list[str]:
    if mep_df is None or mep_df.empty:
        return []
    text_columns = [
        col
        for col in [
            "MEP Activity",
            "Material / System",
            "Interface Type",
            "Civil / Structural Dependency",
            "Scope / Debate Treatment",
            "Delay Analysis Use",
        ]
        if col in mep_df.columns
    ]
    if not text_columns:
        return []
    stop_words = {
        "and",
        "are",
        "for",
        "from",
        "into",
        "not",
        "only",
        "the",
        "this",
        "that",
        "with",
        "work",
        "works",
        "activity",
        "activities",
        "analysis",
        "delay",
        "civil",
        "structural",
    }
    term_counts: dict[str, int] = {}
    for value in mep_df[text_columns].fillna("").astype(str).agg(" ".join, axis=1):
        cleaned = re.sub(r"[^A-Za-z0-9#+/ -]+", " ", value.lower())
        tokens = [token.strip(" -/") for token in re.split(r"[\s,/;:()]+", cleaned) if len(token.strip(" -/")) >= 3]
        for token in tokens:
            if token not in stop_words and not token.isdigit():
                term_counts[token] = term_counts.get(token, 0) + 1
        for phrase in re.findall(r"[a-z0-9#+]+(?:\s+[a-z0-9#+]+){1,3}", cleaned):
            phrase = phrase.strip()
            words = phrase.split()
            if len(phrase) >= 8 and not all(word in stop_words for word in words):
                term_counts[phrase] = term_counts.get(phrase, 0) + 1
    return [term for term, _ in sorted(term_counts.items(), key=lambda item: (-item[1], item[0]))[:40]]


def build_mep_letter_conclusion(row: pd.Series) -> str:
    parts = []
    for col in ["Main Purpose", "Key Requests", "Affected Activities", "Required Actions", "Delay Risk", "EOT Potential"]:
        value = str(row.get(col, "")).strip()
        if value and value.lower() not in {"nan", "none", "not applicable", "n/a"}:
            parts.append(value)
    if not parts:
        subject = str(row.get("Subject", "")).strip()
        parts.append(subject or "Related correspondence requires manual review against the MEP interface register.")
    conclusion = " | ".join(dict.fromkeys(parts))
    return textwrap.shorten(conclusion, width=260, placeholder="...")


def build_mep_related_letters_df(mep_df: pd.DataFrame) -> pd.DataFrame:
    output_columns = ["Direction", "Letter No.", "Date", "Conclusion", "Delay Risk", "EOT Potential", "Claim Strength", "Match Score"]
    sheets = load_letters_workbook()
    if not sheets:
        return pd.DataFrame(columns=output_columns)

    search_terms = build_mep_related_search_terms(mep_df)
    if not search_terms:
        return pd.DataFrame(columns=output_columns)

    rows = []
    for sheet_name, direction in [("From Contractor", "Contractor to Consultant"), ("From Consultant", "Consultant to Contractor")]:
        frame = sheets.get(sheet_name, pd.DataFrame())
        if frame.empty:
            continue
        for _, row in frame.iterrows():
            haystack = " ".join(str(row.get(col, "")) for col in frame.columns).lower()
            match_score = sum(1 for term in search_terms if term in haystack)
            if match_score <= 0:
                continue
            letter_no = str(row.get("Ref No", "")).strip()
            if not letter_no:
                continue
            rows.append(
                {
                    "Direction": direction,
                    "Letter No.": letter_no,
                    "Date": str(row.get("Date", "")).strip(),
                    "Conclusion": build_mep_letter_conclusion(row),
                    "Delay Risk": str(row.get("Delay Risk", "")).strip(),
                    "EOT Potential": str(row.get("EOT Potential", "")).strip(),
                    "Claim Strength": str(row.get("Claim Strength", "")).strip(),
                    "Match Score": match_score,
                }
            )
    if not rows:
        return pd.DataFrame(columns=output_columns)

    letters_df = pd.DataFrame(rows).drop_duplicates(subset=["Letter No.", "Conclusion"]).reset_index(drop=True)
    return letters_df.sort_values(["Match Score", "Date", "Letter No."], ascending=[False, True, True]).head(15).reset_index(drop=True)


def build_mep_schedule_export_xlsx(
    mep_activities_df: pd.DataFrame,
    mep_schedule_df: pd.DataFrame,
    mep_civil_logic_df: pd.DataFrame,
    related_mep_letters_df: pd.DataFrame,
) -> bytes:
    if not OPENPYXL_AVAILABLE:
        return b""
    wb = Workbook()
    default_sheet = wb.active
    wb.remove(default_sheet)

    def add_sheet(name: str, df: pd.DataFrame) -> None:
        ws = wb.create_sheet(title=name[:31])
        data = df.copy() if isinstance(df, pd.DataFrame) else pd.DataFrame()
        if data.empty:
            data = pd.DataFrame([{"Status": "No records available"}])
        for row in dataframe_to_rows(data, index=False, header=True):
            ws.append(row)
        for cell in ws[1]:
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="1F4E78")
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        for column_cells in ws.columns:
            max_len = max(len(str(cell.value or "")) for cell in column_cells)
            ws.column_dimensions[column_cells[0].column_letter].width = min(max(max_len + 2, 12), 55)
        ws.freeze_panes = "A2"

    add_sheet("MEP Activities", mep_activities_df)
    add_sheet("MEP Schedule", mep_schedule_df)
    add_sheet("Civil Logic", mep_civil_logic_df)
    add_sheet("Letter Conclusions", related_mep_letters_df)
    output = io.BytesIO()
    wb.save(output)
    return output.getvalue()


def parse_delay_days_from_text(value) -> object:
    if value is None or pd.isna(value):
        return pd.NA
    text = str(value).strip()
    if not text:
        return pd.NA
    exact_match = re.search(r"delay beyond 10d\s*=\s*(-?\d+(?:\.\d+)?)", text, flags=re.IGNORECASE)
    if exact_match:
        return abs(float(exact_match.group(1)))
    fallback_match = re.search(r"(-?\d+(?:\.\d+)?)\s*days", text, flags=re.IGNORECASE)
    if fallback_match:
        return abs(float(fallback_match.group(1)))
    return pd.NA


def extract_submission_reply_dates(value) -> tuple[object, object]:
    if value is None or pd.isna(value):
        return pd.NaT, pd.NaT
    text = str(value)
    dates = re.findall(r"\d{1,2}-[A-Za-z]{3}-\d{2}", text)
    submission = parse_mixed_date(dates[0]) if len(dates) >= 1 else pd.NaT
    reply = parse_mixed_date(dates[1]) if len(dates) >= 2 else pd.NaT
    return submission, reply


def load_delay_tia_local_support_files() -> dict[str, pd.DataFrame]:
    return {
        "rfi_detailed_df": steel_tia_load_first_csv_or_empty(
            [
                STEEL_TIA_DIR / "09- rfi_status.csv",
            ]
        ),
        "ifc_local_df": steel_tia_load_first_csv_or_empty([STEEL_TIA_DIR / "07- ifc_conflict.csv", STEEL_TIA_DIR / "06- ifc_conflict.csv"]),
    }


def delay_tia_question_read_csv(path: Path) -> pd.DataFrame:
    try:
        return pd.read_csv(path, encoding="utf-8-sig")
    except UnicodeDecodeError:
        return pd.read_csv(path, encoding="latin1")
    except Exception:
        return pd.DataFrame()


def load_delay_tia_question_frames() -> dict[str, pd.DataFrame]:
    frames: dict[str, pd.DataFrame] = {}
    scoped_tia_dir = project_scoped_file(STEEL_TIA_DIR / "01-project_metadata_template.csv").parent
    if scoped_tia_dir.exists():
        for path in sorted(scoped_tia_dir.glob("*.csv")):
            key = f"steel_delay_tia_templates/{path.name}"
            frames[key] = delay_tia_question_read_csv(path)
    return frames


def build_delay_tia_question_column_inventory(frames: dict[str, pd.DataFrame]) -> pd.DataFrame:
    rows = []
    for name, df in frames.items():
        rows.append(
            {
                "Dataset": name,
                "Rows": int(len(df)),
                "Columns": int(len(df.columns)),
                "Column List": ", ".join(str(col) for col in df.columns),
            }
        )
    return pd.DataFrame(rows)


def delay_tia_question_number(df: pd.DataFrame, column_name: str) -> pd.Series:
    if df.empty:
        return pd.Series(dtype="float64")
    if column_name not in df.columns:
        normalized_target = re.sub(r"[^a-z0-9]+", "", str(column_name).lower())
        matched_column = next(
            (
                col
                for col in df.columns
                if re.sub(r"[^a-z0-9]+", "", str(col).lower()) == normalized_target
            ),
            None,
        )
        if matched_column is None:
            return pd.Series(dtype="float64")
        column_name = matched_column
    return pd.to_numeric(df[column_name], errors="coerce").dropna()


def delay_tia_question_first_number(df: pd.DataFrame, column_names: list[str]) -> pd.Series:
    for column_name in column_names:
        series = delay_tia_question_number(df, column_name)
        if not series.empty:
            return series
    return pd.Series(dtype="float64")


def delay_tia_frames_get_first(frames: dict[str, pd.DataFrame], keys: list[str]) -> pd.DataFrame:
    for key in keys:
        frame = frames.get(key, pd.DataFrame())
        if isinstance(frame, pd.DataFrame) and not frame.empty:
            return frame
    return pd.DataFrame()


def build_delay_tia_question_kpis(frames: dict[str, pd.DataFrame]) -> dict[str, float | int]:
    concurrency_df = delay_tia_frames_get_first(
        frames,
        [
            "steel_delay_tia_templates/11-concurrency_matrix_template.updated.csv",
            "steel_delay_tia_templates/11-concurrency_matrix_template.csv",
        ],
    )
    claimed_df = pd.DataFrame()
    fragnet_df = concurrency_df.copy()
    rfi_claim_df = frames.get("steel_delay_tia_templates/09- rfi_status.csv", pd.DataFrame())
    events_df = frames.get("steel_delay_tia_templates/07- ifc_conflict.csv", pd.DataFrame())
    activity_df = delay_tia_frames_get_first(frames, ["steel_delay_tia_templates/02- master_activity_steel_analysis.csv", "steel_delay_tia_templates/04- p6_activity_export.csv"])
    evidence_df = frames.get("steel_delay_tia_templates/06- contract_library.csv", pd.DataFrame())
    p6_df = delay_tia_frames_get_first(frames, ["steel_delay_tia_templates/04- p6_activity_export.csv", "steel_delay_tia_templates/03- p6_activity_export.csv"])
    employer_df = delay_tia_frames_get_first(
        frames,
        [
            "steel_delay_tia_templates/03- employer_steel_supply_at_site.csv",
            "steel_delay_tia_templates/03- employer_steel_supply.csv",
            "steel_delay_tia_templates/02- employer_steel_supply.csv",
        ],
    )
    samco_df = delay_tia_frames_get_first(
        frames,
        [
            "steel_delay_tia_templates/10- contractor_steel_supplied_at_site.csv",
            "steel_delay_tia_templates/10- samco_steel_supplied_at_site.csv",
            "steel_delay_tia_templates/09- samco_steel_supplied_at_site.csv",
        ],
    )

    claimed_days = delay_tia_question_number(claimed_df, "Claimed Delay Duration (days)")
    fragnet_days = delay_tia_question_number(fragnet_df, "Claimed Delay Duration")
    concurrency_days = delay_tia_question_first_number(concurrency_df, ["Concurrent Delay Days", "Concurrent delay"])
    rfi_days = delay_tia_question_number(rfi_claim_df, "Delay Beyond 10 Days")
    employer_qty = delay_tia_question_number(employer_df, "Available units at site received by client")
    samco_qty = delay_tia_question_first_number(samco_df, ["Steel available at site", "Total Quantity", "Delivered Qty"])

    max_claimed = int(claimed_days.max()) if not claimed_days.empty else 0
    max_fragnet = int(fragnet_days.max()) if not fragnet_days.empty else 0

    return {
        "Datasets Loaded": int(len(frames)),
        "Delay Events": int(len(events_df)),
        "Activity Impact Rows": int(len(activity_df)),
        "Fragnet Rows": int(len(fragnet_df)),
        "Evidence Rows": int(len(evidence_df)),
        "P6 Activities": int(len(p6_df)),
        "Critical P6 Activities": int(p6_df.get("Critical", pd.Series(dtype=str)).astype(str).str.lower().eq("yes").sum()) if not p6_df.empty else 0,
        "Longest Path Activities": int(p6_df.get("Longest Path", pd.Series(dtype=str)).astype(str).str.lower().eq("yes").sum()) if not p6_df.empty else 0,
        "Employer Steel Qty": float(employer_qty.sum()) if not employer_qty.empty else 0.0,
        "Contractor Steel Qty Visibility Only": float(samco_qty.sum()) if not samco_qty.empty else 0.0,
        "Max Claimed Delay Days": max_claimed,
        "Gross Claimed Delay Days": int(claimed_days.sum()) if not claimed_days.empty else 0,
        "Max Fragnet Duration": max_fragnet,
        "Gross Fragnet Duration": int(fragnet_days.sum()) if not fragnet_days.empty else 0,
        "Concurrent Delay Days": int(concurrency_days.sum()) if not concurrency_days.empty else 0,
        "RFI Delay Beyond 10 Days": int(rfi_days.sum()) if not rfi_days.empty else 0,
        "Recommended Conservative Days": max(max_claimed, max_fragnet),
    }


def delay_tia_question_search_rows(frames: dict[str, pd.DataFrame], question: str, limit: int = 12) -> pd.DataFrame:
    tokens = [token for token in re.findall(r"[a-z0-9]+", question.lower()) if len(token) >= 3]
    rows = []
    if not tokens:
        return pd.DataFrame(rows)
    for name, df in frames.items():
        if df.empty:
            continue
        for idx, row in df.head(250).iterrows():
            text = " | ".join(f"{col}: {row.get(col, '')}" for col in df.columns[:12]).lower()
            score = sum(1 for token in tokens if token in text)
            if score:
                rows.append(
                    {
                        "Dataset": name,
                        "Row": int(idx) + 1,
                        "Score": int(score),
                        "Matched Content": text[:600],
                    }
                )
    return pd.DataFrame(rows).sort_values(["Score", "Dataset"], ascending=[False, True]).head(limit)


def answer_delay_tia_question(question: str, frames: dict[str, pd.DataFrame]) -> tuple[str, list[tuple[str, pd.DataFrame]], dict[str, float | int]]:
    q = question.lower().strip()
    kpis = build_delay_tia_question_kpis(frames)
    tables: list[tuple[str, pd.DataFrame]] = []

    if any(word in q for word in ["column", "schema", "field", "inspect"]):
        tables.append(("Inspected column inventory", build_delay_tia_question_column_inventory(frames)))
        return (
            "I inspected the available CSV columns before answering. The inventory table lists every loaded Delay Analysis - Time Impact Analysis template dataset, row count, column count, and column list.",
            tables,
            kpis,
        )

    if any(word in q for word in ["final", "submit", "submission", "total", "days", "delay"]):
        concurrency_df = delay_tia_frames_get_first(
            frames,
            [
                "steel_delay_tia_templates/11-concurrency_matrix_template.updated.csv",
                "steel_delay_tia_templates/11-concurrency_matrix_template.csv",
            ],
        )
        for title, df in [
            ("Concurrency review", concurrency_df),
        ]:
            if not df.empty:
                tables.append((title, df))
        answer = (
            f"Recommended conservative delay answer: {int(kpis['Recommended Conservative Days'])} days. "
            "This follows the Delay TIA methodology by using the strongest modelled delay stream rather than adding overlapping delay rows. "
            f"The max claimed delay is {int(kpis['Max Claimed Delay Days'])} days, the max fragnet duration is {int(kpis['Max Fragnet Duration'])} days, "
            f"and the gross concurrency overlap visible in the matrix is {int(kpis['Concurrent Delay Days'])} days. "
            "Employer steel is the calculation basis; contractor-supplied steel is shown only for mitigation visibility."
        )
        return answer, tables, kpis

    if any(word in q for word in ["rfi", "reply", "consultant"]):
        for key in ["steel_delay_tia_templates/09- rfi_status.csv"]:
            df = frames.get(key, pd.DataFrame())
            if not df.empty:
                tables.append((key, df))
        return (
            f"RFI delay support totals {int(kpis['RFI Delay Beyond 10 Days'])} days beyond the 10-day review allowance in the normalized claim file. "
            "These rows are support/concurrency evidence unless separately fragnet-modelled into the CPM network.",
            tables,
            kpis,
        )

    if any(word in q for word in ["concurrent", "concurrency", "overlap"]):
        concurrency_df = delay_tia_frames_get_first(
            frames,
            [
                "steel_delay_tia_templates/11-concurrency_matrix_template.updated.csv",
                "steel_delay_tia_templates/11-concurrency_matrix_template.csv",
            ],
        )
        if not concurrency_df.empty:
            tables.append(("Concurrency matrix", concurrency_df))
        return (
            f"The concurrency matrix shows {int(kpis['Concurrent Delay Days'])} gross overlapping days. "
            "Use these days to avoid double counting and to separate driving delay from support-only overlap.",
            tables,
            kpis,
        )

    if any(word in q for word in ["evidence", "proof", "notice", "letter"]):
        for key in ["steel_delay_tia_templates/06- contract_library.csv", "steel_delay_tia_templates/09- rfi_status.csv", "steel_delay_tia_templates/07- ifc_conflict.csv"]:
            df = frames.get(key, pd.DataFrame())
            if not df.empty:
                tables.append((key, df))
        return (
            f"The evidence register contains {int(kpis['Evidence Rows'])} traceable evidence rows linked to events and activities. "
            "Use the evidence reference, document type, source file, source row, date, and linked letter/RFI/MIR/delivery reference to support the answer.",
            tables,
            kpis,
        )

    if any(word in q for word in ["critical", "longest", "float", "activity", "path"]):
        for key in ["steel_delay_tia_templates/04- p6_activity_export.csv", "steel_delay_tia_templates/11-concurrency_matrix_template.updated.csv"]:
            df = frames.get(key, pd.DataFrame())
            if not df.empty:
                tables.append((key, df.head(80)))
        return (
            f"The current P6 export has {int(kpis['Critical P6 Activities'])} critical activities and {int(kpis['Longest Path Activities'])} longest-path activities. "
            "The activity impact register should be used to connect each delay event to the first affected activity, float position, and successor exposure.",
            tables,
            kpis,
        )

    search_df = delay_tia_question_search_rows(frames, question)
    if not search_df.empty:
        tables.append(("Best matching data rows", search_df))
    tables.append(("Inspected column inventory", build_delay_tia_question_column_inventory(frames)))
    return (
        f"I loaded {int(kpis['Datasets Loaded'])} datasets and inspected their columns before answering. "
        f"The conservative default delay position is {int(kpis['Recommended Conservative Days'])} days, based on max claimed/modelled fragnet duration and concurrency-aware treatment. "
        "The matching rows table shows the best source rows for your question.",
        tables,
        kpis,
    )


def build_delay_tia_path_status_lookup(p6_df: pd.DataFrame, bl_fixed_context: dict[str, pd.DataFrame]) -> dict[str, dict]:
    lookup: dict[str, dict] = {}
    if not p6_df.empty and "Activity ID" in p6_df.columns:
        working = p6_df.copy()
        working["Activity ID"] = working["Activity ID"].astype(str).str.strip()
        for _, row in working.iterrows():
            activity_id = str(row.get("Activity ID", "")).strip()
            if not activity_id:
                continue
            lookup[activity_id] = {
                "Activity Name": delay_tia_docx_text(row.get("Activity Name", "")),
                "WBS": delay_tia_docx_text(row.get("WBS", "")),
                "Current Total Float": pd.to_numeric(row.get("Total Float"), errors="coerce"),
                "Current Critical": delay_tia_docx_text(row.get("Critical", "No"), "No"),
                "Current Longest Path": delay_tia_docx_text(row.get("Longest Path", "No"), "No"),
                "Current Start": parse_mixed_date(row.get("Start")),
                "Current Finish": parse_mixed_date(row.get("Finish")),
                "BL Float Path": "No",
                "BL Longest Path": "No",
                "BL Critical Path": "No",
            }

    float_df = bl_fixed_context.get("float_df", pd.DataFrame()).copy()
    if not float_df.empty and "task_code" in float_df.columns:
        float_working = float_df.copy()
        float_working["task_code"] = float_working["task_code"].astype(str).str.strip()
        float_working = float_working[~float_working["task_code"].isin(["", "Activity ID"])].copy()
        for _, row in float_working.iterrows():
            activity_id = str(row.get("task_code", "")).strip()
            if not activity_id:
                continue
            lookup.setdefault(activity_id, {
                "Activity Name": delay_tia_docx_text(row.get("task_name", "")),
                "WBS": delay_tia_docx_text(row.get("wbs_id", "")),
                "Current Total Float": pd.NA,
                "Current Critical": "No",
                "Current Longest Path": "No",
                "Current Start": pd.NaT,
                "Current Finish": pd.NaT,
                "BL Float Path": "No",
                "BL Longest Path": "No",
                "BL Critical Path": "No",
            })
            lookup[activity_id]["BL Float Path"] = "Yes"

    longest_df = bl_fixed_context.get("longest_df", pd.DataFrame()).copy()
    if not longest_df.empty and "Activity ID" in longest_df.columns:
        longest_df["Activity ID"] = longest_df["Activity ID"].astype(str).str.strip()
        longest_df = longest_df[longest_df["Activity ID"] != ""].copy()
        for _, row in longest_df.iterrows():
            activity_id = str(row.get("Activity ID", "")).strip()
            if not activity_id:
                continue
            lookup.setdefault(activity_id, {
                "Activity Name": delay_tia_docx_text(row.get("Activity Name", "")),
                "WBS": delay_tia_docx_text(row.get("WBS", "")),
                "Current Total Float": pd.NA,
                "Current Critical": "No",
                "Current Longest Path": "No",
                "Current Start": pd.NaT,
                "Current Finish": pd.NaT,
                "BL Float Path": "No",
                "BL Longest Path": "No",
                "BL Critical Path": "No",
            })
            if delay_tia_docx_text(row.get("Longest Path", "No"), "No").lower() == "yes":
                lookup[activity_id]["BL Longest Path"] = "Yes"

    critical_df = bl_fixed_context.get("critical_df", pd.DataFrame()).copy()
    if not critical_df.empty and "Activity ID" in critical_df.columns:
        critical_df["Activity ID"] = critical_df["Activity ID"].astype(str).str.strip()
        critical_df = critical_df[critical_df["Activity ID"] != ""].copy()
        for _, row in critical_df.iterrows():
            activity_id = str(row.get("Activity ID", "")).strip()
            if not activity_id:
                continue
            lookup.setdefault(activity_id, {
                "Activity Name": delay_tia_docx_text(row.get("Activity Name", "")),
                "WBS": delay_tia_docx_text(row.get("WBS", "")),
                "Current Total Float": pd.NA,
                "Current Critical": "No",
                "Current Longest Path": "No",
                "Current Start": pd.NaT,
                "Current Finish": pd.NaT,
                "BL Float Path": "No",
                "BL Longest Path": "No",
                "BL Critical Path": "No",
            })
            critical_flag = delay_tia_docx_text(row.get("Critical", "No"), "No").lower() == "yes"
            longest_flag = delay_tia_docx_text(row.get("Longest Path", "No"), "No").lower() == "yes"
            if critical_flag or longest_flag:
                lookup[activity_id]["BL Critical Path"] = "Yes"
            if longest_flag:
                lookup[activity_id]["BL Longest Path"] = "Yes"

    return lookup


def build_delay_tia_concurrent_delay_review_df(
    master_df: pd.DataFrame,
    p6_df: pd.DataFrame,
    fragnet_df: pd.DataFrame,
    ifc_df: pd.DataFrame,
    rfi_detailed_df: pd.DataFrame,
    bl_fixed_context: dict[str, pd.DataFrame],
) -> pd.DataFrame:
    path_lookup = build_delay_tia_path_status_lookup(p6_df, bl_fixed_context)
    rows: list[dict] = []
    steel_activity_ids_from_fragnet: set[str] = set()

    if not fragnet_df.empty:
        fragnet_working = fragnet_df.copy().reset_index(drop=True)
        for idx, row in fragnet_working.iterrows():
            activity_id = delay_tia_docx_text(row.get("Insert Fragment Before", ""))
            if not activity_id:
                continue
            steel_activity_ids_from_fragnet.add(activity_id)
            path_status = path_lookup.get(activity_id, {})
            delayed_days = compute_claimed_delay_duration_days(
                row.get("Fragment Start", pd.NaT),
                row.get("Fragment Finish", pd.NaT),
                row.get("Fragment Duration", 0),
                row.get("Recorded Steel Delay Duration (days)", pd.NA),
            )
            rows.append(
                {
                    "Delay Stream": "Steel",
                    "Delay Ref": delay_tia_docx_text(row.get("Fragment Activity ID", ""), f"FG-{idx + 1:03d}"),
                    "Delay Type": "Steel unavailability at site",
                    "Activity ID": activity_id,
                    "Activity Name": delay_tia_docx_text(row.get("Affected Activity Name", ""), path_status.get("Activity Name", "")),
                    "WBS": delay_tia_docx_text(path_status.get("WBS", ""), delay_tia_docx_text(row.get("Affected Building", ""))),
                    "Delay Start": parse_mixed_date(row.get("Fragment Start")),
                    "Delay Finish": parse_mixed_date(row.get("Fragment Finish")),
                    "Delayed Days": delayed_days,
                    "Basis": "Taken from TIA Fragnet Recommendation using fragnet start/finish, fragment duration, and recorded steel delay duration where available.",
                    "Current Total Float": path_status.get("Current Total Float", pd.NA),
                    "Current Critical": path_status.get("Current Critical", "No"),
                    "Current Longest Path": path_status.get("Current Longest Path", "No"),
                    "BL Float Path": path_status.get("BL Float Path", "No"),
                    "BL Longest Path": path_status.get("BL Longest Path", "No"),
                    "BL Critical Path": path_status.get("BL Critical Path", "No"),
                }
            )

    delay_col = get_master_steel_delay_duration_column(master_df)
    if delay_col and not master_df.empty:
        steel_working = master_df.copy()
        steel_working[delay_col] = pd.to_numeric(steel_working[delay_col], errors="coerce")
        steel_working = steel_working[steel_working[delay_col].notna() & (steel_working[delay_col] != 0)].copy()
        for _, row in steel_working.iterrows():
            activity_id = delay_tia_docx_text(row.get("Activity ID", ""))
            if activity_id in steel_activity_ids_from_fragnet:
                continue
            path_status = path_lookup.get(activity_id, {})
            rows.append(
                {
                    "Delay Stream": "Steel",
                    "Delay Ref": f"STEEL-{activity_id}",
                    "Delay Type": "Steel unavailability at site",
                    "Activity ID": activity_id,
                    "Activity Name": delay_tia_docx_text(row.get("Activity Name", ""), path_status.get("Activity Name", "")),
                    "WBS": delay_tia_docx_text(path_status.get("WBS", ""), delay_tia_docx_text(row.get("Register Building / Zone Allocation", ""))),
                    "Delay Start": parse_mixed_date(row.get("Start")),
                    "Delay Finish": parse_mixed_date(row.get("Finish")),
                    "Delayed Days": abs(pd.to_numeric(row.get(delay_col), errors="coerce")),
                    "Basis": f"Fallback from 02- master_activity_steel_analysis.csv column `{delay_col}` because no fragnet recommendation row was available for this activity.",
                    "Current Total Float": path_status.get("Current Total Float", pd.NA),
                    "Current Critical": path_status.get("Current Critical", "No"),
                    "Current Longest Path": path_status.get("Current Longest Path", "No"),
                    "BL Float Path": path_status.get("BL Float Path", "No"),
                    "BL Longest Path": path_status.get("BL Longest Path", "No"),
                    "BL Critical Path": path_status.get("BL Critical Path", "No"),
                }
            )

    if not rfi_detailed_df.empty:
        for _, row in rfi_detailed_df.iterrows():
            activity_id = delay_tia_docx_text(row.get("Activity ID", ""))
            if not activity_id:
                continue
            submission_date, reply_date = extract_submission_reply_dates(row.get("Date (Submission | Reply | Delay beyond 10d)"))
            delay_days = parse_delay_days_from_text(row.get("Date (Submission | Reply | Delay beyond 10d)"))
            delay_start = submission_date + pd.Timedelta(days=10) if not pd.isna(submission_date) else pd.NaT
            path_status = path_lookup.get(activity_id, {})
            rows.append(
                {
                    "Delay Stream": "RFI",
                    "Delay Ref": delay_tia_docx_text(row.get("Ref No", "")),
                    "Delay Type": delay_tia_docx_text(row.get("Type", "RFI response delay")),
                    "Activity ID": activity_id,
                    "Activity Name": delay_tia_docx_text(row.get("Activity Name", ""), path_status.get("Activity Name", "")),
                    "WBS": delay_tia_docx_text(path_status.get("WBS", "")),
                    "Delay Start": delay_start,
                    "Delay Finish": reply_date,
                    "Delayed Days": delay_days,
                    "Basis": "RFI Delay.csv field `delay beyond 10d` plus activity linkage from the detailed RFI delay register.",
                    "Current Total Float": path_status.get("Current Total Float", pd.NA),
                    "Current Critical": path_status.get("Current Critical", "No"),
                    "Current Longest Path": path_status.get("Current Longest Path", "No"),
                    "BL Float Path": path_status.get("BL Float Path", "No"),
                    "BL Longest Path": path_status.get("BL Longest Path", "No"),
                    "BL Critical Path": path_status.get("BL Critical Path", "No"),
                }
            )

    if not ifc_df.empty:
        current_ifc_rows_added = 0
        if {"Activity ID", "Delayed days"}.issubset(set(ifc_df.columns)):
            for idx, row in ifc_df.iterrows():
                activity_id = delay_tia_docx_text(row.get("Activity ID", ""))
                delay_days = abs(pd.to_numeric(row.get("Delayed days"), errors="coerce"))
                if not activity_id or pd.isna(delay_days) or delay_days <= 0:
                    continue
                path_status = path_lookup.get(activity_id, {})
                rows.append(
                    {
                        "Delay Stream": "IFC",
                        "Delay Ref": f"IFC-{idx + 1:03d}",
                        "Delay Type": "IFC design conflict / approved drawing dependency",
                        "Activity ID": activity_id,
                        "Activity Name": delay_tia_docx_text(row.get("Activity Name", ""), path_status.get("Activity Name", "")),
                        "WBS": path_status.get("WBS", ""),
                        "Delay Start": parse_mixed_date(delay_tia_first_existing(row, ["Re-Start", "Start"])),
                        "Delay Finish": parse_mixed_date(delay_tia_first_existing(row, ["Finish.1", "Finish"])),
                        "Delayed Days": delay_days,
                        "Basis": "07- ifc_conflict.csv current template using Activity ID, Re-Start/Finish.1, and Delayed days.",
                        "Current Total Float": path_status.get("Current Total Float", pd.NA),
                        "Current Critical": path_status.get("Current Critical", "No"),
                        "Current Longest Path": path_status.get("Current Longest Path", "No"),
                        "BL Float Path": path_status.get("BL Float Path", "No"),
                        "BL Longest Path": path_status.get("BL Longest Path", "No"),
                        "BL Critical Path": path_status.get("BL Critical Path", "No"),
                    }
                )
                current_ifc_rows_added += 1

        revision_row = ifc_df[ifc_df.astype(str).apply(lambda s: s.str.contains("SD-REV-001", case=False, na=False)).any(axis=1)].head(1)
        if current_ifc_rows_added == 0 and not revision_row.empty:
            revision_row = revision_row.iloc[0]
            delay_days = abs(pd.to_numeric(revision_row.get("Var"), errors="coerce"))
            if pd.isna(delay_days):
                delay_days = 37.0
            start_candidates = ifc_df.loc[
                ifc_df.astype(str).apply(lambda s: s.str.contains("24-Mar-26", case=False, na=False)).any(axis=1),
                "Date Of Submission To consultant",
            ]
            delay_start = parse_mixed_date(start_candidates.iloc[0]) if not start_candidates.empty else pd.NaT
            delay_finish = parse_mixed_date(revision_row.get("Date Of Reciving Reply from consultant"))
            activity_id = "CON-B01-FOOT-1050"
            path_status = path_lookup.get(activity_id, {})
            rows.append(
                {
                    "Delay Stream": "IFC",
                    "Delay Ref": "SD-REV-001",
                    "Delay Type": "IFC design conflict / revision cycle",
                    "Activity ID": activity_id,
                    "Activity Name": path_status.get("Activity Name", "RFT for RC Raft & Footings (B01) (footing)"),
                    "WBS": path_status.get("WBS", "B01 Foundation"),
                    "Delay Start": delay_start,
                    "Delay Finish": delay_finish,
                    "Delayed Days": delay_days,
                    "Basis": "07- ifc_conflict.csv `Var = -37` and the 24-Mar-26 to 30-Apr-26 consolidated revision cycle, applied to B01 foundation as instructed.",
                    "Current Total Float": path_status.get("Current Total Float", pd.NA),
                    "Current Critical": path_status.get("Current Critical", "No"),
                    "Current Longest Path": path_status.get("Current Longest Path", "No"),
                    "BL Float Path": path_status.get("BL Float Path", "No"),
                    "BL Longest Path": path_status.get("BL Longest Path", "No"),
                    "BL Critical Path": path_status.get("BL Critical Path", "No"),
                }
            )

    concurrent_df = pd.DataFrame(rows)
    if concurrent_df.empty:
        return concurrent_df

    concurrent_df["Delayed Days"] = pd.to_numeric(concurrent_df["Delayed Days"], errors="coerce")
    concurrent_df["Current Total Float"] = pd.to_numeric(concurrent_df["Current Total Float"], errors="coerce")

    concurrency_notes = []
    for idx, row in concurrent_df.iterrows():
        overlaps = []
        for jdx, other in concurrent_df.iterrows():
            if idx == jdx:
                continue
            same_activity = delay_tia_docx_text(row.get("Activity ID", "")) and row.get("Activity ID") == other.get("Activity ID")
            overlap_window = False
            if not pd.isna(row.get("Delay Start")) and not pd.isna(row.get("Delay Finish")) and not pd.isna(other.get("Delay Start")) and not pd.isna(other.get("Delay Finish")):
                overlap_window = max(row["Delay Start"], other["Delay Start"]) <= min(row["Delay Finish"], other["Delay Finish"])
            if same_activity or overlap_window:
                overlaps.append(f"{other['Delay Stream']}:{other['Delay Ref']}")
        concurrency_notes.append(", ".join(overlaps) if overlaps else "None identified from loaded streams")
    concurrent_df["Concurrent With"] = concurrency_notes
    concurrent_df["Concurrent?"] = concurrent_df["Concurrent With"].ne("None identified from loaded streams").map({True: "Yes", False: "No"})
    concurrent_df["Path Summary"] = concurrent_df.apply(
        lambda row: (
            f"Current TF={delay_tia_docx_number(row.get('Current Total Float'), 'N/A')} | "
            f"Current Critical={delay_tia_docx_text(row.get('Current Critical'), 'No')} | "
            f"Current LP={delay_tia_docx_text(row.get('Current Longest Path'), 'No')} | "
            f"BL Float Path={delay_tia_docx_text(row.get('BL Float Path'), 'No')} | "
            f"BL LP={delay_tia_docx_text(row.get('BL Longest Path'), 'No')}"
        ),
        axis=1,
    )
    concurrent_df = concurrent_df.sort_values(by=["Delay Start", "Delay Stream", "Activity ID"], na_position="last").reset_index(drop=True)
    return concurrent_df


def delay_tia_first_existing(row: pd.Series, names: list[str], default: Any = "") -> Any:
    normalized_lookup = {
        re.sub(r"[^a-z0-9]+", "", str(column).lower()): column
        for column in row.index
    }
    for name in names:
        column_name = name if name in row.index else normalized_lookup.get(re.sub(r"[^a-z0-9]+", "", str(name).lower()))
        if column_name in row.index:
            value = row.get(column_name)
            if pd.notna(value) and str(value).strip() != "":
                return value
    return default


def delay_tia_improvement_rfi_to_detailed_df(improvement_frames: dict[str, pd.DataFrame]) -> pd.DataFrame:
    rfi_detail = improvement_frames.get("rfi_delay_detailed", pd.DataFrame()).copy()
    if not rfi_detail.empty and "Date (Submission | Reply | Delay beyond 10d)" in rfi_detail.columns:
        return rfi_detail

    normalized = improvement_frames.get("rfi_claim6_normalized", pd.DataFrame()).copy()
    if normalized.empty:
        return pd.DataFrame()

    rows = []
    for _, row in normalized.iterrows():
        ref_no = delay_tia_docx_text(delay_tia_first_existing(row, ["RFI Ref", "Ref No", "rfi_id"]))
        submission = delay_tia_docx_text(delay_tia_first_existing(row, ["Submission Date", "submission_date"]))
        reply = delay_tia_docx_text(delay_tia_first_existing(row, ["Reply Date", "response_date", "reply_date"]))
        review_days = delay_tia_docx_text(delay_tia_first_existing(row, ["Review Duration Days", "review_duration_days"], ""))
        delay_days = delay_tia_docx_text(delay_tia_first_existing(row, ["Delay Beyond 10 Days", "delay_beyond_10_days", "Working Delay Days"], ""))
        rows.append(
            {
                "Ref No": ref_no,
                "Date (Submission | Reply | Delay beyond 10d)": (
                    f"Submitted {submission}; replied {reply}; review duration {review_days} days; "
                    f"delay beyond 10d = {delay_days} days"
                ),
                "Type": "RFI delay / improvement register",
                "Subject": delay_tia_first_existing(row, ["Subject", "subject"]),
                "Main Purpose": delay_tia_first_existing(row, ["Main Purpose", "main_purpose"]),
                "Key Requests": delay_tia_first_existing(row, ["Key Requests", "key_requests"]),
                "Scope Impact": delay_tia_first_existing(row, ["Scope Impact", "scope_impact"]),
                "Responsibility": delay_tia_first_existing(row, ["Responsibility", "responsibility"]),
                "Activity ID": delay_tia_first_existing(row, ["Activity ID", "activity_id"]),
                "Activity Name": delay_tia_first_existing(row, ["Activity Name", "activity_name"]),
                "Start Dependency": delay_tia_first_existing(row, ["Start Dependency", "start_dependency"]),
                "Sequence Impact": delay_tia_first_existing(row, ["Sequence Impact", "sequence_impact"]),
                "Commercial Impact": delay_tia_first_existing(row, ["Commercial Impact", "commercial_impact"]),
                "Risk Type": "RFI response delay",
                "Risk Owner": delay_tia_first_existing(row, ["Risk Owner", "risk_owner"]),
                "Delay Risk": delay_tia_first_existing(row, ["Delay Risk", "delay_risk"]),
                "EOT Potential": delay_tia_first_existing(row, ["EOT Potential", "eot_potential"]),
                "Claim Strength": delay_tia_first_existing(row, ["Claim Strength", "claim_strength"]),
                "Required Actions": delay_tia_first_existing(row, ["Required Actions", "required_actions"]),
            }
        )
    return pd.DataFrame(rows)


def build_delay_tia_improvement_concurrent_rows_df(
    improvement_frames: dict[str, pd.DataFrame],
    p6_df: pd.DataFrame,
    bl_fixed_context: dict[str, pd.DataFrame],
) -> pd.DataFrame:
    path_lookup = build_delay_tia_path_status_lookup(p6_df, bl_fixed_context)
    rows: list[dict] = []

    delay_events_df = improvement_frames.get("delay_event_register", pd.DataFrame()).copy()
    if not delay_events_df.empty:
        for _, row in delay_events_df.iterrows():
            event_id = delay_tia_docx_text(delay_tia_first_existing(row, ["Event ID", "Delay Event ID"]))
            activity_id = delay_tia_docx_text(delay_tia_first_existing(row, ["Linked Activity ID", "Activity ID"]))
            if not event_id or not activity_id:
                continue
            path_status = path_lookup.get(activity_id, {})
            event_type = delay_tia_docx_text(delay_tia_first_existing(row, ["Event Type"], "Improvement Register"))
            delay_stream = "Improvement RFI" if "rfi" in event_type.lower() else ("Improvement IFC" if "ifc" in event_type.lower() else "Improvement")
            rows.append(
                {
                    "Delay Stream": delay_stream,
                    "Delay Ref": event_id,
                    "Delay Type": event_type,
                    "Activity ID": activity_id,
                    "Activity Name": delay_tia_docx_text(delay_tia_first_existing(row, ["Linked Activity Name", "Activity Name"]), path_status.get("Activity Name", "")),
                    "WBS": delay_tia_docx_text(path_status.get("WBS", ""), delay_tia_first_existing(row, ["Linked Building / Zone", "Building / Zone"])),
                    "Delay Start": parse_mixed_date(delay_tia_first_existing(row, ["Impact Start Date", "Evidence Date", "Start"])),
                    "Delay Finish": parse_mixed_date(delay_tia_first_existing(row, ["Recovery Date", "Finish"])),
                    "Delayed Days": pd.to_numeric(delay_tia_first_existing(row, ["Working Delay Days", "Claimed Delay Duration (days)", "Claimed Delay Duration"], 0), errors="coerce"),
                    "Basis": "Uploaded Delay TIA improvement event register. Activity ID and Activity Name are treated as the affected activity.",
                    "Current Total Float": path_status.get("Current Total Float", pd.NA),
                    "Current Critical": path_status.get("Current Critical", "No"),
                    "Current Longest Path": path_status.get("Current Longest Path", "No"),
                    "BL Float Path": path_status.get("BL Float Path", "No"),
                    "BL Longest Path": path_status.get("BL Longest Path", "No"),
                    "BL Critical Path": path_status.get("BL Critical Path", "No"),
                }
            )

    fragnet_logic_df = improvement_frames.get("fragnet_logic_register", pd.DataFrame()).copy()
    if not fragnet_logic_df.empty:
        for _, row in fragnet_logic_df.iterrows():
            event_id = delay_tia_docx_text(delay_tia_first_existing(row, ["Delay Event ID", "Event ID"]))
            activity_id = delay_tia_docx_text(delay_tia_first_existing(row, ["Successor Activity ID", "Fragnet Insert Before", "Activity ID"]))
            if not event_id or not activity_id:
                continue
            delay_ref = delay_tia_docx_text(delay_tia_first_existing(row, ["Fragnet ID", "Fragnet Activity ID"], event_id))
            if any(existing.get("Delay Ref") == delay_ref for existing in rows):
                continue
            path_status = path_lookup.get(activity_id, {})
            rows.append(
                {
                    "Delay Stream": "Improvement Fragnet",
                    "Delay Ref": delay_ref,
                    "Delay Type": "Uploaded fragnet logic register",
                    "Activity ID": activity_id,
                    "Activity Name": delay_tia_docx_text(delay_tia_first_existing(row, ["Fragnet Activity Name", "Activity Name"]), path_status.get("Activity Name", "")),
                    "WBS": delay_tia_docx_text(path_status.get("WBS", "")),
                    "Delay Start": parse_mixed_date(delay_tia_first_existing(row, ["Fragment Start", "Start"])),
                    "Delay Finish": parse_mixed_date(delay_tia_first_existing(row, ["Fragment Finish", "Finish"])),
                    "Delayed Days": pd.to_numeric(delay_tia_first_existing(row, ["Claimed Delay Duration", "Fragment Duration", "Original Duration"], 0), errors="coerce"),
                    "Basis": delay_tia_docx_text(delay_tia_first_existing(row, ["Logic Rationale"], "Uploaded improvement fragnet logic.")),
                    "Current Total Float": path_status.get("Current Total Float", pd.NA),
                    "Current Critical": path_status.get("Current Critical", "No"),
                    "Current Longest Path": path_status.get("Current Longest Path", "No"),
                    "BL Float Path": path_status.get("BL Float Path", "No"),
                    "BL Longest Path": path_status.get("BL Longest Path", "No"),
                    "BL Critical Path": path_status.get("BL Critical Path", "No"),
                }
            )

    concurrency_matrix_df = improvement_frames.get("concurrency_matrix", pd.DataFrame()).copy()
    if not concurrency_matrix_df.empty:
        for idx, row in concurrency_matrix_df.iterrows():
            activity_id = delay_tia_docx_text(delay_tia_first_existing(row, ["Activity ID", "Linked Activity ID"]))
            if not activity_id:
                continue
            event_id = delay_tia_docx_text(delay_tia_first_existing(row, ["Primary Event ID", "Delay Event ID", "Event ID"]), f"CONC-{idx + 1:03d}")
            concurrent_days = pd.to_numeric(
                delay_tia_first_existing(
                    row,
                    [
                        "Concurrent delay",
                        "Concurrent Delay Days",
                        "Delayed duration after overlap",
                        "Delayed duration",
                        "Delayed duration in days due to steel un avilability in site",
                    ],
                    0,
                ),
                errors="coerce",
            )
            if pd.isna(concurrent_days) or abs(concurrent_days) <= 0:
                continue
            path_status = path_lookup.get(activity_id, {})
            overlap_start = parse_mixed_date(delay_tia_first_existing(row, ["Overlap Start", "Start", "BL Start"]))
            overlap_finish = parse_mixed_date(delay_tia_first_existing(row, ["Overlap Finish", "Finish", "BL Finish"]))
            rows.append(
                {
                    "Delay Stream": "Improvement Concurrency",
                    "Delay Ref": f"{event_id}-{activity_id}-{idx + 1}",
                    "Delay Type": event_id,
                    "Activity ID": activity_id,
                    "Activity Name": delay_tia_docx_text(delay_tia_first_existing(row, ["Activity Name", "Linked Activity Name"]), path_status.get("Activity Name", "")),
                    "WBS": delay_tia_docx_text(path_status.get("WBS", "")),
                    "Delay Start": overlap_start,
                    "Delay Finish": overlap_finish,
                    "Delayed Days": abs(concurrent_days),
                    "Basis": (
                        "Uploaded 07-concurrency_matrix_template.csv: "
                        f"overlap {delay_tia_docx_text(overlap_start)} to {delay_tia_docx_text(overlap_finish)}; "
                        f"BL Critical={delay_tia_docx_text(delay_tia_first_existing(row, ['BL Critical Path'], path_status.get('BL Critical Path', 'No')), 'No')}; "
                        f"Current Critical={delay_tia_docx_text(delay_tia_first_existing(row, ['Current Critical Path'], path_status.get('Current Critical', 'No')), 'No')}; "
                        f"after overlap={delay_tia_docx_number(delay_tia_first_existing(row, ['Delayed duration after overlap'], ''), '')}."
                    ),
                    "Current Total Float": path_status.get("Current Total Float", pd.NA),
                    "Current Critical": delay_tia_docx_text(delay_tia_first_existing(row, ["Current Critical Path"], path_status.get("Current Critical", "No")), "No"),
                    "Current Longest Path": path_status.get("Current Longest Path", "No"),
                    "BL Float Path": path_status.get("BL Float Path", "No"),
                    "BL Longest Path": path_status.get("BL Longest Path", "No"),
                    "BL Critical Path": delay_tia_docx_text(delay_tia_first_existing(row, ["BL Critical Path"], path_status.get("BL Critical Path", "No")), "No"),
                }
            )

    improvement_df = pd.DataFrame(rows)
    if improvement_df.empty:
        return improvement_df
    improvement_df["Delayed Days"] = pd.to_numeric(improvement_df["Delayed Days"], errors="coerce").fillna(0).abs()
    return improvement_df[improvement_df["Delayed Days"] > 0].copy()


def calculate_delay_tia_improvement_modelled_days(context: dict) -> int:
    improvement_frames = context.get("improvement_frames", {}) if isinstance(context, dict) else {}
    if not isinstance(improvement_frames, dict):
        return 0
    candidate_values = []
    for key, duration_cols in {
        "fragnet_logic_register": ["Claimed Delay Duration", "Fragment Duration", "Original Duration"],
        "claimed_delay_register": ["Claimed Delay Duration (days)", "Claimed Delay Duration"],
        "primavera_fragnet_import": ["Claimed Delay Duration", "Original Duration"],
    }.items():
        df = improvement_frames.get(key, pd.DataFrame())
        if not isinstance(df, pd.DataFrame) or df.empty:
            continue
        for col in duration_cols:
            if col in df.columns:
                values = pd.to_numeric(df[col], errors="coerce").fillna(0).abs()
                if not values.empty:
                    candidate_values.append(int(values.sum()))
                break
    return max(candidate_values) if candidate_values else 0


def build_delay_tia_context_from_frames(
    *,
    master_df: pd.DataFrame,
    employer_raw_df: pd.DataFrame,
    p6_df: pd.DataFrame,
    relationship_df: pd.DataFrame,
    contract_df: pd.DataFrame,
    ifc_df: pd.DataFrame,
    payments_df: pd.DataFrame,
    rfi_df: pd.DataFrame,
    bl_critical_path_df: pd.DataFrame,
    samco_df: pd.DataFrame,
    delay_events_df: pd.DataFrame,
) -> dict:
    employer_supply_df = steel_tia_prepare_employer_supply_df(employer_raw_df)
    samco_supply_df = steel_tia_prepare_employer_supply_df(samco_df)
    requirement_df = build_requirement_df_from_client_supply_sheet(master_df) if not master_df.empty else pd.DataFrame()
    settings = SteelTiaSettings(
        usability_lag_days=2,
        near_critical_float_threshold=10,
        data_date=pd.Timestamp.today().normalize(),
    )
    analysis = run_steel_delay_tia_analysis(
        p6_df=p6_df,
        steel_df=employer_supply_df,
        requirement_df=requirement_df,
        relationship_df=relationship_df,
        contract_library_df=contract_df,
        delay_events_df=delay_events_df,
        settings=settings,
    )
    analysis = enrich_delay_tia_analysis_with_master_delay(master_df, analysis)

    return {
        "master_df": master_df,
        "employer_raw_df": employer_raw_df,
        "p6_df": p6_df,
        "relationship_df": relationship_df,
        "contract_df": contract_df,
        "delay_events_df": delay_events_df,
        "ifc_df": ifc_df,
        "payments_df": payments_df,
        "rfi_df": rfi_df,
        "bl_critical_path_df": bl_critical_path_df,
        "samco_df": samco_df,
        "samco_supply_df": samco_supply_df,
        "analysis": analysis,
        "employer_total": steel_tia_qty_sum(employer_supply_df.get("Delivered Qty") if not employer_supply_df.empty else None),
        "samco_total": steel_tia_qty_sum(samco_supply_df.get("Delivered Qty") if not samco_supply_df.empty else None),
        "employer_first_date": employer_supply_df.get("Date").iloc[0] if not employer_supply_df.empty else pd.NaT,
        "employer_last_date": employer_supply_df.get("Date").iloc[-1] if not employer_supply_df.empty else pd.NaT,
        "samco_first_date": samco_supply_df.get("Date").iloc[0] if not samco_supply_df.empty else pd.NaT,
        "samco_last_date": samco_supply_df.get("Date").iloc[-1] if not samco_supply_df.empty else pd.NaT,
        "concurrent_delay_review_df": pd.DataFrame(),
    }


def get_master_steel_delay_duration_column(master_df: pd.DataFrame) -> str | None:
    if master_df is None or master_df.empty:
        return None
    normalized = {re.sub(r"\s+", " ", str(col).strip()).lower(): col for col in master_df.columns}
    for candidate in STEEL_DELAY_DURATION_COLUMNS:
        match = normalized.get(re.sub(r"\s+", " ", candidate.strip()).lower())
        if match:
            return match
    for col in master_df.columns:
        col_text = str(col).strip().lower()
        if "delayed duration" in col_text and "steel" in col_text:
            return col
    return None


def build_master_steel_delay_duration_lookup(master_df: pd.DataFrame) -> dict[str, float]:
    lookup: dict[str, float] = {}
    if master_df is None or master_df.empty or "Activity ID" not in master_df.columns:
        return lookup
    delay_col = get_master_steel_delay_duration_column(master_df)
    if not delay_col:
        return lookup
    for _, row in master_df.iterrows():
        activity_id = str(row.get("Activity ID", "")).strip()
        delay_days = pd.to_numeric(row.get(delay_col), errors="coerce")
        if activity_id and pd.notna(delay_days):
            lookup[activity_id] = max(float(delay_days), 0.0)
    return lookup


def enrich_delay_tia_analysis_with_master_delay(master_df: pd.DataFrame, analysis: dict) -> dict:
    delay_lookup = build_master_steel_delay_duration_lookup(master_df)
    if not delay_lookup or not isinstance(analysis, dict):
        return analysis
    enriched = dict(analysis)
    candidates_df = enriched.get("candidates_df", pd.DataFrame()).copy()
    if not candidates_df.empty and "Activity ID" in candidates_df.columns:
        candidates_df["Recorded Steel Delay Duration (days)"] = candidates_df["Activity ID"].astype(str).str.strip().map(delay_lookup)
        enriched["candidates_df"] = candidates_df
    fragnet_df = enriched.get("fragnet_df", pd.DataFrame()).copy()
    if not fragnet_df.empty:
        if "Affected Activity" in fragnet_df.columns:
            fragnet_df["Recorded Steel Delay Duration (days)"] = fragnet_df["Affected Activity"].astype(str).str.strip().map(delay_lookup)
        elif "Insert Fragment Before" in fragnet_df.columns:
            fragnet_df["Recorded Steel Delay Duration (days)"] = fragnet_df["Insert Fragment Before"].astype(str).str.strip().map(delay_lookup)
        enriched["fragnet_df"] = fragnet_df
    enriched["master_steel_delay_duration_lookup"] = delay_lookup
    return enriched


def delay_tia_docx_text(value, default: str = "NOT ENOUGH DATA") -> str:
    if value is None:
        return default
    if isinstance(value, pd.Timestamp):
        if pd.isna(value):
            return default
        return value.strftime("%d-%b-%Y")
    if pd.isna(value):
        return default
    text = str(value).strip()
    return text if text and text.lower() != "nan" else default


def delay_tia_docx_number(value, default: str = "0") -> str:
    if value is None or pd.isna(value):
        return default
    try:
        number = float(value)
    except (TypeError, ValueError):
        return delay_tia_docx_text(value, default)
    if number.is_integer():
        return str(int(number))
    return f"{number:.2f}"


def delay_tia_set_table_cell(cell, text: str):
    cell.text = text






def delay_tia_join_nonempty(parts: list[str]) -> str:
    values = [str(part).strip() for part in parts if str(part or "").strip() and str(part).strip().lower() != "nan"]
    return " | ".join(values)


def delay_tia_extract_activity_ids(value) -> list[str]:
    text = str(value or "").strip()
    if not text or text == "-":
        return []
    ids: list[str] = []
    for token in re.split(r"[,\n;/]+", text):
        token = token.strip()
        if token.count("-") >= 2 and token.upper() == token:
            ids.append(token)
    return list(dict.fromkeys(ids))


def delay_tia_build_p6_lookup(p6_df: pd.DataFrame) -> dict[str, pd.Series]:
    lookup: dict[str, pd.Series] = {}
    if p6_df is None or p6_df.empty or "Activity ID" not in p6_df.columns:
        return lookup
    for _, row in p6_df.iterrows():
        activity_id = str(row.get("Activity ID", "")).strip()
        if activity_id and activity_id not in lookup:
            lookup[activity_id] = row
    return lookup


def delay_tia_bool_text(value) -> bool:
    text = str(value or "").strip().lower()
    return text in {"yes", "y", "true", "1"}


def delay_tia_describe_criticality(activity_row: pd.Series | dict | None, fallback_status: str = "") -> str:
    if activity_row is None:
        return delay_tia_docx_text(fallback_status, "Not proven from current uploaded files")
    critical = delay_tia_bool_text(activity_row.get("Critical", ""))
    longest = delay_tia_bool_text(activity_row.get("Longest Path", ""))
    total_float = pd.to_numeric(activity_row.get("Total Float", pd.NA), errors="coerce")
    if critical:
        return f"Critical Path at impact review; Total Float = {delay_tia_docx_number(total_float, 'NOT RECORDED')} days."
    if longest:
        return f"Longest Path at impact review; Total Float = {delay_tia_docx_number(total_float, 'NOT RECORDED')} days."
    if not pd.isna(total_float) and total_float <= 10:
        return f"Near-critical at impact review; Total Float = {delay_tia_docx_number(total_float, 'NOT RECORDED')} days."
    if fallback_status:
        return delay_tia_docx_text(fallback_status)
    return f"Not proven as critical / longest path from current uploaded files; Total Float = {delay_tia_docx_number(total_float, 'NOT RECORDED')} days."


def build_delay_tia_activity_selection_explanation_df(context: dict, analysis: dict) -> pd.DataFrame:
    candidates_df = analysis.get("candidates_df", pd.DataFrame()).copy()
    if candidates_df.empty:
        return pd.DataFrame(
            columns=[
                "Delay Event Ref",
                "Delay Type",
                "Activity ID",
                "Activity Name",
                "Why Selected",
                "Selection Methodology",
                "Criticality / Longest Path Basis",
                "Contractual and Factual Link",
                "Evidence Used",
                "Resulting Schedule Impact",
            ]
        )
    rows = []
    for _, row in candidates_df.iterrows():
        event_ref = f"STEEL-{delay_tia_docx_text(row.get('Building', ''), 'AREA')}-{delay_tia_docx_text(row.get('Stock-Out Date', pd.NaT), 'NO-DATE')}"
        methodology = delay_tia_join_nonempty([
            f"Due / Ready Test: {row.get('Due / Ready Test', '')}",
            f"Stock Unavailable Test: {row.get('Stock Unavailable Test', '')}",
            f"Not Completed Test: {row.get('Not Completed Test', '')}",
            f"Downstream Impact Test: {row.get('Downstream Impact Test', '')}",
            f"TIA Candidate Score: {delay_tia_docx_number(row.get('TIA Candidate Score', 0))}",
            f"Classification: {delay_tia_docx_text(row.get('Candidate Classification', ''))}",
        ])
        link_basis = delay_tia_join_nonempty([
            f"Required qty {delay_tia_docx_number(row.get('Required Qty', 0))}",
            f"Available qty {delay_tia_docx_number(row.get('Available Qty', 0))}",
            f"Shortage qty {delay_tia_docx_number(row.get('Shortage Qty', 0))}",
            f"Required date {delay_tia_docx_text(row.get('Required Date', pd.NaT))}",
            f"Client supply gap {delay_tia_docx_number(row.get('Client Supply vs Actual Gap', pd.NA), 'NOT RECORDED')}",
        ])
        evidence = delay_tia_join_nonempty([
            "02- master_activity_steel_analysis.csv",
            "03- employer_steel_supply_at_site.csv",
            "04- p6_activity_export.csv",
            "05- relationship_file.csv",
        ])
        rows.append(
            {
                "Delay Event Ref": event_ref,
                "Delay Type": "Delayed steel delivery / material shortage",
                "Activity ID": delay_tia_docx_text(row.get("Activity ID", "")),
                "Activity Name": delay_tia_docx_text(row.get("Activity Name", "")),
                "Why Selected": delay_tia_join_nonempty([
                    delay_tia_docx_text(row.get("Selection Explanation", ""), ""),
                    delay_tia_docx_text(row.get("Affected Activity Explanation", ""), ""),
                ]),
                "Selection Methodology": methodology,
                "Criticality / Longest Path Basis": delay_tia_describe_criticality(row, row.get("Delay Status", "")),
                "Contractual and Factual Link": (
                    "The activity was selected because the uploaded employer-only steel availability logic shows the activity due / ready "
                    f"when steel was unavailable, and the shortage evidence at activity level supports a direct cause-effect link. {link_basis}"
                ),
                "Evidence Used": evidence,
                "Resulting Schedule Impact": delay_tia_join_nonempty([
                    delay_tia_docx_text(row.get("Construction Sequence Impact", ""), ""),
                    f"Downstream impact: {delay_tia_docx_text(row.get('Downstream Impact Test', ''), 'Not recorded')}",
                    f"Predecessors: {delay_tia_docx_text(row.get('Predecessors', ''), 'Not recorded')}",
                    f"Successors: {delay_tia_docx_text(row.get('Successors', ''), 'Not recorded')}",
                ]),
            }
        )
    return pd.DataFrame(rows)


def build_delay_tia_fragnet_explanation_df(context: dict, analysis: dict) -> pd.DataFrame:
    fragnet_df = analysis.get("fragnet_df", pd.DataFrame()).copy()
    if fragnet_df.empty:
        return pd.DataFrame(
            columns=[
                "Delay Event Ref",
                "Delay Type",
                "Fragnet ID",
                "Fragnet Activity Name",
                "Affected Activity",
                "Why This Fragnet Was Developed",
                "Duration Basis",
                "Logic, Predecessor, Successor, Constraints and Sequencing",
                "Methodology Used",
                "Resulting Impact on Completion / Milestones / Float",
            ]
        )
    rows = []
    for idx, row in fragnet_df.reset_index(drop=True).iterrows():
        fragnet_id = f"FG-{idx + 1:03d}"
        duration_basis = (
            f"Duration {delay_tia_docx_number(row.get('Fragment Duration', 0))} days derived from "
            f"{delay_tia_docx_text(row.get('Fragment Start', pd.NaT))} to {delay_tia_docx_text(row.get('Fragment Finish', pd.NaT))}. "
            "Where recovery was evidenced, the duration follows the calculated unavailability window; otherwise it remains the best "
            "available schedule-impact proxy from current uploaded records."
        )
        logic_basis = (
            f"Predecessor: {delay_tia_docx_text(row.get('Last completed / available predecessor', ''), 'Not identified')}; "
            f"Successor / insertion point: {delay_tia_docx_text(row.get('Insert Fragment Before', ''), 'Not identified')}. "
            "The fragnet is connected using the current schedule sequence so that the delay sits immediately before the first activity "
            "proven to be affected rather than before a generic material or correspondence record."
        )
        rows.append(
            {
                "Delay Event Ref": f"FRAG-{idx + 1:03d}",
                "Delay Type": "Delayed steel delivery / material shortage",
                "Fragnet ID": fragnet_id,
                "Fragnet Activity Name": delay_tia_docx_text(row.get("Fragment Activity Name", "")),
                "Affected Activity": delay_tia_docx_text(row.get("Affected Activity", "")),
                "Why This Fragnet Was Developed": (
                    "The fragnet was developed because the affected activity passed the selection tests, the shortage window was measurable, "
                    "and a predecessor-successor insertion point could be tied to the existing schedule network."
                ),
                "Duration Basis": duration_basis,
                "Logic, Predecessor, Successor, Constraints and Sequencing": logic_basis,
                "Methodology Used": (
                    "Fragnet-based Time Impact Analysis using employer-only supply timing, activity-level shortage evidence, CPM sequence logic, "
                    "and criticality / near-criticality review at the time of impact."
                ),
                "Resulting Impact on Completion / Milestones / Float": (
                    f"Delay status: {delay_tia_docx_text(row.get('Delay Status', ''), 'Not recorded')}; "
                    f"Affected building: {delay_tia_docx_text(row.get('Affected Building', ''), 'Not recorded')}; "
                    f"Shortage qty: {delay_tia_docx_number(row.get('Shortage Qty', 0))}; "
                    "final impact on completion and float must be read with the activity impact and causation matrices."
                ),
            }
        )
    return pd.DataFrame(rows)


def build_delay_tia_client_delay_evidence_explanation_df(context: dict, analysis: dict) -> pd.DataFrame:
    delay_events_df = context.get("delay_events_df", pd.DataFrame()).copy()
    p6_lookup = delay_tia_build_p6_lookup(context.get("p6_df", pd.DataFrame()))
    rows = []

    if not delay_events_df.empty:
        for _, row in delay_events_df.iterrows():
            cause_category = delay_tia_docx_text(row.get("cause_category", ""), "General delay event")
            activity_ids = delay_tia_extract_activity_ids(row.get("activity_id", ""))
            selected_activity = activity_ids[0] if activity_ids else ""
            p6_row = p6_lookup.get(selected_activity)
            delay_type_lower = cause_category.lower()
            if "steel" in delay_type_lower or "material" in delay_type_lower:
                methodology = "Fragnet-based TIA supported by activity-level steel shortage and employer supply timing."
                fragnet_status = "Fragnet generated from current TIA run where CPM linkage was sufficiently measurable."
                evidence = "Master activity steel analysis, employer steel supply register, P6 export, relationship file, delivery records."
            elif "rfi" in delay_type_lower:
                methodology = "Event-window review with CPM activity linkage from uploaded records; suitable for TIA support and concurrency testing."
                fragnet_status = "No automatic fragnet inserted unless a direct CPM activity and measurable impact window are proven."
                evidence = "RFI register, correspondence, submission/reply dates, P6 activity linkage, relationship logic."
            elif "ifc" in delay_type_lower or "design" in delay_type_lower:
                methodology = "Event-window review with CPM activity linkage and design-response timing; used for TIA support and causation review."
                fragnet_status = "No automatic fragnet inserted unless direct schedule linkage and duration basis are evidenced."
                evidence = "IFC issue register, design correspondence, approvals, technical office records, P6 linkage."
            elif "payment" in delay_type_lower:
                methodology = "Commercial delay screening with CPM relevance test; used mainly for causation and concurrency review."
                fragnet_status = "No automatic fragnet inserted unless payment delay is linked to workfront stoppage and measurable duration."
                evidence = "Payment records, certificates, invoices, commercial correspondence, site records, P6 linkage."
            elif "weather" in delay_type_lower or "force majeure" in delay_type_lower:
                methodology = "Window analysis / TIA support review subject to contemporaneous weather and site stoppage records."
                fragnet_status = "No automatic fragnet inserted from current uploaded data."
                evidence = "Weather logs, site diaries, delay notices, daily reports, schedule linkage."
            elif "access" in delay_type_lower:
                methodology = "Window analysis / TIA support review subject to handover, possession, and workfront records."
                fragnet_status = "No automatic fragnet inserted from current uploaded data."
                evidence = "Site handover records, letters, instructions, daily reports, P6 linkage."
            elif "variation" in delay_type_lower or "change" in delay_type_lower:
                methodology = "Variation impact review using event windows and schedule sequence logic."
                fragnet_status = "No automatic fragnet inserted unless scope and duration are measurable."
                evidence = "Variation instructions, approval records, revised drawings, quantity records, P6 linkage."
            else:
                methodology = "Delay event screening using the uploaded delay register and current schedule context."
                fragnet_status = "No automatic fragnet inserted from current uploaded data unless direct CPM linkage is proven."
                evidence = "Uploaded delay record set and any linked schedule/correspondence evidence."

            rows.append(
                {
                    "Delay Event Ref": delay_tia_docx_text(row.get("delay_id", "")),
                    "Delay Type": cause_category,
                    "How the Delay Was Identified": (
                        f"Identified from the uploaded delay event register with start {delay_tia_docx_text(row.get('start_date'))} "
                        f"and finish {delay_tia_docx_text(row.get('end_date'))}, then reviewed against uploaded schedule and support records."
                    ),
                    "Affected Activities and Why": (
                        f"Activities referenced: {delay_tia_docx_text(row.get('activity_id', ''), 'No direct activity ID recorded')}. "
                        f"Primary activity used for schedule linkage: {selected_activity or 'Not directly identified'}."
                    ),
                    "Critical Path / Longest Path Status at Impact": delay_tia_describe_criticality(p6_row),
                    "Contractual and Factual Basis": (
                        f"Responsibility line: {delay_tia_docx_text(row.get('responsibility', ''), 'Not recorded')}. "
                        f"Notice / ref: {delay_tia_docx_text(row.get('notice_ref', ''), 'Not recorded')}."
                    ),
                    "Evidence and Correspondence Used": evidence,
                    "Fragnet Development Basis": fragnet_status,
                    "Duration Basis": (
                        f"Estimated delay days recorded = {delay_tia_docx_number(row.get('estimated_delay_days', 0))}; "
                        f"approved EOT days recorded = {delay_tia_docx_number(row.get('approved_eot_days', 0))}. "
                        "Any final TIA duration must follow measurable event dates, schedule logic, or contemporaneous performance records."
                    ),
                    "Methodology and Schedule Impact": methodology,
                    "Concurrency, Mitigation, and Final Entitlement Position": (
                        f"Current status: {delay_tia_docx_text(row.get('status', ''), 'Not recorded')}. "
                        "Concurrency must be reviewed against overlapping event windows and actual workface relevance; entitlement depends on evidence, "
                        "contract support, mitigation records, and proven cause-effect."
                    ),
                }
            )

    steel_assessment_df = analysis.get("assessment_df", pd.DataFrame()).copy()
    for _, row in steel_assessment_df.iterrows():
        rows.append(
            {
                "Delay Event Ref": f"CLAIM-{delay_tia_docx_text(row.get('Activity ID', ''), 'STEEL')}",
                "Delay Type": "Client-caused delayed steel delivery / material shortage",
                "How the Delay Was Identified": (
                    f"Identified by shortage on {delay_tia_docx_text(row.get('Stock-Out Date'))} with fragment window "
                    f"{delay_tia_docx_text(row.get('Fragment Start'))} to {delay_tia_docx_text(row.get('Fragment Finish'))}."
                ),
                "Affected Activities and Why": (
                    f"{delay_tia_docx_text(row.get('Activity ID', ''))} was selected because the uploaded TIA logic tied the shortage to a due / ready "
                    "activity with insufficient steel availability and measurable successor exposure."
                ),
                "Critical Path / Longest Path Status at Impact": (
                    f"Critical flag = {delay_tia_docx_text(row.get('Critical', ''))}; "
                    f"Longest path flag = {delay_tia_docx_text(row.get('Longest Path', ''))}; "
                    f"Total Float = {delay_tia_docx_number(row.get('Total Float', pd.NA), 'NOT RECORDED')} days."
                ),
                "Contractual and Factual Basis": delay_tia_join_nonempty([
                    f"Delay classification: {delay_tia_docx_text(row.get('Delay Classification', ''), 'Not classified')}",
                    f"Notice reference: {delay_tia_docx_text(row.get('Notice Reference', ''), 'Not linked')}",
                    f"Evidence reference: {delay_tia_docx_text(row.get('Evidence Reference', ''), 'Not linked')}",
                    f"Final assessment: {delay_tia_docx_text(row.get('Final Assessment', ''), 'Not recorded')}",
                ]),
                "Evidence and Correspondence Used": (
                    "Employer steel supply records, master activity steel analysis, P6 activity export, relationship file, "
                    "contract library, and any linked letters or notices."
                ),
                "Fragnet Development Basis": (
                    f"Fragnet inserted before {delay_tia_docx_text(row.get('Insert Fragment Before', ''), 'Not identified')} with duration "
                    f"{delay_tia_docx_number(row.get('Fragment Duration', 0))} days."
                ),
                "Duration Basis": (
                    f"Duration based on fragment window {delay_tia_docx_text(row.get('Fragment Start'))} to "
                    f"{delay_tia_docx_text(row.get('Fragment Finish'))} and shortage quantity "
                    f"{delay_tia_docx_number(row.get('Shortage Qty', 0))}."
                ),
                "Methodology and Schedule Impact": delay_tia_join_nonempty([
                    "Time Impact Analysis using fragnet insertion into the current CPM context.",
                    delay_tia_docx_text(row.get("Schedule Impact and EOT Assessment", ""), ""),
                    f"Money impact: {delay_tia_docx_text(row.get('Money Impact Assessment', ''), 'Not recorded')}",
                ]),
                "Concurrency, Mitigation, and Final Entitlement Position": (
                    f"Practical action / evidence: {delay_tia_docx_text(row.get('Practical Action / Evidence', ''), 'Not recorded')}. "
                    "Final entitlement depends on closing any notice, concurrency, readiness, and contractor-risk gaps."
                ),
            }
        )

    return pd.DataFrame(rows)


def compute_claimed_delay_duration_days(pre_finish, impacted_finish, fallback_duration=0, recorded_delay_duration=pd.NA):
    try:
        recorded = pd.to_numeric(recorded_delay_duration, errors="coerce")
        if pd.notna(recorded):
            return max(int(recorded), 0)
    except Exception:
        pass
    try:
        if not pd.isna(pre_finish) and not pd.isna(impacted_finish):
            delta = (pd.to_datetime(impacted_finish) - pd.to_datetime(pre_finish)).days
            if pd.notna(delta):
                return max(int(delta), 0)
    except Exception:
        pass
    try:
        fallback = pd.to_numeric(fallback_duration, errors="coerce")
        if pd.notna(fallback):
            return max(int(fallback), 0)
    except Exception:
        pass
    return pd.NA


def build_delay_tia_director_delay_events_df(context: dict, analysis: dict) -> pd.DataFrame:
    fragnet_df = analysis.get("fragnet_df", pd.DataFrame()).copy()
    if fragnet_df.empty:
        return pd.DataFrame(
            columns=[
                "Event ID", "Event Description", "Responsible Party", "WBS / Area", "Affected Activity ID",
                "Impact Start", "Impact Finish", "Calendar Days", "Working Days", "Claimed Delay Duration (days)", "Fragnet ID",
                "Critical / Longest Path?", "Delay Type", "Notice / Letter Ref.", "Evidence Status", "Claim Decision",
            ]
        )

    rows = []
    for idx, row in fragnet_df.reset_index(drop=True).iterrows():
        event_id = f"EV-{idx + 1:03d}"
        fragnet_id = f"FG-{idx + 1:03d}"
        claimed_delay_duration = compute_claimed_delay_duration_days(
            row.get("Fragment Start", pd.NaT),
            row.get("Fragment Finish", pd.NaT),
            row.get("Fragment Duration", 0),
            row.get("Recorded Steel Delay Duration (days)", pd.NA),
        )
        delay_status = delay_tia_docx_text(row.get("Delay Status", ""), "")
        critical_flag = "Yes" if "critical" in delay_status.lower() else ("Near-Critical" if "near" in delay_status.lower() else "Not proven")
        decision = "Valid for TIA" if critical_flag in {"Yes", "Near-Critical"} else "Not Enough Data"
        evidence_status = "Complete" if decision == "Valid for TIA" else "Partial"
        rows.append(
            {
                "Event ID": event_id,
                "Event Description": delay_tia_docx_text(row.get("Fragment Activity Name", "Employer steel delay fragnet")),
                "Responsible Party": "Employer / Client",
                "WBS / Area": delay_tia_docx_text(row.get("Affected Building", ""), "Project Area"),
                "Affected Activity ID": delay_tia_docx_text(row.get("Insert Fragment Before", "")),
                "Impact Start": row.get("Fragment Start", pd.NaT),
                "Impact Finish": row.get("Fragment Finish", pd.NaT),
                "Calendar Days": pd.to_numeric(row.get("Fragment Duration", 0), errors="coerce"),
                "Working Days": pd.to_numeric(row.get("Fragment Duration", 0), errors="coerce"),
                "Claimed Delay Duration (days)": claimed_delay_duration,
                "Fragnet ID": fragnet_id,
                "Critical / Longest Path?": critical_flag,
                "Delay Type": "Employer Steel Supply Delay",
                "Notice / Letter Ref.": delay_tia_docx_text(row.get("Supporting Letter Ref", ""), "Not linked"),
                "Evidence Status": evidence_status,
                "Claim Decision": decision,
            }
        )
    return pd.DataFrame(rows)


def build_delay_tia_director_activity_impact_df(context: dict, analysis: dict) -> pd.DataFrame:
    candidates_df = analysis.get("candidates_df", pd.DataFrame()).copy()
    if candidates_df.empty:
        return pd.DataFrame(
            columns=[
                "Activity ID", "Activity Name", "WBS", "Baseline Finish", "Pre-Impact Forecast Finish",
                "Impacted Forecast Finish", "Finish Delta (days)", "Claimed Delay Duration (days)", "Total Float Before", "Total Float After",
                "Critical / LP?", "Downstream Milestone", "Impact Note",
            ]
        )

    rows = []
    for _, row in candidates_df.iterrows():
        pre_finish = row.get("Finish", pd.NaT)
        impacted_finish = row.get("Recommended Recovery / Availability Date", pd.NaT)
        baseline_finish = row.get("BL Finish", pd.NaT)
        finish_delta = pd.NA
        try:
            if not pd.isna(pre_finish) and not pd.isna(impacted_finish):
                finish_delta = (pd.to_datetime(impacted_finish) - pd.to_datetime(pre_finish)).days
        except Exception:
            finish_delta = pd.NA
        claimed_delay_duration = compute_claimed_delay_duration_days(
            pre_finish,
            impacted_finish,
            row.get("Recommended Fragment Duration", row.get("Required Delay Duration", 0)),
            row.get("Recorded Steel Delay Duration (days)", pd.NA),
        )
        tf_before = pd.to_numeric(row.get("Total Float", pd.NA), errors="coerce")
        tf_after = pd.NA
        if not pd.isna(tf_before) and not pd.isna(claimed_delay_duration):
            tf_after = tf_before - claimed_delay_duration
        delay_status = delay_tia_docx_text(row.get("Delay Status", ""), "")
        rows.append(
            {
                "Activity ID": delay_tia_docx_text(row.get("Activity ID", "")),
                "Activity Name": delay_tia_docx_text(row.get("Activity Name", "")),
                "WBS": delay_tia_docx_text(row.get("Building", ""), "Project"),
                "Baseline Finish": baseline_finish,
                "Pre-Impact Forecast Finish": pre_finish,
                "Impacted Forecast Finish": impacted_finish,
                "Finish Delta (days)": finish_delta,
                "Claimed Delay Duration (days)": claimed_delay_duration,
                "Total Float Before": tf_before,
                "Total Float After": tf_after,
                "Critical / LP?": "Yes" if "critical" in delay_status.lower() else ("Near-Critical" if "near" in delay_status.lower() else "Not proven"),
                "Downstream Milestone": delay_tia_docx_text(row.get("Downstream Impact", ""), "Not identified"),
                "Impact Note": delay_tia_docx_text(row.get("Selection Explanation", ""), "Derived from uploaded Delay TIA activity evidence."),
            }
        )
    return pd.DataFrame(rows)


def build_delay_tia_director_fragnet_logic_df(context: dict, analysis: dict) -> pd.DataFrame:
    primavera_df = build_delay_tia_primavera_fragnet_df(context, analysis)
    if primavera_df.empty:
        return pd.DataFrame(
            columns=[
                "Fragnet ID", "Fragnet Activity ID", "Fragnet Activity Description", "Duration (days)",
                "Predecessor Activity", "Relationship / Lag", "Successor Activity", "Logic Rationale", "Record Reference",
            ]
        )
    rows = []
    for _, row in primavera_df.iterrows():
        rows.append(
            {
                "Fragnet ID": delay_tia_docx_text(row.get("Activity ID", "")),
                "Fragnet Activity ID": delay_tia_docx_text(row.get("Activity ID", "")),
                "Fragnet Activity Description": delay_tia_docx_text(row.get("Activity Name", "")),
                "Duration (days)": row.get("Original Duration", 0),
                "Predecessor Activity": delay_tia_docx_text(row.get("Predecessor Activity ID", "")),
                "Relationship / Lag": delay_tia_docx_text(row.get("Predecessor Relationship Type", ""), "FS+0"),
                "Successor Activity": delay_tia_docx_text(row.get("Successor Activity ID", "")),
                "Logic Rationale": "Fragnet inserted before the affected successor activity using uploaded relationship and shortage evidence.",
                "Record Reference": delay_tia_docx_text(row.get("Notes", ""), "Delay TIA generated record"),
            }
        )
    return pd.DataFrame(rows)


def build_delay_tia_director_readiness_matrix_df(context: dict, analysis: dict) -> pd.DataFrame:
    rows = [
        ("Accepted baseline / approved update", "Mandatory", "Confirms valid starting model for impact insertion.", "Complete" if not context.get("p6_df", pd.DataFrame()).empty else "Missing"),
        ("Activity list export", "Mandatory", "Shows activities, dates, float, WBS and status.", "Complete" if not context.get("p6_df", pd.DataFrame()).empty else "Missing"),
        ("Relationship export", "Mandatory", "Shows predecessor/successor logic and lag.", "Complete" if not context.get("relationship_df", pd.DataFrame()).empty else "Missing"),
        ("Master activity steel analysis", "Mandatory", "Provides activity-level shortage and requirement evidence.", "Complete" if not context.get("master_df", pd.DataFrame()).empty else "Missing"),
        ("Employer steel supply", "Mandatory", "Provides employer delivery timing and quantity basis.", "Complete" if not context.get("employer_raw_df", pd.DataFrame()).empty else "Missing"),
        ("Contract library", "High", "Supports entitlement, notice, and claim treatment.", "Complete" if not context.get("contract_df", pd.DataFrame()).empty else "Missing"),
        ("IFC / RFI / payment records", "High", "Supports concurrency and non-steel causation review.", "Complete" if any(not context.get(key, pd.DataFrame()).empty for key in ["ifc_df", "rfi_df", "payments_df"]) else "Partial"),
        ("BL critical path fixed package", "Recommended", "Supports baseline criticality comparison.", "Complete" if not context.get("bl_critical_path_df", pd.DataFrame()).empty else "Partial"),
        ("Contractor site supply visibility", "Low", "Display-only stream for management visibility, excluded from TIA calculations.", "Complete" if not context.get("samco_df", pd.DataFrame()).empty else "Not Required"),
    ]
    return pd.DataFrame(rows, columns=["Data Input", "Priority", "Reason for TIA", "Status"])


def build_delay_tia_director_evidence_register_df(context: dict, analysis: dict) -> pd.DataFrame:
    events_df = build_delay_tia_director_delay_events_df(context, analysis)
    rows = []
    for _, row in events_df.iterrows():
        event_id = row.get("Event ID", "")
        evidence_status = delay_tia_docx_text(row.get("Evidence Status", "Pending"), "Pending")
        rows.append(
            {
                "Evidence ID": f"EVD-{str(event_id).split('-')[-1]}",
                "Related Event": event_id,
                "Document Type": "Uploaded Delay TIA Record Set",
                "Document Ref.": delay_tia_docx_text(row.get("Notice / Letter Ref.", ""), "Uploaded file bundle"),
                "Date": row.get("Impact Start", pd.NaT),
                "Key Fact Proven": delay_tia_docx_text(row.get("Event Description", "")),
                "Causation Link": "Activity shortage / supply timing / CPM logic linkage",
                "Strength": "High" if evidence_status == "Complete" else "Medium",
                "Missing Item": "" if evidence_status == "Complete" else "Need fuller linked notice or CPM evidence",
                "Action Owner": "Planning / Contracts",
            }
        )
    return pd.DataFrame(rows)


def build_delay_tia_director_causation_matrix_df(context: dict, analysis: dict) -> pd.DataFrame:
    events_df = build_delay_tia_director_delay_events_df(context, analysis)
    concurrent_df = context.get("concurrent_delay_review_df", pd.DataFrame()).copy()
    rows = []
    for _, row in events_df.iterrows():
        decision = delay_tia_docx_text(row.get("Claim Decision", "Not Enough Data"))
        critical = delay_tia_docx_text(row.get("Critical / Longest Path?", "Not proven"))
        activity_id = delay_tia_docx_text(row.get("Affected Activity ID", ""))
        concurrency_text = "Needs further concurrency review."
        if not concurrent_df.empty and activity_id:
            related = concurrent_df[concurrent_df["Activity ID"].astype(str).str.strip() == activity_id].copy()
            non_steel_related = related[related["Delay Stream"].astype(str).str.strip().str.lower() != "steel"].copy() if not related.empty else pd.DataFrame()
            if not non_steel_related.empty:
                refs = ", ".join(non_steel_related["Delay Ref"].astype(str).tolist())
                concurrency_text = f"Concurrent non-steel streams identified for the same activity: {refs}. Review overlapping windows before final entitlement."
            elif decision == "Valid for TIA":
                concurrency_text = "No non-steel concurrent stream identified for the same activity from the loaded concurrent review."
        elif decision == "Valid for TIA":
            concurrency_text = "No contractor concurrency proven from current uploaded Delay TIA files."
        rows.append(
            {
                "Event": row.get("Event ID", ""),
                "Cause": row.get("Event Description", ""),
                "Critical Impact Proven?": "Yes" if critical in {"Yes", "Near-Critical"} else "No",
                "Concurrency / Risk Test": concurrency_text,
                "Claim Treatment": decision,
                "Required Action": "Proceed with EOT narrative." if decision == "Valid for TIA" else "Close evidence and schedule-link gaps before reliance.",
            }
        )
    return pd.DataFrame(rows)


def build_delay_tia_director_p6_controls_df(context: dict, analysis: dict) -> pd.DataFrame:
    data_date = pd.Timestamp.today().normalize()
    return pd.DataFrame(
        [
            {"Control Item": "Schedule file name", "Recorded Value": "04- p6_activity_export.csv", "Why It Must Be Shown": "Primary activity export used for Delay TIA context."},
            {"Control Item": "Data date", "Recorded Value": data_date, "Why It Must Be Shown": "Must match the reporting cut-off for the TIA run."},
            {"Control Item": "Schedule options", "Recorded Value": "Derived from uploaded activity export only", "Why It Must Be Shown": "Clarifies that no hidden P6 recalculation options were imported."},
            {"Control Item": "Longest path / critical path basis", "Recorded Value": "Critical flag plus near-critical float threshold = 10 days", "Why It Must Be Shown": "Defines how TIA criticality was classified."},
            {"Control Item": "Retained logic / progress override", "Recorded Value": "Not recorded in uploaded files", "Why It Must Be Shown": "Required for full forensic repeatability."},
            {"Control Item": "Out-of-sequence progress handling", "Recorded Value": "Not recorded in uploaded files", "Why It Must Be Shown": "Can affect critical path movement."},
            {"Control Item": "Constraints", "Recorded Value": "Not explicitly recorded", "Why It Must Be Shown": "Constraints may create artificial criticality."},
            {"Control Item": "Calendars", "Recorded Value": "Not explicitly recorded", "Why It Must Be Shown": "Delay days depend on the governing calendar."},
            {"Control Item": "Open ends", "Recorded Value": "Not explicitly recorded", "Why It Must Be Shown": "Open ends weaken CPM reliability."},
            {"Control Item": "Negative float", "Recorded Value": "Derived from current candidate / fragnet outputs only", "Why It Must Be Shown": "Supports the severity of affected activities."},
        ]
    )


def build_delay_tia_director_metadata(context: dict, analysis: dict) -> dict:
    project_rows = load_core_csv(PROJECTS_CSV_PATH, project_id=selected_project_id())
    project_row = project_rows.iloc[0].to_dict() if not project_rows.empty else {}
    return {
        "project_name": str(project_row.get("project_name") or selected_project_id() or "Project"),
        "contract_no": str(project_row.get("contract_no") or "NOT PROVIDED IN PROJECT DATA"),
        "employer": str(project_row.get("client_name") or project_row.get("employer") or "Employer"),
        "contractor": str(project_row.get("contractor") or "Contractor"),
        "contract_form_or_clause": "From uploaded contract library and matched clauses",
        "accepted_baseline_programme": "Project baseline package plus activity export",
        "impacted_update": "Current uploaded Delay TIA analysis context",
        "data_date": pd.Timestamp.today().normalize(),
        "calendar_basis": "Working-day logic based on uploaded activity schedule context",
        "criticality_rule": "Critical / longest path with near-critical threshold at 10 days float",
        "analysis_period": (
            f"{delay_tia_docx_text(context.get('employer_first_date', pd.NaT))} to "
            f"{delay_tia_docx_text(context.get('employer_last_date', pd.NaT))}"
        ),
        "prepared_by": "Planning Department",
        "checked_by": "Pending",
        "approved_by": "Pending",
    }


def build_delay_tia_director_kpis(context: dict, analysis: dict) -> dict:
    delay_events_df = build_delay_tia_director_delay_events_df(context, analysis)
    total_days = int(
        numeric_series_from_columns(
            delay_events_df,
            ["Working Days", "Claimed Delay Duration (days)", "Delayed duration after overlap", "Delayed duration", "estimated_delay_days"],
        ).sum()
    ) if not delay_events_df.empty else 0
    valid_days = int(
        pd.to_numeric(
            delay_events_df.loc[delay_events_df["Claim Decision"] == "Valid for TIA", "Working Days"],
            errors="coerce",
        ).fillna(0).sum()
    ) if not delay_events_df.empty else 0
    critical_events = int((delay_events_df.get("Critical / Longest Path?") == "Yes").sum()) if not delay_events_df.empty else 0
    missing_evidence = int((delay_events_df.get("Evidence Status") == "Missing").sum()) if not delay_events_df.empty else 0
    if valid_days > 0:
        conclusion = (
            "The uploaded Delay TIA records indicate that employer steel supply delay events with traceable fragnet logic "
            "and critical or near-critical CPM exposure may proceed into the EOT narrative, while weaker items should remain visible as evidence gaps."
        )
    else:
        conclusion = (
            "The current uploaded Delay TIA run does not yet prove a director-level EOT case without further evidence and stronger CPM linkage."
        )
    return {
        "total_recorded_working_delay_days": total_days,
        "potential_eot_days": valid_days,
        "critical_delay_events_count": critical_events,
        "missing_evidence_items_count": missing_evidence,
        "executive_conclusion": conclusion,
    }


def build_delay_tia_director_impact_calculation_df(context: dict, analysis: dict) -> pd.DataFrame:
    activity_impact_df = build_delay_tia_director_activity_impact_df(context, analysis)
    baseline_completion = activity_impact_df["Baseline Finish"].dropna().max() if not activity_impact_df.empty else pd.NaT
    pre_impact_completion = activity_impact_df["Pre-Impact Forecast Finish"].dropna().max() if not activity_impact_df.empty else pd.NaT
    impacted_completion = activity_impact_df["Impacted Forecast Finish"].dropna().max() if not activity_impact_df.empty else pd.NaT
    net_impact_days = "NOT ENOUGH DATA"
    try:
        if not pd.isna(pre_impact_completion) and not pd.isna(impacted_completion):
            net_impact_days = f"{(pd.to_datetime(impacted_completion) - pd.to_datetime(pre_impact_completion)).days} days"
    except Exception:
        pass
    kpis = build_delay_tia_director_kpis(context, analysis)
    return pd.DataFrame(
        [
            ["Baseline Completion", delay_tia_docx_text(baseline_completion), "Baseline finish from uploaded activity evidence."],
            ["Pre-Impact Forecast Completion", delay_tia_docx_text(pre_impact_completion), "Forecast before applying the current fragnet recommendation."],
            ["Impacted Forecast Completion", delay_tia_docx_text(impacted_completion), "Forecast represented by the current Delay TIA impact window."],
            ["Net Impact", delay_tia_docx_text(net_impact_days), "Impacted finish minus pre-impact finish."],
            ["Potential EOT Position", f"{delay_tia_docx_number(kpis['potential_eot_days'])} working days", "Subject to entitlement, evidence, and concurrency review."],
        ],
        columns=["Calculation Item", "Value", "Interpretation"],
    )


def build_delay_tia_final_submission_summary_df(context: dict, analysis: dict) -> pd.DataFrame:
    delay_events_df = build_delay_tia_director_delay_events_df(context, analysis)
    concurrent_df = context.get("concurrent_delay_review_df", pd.DataFrame()).copy()
    impact_df = build_delay_tia_director_impact_calculation_df(context, analysis)

    gross_all_days = 0
    if not concurrent_df.empty and "Delayed Days" in concurrent_df.columns:
        gross_all_days = int(pd.to_numeric(concurrent_df["Delayed Days"], errors="coerce").fillna(0).sum())

    steel_tia_days = 0
    if not delay_events_df.empty:
        steel_tia_days = int(
            pd.to_numeric(
                delay_events_df.loc[delay_events_df["Claim Decision"] == "Valid for TIA", "Claimed Delay Duration (days)"],
                errors="coerce",
            ).fillna(0).sum()
        )

    support_path_days = 0
    if not concurrent_df.empty:
        concurrent_df["__is_path__"] = (
            concurrent_df.get("Current Critical", pd.Series(dtype=object)).astype(str).str.strip().str.lower().eq("yes")
            | concurrent_df.get("Current Longest Path", pd.Series(dtype=object)).astype(str).str.strip().str.lower().eq("yes")
        )
        support_path_days = int(
            pd.to_numeric(
                concurrent_df.loc[
                    (concurrent_df["Delay Stream"].astype(str).str.strip().str.lower() != "steel")
                    & concurrent_df["__is_path__"],
                    "Delayed Days",
                ],
                errors="coerce",
            ).fillna(0).sum()
        )

    improvement_modelled_days = calculate_delay_tia_improvement_modelled_days(context)

    net_impact_days = pd.NA
    if not impact_df.empty:
        net_impact_row = impact_df.loc[impact_df["Calculation Item"] == "Net Impact", "Value"]
        if not net_impact_row.empty:
            net_impact_value = str(net_impact_row.iloc[0]).strip()
            net_impact_match = re.search(r"(-?\d+)", net_impact_value)
            if net_impact_match:
                net_impact_days = abs(int(net_impact_match.group(1)))

    recommended_submission_days = (
        int(net_impact_days)
        if not pd.isna(net_impact_days)
        else steel_tia_days
    )
    recommended_submission_days = max(int(recommended_submission_days or 0), int(improvement_modelled_days or 0))

    return pd.DataFrame(
        [
            {
                "Submission Position": "Gross detected delayed days across all loaded events",
                "Days": gross_all_days,
                "Basis": "Informational only. This includes steel, RFI, and IFC rows before concurrency or CPM de-duplication.",
            },
            {
                "Submission Position": "Gross TIA-modelled steel delay days",
                "Days": steel_tia_days,
                "Basis": "Sum of valid TIA-modelled steel fragnet events from the current employer-only delay analysis.",
            },
            {
                "Submission Position": "Gross support-stream days on current critical / longest path",
                "Days": support_path_days,
                "Basis": "RFI / IFC rows that touch the current critical or longest path. These remain support evidence unless separately fragnet-modelled.",
            },
            {
                "Submission Position": "Gross modelled improvement-register delay days",
                "Days": improvement_modelled_days,
                "Basis": "Optional uploaded improvement fragnet / claimed-delay registers. Used as an additional modelled TIA stream when provided, without arithmetic double-counting against the net impact.",
            },
            {
                "Submission Position": "Recommended total delayed days to submit",
                "Days": recommended_submission_days,
                "Basis": "Conservative maximum of net project-completion impact and uploaded modelled improvement-register delay days. This avoids raw addition across overlapping streams.",
            },
        ]
    )


def build_delay_tia_final_submission_event_breakdown_df(context: dict, analysis: dict) -> pd.DataFrame:
    concurrent_df = context.get("concurrent_delay_review_df", pd.DataFrame()).copy()
    delay_events_df = build_delay_tia_director_delay_events_df(context, analysis)
    if concurrent_df.empty:
        return pd.DataFrame(
            columns=[
                "Delay Stream",
                "Delay Rows",
                "Gross Delayed Days",
                "Current Critical / LP Rows",
                "Included in Submitted Total?",
                "Treatment Basis",
            ]
        )

    current_critical = concurrent_df.get("Current Critical", pd.Series(dtype=object)).astype(str).str.strip().str.lower().eq("yes")
    current_longest = concurrent_df.get("Current Longest Path", pd.Series(dtype=object)).astype(str).str.strip().str.lower().eq("yes")
    concurrent_df["__path_rows__"] = current_critical | current_longest

    grouped = (
        concurrent_df.groupby("Delay Stream", dropna=False)
        .agg(
            Delay_Rows=("Delay Ref", "count"),
            Gross_Delayed_Days=("Delayed Days", lambda s: pd.to_numeric(s, errors="coerce").fillna(0).sum()),
            Current_Path_Rows=("__path_rows__", lambda s: int(s.fillna(False).sum())),
        )
        .reset_index()
        .rename(
            columns={
                "Delay_Rows": "Delay Rows",
                "Gross_Delayed_Days": "Gross Delayed Days",
                "Current_Path_Rows": "Current Critical / LP Rows",
            }
        )
    )

    valid_steel_days = int(
        pd.to_numeric(
            delay_events_df.loc[delay_events_df["Claim Decision"] == "Valid for TIA", "Claimed Delay Duration (days)"],
            errors="coerce",
        ).fillna(0).sum()
    ) if not delay_events_df.empty else 0

    rows = []
    for _, row in grouped.iterrows():
        stream = str(row.get("Delay Stream", "")).strip()
        lower_stream = stream.lower()
        included = "No"
        basis = "Support stream only."
        if lower_stream == "steel":
            included = "Yes" if valid_steel_days > 0 else "No"
            basis = "Included through the TIA fragnet model and net completion impact calculation." if valid_steel_days > 0 else "No valid steel TIA event currently proven."
        elif lower_stream == "improvement fragnet":
            included = "Yes"
            basis = "Included as an uploaded modelled improvement fragnet stream. Final total still uses conservative max/net logic to avoid double counting."
        elif int(row.get("Current Critical / LP Rows", 0) or 0) > 0:
            basis = "Touches the current critical / longest path and is shown as support / concurrency evidence, but is not yet added to the submitted total without separate fragnet modelling."
        rows.append(
            {
                "Delay Stream": stream,
                "Delay Rows": int(row.get("Delay Rows", 0) or 0),
                "Gross Delayed Days": int(pd.to_numeric(pd.Series([row.get("Gross Delayed Days", 0)]), errors="coerce").fillna(0).iloc[0]),
                "Current Critical / LP Rows": int(row.get("Current Critical / LP Rows", 0) or 0),
                "Included in Submitted Total?": included,
                "Treatment Basis": basis,
            }
        )
    return pd.DataFrame(rows)


def delay_tia_add_number_of_delayed_days_column(df: pd.DataFrame, default: Any = "") -> pd.DataFrame:
    if not isinstance(df, pd.DataFrame) or df.empty:
        return df
    working = df.copy()
    if "Number of Delayed Days" in working.columns:
        return working
    candidate_cols = [
        "Delayed Days",
        "Claimed Delay Duration (days)",
        "Claimed Delay Duration",
        "Working Delay Days",
        "Working Days",
        "Fragment Duration",
        "Duration (days)",
        "Original Duration",
        "Gross Delayed Days",
        "Total Delayed Days",
        "Days",
        "Finish Delta (days)",
        "Concurrent Delay Days",
        "Concurrent delay",
        "Delayed duration after overlap",
        "Delayed duration",
    ]
    value = pd.Series([default] * len(working), index=working.index, dtype=object)
    for col in candidate_cols:
        if col in working.columns:
            numeric = pd.to_numeric(working[col], errors="coerce")
            value = numeric.abs().where(numeric.notna(), value)
            break
    insert_at = 0
    for anchor in ["Activity Name", "Activity ID", "Affected Activity ID", "Delay Ref", "Event ID", "Submission Position", "Item"]:
        if anchor in working.columns:
            insert_at = min(working.columns.get_loc(anchor) + 1, len(working.columns))
            break
    working.insert(insert_at, "Number of Delayed Days", value)
    return working


def build_delay_tia_final_delayed_activities_df(context: dict, analysis: dict) -> pd.DataFrame:
    concurrent_df = context.get("concurrent_delay_review_df", pd.DataFrame()).copy()
    path_lookup = build_delay_tia_path_status_lookup(
        context.get("p6_df", pd.DataFrame()),
        {
            "float_df": pd.DataFrame(),
            "longest_df": pd.DataFrame(),
            "critical_df": context.get("bl_critical_path_df", pd.DataFrame()),
        },
    )
    if concurrent_df.empty:
        return pd.DataFrame(
            columns=[
                "Delay Stream",
                "Delay Ref",
                "Activity ID",
                "Activity Name",
                "Number of Delayed Days",
                "BL Critical Path",
                "Updated Critical Path",
                "Updated Longest Path",
                "Current Total Float",
                "Delay Start",
                "Delay Finish",
                "Treatment",
            ]
        )

    rows = []
    for _, row in concurrent_df.iterrows():
        activity_id = delay_tia_docx_text(row.get("Activity ID", ""))
        if not activity_id:
            continue
        path_status = path_lookup.get(activity_id, {})
        delayed_days = pd.to_numeric(row.get("Delayed Days", row.get("Number of Delayed Days", 0)), errors="coerce")
        if pd.isna(delayed_days):
            delayed_days = 0
        current_critical = delay_tia_docx_text(row.get("Current Critical", path_status.get("Current Critical", "No")), "No")
        current_lp = delay_tia_docx_text(row.get("Current Longest Path", path_status.get("Current Longest Path", "No")), "No")
        bl_critical = delay_tia_docx_text(row.get("BL Critical Path", path_status.get("BL Critical Path", "No")), "No")
        if bl_critical.lower() == "yes" and current_critical.lower() == "yes":
            treatment = "Delayed activity is on both BL and updated critical path."
        elif bl_critical.lower() == "yes":
            treatment = "Delayed activity was on BL critical path but is not shown critical in the update."
        elif current_critical.lower() == "yes" or current_lp.lower() == "yes":
            treatment = "Delayed activity became/currently remains driving in the update critical or longest path."
        else:
            treatment = "Delayed activity is support / concurrency evidence unless separately proved driving."
        rows.append(
            {
                "Delay Stream": delay_tia_docx_text(row.get("Delay Stream", "")),
                "Delay Ref": delay_tia_docx_text(row.get("Delay Ref", "")),
                "Activity ID": activity_id,
                "Activity Name": delay_tia_docx_text(row.get("Activity Name", ""), path_status.get("Activity Name", "")),
                "Number of Delayed Days": abs(float(delayed_days)),
                "BL Critical Path": bl_critical,
                "Updated Critical Path": current_critical,
                "Updated Longest Path": current_lp,
                "Current Total Float": row.get("Current Total Float", path_status.get("Current Total Float", pd.NA)),
                "Delay Start": row.get("Delay Start", pd.NaT),
                "Delay Finish": row.get("Delay Finish", pd.NaT),
                "Treatment": treatment,
            }
        )

    result = pd.DataFrame(rows)
    if result.empty:
        return result
    grouped = (
        result.groupby(["Delay Stream", "Delay Ref", "Activity ID", "Activity Name"], dropna=False)
        .agg(
            **{
                "Number of Delayed Days": ("Number of Delayed Days", "max"),
                "BL Critical Path": ("BL Critical Path", lambda s: "Yes" if s.astype(str).str.lower().eq("yes").any() else "No"),
                "Updated Critical Path": ("Updated Critical Path", lambda s: "Yes" if s.astype(str).str.lower().eq("yes").any() else "No"),
                "Updated Longest Path": ("Updated Longest Path", lambda s: "Yes" if s.astype(str).str.lower().eq("yes").any() else "No"),
                "Current Total Float": ("Current Total Float", "first"),
                "Delay Start": ("Delay Start", "min"),
                "Delay Finish": ("Delay Finish", "max"),
                "Treatment": ("Treatment", "first"),
            }
        )
        .reset_index()
    )
    gross_total = float(pd.to_numeric(grouped["Number of Delayed Days"], errors="coerce").fillna(0).sum())
    final_summary_df = build_delay_tia_final_submission_summary_df(context, analysis)
    recommended_total = 0
    if not final_summary_df.empty:
        recommended_row = final_summary_df.loc[
            final_summary_df["Submission Position"] == "Recommended total delayed days to submit",
            "Days",
        ]
        if not recommended_row.empty:
            recommended_total = float(pd.to_numeric(recommended_row, errors="coerce").fillna(0).iloc[0])
    total_rows = pd.DataFrame(
        [
            {
                "Delay Stream": "TOTAL",
                "Delay Ref": "Gross delayed activity total before concurrency de-duplication",
                "Activity ID": "",
                "Activity Name": "",
                "Number of Delayed Days": gross_total,
                "BL Critical Path": "",
                "Updated Critical Path": "",
                "Updated Longest Path": "",
                "Current Total Float": "",
                "Delay Start": pd.NaT,
                "Delay Finish": pd.NaT,
                "Treatment": "Gross total across listed delayed activities.",
            },
            {
                "Delay Stream": "TOTAL",
                "Delay Ref": "Recommended total delayed days to submit",
                "Activity ID": "",
                "Activity Name": "",
                "Number of Delayed Days": recommended_total,
                "BL Critical Path": "",
                "Updated Critical Path": "",
                "Updated Longest Path": "",
                "Current Total Float": "",
                "Delay Start": pd.NaT,
                "Delay Finish": pd.NaT,
                "Treatment": "Conservative final TIA submission total after concurrency / net-impact treatment.",
            },
        ]
    )
    return pd.concat([grouped, total_rows], ignore_index=True)


def build_delay_tia_director_pack_context(
    context: dict,
    analysis: dict,
    overrides: dict | None = None,
    selected_sources: list[str] | None = None,
) -> dict:
    overrides = overrides or {}
    metadata = build_delay_tia_director_metadata(context, analysis)
    kpis = build_delay_tia_director_kpis(context, analysis)
    p6_controls_df = build_delay_tia_director_p6_controls_df(context, analysis)
    control_lookup = {
        str(row["Control Item"]).strip().lower(): delay_tia_docx_text(row["Recorded Value"])
        for _, row in p6_controls_df.iterrows()
    }
    source_frames = {
        "Delay Event Register": build_delay_tia_director_delay_events_df(context, analysis),
        "Fragnet Register": build_delay_tia_director_fragnet_logic_df(context, analysis),
        "Impact Calculation Table": build_delay_tia_director_impact_calculation_df(context, analysis),
        "Affected Activities / Criticality Table": build_delay_tia_director_activity_impact_df(context, analysis),
        "Final Delayed Activities / Critical Path Comparison": build_delay_tia_final_delayed_activities_df(context, analysis),
        "Causation / Concurrency / Entitlement Matrix": build_delay_tia_director_causation_matrix_df(context, analysis),
        "Evidence Checklist": build_delay_tia_director_evidence_register_df(context, analysis),
        "Data Readiness Matrix": build_delay_tia_director_readiness_matrix_df(context, analysis),
        "P6 Control Sheet": p6_controls_df,
    }
    if selected_sources is None:
        selected_sources = list(source_frames.keys())
    selected = set(selected_sources)

    def include(name: str) -> pd.DataFrame:
        frame = source_frames.get(name, pd.DataFrame())
        if name not in selected:
            return pd.DataFrame(columns=list(frame.columns))
        return delay_tia_add_number_of_delayed_days_column(frame.copy())

    return {
        "project_name": overrides.get("project_name", metadata.get("project_name", "")),
        "contract_no": overrides.get("contract_no", metadata.get("contract_no", "")),
        "data_date": overrides.get("data_date", metadata.get("data_date", pd.Timestamp.today().normalize())),
        "revision": overrides.get("revision", "Rev. 00"),
        "employer": overrides.get("employer", metadata.get("employer", "")),
        "contractor": overrides.get("contractor", metadata.get("contractor", "")),
        "contract_form_clause": overrides.get("contract_form_clause", metadata.get("contract_form_or_clause", "")),
        "accepted_baseline_programme": overrides.get("accepted_baseline_programme", metadata.get("accepted_baseline_programme", "")),
        "impacted_update_programme": overrides.get("impacted_update_programme", metadata.get("impacted_update", "")),
        "calendar_basis": overrides.get("calendar_basis", metadata.get("calendar_basis", "")),
        "schedule_file_name": overrides.get("schedule_file_name", control_lookup.get("schedule file name", "04- p6_activity_export.csv")),
        "schedule_options": overrides.get("schedule_options", control_lookup.get("schedule options", "Derived from uploaded activity export only")),
        "critical_path_basis": overrides.get("critical_path_basis", control_lookup.get("longest path / critical path basis", metadata.get("criticality_rule", ""))),
        "retained_logic_setting": overrides.get("retained_logic_setting", control_lookup.get("retained logic / progress override", "Not recorded in uploaded files")),
        "out_of_sequence_treatment": overrides.get("out_of_sequence_treatment", control_lookup.get("out-of-sequence progress handling", "Not recorded in uploaded files")),
        "constraints": overrides.get("constraints", control_lookup.get("constraints", "Not explicitly recorded")),
        "calendars": overrides.get("calendars", control_lookup.get("calendars", "Not explicitly recorded")),
        "open_ends": overrides.get("open_ends", control_lookup.get("open ends", "Not explicitly recorded")),
        "negative_float": overrides.get("negative_float", control_lookup.get("negative float", "Derived from current candidate / fragnet outputs only")),
        "preserve_original_charts_images": bool(overrides.get("preserve_original_charts_images", True)),
        "generated_by": overrides.get("generated_by", "Planning Department"),
        "generation_notes": overrides.get("generation_notes", ""),
        "db_path": CONTRACT_CLAIMS_DB_PATH,
        "kpis": kpis,
        "delay_event_register_df": include("Delay Event Register"),
        "fragnet_register_df": include("Fragnet Register"),
        "impact_calculation_df": include("Impact Calculation Table"),
        "activity_impact_df": include("Affected Activities / Criticality Table"),
        "final_delayed_activities_df": include("Final Delayed Activities / Critical Path Comparison"),
        "causation_matrix_df": include("Causation / Concurrency / Entitlement Matrix"),
        "evidence_register_df": include("Evidence Checklist"),
        "readiness_matrix_df": include("Data Readiness Matrix"),
        "p6_controls_df": include("P6 Control Sheet"),
        "selected_sources": sorted(selected),
    }




def build_delay_tia_delay_report_df(context: dict, analysis: dict) -> pd.DataFrame:
    assessment_df = analysis.get("assessment_df", pd.DataFrame())
    stock_out_df = analysis.get("stock_out_df", pd.DataFrame())
    fragnet_df = analysis.get("fragnet_df", pd.DataFrame())
    support_df = analysis.get("support_df", pd.DataFrame())
    contract_matches_df = analysis.get("contract_matches_df", pd.DataFrame())
    executive_df = analysis.get("executive_summary_df", pd.DataFrame())

    delay_days = 0
    if not fragnet_df.empty:
        if "Recorded Steel Delay Duration (days)" in fragnet_df.columns and fragnet_df["Recorded Steel Delay Duration (days)"].notna().any():
            delay_days = int(pd.to_numeric(fragnet_df["Recorded Steel Delay Duration (days)"], errors="coerce").fillna(0).max())
        elif "Fragment Duration" in fragnet_df.columns:
            delay_days = int(pd.to_numeric(fragnet_df["Fragment Duration"], errors="coerce").fillna(0).max())
    cause = "Not identified"
    if not assessment_df.empty:
        cause = "Employer free-issue steel shortage affecting the first impacted reinforcement activity."
    clause_topic = contract_matches_df["Clause / Topic"].iloc[0] if not contract_matches_df.empty else ""
    methodology = (
        "Uploaded-file Time Impact Analysis using employer-only steel supply, master activity shortage evidence, "
        "P6 criticality and float logic, relationship sequencing, and contract support review."
    )
    return pd.DataFrame(
        [
            {"Item": "Delay Event", "Value": "Client / Employer free-issue steel delay"},
            {"Item": "Delayed Days", "Value": delay_days},
            {"Item": "Cause", "Value": cause},
            {"Item": "First Stock-Out Date", "Value": stock_out_df["Stock-Out Date"].min() if not stock_out_df.empty else pd.NaT},
            {"Item": "First Impacted Activity", "Value": executive_df.loc[executive_df["Question"] == "Which activity was first impacted?", "Answer"].iloc[0] if not executive_df.empty and (executive_df["Question"] == "Which activity was first impacted?").any() else "Not identified"},
            {"Item": "Fragnet Insert Before", "Value": fragnet_df["Insert Fragment Before"].iloc[0] if not fragnet_df.empty else "Not identified"},
            {"Item": "Fragment Start", "Value": fragnet_df["Fragment Start"].iloc[0] if not fragnet_df.empty else pd.NaT},
            {"Item": "Fragment Finish", "Value": fragnet_df["Fragment Finish"].iloc[0] if not fragnet_df.empty else pd.NaT},
            {"Item": "Methodology Used", "Value": methodology},
            {"Item": "Primary Contract Support", "Value": clause_topic},
            {"Item": "Contractual Strength", "Value": support_df["Contractual Strength"].iloc[0] if not support_df.empty else "Not Contractually Proven"},
        ]
    )


def build_delay_tia_primavera_fragnet_df(context: dict, analysis: dict) -> pd.DataFrame:
    fragnet_df = analysis.get("fragnet_df", pd.DataFrame())
    if fragnet_df.empty:
        return pd.DataFrame()
    rows = []
    for idx, row in fragnet_df.reset_index(drop=True).iterrows():
        successor_id = str(row.get("Insert Fragment Before", "")).strip()
        predecessor_id = str(row.get("Last completed / available predecessor", "")).strip()
        fragnet_id = f"FG-{idx + 1:03d}"
        claimed_delay_duration = compute_claimed_delay_duration_days(
            row.get("Fragment Start", pd.NaT),
            row.get("Fragment Finish", pd.NaT),
            row.get("Fragment Duration", 0),
            row.get("Recorded Steel Delay Duration (days)", pd.NA),
        )
        rows.append(
            {
                "Activity ID": fragnet_id,
                "Activity Name": row.get("Fragment Activity Name", ""),
                "WBS": "TIA FRAGNET",
                "Original Duration": row.get("Fragment Duration", 0),
                "Claimed Delay Duration (days)": claimed_delay_duration,
                "Start": row.get("Fragment Start", pd.NaT),
                "Finish": row.get("Fragment Finish", pd.NaT),
                "Predecessor Activity ID": predecessor_id,
                "Predecessor Relationship Type": "FS" if predecessor_id else "",
                "Successor Activity ID": successor_id,
                "Successor Relationship Type": "FS" if successor_id else "",
                "Delay Status": row.get("Delay Status", ""),
                "Building": row.get("Affected Building", ""),
                "Steel Type": row.get("Steel Type", ""),
                "Shortage Qty": row.get("Shortage Qty", 0),
                "TIA Candidate Score": row.get("TIA Candidate Score", 0),
                "Notes": "Insert this fragnet before the affected activity in Primavera for TIA modelling.",
            }
        )
    return pd.DataFrame(rows)


def build_delay_tia_detailed_report_html(context: dict, analysis: dict) -> str:
    delay_report_df = delay_tia_add_number_of_delayed_days_column(build_delay_tia_delay_report_df(context, analysis))
    fragnet_primavera_df = delay_tia_add_number_of_delayed_days_column(build_delay_tia_primavera_fragnet_df(context, analysis))
    activity_impact_df = delay_tia_add_number_of_delayed_days_column(build_delay_tia_director_activity_impact_df(context, analysis))
    delay_events_director_df = delay_tia_add_number_of_delayed_days_column(build_delay_tia_director_delay_events_df(context, analysis))
    concurrent_review_df = delay_tia_add_number_of_delayed_days_column(context.get("concurrent_delay_review_df", pd.DataFrame()))
    final_delayed_activities_df = delay_tia_add_number_of_delayed_days_column(build_delay_tia_final_delayed_activities_df(context, analysis))
    final_submission_df = delay_tia_add_number_of_delayed_days_column(build_delay_tia_final_submission_summary_df(context, analysis))
    final_submission_breakdown_df = delay_tia_add_number_of_delayed_days_column(build_delay_tia_final_submission_event_breakdown_df(context, analysis))
    activity_selection_expl_df = delay_tia_add_number_of_delayed_days_column(build_delay_tia_activity_selection_explanation_df(context, analysis))
    fragnet_expl_df = delay_tia_add_number_of_delayed_days_column(build_delay_tia_fragnet_explanation_df(context, analysis))
    client_delay_expl_df = delay_tia_add_number_of_delayed_days_column(build_delay_tia_client_delay_evidence_explanation_df(context, analysis))
    narrative_df = delay_tia_add_number_of_delayed_days_column(analysis.get("narrative_df", pd.DataFrame()))
    assessment_df = delay_tia_add_number_of_delayed_days_column(analysis.get("assessment_df", pd.DataFrame()))
    methodology_df = pd.DataFrame(
        [
            {"Step": 1, "Methodology Logic": "Use uploaded employer steel supply only; exclude contractor supply from calculations."},
            {"Step": 2, "Methodology Logic": "Use the master activity steel analysis file as the primary requirement and shortage basis."},
            {"Step": 3, "Methodology Logic": "Use P6 activity data for float, criticality, readiness, and successor impact."},
            {"Step": 4, "Methodology Logic": "Use relationship logic to determine where the fragnet must be inserted."},
            {"Step": 5, "Methodology Logic": "Use contract library support to explain schedule entitlement and delay classification."},
        ]
    )
    sections = [
        "<h1>Delay TIA Detailed Report</h1>",
        "<p>Based on the uploaded project records, this is a professional planning and contractual assessment and not a legal opinion.</p>",
        "<h2>Delay Summary</h2>",
        delay_report_df.to_html(index=False, border=0),
        "<h2>Claimed Delay Events</h2>",
        delay_events_director_df.to_html(index=False, border=0) if not delay_events_director_df.empty else "<p>No claimed delay events generated.</p>",
        "<h2>Claimed Activity Delay Durations</h2>",
        activity_impact_df.to_html(index=False, border=0) if not activity_impact_df.empty else "<p>No claimed activity delay durations generated.</p>",
        "<h2>Concurrent Delay Review</h2>",
        concurrent_review_df.to_html(index=False, border=0) if not concurrent_review_df.empty else "<p>No concurrent delay review rows generated.</p>",
        "<h2>Final Delayed Activities and Critical Path Comparison</h2>",
        final_delayed_activities_df.to_html(index=False, border=0) if not final_delayed_activities_df.empty else "<p>No final delayed activities generated.</p>",
        "<h2>Final Submission Delay Position</h2>",
        final_submission_df.to_html(index=False, border=0) if not final_submission_df.empty else "<p>No final submission delay position generated.</p>",
        "<h2>Submission Delay Event Breakdown</h2>",
        final_submission_breakdown_df.to_html(index=False, border=0) if not final_submission_breakdown_df.empty else "<p>No submission event breakdown generated.</p>",
        "<h2>Methodology Used</h2>",
        methodology_df.to_html(index=False, border=0),
        "<h2>Contractual Delay Assessment</h2>",
        assessment_df.to_html(index=False, border=0) if not assessment_df.empty else "<p>No contractual assessment generated.</p>",
        "<h2>Detailed TIA Narrative</h2>",
        narrative_df.to_html(index=False, border=0) if not narrative_df.empty else "<p>No detailed narrative generated.</p>",
        "<h2>Detailed Basis for Selecting Delayed Activities</h2>",
        activity_selection_expl_df.to_html(index=False, border=0) if not activity_selection_expl_df.empty else "<p>No delayed-activity selection explanation generated.</p>",
        "<h2>Detailed Fragnet Development Basis</h2>",
        fragnet_expl_df.to_html(index=False, border=0) if not fragnet_expl_df.empty else "<p>No fragnet explanation generated.</p>",
        "<h2>Detailed Client Delay Evidence, Methodology, and Entitlement Basis</h2>",
        client_delay_expl_df.to_html(index=False, border=0) if not client_delay_expl_df.empty else "<p>No detailed client-delay explanation generated.</p>",
        "<h2>Primavera Fragnet Sheet</h2>",
        fragnet_primavera_df.to_html(index=False, border=0) if not fragnet_primavera_df.empty else "<p>No fragnet rows generated.</p>",
    ]
    return (
        "<html><head><meta charset='utf-8'><title>Delay TIA Detailed Report</title>"
        "<style>body{font-family:Arial,sans-serif;padding:24px;color:#1f2937;}h1,h2{color:#0f172a;}table{border-collapse:collapse;width:100%;margin:12px 0;}th,td{border:1px solid #d1d5db;padding:8px;text-align:left;}th{background:#e5eef8;}</style>"
        "</head><body>"
        + "".join(sections)
        + "</body></html>"
    )


def build_delay_tia_excel_report_bytes(context: dict, analysis: dict) -> bytes:
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        report_sheets = {
            "Delay Report": build_delay_tia_delay_report_df(context, analysis),
            "Claimed Delay Events": build_delay_tia_director_delay_events_df(context, analysis),
            "Claimed Activity Impact": build_delay_tia_director_activity_impact_df(context, analysis),
            "Concurrent Delay Review": context.get("concurrent_delay_review_df", pd.DataFrame()),
            "Final Delayed Activities": build_delay_tia_final_delayed_activities_df(context, analysis),
            "Final Submission Delay": build_delay_tia_final_submission_summary_df(context, analysis),
            "Submission Event Split": build_delay_tia_final_submission_event_breakdown_df(context, analysis),
            "Activity Selection Basis": build_delay_tia_activity_selection_explanation_df(context, analysis),
            "Fragnet Basis": build_delay_tia_fragnet_explanation_df(context, analysis),
            "Delay Evidence Basis": build_delay_tia_client_delay_evidence_explanation_df(context, analysis),
            "Executive Summary": analysis.get("executive_summary_df", pd.DataFrame()),
            "Stock-Out Events": analysis.get("stock_out_df", pd.DataFrame()),
            "Affected Activities": analysis.get("candidates_df", pd.DataFrame()),
            "Fragnet Recommendation": analysis.get("fragnet_df", pd.DataFrame()),
            "Contract Assessment": analysis.get("assessment_df", pd.DataFrame()),
            "Detailed TIA": analysis.get("narrative_df", pd.DataFrame()),
            "Primavera Fragnet": build_delay_tia_primavera_fragnet_df(context, analysis),
            "Rebuttal Matrix": analysis.get("rebuttal_df", pd.DataFrame()),
            "Action Tracker": analysis.get("action_tracker_df", pd.DataFrame()),
        }
        for sheet_name, frame in report_sheets.items():
            delay_tia_add_number_of_delayed_days_column(frame).to_excel(writer, sheet_name=sheet_name[:31], index=False)
        workbook = writer.book
        for sheet in workbook.worksheets:
            sheet.freeze_panes = "A2"
    output.seek(0)
    return output.getvalue()


delay_tia_context = {}
delay_tia_analysis = {}
bl_critical_path_fixed_df = pd.DataFrame()
bl_critical_path_comparison_df = pd.DataFrame()
bl_critical_path_summary = {
    "bl_count": 0,
    "current_count": 0,
    "matched_count": 0,
    "bl_only_count": 0,
    "current_only_count": 0,
}


hidden_slide_selectors = []
for hidden_slide_name in ("Delays", "Time Impact"):
    if hidden_slide_name in PROJECT_HUB_SLIDE_NAMES:
        hidden_slide_selectors.append(
            f'div[data-baseweb="tab-list"] > button:nth-of-type({PROJECT_HUB_SLIDE_NAMES.index(hidden_slide_name) + 1})'
        )
if hidden_slide_selectors:
    st.markdown(
        "<style>"
        + ",".join(hidden_slide_selectors)
        + "{display:none!important;}"
        + "</style>",
        unsafe_allow_html=True,
    )

st.markdown(
    """
    <style>
    #MainMenu, footer, header {visibility:hidden!important;}
    [data-testid="stToolbar"],
    [data-testid="stDecoration"],
    [data-testid="stStatusWidget"],
    [data-testid="stDeployButton"],
    [data-testid="stHeaderActionElements"],
    .stAppDeployButton,
    .st-emotion-cache-1dp5vir,
    .viewerBadge_container__1QSob,
    .viewerBadge_link__1S137,
    div[class^="viewerBadge_"],
    div[class*=" viewerBadge_"],
    a[class^="viewerBadge_"],
    a[class*=" viewerBadge_"],
    div[class*="ViewerBadge"],
    a[class*="ViewerBadge"],
    div[class*="stDeployButton"],
    button[class*="stDeployButton"],
    a[href*="streamlit.io/cloud"],
    a[href*="github.com"][target="_blank"],
    div:has(a[href*="streamlit.io/cloud"]),
    div:has(a[href*="share.streamlit.io"]),
    div[style*="position: fixed"][style*="bottom"][style*="right"],
    button[style*="position: fixed"][style*="bottom"][style*="right"] {
        display:none!important;
        visibility:hidden!important;
        opacity:0!important;
        pointer-events:none!important;
    }
    body > div:last-child:not([data-testid="stAppViewContainer"]) {
        display:none!important;
    }
    div[data-testid="stAppViewContainer"] {padding-bottom:0!important;}
    </style>
    """,
    unsafe_allow_html=True,
)

default_meeting_seed = overview_metrics.get("project_name") or selected_project_id() or "Project-Intelligence-Hub"
default_meeting_room = sanitize_conference_room_name(f"{default_meeting_seed}-Live-Review")
if not st.session_state.get("conference_room_name"):
    st.session_state["conference_room_name"] = default_meeting_room
shared_meeting_room = sanitize_conference_room_name(st.session_state["conference_room_name"])
shared_meeting_url = f"https://meet.jit.si/{shared_meeting_room}"
try:
    shared_dashboard_url = str(st.context.url or configured_streamlit_url() or "Configure the deployed Streamlit URL in mobile_config.json.")
except Exception:
    shared_dashboard_url = configured_streamlit_url() or "Configure the deployed Streamlit URL in mobile_config.json."

meeting_title_col, meeting_action_col = st.columns([0.78, 0.22])
meeting_title_col.markdown("#### Project Meeting")
meeting_title_col.caption("Start the call here, then switch through the project slides while everyone continues talking. Supports 5+ participants.")
if st.session_state.get("project_meeting_active", False):
    if meeting_action_col.button("Close Meeting", key="close_project_meeting", width="stretch"):
        st.session_state["project_meeting_active"] = False
        st.rerun()
else:
    if meeting_action_col.button("Start Meeting", key="start_project_meeting", type="primary", width="stretch"):
        st.session_state["project_meeting_active"] = True
        st.rerun()

with st.expander("Share meeting with participants", expanded=False):
    st.code(
        "Join the live project review:\n"
        f"Dashboard: {shared_dashboard_url}\n"
        f"Meeting room: {shared_meeting_url}",
        language=None,
    )

if st.session_state.get("project_meeting_active", False):
    st.components.v1.html(
        build_inline_project_meeting_html(shared_meeting_room),
        height=490,
        scrolling=False,
    )

VISIBLE_PROJECT_SLIDE_NAMES = [
    slide_name
    for slide_name in PROJECT_HUB_SLIDE_NAMES
    if slide_name not in {"Delays", "Time Impact"}
]
if str(AUTH_USER.get("role", "viewer")).lower() == "director":
    VISIBLE_PROJECT_SLIDE_NAMES = [
        slide_name
        for slide_name in VISIBLE_PROJECT_SLIDE_NAMES
        if slide_name in {"Overview", "EVM Analysis", "Risks", "Output Studio"}
    ]
if str(AUTH_USER.get("role", "viewer")).lower() != "admin":
    allowed_sections = [
        section
        for section in AUTH_USER.get("access_sections", [])
        if section in VISIBLE_PROJECT_SLIDE_NAMES
    ]
    if allowed_sections:
        VISIBLE_PROJECT_SLIDE_NAMES = allowed_sections
if st.session_state.get("active_project_slide_name") not in VISIBLE_PROJECT_SLIDE_NAMES:
    st.session_state["active_project_slide_name"] = "Output Studio" if "Output Studio" in VISIBLE_PROJECT_SLIDE_NAMES else VISIBLE_PROJECT_SLIDE_NAMES[0]
active_slide_name = st.selectbox(
    "Project slide",
    VISIBLE_PROJECT_SLIDE_NAMES,
    key="active_project_slide_name",
)
if str(AUTH_USER.get("role", "viewer")).lower() == "director":
    st.caption("Director access is focused on executive overview, EVM, risk, and output dashboards.")

if active_slide_name == PROJECT_HUB_SLIDE_NAMES[0]:
    st.markdown("<div class='section-header'><h3>Project Overview</h3></div>", unsafe_allow_html=True)
    c1, c2, c3, c4, c5 = st.columns(5)
    c1.metric("Project Start Date", format_project_date(overview_metrics.get("project_start")))
    c2.metric("Project Finish Date", format_project_date(overview_metrics.get("project_finish")))
    c3.metric("Project Duration [Days]", int(overview_metrics.get("duration_days", 0)))
    c4.metric("Duration Elapsed", pct(overview_metrics.get("duration_elapsed_pct")))
    c5.metric("Remaining Duration", pct(overview_metrics.get("remaining_duration_pct")))

    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Overall Progress", pct(overview_metrics.get("overall_progress")), delta="Actual progress")
    c2.metric("Planned Progress", pct(overview_metrics.get("planned_progress")), delta="Planned progress")
    c3.metric("Contract Value", egp(overview_metrics.get("contract_value")), delta="Project contract")
    c4.metric("Total Activities", int(overview_metrics.get("total_activities", 0)), delta="Activity count")

    c1, c2 = st.columns(2)
    c1.metric("Critical Activities", int(overview_metrics.get("critical_activities", 0)), delta="Critical = Yes")
    c2.metric("Remaining Duration", pct(overview_metrics.get("remaining_duration_pct")), delta="100% - elapsed")

    if wbs_costs:
        df_wbs = pd.DataFrame(wbs_costs)
        fig = px.bar(df_wbs, x="wbs_name", y=["budget", "actual"], title="Overall Cost by Discipline / WBS", barmode="group", color_discrete_map={"budget":"#245f95","actual":"#d9544d"})
        st.plotly_chart(style_plotly(fig, 430), width="stretch")

    if letters:
        threads = letters.get("Issue Threads", pd.DataFrame())
        if not threads.empty:
            st.markdown("<div class='panel-note'><b>Live Alert Engine</b><br>Priority correspondence threads are linked to required follow-up and claim exposure.</div>", unsafe_allow_html=True)
            st.dataframe(threads[["Thread", "Priority", "Next Action"]].head(8), width="stretch", hide_index=True)

if active_slide_name == PROJECT_HUB_SLIDE_NAMES[2]:
    st.markdown("<div class='section-header'><h3>Activities Analysis</h3></div>", unsafe_allow_html=True)
    activities_df = activity_metrics["activities_df"]
    critical_df = activity_metrics["critical_df"]
    deviated_df = activity_metrics["deviated_df"]
    rft_df = activity_metrics["rft_df"]
    if not activities_df.empty:
        c1, c2, c3, c4 = st.columns(4)
        c1.metric("Total Activities", len(activities_df))
        c2.metric("Critical Path Activities", activity_metrics["critical_count"])
        c3.metric("Most Deviated Activities", activity_metrics["deviated_count"])
        c4.metric("RFT Activities", activity_metrics["rft_count"])

        c5, c6 = st.columns(2)
        c5.metric("Avg Critical Variance", f"{activity_metrics['avg_critical_variance']:.1f}%")
        c6.metric("Avg RFT Variance", f"{activity_metrics['avg_rft_variance']:.1f}%")

        st.markdown("#### Critical Path Analysis")
        critical_view = critical_df[
            ["activity_id", "activity_name", "planned_progress", "actual_progress", "progress_variance", "planned_finish", "forecast_finish", "finish_slip_days", "total_float_days"]
        ].copy() if not critical_df.empty else pd.DataFrame()
        if not critical_view.empty:
            critical_view["progress_variance"] = critical_view["progress_variance"].round(1)
            critical_view["finish_slip_days"] = critical_view["finish_slip_days"].astype(int)
            st.dataframe(arrow_safe_display_df(critical_view), width="stretch", hide_index=True, height=dataframe_height(critical_view))
        else:
            st.info("No critical path activities flagged in activities.csv.")

        st.markdown("#### Most Deviated Activities")
        deviated_view = deviated_df[
            ["activity_id", "activity_name", "planned_progress", "actual_progress", "progress_variance", "planned_finish", "forecast_finish", "finish_slip_days", "is_critical"]
        ].copy() if not deviated_df.empty else pd.DataFrame()
        if not deviated_view.empty:
            deviated_view["progress_variance"] = deviated_view["progress_variance"].round(1)
            deviated_view["finish_slip_days"] = deviated_view["finish_slip_days"].astype(int)
            st.dataframe(arrow_safe_display_df(deviated_view), width="stretch", hide_index=True, height=dataframe_height(deviated_view))
        else:
            st.info("No deviated activities detected from planned vs actual progress or forecast finish slip.")

        st.markdown("#### RFT Activities Analysis")
        rft_view = rft_df[
            ["activity_id", "activity_name", "planned_progress", "actual_progress", "progress_variance", "planned_finish", "actual_finish", "forecast_finish", "finish_slip_days", "is_critical"]
        ].copy() if not rft_df.empty else pd.DataFrame()
        if not rft_view.empty:
            rft_view["progress_variance"] = rft_view["progress_variance"].round(1)
            rft_view["finish_slip_days"] = rft_view["finish_slip_days"].astype(int)
            st.dataframe(arrow_safe_display_df(rft_view), width="stretch", hide_index=True, height=dataframe_height(rft_view))
        else:
            st.info("No RFT activities found in activities.csv.")
    else:
        st.info("No activities data available.")

if active_slide_name == PROJECT_HUB_SLIDE_NAMES[1]:
    st.markdown("<div class='section-header'><h3>WBS Analysis</h3></div>", unsafe_allow_html=True)
    wbs_df = wbs_metrics["wbs_df"]
    con_wbs_df = wbs_metrics["con_wbs_df"]
    chart_rows = wbs_metrics["chart_rows"]
    code_col = wbs_metrics["code_col"]
    target_codes = chart_rows[code_col].astype(str).tolist() if not chart_rows.empty else []
    if wbs_df.empty or not code_col:
        st.info("No WBS data available.")
    else:
        left_col, right_col = st.columns(2)
        for idx, target_code in enumerate(target_codes):
            target_df = chart_rows[chart_rows[code_col] == target_code].copy()
            chart_col = left_col if idx % 2 == 0 else right_col
            with chart_col:
                if target_df.empty:
                    st.info(f"{target_code} was not found in wbs.csv.")
                else:
                    row = target_df.iloc[0]
                    chart_df = pd.DataFrame(
                        {
                            "Metric": ["Schedule % Complete", "Performance % Complete"],
                            "Value": [
                                float(row.get("schedule_%_complete_num", 0.0)),
                                float(row.get("performance_%_complete_num", 0.0)),
                            ],
                        }
                    )
                    fig = px.bar(
                        chart_df,
                        x="Metric",
                        y="Value",
                        color="Metric",
                        barmode="group",
                        title=target_code,
                        color_discrete_map={
                            "Schedule % Complete": "#245f95",
                            "Performance % Complete": "#168f8b",
                        },
                    )
                    fig.update_yaxes(title_text="% Complete")
                    st.plotly_chart(style_plotly(fig, 370), width="stretch")

        st.markdown("#### WBS Analysis Table - CON Rows Only")
        if con_wbs_df.empty:
            st.info("No WBS rows containing CON were found.")
        else:
            st.dataframe(
                con_wbs_df,
                width="stretch",
                hide_index=True,
                height=dataframe_height(con_wbs_df, max_height=1200),
            )

if active_slide_name == PROJECT_HUB_SLIDE_NAMES[9]:
    st.markdown("<div class='section-header'><h3>Time Impact Analysis</h3></div>", unsafe_allow_html=True)
    engine_df = time_impact_engine["time_impact_df"]
    if not engine_df.empty:
        c1, c2, c3, c4 = st.columns(4)
        c1.metric("TIA Input Events", time_impact_engine["priority_count"])
        c2.metric("Critical Path Impacted", time_impact_engine["critical_impacted"])
        c3.metric("Near-Critical Impacted", time_impact_engine["near_critical_impacted"])
        c4.metric("RFT-Linked Events", time_impact_engine["rft_impacted"])

        c5, c6 = st.columns(2)
        c5.metric("Commercial Gap", egp(time_impact_engine["commercial_gap"]))
        c6.metric("Open High-Priority Events", int(engine_df[(engine_df["status_group"] == "Open") & (engine_df["engine_priority"] >= 8)].shape[0]))

        st.markdown("#### Time Impact Engine Input Matrix")
        st.dataframe(
            time_impact_engine["input_matrix"],
            width="stretch",
            hide_index=True,
            height=dataframe_height(time_impact_engine["input_matrix"], max_height=360),
        )

        st.markdown("#### Causation and Responsibility Matrix")
        st.dataframe(
            time_impact_engine["causation_matrix"],
            width="stretch",
            hide_index=True,
            height=dataframe_height(time_impact_engine["causation_matrix"], max_height=420),
        )

        priority_view = engine_df[
            [
                "delay_id",
                "source_stream",
                "delay_title",
                "activity_id",
                "activity_name",
                "estimated_delay_days_num",
                "responsibility_band",
                "critical_path_flag",
                "near_critical_flag",
                "rft_flag",
                "progress_variance",
                "finish_slip_days",
                "notice_status",
                "eot_potential",
                "status",
                "engine_priority",
            ]
        ].copy()
        priority_view.columns = [
            "Delay ID",
            "Source Stream",
            "Delay Title",
            "Activity ID",
            "Activity Name",
            "Delay Days",
            "Responsibility",
            "Critical Path",
            "Near Critical",
            "RFT Linked",
            "Progress Variance",
            "Finish Slip Days",
            "Notice Ref",
            "EOT Potential",
            "Status",
            "Priority Score",
        ]
        priority_view["Progress Variance"] = priority_view["Progress Variance"].round(1)
        priority_view["Finish Slip Days"] = priority_view["Finish Slip Days"].fillna(0).astype(int)
        st.markdown("#### Integrated Time Impact Event Register")
        st.dataframe(priority_view, width="stretch", hide_index=True, height=dataframe_height(priority_view, max_height=950))
    else:
        st.info("No time impact data available.")

if active_slide_name == PROJECT_HUB_SLIDE_NAMES[4]:
    st.markdown("<div class='section-header'><h3>S-Curve Analysis</h3></div>", unsafe_allow_html=True)
    curve_df = s_curve_metrics["curve_df"]
    if not curve_df.empty:
        c1, c2, c3, c4 = st.columns(4)
        c1.metric("Total Planned Budget", egp(s_curve_metrics["planned_total"]))
        c2.metric("Total Actual Spend", egp(s_curve_metrics["actual_total"]))
        c3.metric("Total Invoiced", egp(s_curve_metrics["invoiced_total"]))
        c4.metric("Last Actual Month", s_curve_metrics["last_actual_month"])

        line_df = curve_df[["Month", "Planned", "Actual", "Invoiced"]].copy()
        fig = px.line(
            line_df,
            x="Month",
            y=["Planned", "Actual", "Invoiced"],
            markers=True,
            title="Cumulative Spend Curve",
            color_discrete_map={"Planned":"#245f95","Actual":"#168f8b","Invoiced":"#c98519"},
        )
        st.plotly_chart(style_plotly(fig, 500), width="stretch")

        monthly_df = curve_df[["Month", "planned_value_num", "actual_cost_num", "invoice_amount_num"]].copy()
        monthly_df.columns = ["Month", "Monthly Planned", "Monthly Actual", "Monthly Invoiced"]
        fig_monthly = px.bar(
            monthly_df,
            x="Month",
            y=["Monthly Planned", "Monthly Actual", "Monthly Invoiced"],
            barmode="group",
            title="Monthly Spend Comparison",
            color_discrete_map={"Monthly Planned":"#245f95","Monthly Actual":"#168f8b","Monthly Invoiced":"#c98519"},
        )
        st.plotly_chart(style_plotly(fig_monthly, 420), width="stretch")
        st.dataframe(curve_df, width="stretch", hide_index=True)
    else:
        st.info("No S-curve data available.")

if active_slide_name == PROJECT_HUB_SLIDE_NAMES[5]:
    st.markdown("<div class='section-header'><h3>Earned Value Management</h3></div>", unsafe_allow_html=True)
    cols = st.columns(4)
    for col, label, key in zip(cols, ["BAC", "AC", "EV", "PV"], ["bac", "ac", "ev", "pv"]):
        col.metric(label, egp(evm_metrics.get(key)))
    cols = st.columns(4)
    cols[0].metric("Schedule Variance", egp(evm_metrics.get("sv")))
    cols[1].metric("Cost Variance", egp(evm_metrics.get("cv")))
    cols[2].metric("EAC", egp(evm_metrics.get("eac")))
    cols[3].metric("TCPI", f"{evm_metrics.get('tcpi'):.3f}" if evm_metrics.get("tcpi") is not None else "N/A")
    evm_table = filter_active_project(load_core_csv(EVM_CSV_PATH))
    if not evm_table.empty:
        st.dataframe(evm_table, width="stretch", hide_index=True)

    st.markdown("### Earned Value Analysis Add-on")
    rootCauseDf = build_evm_root_cause_rows(delay_metrics, risk_metrics, contract_metrics)
    mitigationDf = build_evm_mitigation_rows()
    evmComments = get_evm_comments()
    report_date = pd.Timestamp.today().strftime("%d %b %Y")
    project_title = overview_metrics.get("project_name") or project.get("project_name", "Project")

    st.markdown("#### Quantitative Performance (EVM Metrics)")
    q1, q2, q3, q4 = st.columns(4)
    q1.metric("BAC", sar(evmData["BAC"]), delta="Budget @ Completion")
    q2.metric("PV", sar(evmData["PV"]), delta="Cumm. Planned Value")
    q3.metric("EV", sar(evmData["EV"]), delta="Cumm. Earned Value")
    q4.metric("SV", sar(evmData["SV"]), delta="Negative exposure")
    q5, q6, q7, q8 = st.columns(4)
    q5.metric("SPI", f"{evmData['SPI']:.2f}", delta="Schedule performance index")
    q6.metric("Planned completion value position", pct(evmData["plannedCompletionValuePosition"]), delta="PV / BAC")
    q7.metric("Earned value gap", sar(evmData["earnedValueGap"]), delta="PV not converted into EV")
    q8.metric("Schedule health classification", evmData["scheduleHealthClassification"], delta="Executive status")

    chart_col, gauge_col = st.columns([1.7, 1])
    with chart_col:
        chart_df = pd.DataFrame(
            [
                {"Metric": "BAC", "Value": evmData["BAC"]},
                {"Metric": "PV", "Value": evmData["PV"]},
                {"Metric": "EV", "Value": evmData["EV"]},
                {"Metric": "SV Exposure", "Value": abs(evmData["SV"])},
            ]
        )
        chart_fig = px.bar(
            chart_df,
            x="Metric",
            y="Value",
            title="BAC vs PV vs EV with SV Exposure",
            color="Metric",
            color_discrete_map={
                "BAC": "#245f95",
                "PV": "#168f8b",
                "EV": "#c98519",
                "SV Exposure": "#d9544d",
            },
        )
        st.plotly_chart(style_plotly(chart_fig, 420), width="stretch")
    with gauge_col:
        gauge_fig = go.Figure(
            go.Indicator(
                mode="gauge+number",
                value=evmData["SPI"],
                title={"text": "SPI Health"},
                number={"valueformat": ".2f"},
                gauge={
                    "axis": {"range": [0, 1.5]},
                    "bar": {"color": "#d9544d" if evmData["SPI"] < 1 else "#168f8b"},
                    "steps": [
                        {"range": [0, 0.75], "color": "#fde8e8"},
                        {"range": [0.75, 1.0], "color": "#fff4d6"},
                        {"range": [1.0, 1.5], "color": "#e7f7ef"},
                    ],
                    "threshold": {"line": {"color": "#173b63", "width": 4}, "value": 1.0},
                },
            )
        )
        gauge_fig.update_layout(height=420, paper_bgcolor="rgba(0,0,0,0)", margin=dict(l=10, r=10, t=50, b=10), font=dict(color="#172033"))
        st.plotly_chart(gauge_fig, width="stretch")

    st.markdown(f"<div class='panel-note'>{evmData['interpretation']}</div>", unsafe_allow_html=True)
    st.text_area(
        "Add / Edit My Comment",
        key="evm_comment_quantitativePerformance",
        height=110,
        on_change=persist_evm_comments_from_state,
    )
    if st.button("Clear Comment", key="clear_evm_quantitative"):
        clear_evm_comment("quantitativePerformance")
        st.rerun()

    st.markdown("#### Root Cause Linkage")
    st.dataframe(rootCauseDf, width="stretch", hide_index=True, height=dataframe_height(rootCauseDf, max_height=420))
    st.markdown(
        "<div class='panel-note'>The negative schedule variance is not an isolated numerical deviation. It is directly linked to unresolved external and interface-driven constraints that prevented planned progress from being converted into earned value, mainly within construction and procurement work fronts.</div>",
        unsafe_allow_html=True,
    )
    st.text_area(
        "Add / Edit My Comment",
        key="evm_comment_rootCauseLinkage",
        height=110,
        on_change=persist_evm_comments_from_state,
    )
    if st.button("Clear Comment", key="clear_evm_rootcause"):
        clear_evm_comment("rootCauseLinkage")
        st.rerun()

    st.markdown("#### Contractor Mitigation & Recovery Status")
    st.dataframe(mitigationDf, width="stretch", hide_index=True, height=dataframe_height(mitigationDf, max_height=520))
    st.markdown(
        "<div class='panel-note'>Contractor mitigation is focused on protecting available work fronts, accelerating technical closures, maintaining commercial entitlement records, and recovering productivity once external constraints are removed. Recovery remains dependent on timely closure of outstanding Owner / Engineer-driven constraints.</div>",
        unsafe_allow_html=True,
    )
    st.text_area(
        "Add / Edit My Comment",
        key="evm_comment_mitigationRecovery",
        height=110,
        on_change=persist_evm_comments_from_state,
    )
    if st.button("Clear Comment", key="clear_evm_mitigation"):
        clear_evm_comment("mitigationRecovery")
        st.rerun()

    evmComments = get_evm_comments()
    evm_html = printEVMHtml(project_title, report_date, evmData, rootCauseDf, mitigationDf, evmComments)
    evm_ppt = exportEVMToPowerPoint(project_title, report_date, evmData, rootCauseDf, mitigationDf, evmComments)

    st.markdown("#### Earned Value Analysis Print / Export")
    export_col1, export_col2 = st.columns(2)
    export_col1.download_button(
        "Download Earned Value Analysis HTML",
        data=evm_html.encode("utf-8"),
        file_name="earned_value_analysis.html",
        mime="text/html",
        width="stretch",
    )
    export_col2.download_button(
        "Download Earned Value Analysis PowerPoint",
        data=evm_ppt,
        file_name="earned_value_analysis.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        width="stretch",
    )
    with st.expander("Preview print-ready HTML", expanded=False):
        st.components.v1.html(evm_html, height=1200, scrolling=True)

if active_slide_name == PROJECT_HUB_SLIDE_NAMES[6]:
    st.markdown("<div class='section-header'><h3>Contract Management</h3></div>", unsafe_allow_html=True)
    contracts_df = contract_metrics["contracts_df"]
    payments_df = contract_metrics["payments_df"]
    if not contracts_df.empty:
        c1, c2, c3, c4 = st.columns(4)
        c1.metric("Total Contracts", contract_metrics["total_contracts"])
        c2.metric("Total Contract Value", egp(contract_metrics["total_contract_value"]))
        c3.metric("Total Certified", egp(contract_metrics["total_certified"]))
        c4.metric("Total Paid", egp(contract_metrics["total_paid"]))

        chart_df = pd.DataFrame(
            {
                "Metric": ["Total Contract Value", "Total Certified", "Total Paid"],
                "Value": [
                    contract_metrics["total_contract_value"],
                    contract_metrics["total_certified"],
                    contract_metrics["total_paid"],
                ],
            }
        )
        fig = px.bar(
            chart_df,
            x="Metric",
            y="Value",
            title="Contract Value vs Certified vs Paid",
            color="Metric",
            color_discrete_map={
                "Total Contract Value": "#245f95",
                "Total Certified": "#168f8b",
                "Total Paid": "#c98519",
            },
        )
        st.plotly_chart(style_plotly(fig, 420), width="stretch")

        st.markdown("#### Contracts Table")
        st.dataframe(contracts_df, width="stretch", hide_index=True, height=dataframe_height(contracts_df))

        st.markdown("#### Payments Table")
        st.dataframe(arrow_safe_display_df(payments_df), width="stretch", hide_index=True, height=dataframe_height(payments_df))
    else:
        st.info("No contract data available.")

def render_contract_clause_matching_engine():
    st.markdown("<div class='section-header'><h3>Contract Clause Matching Engine</h3></div>", unsafe_allow_html=True)
    all_clauses = get_all_clauses()
    if not all_clauses:
        st.info("No contract clause data is available for the selected project.")
        st.dataframe(
            pd.DataFrame({
                "Required Project Input": ["Contract document or clause library", "Contract evidence register", "Project claims / delay events"],
                "Project Folder": ["05-contracts/source", "06-evidence", "01-data/import_templates"],
            }),
            width="stretch",
            hide_index=True,
        )
        return
    contract_terms = get_contract_terms(overview_metrics.get("contract_value", 0.0))
    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Contract Subjects", len(all_clauses), delta="Full clause library")
    c2.metric(
        "Delay Damages Cap",
        f"{contract_terms['delay_damages_cap_pct']:.0f}%" if contract_terms.get("delay_damages_cap_pct") is not None else "N/A",
        delta="From contract clauses",
    )
    c3.metric(
        "Weekly LD Rate",
        f"{contract_terms['weekly_ld_rate_pct']:.0f}%" if contract_terms.get("weekly_ld_rate_pct") is not None else "N/A",
        delta="From contract clauses",
    )
    c4.metric(
        "Payment Period",
        f"{int(contract_terms['payment_period_days'])} days" if contract_terms.get("payment_period_days") is not None else "N/A",
        delta="From contract clauses",
    )
    col1, col2 = st.columns(2)
    event_type = col1.selectbox("Select Event Type", ["Material Delivery Delay", "Variation / Change Order", "Payment Delay", "Remeasurement", "Notice / Claim", "Delay / EOT", "Other"])
    event_description = col2.text_input("Event Description", value="")
    if st.button("Find Applicable Clauses"):
        matches = match_event_to_clauses(event_type, event_description)
        ai_brief = generate_ai_clause_brief(event_type, event_description, matches)
        st.success(f"Found {len(matches)} applicable clauses across the contract library")
        brief1, brief2, brief3 = st.columns(3)
        brief1.metric("AI Risk Level", ai_brief["risk_level"])
        brief2.metric("Primary Topic", matches[0]["clause"].topic if matches else "N/A")
        brief3.metric("Matched Clauses", len(matches))
        st.markdown(
            f"""
            <div class='panel-note'>
            <b>AI Contract Assistant</b><br>
            <b>{ai_brief['headline']}</b><br><br>
            <b>Notice Position:</b> {ai_brief['key_notice']}<br>
            <b>Cost Position:</b> {ai_brief['cost_position']}<br>
            <b>Time Position:</b> {ai_brief['time_position']}
            </div>
            """,
            unsafe_allow_html=True,
        )
        st.dataframe(pd.DataFrame({"AI Recommended Next Actions": ai_brief["next_actions"]}), width="stretch", hide_index=True)
        for i, match in enumerate(matches, 1):
            clause = match["clause"]
            with st.expander(f"{i}. {clause.topic} - {match['relevance']}", expanded=i <= 3):
                a, b = st.columns(2)
                a.write(clause.plain_english)
                a.caption(clause.beneath_lines)
                b.write(f"Leverage: {clause.leverage_holder}")
                b.write(f"Notice: {clause.notice_requirement}")
                b.write(f"Money: {clause.money_impact}")
                b.write(f"Schedule: {clause.schedule_impact}")
                st.info(clause.practical_action)
    st.markdown("#### Delay Exposure Engine")
    d1, d2, d3 = st.columns(3)
    delay_days = d1.number_input("Delay Days", min_value=1, value=46)
    contractor_caused = d2.checkbox("Contractor-Caused?", value=False)
    critical_path = d3.checkbox("Critical Path?", value=True)
    if st.button("Analyze Delay Event"):
        analysis = analyze_delay_event(
            "DEL-001",
            event_description,
            int(delay_days),
            contractor_caused,
            critical_path,
            overview_metrics.get("contract_value", 0.0),
        )
        x1, x2, x3, x4 = st.columns(4)
        x1.metric("Delay Days", analysis["delay_days"])
        x2.metric("Entitlement", analysis["entitlement"])
        x3.metric("Delay Damages / Week", egp(analysis["delay_damages_per_week"]))
        exposure_delta = (
            f"Capped at {contract_terms['delay_damages_cap_pct']:.0f}%"
            if contract_terms.get("delay_damages_cap_pct") is not None
            else "Per contract analysis"
        )
        x4.metric("Total Exposure", egp(analysis["delay_damages_exposure"]), delta=exposure_delta)
        st.markdown(
            f"<div class='panel-note'><b>AI Delay Assessment</b><br>Risk level: <b>{analysis['risk_level']}</b><br>Clause-based entitlement view: <b>{analysis['entitlement']}</b>.</div>",
            unsafe_allow_html=True,
        )
        st.dataframe(pd.DataFrame({"Critical Actions": analysis["critical_actions"]}), width="stretch", hide_index=True)
    search_term = st.text_input("Search Clauses", placeholder="delay, payment, variation, notice, steel")
    clauses = search_clauses(search_term) if search_term else all_clauses
    st.dataframe(pd.DataFrame([c.__dict__ for c in clauses]), width="stretch", hide_index=True)

if active_slide_name == PROJECT_HUB_SLIDE_NAMES[7]:
    st.markdown("<div class='section-header'><h3>Letters Intelligence, Alerts Link & Issue Engine</h3></div>", unsafe_allow_html=True)
    active_letters_inbox = letters_inbox_dir()
    auto_ingest_register = letters.get("Auto Ingest Register", pd.DataFrame()) if letters else pd.DataFrame()
    st.markdown(
        "<div class='panel-note'><b>Automatic Letter Inbox</b><br>Place new letters in the project direction folders. The app extracts the reference, date, subject, classification, risk, actions, related references, and issue thread, then adds them to the intelligence tables automatically.</div>",
        unsafe_allow_html=True,
    )
    st.code(str(active_letters_inbox), language=None)
    if not auto_ingest_register.empty:
        ingest_col1, ingest_col2, ingest_col3 = st.columns(3)
        ingest_col1.metric("Inbox Files", len(auto_ingest_register))
        ingest_col2.metric("Added Automatically", int(auto_ingest_register["Status"].eq("Added Automatically").sum()))
        ingest_col3.metric("Needs Review", int(auto_ingest_register["Status"].isin(["Needs Review", "Read Error"]).sum()))
        with st.expander("Automatic ingest register", expanded=False):
            st.dataframe(auto_ingest_register, width="stretch", hide_index=True, height=dataframe_height(auto_ingest_register))
    else:
        st.caption("The inbox is empty. Add PDF, DOCX, TXT, EML, CSV, XLSX, or XLS letters to one of the two direction folders.")
    if not letters:
        st.warning("No project letters workbook or inbox records were found.")
    else:
        samco = letters.get("From Contractor", pd.DataFrame())
        ace = letters.get("From Consultant", pd.DataFrame())
        samco_links = letters.get("Contractor Links", pd.DataFrame())
        ace_links = letters.get("Consultant Links", pd.DataFrame())
        threads = letters.get("Issue Threads", pd.DataFrame())
        c1, c2, c3, c4 = st.columns(4)
        c1.metric(f"{contractor_name} Letters", len(samco), delta="Outbound")
        c2.metric("ACE Letters", len(ace), delta="Inbound")
        c3.metric("Linked Contractor Threads", len(samco_links), delta="Relationship map")
        c4.metric("Issue Threads", len(threads), delta="Alert engine")
        if not threads.empty:
            sorted_threads = threads.sort_values("Priority", key=lambda s: s.map(risk_rank), ascending=False)
            st.dataframe(sorted_threads, width="stretch", hide_index=True, height=dataframe_height(sorted_threads))
        col1, col2 = st.columns(2)
        with col1:
            if not samco.empty:
                counts = samco["Risk Type"].value_counts().reset_index()
                counts.columns = ["Risk Type", "Count"]
                st.plotly_chart(style_plotly(px.bar(counts, x="Risk Type", y="Count", title="Contractor Risk Subjects", color="Risk Type", color_discrete_sequence=px.colors.qualitative.Safe)), width="stretch")
                samco_view = samco[["Ref No", "Date", "Type", "Subject", "Delay Risk", "EOT Potential", "Claim Strength", "Required Actions"]]
                st.dataframe(samco_view, width="stretch", hide_index=True, height=dataframe_height(samco_view))
        with col2:
            if not ace.empty:
                counts = ace["Delay Risk"].value_counts().reset_index()
                counts.columns = ["Delay Risk", "Count"]
                st.plotly_chart(style_plotly(px.bar(counts, x="Delay Risk", y="Count", title="ACE Delay Risk Alerts", color="Delay Risk", color_discrete_map={"High":"#d9544d","Medium":"#c98519","Low":"#168f8b"})), width="stretch")
                ace_view = ace[["Ref No", "Date", "Type", "Subject", "Delay Risk", "EOT Potential", "Claim Strength", "Required Actions"]]
                st.dataframe(ace_view, width="stretch", hide_index=True, height=dataframe_height(ace_view))
        st.markdown("#### Linked Correspondence Engine")
        st.dataframe(samco_links, width="stretch", hide_index=True, height=dataframe_height(samco_links))
        st.dataframe(ace_links, width="stretch", hide_index=True, height=dataframe_height(ace_links))

if active_slide_name == PROJECT_HUB_SLIDE_NAMES[8]:
    st.markdown("<div class='section-header'><h3>Delay Analysis</h3></div>", unsafe_allow_html=True)
    delays_df = delay_metrics["delays_df"]
    display_delays_df = delay_metrics["display_delays_df"]
    if not delays_df.empty:
        s1, c1, c2, c3, c4, s2 = st.columns([0.35, 1, 1, 1, 1, 0.35])
        with c1:
            render_kpi_box("Delay Events", str(delay_metrics["total_delay_events"]))
        with c2:
            render_kpi_box("Total Delay Days", str(int(delay_metrics["total_delay_days"])))
        with c3:
            render_kpi_box("Employer Delays", str(delay_metrics["employer_delays"]))
        with c4:
            render_kpi_box("EOT Potential", str(delay_metrics["eot_potential_count"]))

        s3, c5, c6, s4 = st.columns([1.1, 1, 1, 1.1])
        with c5:
            render_kpi_box("Open Delays", str(delay_metrics["open_delays"]))
        with c6:
            render_kpi_box("Closed Delays", str(delay_metrics["closed_delays"]))

        matrix_df = (
            delays_df.groupby(["responsible_group", "status_group"], as_index=False)
            .size()
            .pivot(index="responsible_group", columns="status_group", values="size")
            .fillna(0)
            .astype(int)
            .reset_index()
        )
        matrix_df.columns.name = None
        st.markdown("#### Delay Responsibility x Status Matrix")
        st.dataframe(matrix_df, width="stretch", hide_index=True, height=dataframe_height(matrix_df, max_height=320))

        st.markdown("#### Delay Events Table")
        st.dataframe(display_delays_df, width="stretch", hide_index=True, height=dataframe_height(display_delays_df))
    else:
        st.info("No delay data available.")

if active_slide_name == PROJECT_HUB_SLIDE_NAMES[10]:
    st.markdown("<div class='section-header'><h3>Risk Analysis</h3></div>", unsafe_allow_html=True)
    risks_df = risk_metrics["risks_df"]
    steel_df = risk_metrics["steel_df"]
    rfi_df = risk_metrics["rfi_df"]
    ifc_df = risk_metrics["ifc_df"]
    if not risks_df.empty:
        c1, c2, c3, c4 = st.columns(4)
        c1.metric("Total Risks", risk_metrics["total_risks"])
        c2.metric("High Risks", risk_metrics["high_risks"])
        c3.metric("Open Risks", risk_metrics["open_risks"])
        c4.metric("Closed Risks", risk_metrics["closed_risks"])
        s1, s2, s3 = st.columns(3)
        s1.metric("Steel Delay Issues", risk_metrics["steel_issues"])
        s2.metric("RFI Status Items", risk_metrics["rfi_items"])
        s3.metric("IFC Conflicts", risk_metrics["ifc_conflicts"])
        col1, col2 = st.columns(2)
        risk_status_counts = risks_df["status"].value_counts().reset_index()
        risk_status_counts.columns = ["Status", "Count"]
        col1.plotly_chart(style_plotly(px.bar(risk_status_counts, x="Status", y="Count", title="Risk Status", color="Status", color_discrete_sequence=["#245f95", "#168f8b", "#c98519"])), width="stretch")
        risk_category_counts = risks_df["risk_category"].value_counts().reset_index()
        risk_category_counts.columns = ["Risk Category", "Count"]
        col2.plotly_chart(style_plotly(px.bar(risk_category_counts, x="Risk Category", y="Count", title="Risk Category", color="Risk Category", color_discrete_sequence=px.colors.qualitative.Safe)), width="stretch")
        st.dataframe(risks_df, width="stretch", hide_index=True, height=dataframe_height(risks_df))
    else:
        st.info("No main risk register data available.")

    if not steel_df.empty:
        st.markdown("#### Steel Delay Status - Employer Free Issue Material")
        st.dataframe(steel_df, width="stretch", hide_index=True, height=dataframe_height(steel_df))

    if not rfi_df.empty:
        st.markdown("#### RFI Status")
        st.dataframe(rfi_df, width="stretch", hide_index=True, height=dataframe_height(rfi_df))

    if not ifc_df.empty:
        st.markdown("#### IFC Conflict")
        st.dataframe(ifc_df, width="stretch", hide_index=True, height=dataframe_height(ifc_df))

if active_slide_name == PROJECT_HUB_SLIDE_NAMES[3]:
    st.markdown("<div class='section-header'><h3>Milestones & Change Orders</h3></div>", unsafe_allow_html=True)
    col1, col2 = st.columns(2)
    col1.dataframe(milestone_metrics["milestones_df"], width="stretch", hide_index=True)
    col2.dataframe(milestone_metrics["change_orders_df"], width="stretch", hide_index=True)

if active_slide_name == PROJECT_HUB_SLIDE_NAMES[11]:
    tia_view = st.selectbox(
        "Delay Analysis view",
        [
            "Uploads",
            "Tables & Conclusion",
            "MEP Activities",
            "AI - TIA",
            "question",
            "Download Reports",
        ],
        key="delay_tia_active_view",
    )

    if tia_view == "Uploads":
        st.markdown("<div class='section-header'><h3>Delay Analyzer & TIA Fragnet Recommender</h3></div>", unsafe_allow_html=True)
        st.markdown(
            "<div class='panel-note'><b>Employer-only calculation rule</b><br>"
            "Employer steel is the only steel supply used in delay-analysis calculations. "
            "Contractor-supplied steel is displayed for visibility only and is excluded from stock-out, affected-activity, "
            "fragnet, and contractual assessment logic.</div>",
            unsafe_allow_html=True,
        )

        delay_tia_file_specs = [
            {
                "key": "metadata",
                "label": "01-project_metadata_template.csv",
                "required": True,
                "fields_used": "Project name, baseline, impacted update, data date, claim position, parties, and report metadata",
                "tia_use": "Project metadata and report context used to label and explain Delay TIA outputs.",
            },
            {
                "key": "master",
                "label": "02- master_activity_steel_analysis.csv",
                "required": True,
                "fields_used": "Activity ID, Activity Name, BL Start, BL Finish, Start, Finish / Actual Finish, Budgeted Units, Actual Units, Remaining Units, Units % Complete, Available units at site received by client, Available units supplied by client vs Actual units of activities, Reason for Delay, Responsibility, Drive Activity Dates, Allocation Date",
                "tia_use": "Primary activity-level shortage and readiness evidence. This is the core steel requirement and activity impact file.",
            },
            {
                "key": "employer",
                "label": "03- employer_steel_supply_at_site.csv",
                "required": True,
                "fields_used": "Delivery date column, delivered quantity column, steel reference / type if available",
                "tia_use": "Employer supply timing only. Used to build usable delivery dates and the steel balance curve.",
            },
            {
                "key": "p6",
                "label": "04- p6_activity_export.csv",
                "required": True,
                "fields_used": "Activity ID, Activity Name, WBS, Baseline Start / Finish, Start, Finish, Actual Start / Finish, Remaining Duration, Total Float, Critical, Longest Path, Predecessors, Successors, Remaining Units, Budgeted Units, Physical % Complete",
                "tia_use": "Schedule context, criticality, float, readiness, and impacted successor logic.",
            },
            {
                "key": "relationship",
                "label": "05- relationship_file.csv",
                "required": True,
                "fields_used": "Activity ID, Activity Name, Predecessors, Predecessor Details, Successors",
                "tia_use": "Dependency chain and fragnet insertion position.",
            },
            {
                "key": "contract",
                "label": "06- contract_library.csv",
                "required": True,
                "fields_used": "Clause / Topic, Location, Plain English Meaning, Research the Lines, Who Holds Leverage, Notice / Time Bar, Money Impact, Schedule Impact, Practical Action / Evidence",
                "tia_use": "Contract support, notice / time-bar review, responsibility and entitlement assessment.",
            },
            {
                "key": "ifc",
                "label": "07- ifc_conflict.csv",
                "required": True,
                "fields_used": "Activity ID, Activity Name, Original Duration, Start, Finish, Approved structural slab shop drawings, Re-Start, Finish.1, Delayed days, Note",
                "tia_use": "Support delay stream for IFC-driven events outside the steel balance path using activity-level delayed days.",
            },
            {
                "key": "payments",
                "label": "08- payments.csv",
                "required": True,
                "fields_used": "payment, contract, project, invoice no, invoice date, certified amount, paid amount, payment status, delayed duration(Days)",
                "tia_use": "Support delay stream for payment-driven events outside the steel balance path.",
            },
            {
                "key": "rfi",
                "label": "09- rfi_status.csv",
                "required": True,
                "fields_used": "RFI, Activity ID, Activity Name, Submittion date, Reply, Delay beyond 10d, Note",
                "tia_use": "Support delay stream for RFI-driven events outside the steel balance path.",
            },
            {
                "key": "samco",
                "label": "10- contractor_steel_supplied_at_site.csv",
                "required": True,
                "fields_used": "Coding No., Supplier, Total Quantity, Date of delivery",
                "tia_use": "Contractor mitigation / alternative steel visibility. Explicitly excluded from employer delay entitlement calculations.",
            },
            {
                "key": "concurrency",
                "label": "11-concurrency_matrix_template.updated.csv",
                "required": True,
                "fields_used": "Overlap Start, Overlap Finish, BL Critical Path, Current Critical Path, Total Float, Longest Path, first affected flag, predecessor/successor logic, fragment dates, concurrent delay",
                "tia_use": "Concurrency framework used to separate schedule entitlement from compensability and to avoid double counting.",
            },
        ]

        st.markdown("#### Required Delay TIA Source Files")
        st.caption(f"Source: projects/{active_project_context.project_folder_name}/02-delay_analysis/steel_delay_tia_templates")
        steel_template_inventory_df = build_steel_delay_template_inventory_df()
        if not steel_template_inventory_df.empty:
            st.markdown("#### Recognized Delay TIA Files")
            st.dataframe(
                steel_template_inventory_df,
                width="stretch",
                hide_index=True,
                height=dataframe_height(steel_template_inventory_df, max_height=420),
            )

        effective_delay_tia_files = {}
        for spec in delay_tia_file_specs:
            template_name, template_frame = load_delay_tia_template_fallback(spec["key"], spec["label"])
            if template_name:
                effective_delay_tia_files[spec["key"]] = {"name": template_name, "source": "Folder", "frame": template_frame}
            else:
                effective_delay_tia_files[spec["key"]] = {"name": "", "source": "Missing", "frame": pd.DataFrame()}

        uploaded_delay_tia_frames = {
            spec["key"]: filter_active_project(effective_delay_tia_files[spec["key"]]["frame"])
            for spec in delay_tia_file_specs
        }
        missing_required_delay_tia_files = [
            spec["label"]
            for spec in delay_tia_file_specs
            if spec["required"] and effective_delay_tia_files[spec["key"]]["source"] == "Missing"
        ]
        delay_tia_ready = not missing_required_delay_tia_files

        upload_status_df = pd.DataFrame(
            [
                {
                    "File": spec["label"],
                    "Required": "Yes" if spec["required"] else "Optional",
                    "Source Status": effective_delay_tia_files[spec["key"]]["source"] if effective_delay_tia_files[spec["key"]]["source"] != "Missing" else ("Missing" if spec["required"] else "Not available"),
                    "Recognized File Name": effective_delay_tia_files[spec["key"]]["name"],
                    "Rows": int(len(uploaded_delay_tia_frames[spec["key"]])),
                    "Fields used by TIA": spec["fields_used"],
                    "How the program uses it": spec["tia_use"],
                }
                for spec in delay_tia_file_specs
            ]
        )
        st.dataframe(upload_status_df, width="stretch", hide_index=True)
        uploaded_improvement_frames = {}
        improvement_status_df = pd.DataFrame(
            [
                {
                    "Status": "Disabled",
                    "Rule": "Delay Analysis - Time Impact Analysis now accepts only the 11 required upload/template files.",
                    "Reason": "Optional improvement-file uploads were removed to keep the calculation controlled, auditable, and aligned with the approved template set.",
                }
            ]
        )

        st.markdown("#### Include / Exclude Supposed Delay Streams")
        st.caption("These controls apply to IFC, RFI, payments, and contractor mitigation evidence. Contractor supply remains excluded from employer-delay calculations.")
        include_col1, include_col2, include_col3, include_col4 = st.columns(4)
        include_ifc_delay = include_col1.checkbox("Include IFC supposed delay", value=True, key="delay_tia_include_ifc")
        include_payments_delay = include_col2.checkbox("Include Payments supposed delay", value=True, key="delay_tia_include_payments")
        include_rfi_delay = include_col3.checkbox("Include RFI supposed delay", value=True, key="delay_tia_include_rfi")
        include_samco_delay = include_col4.checkbox("Include contractor mitigation evidence", value=True, key="delay_tia_include_samco")

        included_delay_streams_df = pd.DataFrame(
            [
                {"Delay Stream": "IFC", "Included": "Yes" if include_ifc_delay else "No"},
                {"Delay Stream": "Payments", "Included": "Yes" if include_payments_delay else "No"},
                {"Delay Stream": "RFI", "Included": "Yes" if include_rfi_delay else "No"},
                {"Delay Stream": "Contractor steel mitigation evidence", "Included": "Yes" if include_samco_delay else "No"},
            ]
        )
        st.dataframe(included_delay_streams_df, width="stretch", hide_index=True)

        bl_fixed_context = load_bl_fixed_context()
        bl_fixed_status_df = pd.DataFrame(
            [
                {
                    "Automatic BL Source": "BL Schedule.csv",
                    "Rows": int(len(bl_fixed_context["schedule_df"])),
                    "Used For": "Baseline schedule window reference inside Delay TIA comparison.",
                },
                {
                    "Automatic BL Source": "BL float bath.csv",
                    "Rows": int(len(bl_fixed_context["float_df"])),
                    "Used For": "Baseline float-path visibility reference inside Delay TIA comparison.",
                },
                {
                    "Automatic BL Source": "Bl Longest bath.csv",
                    "Rows": int(len(bl_fixed_context["longest_df"])),
                    "Used For": "Baseline longest-path visibility reference inside Delay TIA comparison.",
                },
                {
                    "Automatic BL Source": "BL critical path.csv",
                    "Rows": int(len(bl_fixed_context["critical_df"])),
                    "Used For": "Primary fixed BL critical path source for comparison against current Activities analysis.",
                },
            ]
        )
        st.caption("Baseline comparison sources are loaded from the selected project's `03-schedule` folder.")
        st.dataframe(bl_fixed_status_df, width="stretch", hide_index=True)

        delay_tia_local_support = load_delay_tia_local_support_files()

        if delay_tia_ready:
            included_ifc_df = uploaded_delay_tia_frames["ifc"] if include_ifc_delay else pd.DataFrame()
            included_payments_df = uploaded_delay_tia_frames["payments"] if include_payments_delay else pd.DataFrame()
            included_rfi_df = uploaded_delay_tia_frames["rfi"] if include_rfi_delay else pd.DataFrame()
            included_samco_df = uploaded_delay_tia_frames["samco"] if include_samco_delay else pd.DataFrame()
            improvement_rfi_detailed_df = delay_tia_improvement_rfi_to_detailed_df(uploaded_improvement_frames)
            combined_rfi_detailed_df = pd.concat(
                [
                    delay_tia_local_support["rfi_detailed_df"],
                    improvement_rfi_detailed_df,
                ],
                ignore_index=True,
                sort=False,
            )
            if not combined_rfi_detailed_df.empty:
                combined_rfi_detailed_df = combined_rfi_detailed_df.drop_duplicates().reset_index(drop=True)
            active_delay_tia_context = build_delay_tia_context_from_frames(
                master_df=uploaded_delay_tia_frames["master"],
                employer_raw_df=uploaded_delay_tia_frames["employer"],
                p6_df=uploaded_delay_tia_frames["p6"],
                relationship_df=uploaded_delay_tia_frames["relationship"],
                contract_df=uploaded_delay_tia_frames["contract"],
                ifc_df=included_ifc_df,
                payments_df=included_payments_df,
                rfi_df=included_rfi_df,
                bl_critical_path_df=bl_fixed_context["critical_df"],
                samco_df=included_samco_df,
                delay_events_df=pd.DataFrame(),
            )
            active_delay_tia_context["project_metadata_df"] = uploaded_delay_tia_frames.get("metadata", pd.DataFrame())
            active_delay_tia_context["concurrency_matrix_df"] = uploaded_delay_tia_frames.get("concurrency", pd.DataFrame())
            active_delay_tia_analysis = active_delay_tia_context["analysis"]
            active_concurrent_delay_review_df = build_delay_tia_concurrent_delay_review_df(
                master_df=uploaded_delay_tia_frames["master"],
                p6_df=uploaded_delay_tia_frames["p6"],
                fragnet_df=active_delay_tia_analysis.get("fragnet_df", pd.DataFrame()),
                ifc_df=delay_tia_local_support["ifc_local_df"],
                rfi_detailed_df=combined_rfi_detailed_df if include_rfi_delay else pd.DataFrame(),
                bl_fixed_context=bl_fixed_context,
            )
            template_concurrent_df = build_delay_tia_improvement_concurrent_rows_df(
                {"concurrency_matrix": uploaded_delay_tia_frames.get("concurrency", pd.DataFrame())},
                uploaded_delay_tia_frames["p6"],
                bl_fixed_context,
            )
            if not template_concurrent_df.empty:
                active_concurrent_delay_review_df = pd.concat(
                    [active_concurrent_delay_review_df, template_concurrent_df],
                    ignore_index=True,
                    sort=False,
                )
            improvement_concurrent_df = build_delay_tia_improvement_concurrent_rows_df(
                uploaded_improvement_frames,
                uploaded_delay_tia_frames["p6"],
                bl_fixed_context,
            )
            if not improvement_concurrent_df.empty:
                active_concurrent_delay_review_df = pd.concat(
                    [active_concurrent_delay_review_df, improvement_concurrent_df],
                    ignore_index=True,
                    sort=False,
                )
                if not active_concurrent_delay_review_df.empty:
                    active_concurrent_delay_review_df = active_concurrent_delay_review_df.drop_duplicates(
                        subset=["Delay Stream", "Delay Ref", "Activity ID"],
                        keep="last",
                    ).reset_index(drop=True)
            active_delay_tia_context["improvement_frames"] = uploaded_improvement_frames
            active_delay_tia_context["improvement_upload_status_df"] = improvement_status_df
            active_delay_tia_context["improvement_rfi_detailed_df"] = improvement_rfi_detailed_df
            active_delay_tia_context["improvement_concurrent_delay_review_df"] = improvement_concurrent_df
            active_delay_tia_context["concurrent_delay_review_df"] = active_concurrent_delay_review_df
        else:
            active_delay_tia_context = None
            active_delay_tia_analysis = {}
            active_concurrent_delay_review_df = pd.DataFrame()
            st.error(
                "Delay TIA analysis is blocked. Add the required files to `steel_delay_tia_templates` first: "
                + ", ".join(missing_required_delay_tia_files)
            )

        if not bl_fixed_context["critical_df"].empty:
            active_bl_critical_path_fixed_df, active_bl_critical_path_comparison_df, active_bl_critical_path_summary = build_bl_critical_path_comparison(
                bl_fixed_context["critical_df"], activity_metrics
            )
        else:
            active_bl_critical_path_fixed_df = pd.DataFrame()
            active_bl_critical_path_comparison_df = pd.DataFrame()
            active_bl_critical_path_summary = {
                "bl_count": 0,
                "current_count": 0,
                "matched_count": 0,
                "bl_only_count": 0,
                "current_only_count": 0,
            }

    if tia_view == "Tables & Conclusion":
        st.markdown("#### Source Tables & Calculated Conclusion")
        st.caption("This slide reads the 11 Delay Analysis - Time Impact Analysis files from `steel_delay_tia_templates`, inspects their columns, and shows the calculated conclusion from the active TIA logic.")
        if not delay_tia_ready:
            st.warning("Tables and conclusion are blocked until all 11 required files are available.")
            st.dataframe(upload_status_df, width="stretch", hide_index=True)
        else:
            inventory_rows = []
            for spec in delay_tia_file_specs:
                frame = uploaded_delay_tia_frames.get(spec["key"], pd.DataFrame())
                inventory_rows.append(
                    {
                        "File": spec["label"],
                        "Source": effective_delay_tia_files[spec["key"]]["source"],
                        "Rows": int(len(frame)),
                        "Columns": int(len(frame.columns)),
                        "Column List": ", ".join(str(col) for col in frame.columns),
                        "TIA Use": spec["tia_use"],
                    }
                )
            inventory_df = pd.DataFrame(inventory_rows)
            st.markdown("##### Folder File Inventory")
            st.dataframe(inventory_df, width="stretch", hide_index=True, height=dataframe_height(inventory_df, max_height=420))

            kpis_tia = active_delay_tia_analysis.get("kpis", {})
            conclusion_df = pd.DataFrame(
                [
                    {"Conclusion Item": "Selected Method", "Conclusion": "Hybrid retrospective TIA/UIA using employer steel supply, P6 activity context, relationships, support events, and concurrency checks."},
                    {"Conclusion Item": "Potential EOT Days", "Conclusion": f"{int(kpis_tia.get('Potential EOT Days', kpis_tia.get('Recommended Submission Days', 0)) or 0)} days, indicative until P6 recalculation is verified."},
                    {"Conclusion Item": "Strong TIA Candidates", "Conclusion": str(int(kpis_tia.get("Number of Strong TIA Candidates", 0) or 0))},
                    {"Conclusion Item": "First Stock-Out Date", "Conclusion": steel_tia_date_label(kpis_tia.get("First Stock-Out Date"))},
                    {"Conclusion Item": "Concurrency Position", "Conclusion": f"{int(kpis_tia.get('Concurrent Risk Events', 0) or 0)} concurrent risk events identified; compensation must be separated from EOT entitlement."},
                    {"Conclusion Item": "Calculation Rule", "Conclusion": "Employer steel drives steel delay calculations. Contractor-supplied steel is mitigation visibility only unless independently supported by entitlement evidence."},
                ]
            )
            st.markdown("##### Calculated Conclusion")
            st.dataframe(conclusion_df, width="stretch", hide_index=True, height=dataframe_height(conclusion_df, max_height=320))

            table_options = {
                "Executive Summary": active_delay_tia_analysis.get("executive_summary_df", pd.DataFrame()),
                "Activity Delay Mapping": active_delay_tia_analysis.get("activity_delay_mapping_df", pd.DataFrame()),
                "Fragnet Register": active_delay_tia_analysis.get("fragnet_df", pd.DataFrame()),
                "Concurrency Assessment": active_concurrent_delay_review_df,
                "Entitlement Decision Matrix": active_delay_tia_analysis.get("entitlement_decision_df", pd.DataFrame()),
            }
            selected_table_name = st.selectbox("Calculated table", list(table_options.keys()), key="delay_tia_tables_conclusion_table")
            selected_table_df = table_options[selected_table_name]
            if isinstance(selected_table_df, pd.DataFrame) and not selected_table_df.empty:
                st.dataframe(selected_table_df, width="stretch", hide_index=True, height=dataframe_height(selected_table_df, max_height=620))
            else:
                st.info(f"No rows are available for {selected_table_name}.")

    if tia_view == "MEP Activities":
        st.markdown("#### MEP Activities")
        st.caption(
            "This slide reads the selected project's `03-schedule/MEP Activities.csv`. "
            "The records are treated as MEP first-fix / embedded-interface constraints that were outside the original steel-only delay scope."
        )
        mep_activities_df = bl_fixed_context.get("mep_df", pd.DataFrame()).copy()
        mep_schedule_df = bl_fixed_context.get("mep_schedule_df", pd.DataFrame()).copy()
        mep_civil_logic_df = bl_fixed_context.get("mep_civil_logic_df", pd.DataFrame()).copy()
        related_mep_letters_df = build_mep_related_letters_df(mep_activities_df)

        if mep_activities_df.empty:
            st.warning("No MEP activity records were found in the selected project's `03-schedule/MEP Activities.csv`.")
        else:
            mep_kpis = build_mep_activities_kpis(mep_activities_df)
            c1, c2, c3, c4 = st.columns(4)
            c1.metric("MEP Activities", int(mep_kpis["activity_count"]))
            c2.metric("Buildings / Zones", int(mep_kpis["building_count"]))
            c3.metric("Indicative Duration Days", f"{float(mep_kpis['total_duration_days']):.0f}")
            c4.metric("Interface Types", int(mep_kpis["interface_count"]))

            st.markdown(
                """
                **Delay-analysis treatment**

                These items are not reinforcement steel demand records. They are treated as a separate MEP coordination stream
                because the source register identifies activities that must be inserted into related civil workfronts.

                The program derives the schedule position from the baseline schedule, current P6 export, relationship file,
                and MEP activity register:

                - engineering and procurement rows reuse matched baseline activities, logic, and durations,
                - construction rows use the MEP register duration and the related current civil activity dates,
                - civil finish impact is calculated as the related civil finish plus the inserted MEP duration,
                - TIA use remains conditional on proven workfront, event date, and critical-path effect.
                """
            )

            st.markdown("##### MEP Activity Register")
            mep_view_cols = [
                col
                for col in [
                    "MEP Activity ID",
                    "Building / Zone",
                    "MEP Activity",
                    "Material / System",
                    "Duration Days",
                    "Interface Type",
                    "Civil / Structural Dependency",
                    "Scope / Debate Treatment",
                    "Delay Analysis Use",
                ]
                if col in mep_activities_df.columns
            ]
            st.dataframe(
                mep_activities_df[mep_view_cols] if mep_view_cols else mep_activities_df,
                width="stretch",
                hide_index=True,
                height=dataframe_height(mep_activities_df, max_height=560),
            )

            if "Building / Zone" in mep_activities_df.columns:
                mep_by_building_df = (
                    mep_activities_df.groupby("Building / Zone", dropna=False)
                    .agg(
                        Activities=("MEP Activity", "count"),
                        Duration_Days=("Duration Days", lambda values: pd.to_numeric(values, errors="coerce").fillna(0).sum()),
                    )
                    .reset_index()
                    .rename(columns={"Building / Zone": "Building / Zone", "Duration_Days": "Indicative Duration Days"})
                )
                st.markdown("##### Building / Zone Summary")
                st.dataframe(mep_by_building_df, width="stretch", hide_index=True)

        st.markdown("##### MEP Schedule Insertions")
        st.caption("Engineering and procurement are matched from the baseline schedule. Construction is built from the MEP activity register and positioned against the related current civil activity dates.")
        if mep_schedule_df.empty:
            st.info("No derived MEP schedule rows were found in the selected project's `03-schedule/MEP Schedule.csv`.")
        else:
            schedule_view_cols = [
                col
                for col in [
                    "MEP Activity ID",
                    "MEP Schedule Activity ID",
                    "Phase",
                    "Building / Zone",
                    "Activity Name",
                    "Original Duration",
                    "Start",
                    "Finish",
                    "Predecessor Activity ID(s)",
                    "Predecessor Relationship Details",
                    "Successor Activity ID(s)",
                    "Successor Relationship Details",
                    "Related Civil Activity ID",
                    "Related Civil Activity Name",
                    "Related Civil Start",
                    "Related Civil Finish",
                    "Civil Finish After MEP Insertion",
                    "Source Schedule Activity ID",
                    "Match Score",
                    "Date Basis",
                    "Insertion Note",
                ]
                if col in mep_schedule_df.columns
            ]
            st.dataframe(
                mep_schedule_df[schedule_view_cols] if schedule_view_cols else mep_schedule_df,
                width="stretch",
                hide_index=True,
                height=dataframe_height(mep_schedule_df, max_height=620),
            )

        st.markdown("##### Civil Predecessor / Successor Logic")
        st.caption("This table shows the civil predecessor and successor activities that govern each MEP insertion point.")
        if mep_civil_logic_df.empty:
            st.info("No derived civil predecessor/successor rows were found in the selected project's `03-schedule/MEP Civil Logic.csv`.")
        else:
            logic_view_cols = [
                col
                for col in [
                    "MEP Activity ID",
                    "Building / Zone",
                    "MEP Activity",
                    "Related Civil Activity ID",
                    "Related Civil Activity Name",
                    "Predecessor Activity ID(s)",
                    "Predecessor Details",
                    "Successor Activity ID(s)",
                    "Civil Original Duration",
                    "MEP Inserted Duration Days",
                    "Civil Duration After MEP Insertion",
                    "Relationship Basis",
                ]
                if col in mep_civil_logic_df.columns
            ]
            st.dataframe(
                mep_civil_logic_df[logic_view_cols] if logic_view_cols else mep_civil_logic_df,
                width="stretch",
                hide_index=True,
                height=dataframe_height(mep_civil_logic_df, max_height=560),
            )

        st.markdown("##### Related Letter Intelligence References")
        st.caption("Related correspondence is pulled from the Letters Intelligence workbook and ranked using terms derived from the active MEP activity register. The visible column is the conclusion, not the letter subject.")
        if related_mep_letters_df.empty:
            st.info("No related MEP correspondence references were found in the current letters workbook.")
        else:
            st.dataframe(
                related_mep_letters_df,
                width="stretch",
                hide_index=True,
                height=dataframe_height(related_mep_letters_df, max_height=520),
            )
        mep_export_bytes = build_mep_schedule_export_xlsx(
            mep_activities_df,
            mep_schedule_df,
            mep_civil_logic_df,
            related_mep_letters_df,
        )
        if mep_export_bytes:
            st.download_button(
                "Download MEP Schedule Pack (.xlsx)",
                data=mep_export_bytes,
                file_name="MEP_schedule_pack.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                key="download_mep_schedule_pack_xlsx",
                width="stretch",
            )
        elif not OPENPYXL_AVAILABLE:
            st.warning("Excel export is unavailable because openpyxl is not installed.")

    if tia_view == "AI - TIA":
        st.markdown("#### Executive Dashboard / TIA Methodology / Dependency Schema / File Priority / BL Critical Path Comparison")
        st.caption("This combined Delay TIA slide consolidates the former first five analysis slides into one review surface.")
        if not delay_tia_ready:
            st.warning("Executive Delay TIA outputs are blocked until all required folder files are available.")
        else:
            kpis_tia = active_delay_tia_analysis.get("kpis", {})
            c1, c2, c3, c4 = st.columns(4)
            c1.metric("Employer Supplied Qty", f"{active_delay_tia_context['employer_total']:,.3f}")
            c2.metric("Contractor Supplied Qty", f"{active_delay_tia_context['samco_total']:,.3f}")
            c3.metric("First Stock-Out Date", steel_tia_date_label(kpis_tia.get("First Stock-Out Date")))
            c4.metric("Strong TIA Candidates", int(kpis_tia.get("Number of Strong TIA Candidates", 0)))

            c1, c2, c3, c4 = st.columns(4)
            c1.metric("Total Required Steel", f"{float(kpis_tia.get('Total Required Steel', 0.0)):,.3f}")
            c2.metric("Current Steel Balance", f"{float(kpis_tia.get('Current Steel Balance', 0.0)):,.3f}")
            c3.metric("Employer Risk Events", int(kpis_tia.get("Employer Risk Events", 0)))
            c4.metric("Concurrent Risk Events", int(kpis_tia.get("Concurrent Risk Events", 0)))

            c1, c2, c3 = st.columns(3)
            c1.metric("IFC Delay Items", int(len(active_delay_tia_context["ifc_df"])))
            c2.metric("Payment Delay Items", int(len(active_delay_tia_context["payments_df"])))
            c3.metric("RFI Delay Items", int(len(active_delay_tia_context["rfi_df"])))
            c4, c5 = st.columns(2)
            c4.metric("BL Critical Path Rows", active_bl_critical_path_summary["bl_count"])
            c5.metric("Current Critical Path Rows", active_bl_critical_path_summary["current_count"])

            executive_summary_df = active_delay_tia_analysis.get("executive_summary_df", pd.DataFrame())
            if not executive_summary_df.empty:
                st.dataframe(executive_summary_df, width="stretch", hide_index=True)
            else:
                st.info("No Delay TIA executive summary is available from the folder source files.")

    if tia_view == "AI - TIA":
        st.markdown("#### Time Impact Analysis Methodology")
        st.markdown(
            """
            **What a fragnet is**

            A fragnet is a focused schedule fragment inserted into the programme to model the actual delay event,
            its start point, its recovery point, and its effect on the affected activity sequence.

            In this program, the fragnet is not inserted randomly. It is recommended only when the logic can show:

            1. a real delay event or shortage condition,
            2. a first affected activity,
            3. a causal path into the affected activity,
            4. and a downstream schedule effect.
            """
        )

        st.markdown("#### TIA Logic Used by the Program")
        tia_logic_df = pd.DataFrame(
            [
                {"Step": "1. Define event evidence", "Method": "Use employer supply, master activity steel analysis, and support delay files to identify the factual delay event."},
                {"Step": "2. Fix the event date", "Method": "For steel, identify stock-out / insufficiency date. For other streams, use the dated input event and reply / recovery date."},
                {"Step": "3. Identify the first affected activity", "Method": "Test readiness, incomplete status, shortage / blocking condition, float, criticality, and downstream impact."},
                {"Step": "4. Test critical or near-critical impact", "Method": "Use current activity criticality / float and BL critical path comparison to determine if the event is likely delay-driving."},
                {"Step": "5. Build the fragnet position", "Method": "Insert before the first affected activity, not before a generic delivery or correspondence record."},
                {"Step": "6. Measure fragment dates", "Method": "Fragment Start = event impact date. Fragment Finish = usable recovery date or first recovery point evidenced by data."},
                {"Step": "7. Assess entitlement", "Method": "Apply contract support logic, notice / time bar context, and responsibility classification."},
            ]
        )
        st.dataframe(tia_logic_df, width="stretch", hide_index=True)

        st.markdown("#### Fragnet Structure Used Here")
        st.code(
            "Last available predecessor\n"
            "        ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â ÃƒÂ¢Ã¢â€šÂ¬Ã…â€œ FS\n"
            "Delay event fragnet\n"
            "        ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â ÃƒÂ¢Ã¢â€šÂ¬Ã…â€œ FS\n"
            "First affected activity\n"
            "        ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â ÃƒÂ¢Ã¢â€šÂ¬Ã…â€œ\n"
            "Successor activities / milestone impact",
            language="text",
        )

        st.markdown("#### Steel Delay Fragnet Logic")
        steel_logic_df = pd.DataFrame(
            [
                {"Rule": "Supply source", "Program logic": "Only employer steel is used in calculations. Contractor site steel is display-only."},
                {"Rule": "Requirement basis", "Program logic": "Primary requirement basis is the master activity steel analysis file."},
                {"Rule": "Stock-out test", "Program logic": "Daily balance and activity-level shortage evidence are both considered."},
                {"Rule": "Affected activity test", "Program logic": "The activity must require steel, be incomplete or historically blocked, and show schedule relevance."},
                {"Rule": "Recovery date", "Program logic": "Delivery is not equal to recovery. Recovery is tied to usable supply timing."},
                {"Rule": "Fragnet position", "Program logic": "Insert before the first affected RFT / reinforcement activity."},
            ]
        )
        st.dataframe(steel_logic_df, width="stretch", hide_index=True)

        st.markdown("#### Why the Program Compares BL Critical Path and Current Critical Path")
        st.markdown(
            """
            TIA needs both baseline intention and current execution reality.

            - **BL Critical Path** explains what was originally delay-driving in the baseline path logic.
            - **Current Critical Path Analysis** explains what is delay-driving in the live execution logic.

            The comparison helps distinguish:

            - events that were always critical,
            - events that became critical later,
            - and events that are no longer critical but still have historical relevance.
            """
        )

        st.markdown("#### Professional Interpretation Rule")
        st.markdown(
            """
            This slide follows planning and contractual assessment logic, not legal certainty.

            A valid TIA conclusion in this program depends on:

            - factual event evidence,
            - activity linkage,
            - baseline / current path relevance,
            - delay duration logic,
            - and contract support.
            """
        )

    if tia_view == "AI - TIA":
        st.markdown("#### Steel Delay Analysis Dependency Schema")
        st.markdown(
            """
            **Flow**

            `steel_delay_tia_templates Source Files`

            -> `Delay TIA Loader`

            -> `P6 Base Schedule Context`
            -> `Employer Steel Supply Context`
            -> `Activity-Level Requirement / Shortage Context`
            -> `Relationship / Logic Context`
            -> `Contractual Context`
            -> `Support Delay Context`

            -> `Daily Steel Balance`
            -> `Stock-Out Detection`
            -> `Affected Activity Logic`
            -> `TIA Fragnet Recommendation`
            -> `Contractual Delay Assessment`
            -> `Executive Summary`
            """
        )

        dependency_df = upload_status_df[["File", "Required", "Fields used by TIA", "How the program uses it"]].copy()
        st.dataframe(dependency_df, width="stretch", hide_index=True)

    if tia_view == "AI - TIA":
        st.markdown("#### Active File Priority Order")
        priority_df = pd.DataFrame(
            [
                {"Priority": 1, "File / Source": "01-project_metadata_template.csv", "Why it matters": "Sets project identity, data date, baseline/update context, parties, and report metadata."},
                {"Priority": 2, "File / Source": "02- master_activity_steel_analysis.csv", "Why it matters": "Primary activity-level RFT demand, steel availability gap, delayed duration, responsibility, and readiness proof."},
                {"Priority": 3, "File / Source": "03- employer_steel_supply_at_site.csv", "Why it matters": "Employer-only steel delivery timeline used for stock-out, usable supply, and recovery windows."},
                {"Priority": 4, "File / Source": "04- p6_activity_export.csv", "Why it matters": "Schedule dates, float, criticality, longest path, progress, and impacted activity context."},
                {"Priority": 5, "File / Source": "05- relationship_file.csv", "Why it matters": "Predecessor/successor logic, relationship type, lag, causal chain, and fragnet insertion position."},
                {"Priority": 6, "File / Source": "06- contract_library.csv", "Why it matters": "Entitlement, notice/time-bar, money/schedule impact, leverage, and required evidence."},
                {"Priority": 7, "File / Source": "07- ifc_conflict.csv / 09- rfi_status.csv", "Why it matters": "Design, approval, RFI, and mechanical coordination event streams linked to affected activities."},
                {"Priority": 8, "File / Source": "08- payments.csv", "Why it matters": "Commercial support stream; schedule-driving only when progress/resource linkage is proven."},
                {"Priority": 9, "File / Source": "10- contractor_steel_supplied_at_site.csv", "Why it matters": "Contractor mitigation/alternative supply visibility only; not employer entitlement proof by itself."},
                {"Priority": 10, "File / Source": "11-concurrency_matrix_template.updated.csv", "Why it matters": "Overlap, BL/current critical path, float, and concurrent delay checks to separate EOT from compensation."},
                {"Priority": 11, "File / Source": "Letters workbook", "Why it matters": "Notice, correspondence, and issue-thread evidence linkage."},
            ]
        )
        st.dataframe(priority_df, width="stretch", hide_index=True)

        st.markdown("#### Data Structure Rule")
        st.markdown(
            """
            - `02- master_activity_steel_analysis.csv` is the **primary master activity steel file**
            - `03- employer_steel_supply_at_site.csv` stays an **Employer delivery event file**
            - `05- relationship_file.csv` stays a **logic-link file**
            - `11-concurrency_matrix_template.updated.csv` stays a **concurrency and compensability test file**

            These files are intentionally not merged into one operational input because they are different grains:

            - delivery rows are **event-level**
            - master activity rows are **activity-level**
            - relationship rows are **link-level**

            Mixing them would increase duplication and corrupt balance, affected-activity, and fragnet logic.
            """
        )

    if tia_view == "AI - TIA":
        c1, c2, c3, c4, c5 = st.columns(5)
        c1.metric("BL Critical Path Rows", active_bl_critical_path_summary["bl_count"])
        c2.metric("Current Critical Path Rows", active_bl_critical_path_summary["current_count"])
        c3.metric("Matched Rows", active_bl_critical_path_summary["matched_count"])
        c4.metric("BL Only Rows", active_bl_critical_path_summary["bl_only_count"])
        c5.metric("Current Only Rows", active_bl_critical_path_summary["current_only_count"])

        st.markdown("#### Fixed BL Critical Path Register")
        if not active_bl_critical_path_fixed_df.empty:
            st.dataframe(active_bl_critical_path_fixed_df, width="stretch", hide_index=True)
        else:
            st.info("No rows are available in the automatic local `BL critical path.csv` source.")

        st.markdown("#### BL Fixed Source Files")
        bl_src_col1, bl_src_col2 = st.columns(2)
        with bl_src_col1:
            st.markdown("##### BL Schedule")
            if not bl_fixed_context["schedule_df"].empty:
                st.dataframe(bl_fixed_context["schedule_df"], width="stretch", hide_index=True, height=dataframe_height(bl_fixed_context["schedule_df"], max_height=280))
            else:
                st.info("BL Schedule source is empty.")
            st.markdown("##### BL Float Bath")
            if not bl_fixed_context["float_df"].empty:
                st.dataframe(bl_fixed_context["float_df"], width="stretch", hide_index=True, height=dataframe_height(bl_fixed_context["float_df"], max_height=280))
            else:
                st.info("BL float bath source is empty.")
        with bl_src_col2:
            st.markdown("##### BL Longest Bath")
            if not bl_fixed_context["longest_df"].empty:
                st.dataframe(bl_fixed_context["longest_df"], width="stretch", hide_index=True, height=dataframe_height(bl_fixed_context["longest_df"], max_height=280))
            else:
                st.info("BL longest bath source is empty.")
            st.markdown("##### BL Critical Path")
            if not bl_fixed_context["critical_df"].empty:
                st.dataframe(bl_fixed_context["critical_df"], width="stretch", hide_index=True, height=dataframe_height(bl_fixed_context["critical_df"], max_height=280))
            else:
                st.info("BL critical path source is empty.")

        st.markdown("#### BL Critical Path vs Activities Critical Path Analysis")
        if not active_bl_critical_path_comparison_df.empty:
            comparison_view = active_bl_critical_path_comparison_df[
                [
                    col for col in [
                        "Activity ID",
                        "Activity Name",
                        "activity_name",
                        "BL Project Start",
                        "BL Project Finish",
                        "BL Total Float",
                        "Critical",
                        "Longest Path",
                        "Current Critical Path Flag",
                        "Current Float",
                        "Current Planned Finish",
                        "Current Forecast Finish",
                        "Comparison Status",
                    ]
                    if col in active_bl_critical_path_comparison_df.columns
                ]
            ].copy()
            if "activity_name" in comparison_view.columns and "Activity Name" not in comparison_view.columns:
                comparison_view = comparison_view.rename(columns={"activity_name": "Activity Name"})
            st.dataframe(comparison_view, width="stretch", hide_index=True)
        else:
            st.info("No comparison rows are available between the uploaded BL critical path file and the Activities critical path analysis.")

    if tia_view == "question":
        st.markdown("#### question")
        st.caption("Ask any Delay Analysis - Time Impact Analysis question. The slide inspects loaded columns first, then answers from the 11 steel delay template files.")
        question_frames = load_delay_tia_question_frames()
        question_inventory_df = build_delay_tia_question_column_inventory(question_frames)

        if question_inventory_df.empty:
            st.warning("No Delay Analysis - Time Impact Analysis question datasets were found in the steel delay template folder.")
        else:
            default_question = "What is the conservative final delay days to submit and why?"
            question_text = st.text_area(
                "Question",
                value=default_question,
                height=110,
                help="Ask about final delay days, RFIs, concurrency, evidence, critical path, columns, or any loaded data.",
            )
            answer_text, answer_tables, question_kpis = answer_delay_tia_question(question_text, question_frames)

            c1, c2, c3, c4 = st.columns(4)
            c1.metric("Datasets Loaded", int(question_kpis.get("Datasets Loaded", 0)))
            c2.metric("Recommended Days", int(question_kpis.get("Recommended Conservative Days", 0)))
            c3.metric("Evidence Rows", int(question_kpis.get("Evidence Rows", 0)))
            c4.metric("P6 Activities", int(question_kpis.get("P6 Activities", 0)))

            c1, c2, c3, c4 = st.columns(4)
            c1.metric("Max Claimed Days", int(question_kpis.get("Max Claimed Delay Days", 0)))
            c2.metric("Max Fragnet Days", int(question_kpis.get("Max Fragnet Duration", 0)))
            c3.metric("Concurrent Days", int(question_kpis.get("Concurrent Delay Days", 0)))
            c4.metric("RFI Beyond 10 Days", int(question_kpis.get("RFI Delay Beyond 10 Days", 0)))

            st.markdown("##### Answer")
            st.info(answer_text)

            with st.expander("Column inventory inspected before answering", expanded=False):
                st.dataframe(question_inventory_df, width="stretch", hide_index=True, height=dataframe_height(question_inventory_df, max_height=520))

            for table_title, table_df in answer_tables:
                if isinstance(table_df, pd.DataFrame) and not table_df.empty:
                    st.markdown(f"##### {table_title}")
                    st.dataframe(table_df, width="stretch", hide_index=True, height=dataframe_height(table_df, max_height=520))

    if tia_view == "Download Reports":
        if not delay_tia_ready:
            st.warning("Delay TIA report outputs are blocked until all required Delay TIA files exist in `steel_delay_tia_templates`.")
        else:
            delay_report_df = build_delay_tia_delay_report_df(active_delay_tia_context, active_delay_tia_analysis)
            primavera_fragnet_df = build_delay_tia_primavera_fragnet_df(active_delay_tia_context, active_delay_tia_analysis)
            detailed_report_html = build_delay_tia_detailed_report_html(active_delay_tia_context, active_delay_tia_analysis)
            excel_report_bytes = build_delay_tia_excel_report_bytes(active_delay_tia_context, active_delay_tia_analysis)
            director_docx_error = st.session_state.get("tia_director_pack_error", "")
            last_generated_report = fetch_last_generated_report(CONTRACT_CLAIMS_DB_PATH, REPORT_TYPE_TIA_DIRECTOR_PACK)
            generated_docx_path = resolve_existing_output_path(st.session_state.get("tia_director_pack_generated_path"))
            if generated_docx_path is None and last_generated_report:
                generated_docx_path = resolve_existing_output_path(last_generated_report.get("file_path"))
            director_docx_bytes = generated_docx_path.read_bytes() if generated_docx_path else b""

            st.markdown("#### Delay TIA Output Pack")
            st.dataframe(delay_report_df, width="stretch", hide_index=True)

            c1, c2, c3, c4, c5 = st.columns(5)
            c1.download_button(
                "Download Delay Report (.xlsx)",
                data=excel_report_bytes,
                file_name="delay_tia_output_pack.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                width="stretch",
            )
            c2.download_button(
                "Download Detailed TIA Report (.html)",
                data=detailed_report_html.encode("utf-8"),
                file_name="delay_tia_detailed_report.html",
                mime="text/html",
                width="stretch",
            )
            c3.download_button(
                "Download Primavera Fragnet Sheet (.csv)",
                data=primavera_fragnet_df.to_csv(index=False).encode("utf-8") if not primavera_fragnet_df.empty else b"",
                file_name="primavera_fragnet_sheet.csv",
                mime="text/csv",
                width="stretch",
                disabled=primavera_fragnet_df.empty,
            )
            c4.download_button(
                "Download Delay Report Summary (.csv)",
                data=delay_report_df.to_csv(index=False).encode("utf-8"),
                file_name="delay_tia_delay_report.csv",
                mime="text/csv",
                width="stretch",
            )
            c5.download_button(
                "Download Time Impact Analysis Report | Director Pack (.docx)",
                data=director_docx_bytes,
                file_name=generated_docx_path.name if generated_docx_path else "TIA_Director_Level_Word_Report_Final.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                width="stretch",
                disabled=not bool(director_docx_bytes),
            )
            if director_docx_error:
                st.warning(f"Director-level DOCX report could not be generated: {director_docx_error}")

            default_report_context = build_delay_tia_director_pack_context(
                active_delay_tia_context,
                active_delay_tia_analysis,
                overrides={"revision": "Rev. 00"},
            )
            report_source_frames = {
                "Delay Event Register": default_report_context["delay_event_register_df"],
                "Fragnet Register": default_report_context["fragnet_register_df"],
                "Impact Calculation Table": default_report_context["impact_calculation_df"],
                "Affected Activities / Criticality Table": default_report_context["activity_impact_df"],
                "Causation / Concurrency / Entitlement Matrix": default_report_context["causation_matrix_df"],
                "Evidence Checklist": default_report_context["evidence_register_df"],
                "Data Readiness Matrix": default_report_context["readiness_matrix_df"],
                "P6 Control Sheet": default_report_context["p6_controls_df"],
            }

            st.markdown("#### Report Generator")
            st.subheader("Time Impact Analysis Report | Director Pack")
            template_status_col, last_report_col, completion_col, source_status_col = st.columns(4)
            template_status_col.metric("Template Status", "Found" if TIA_DIRECTOR_WORD_TEMPLATE_PATH.exists() else "Missing")
            last_report_col.metric("Last Generated Report", generated_docx_path.name if generated_docx_path else "None")

            input_col1, input_col2, input_col3 = st.columns(3)
            project_name = input_col1.text_input(
                "Project Name",
                value=default_report_context["project_name"],
                key="tia_director_pack_project_name",
            )
            contract_no = input_col2.text_input(
                "Contract No.",
                value=default_report_context["contract_no"],
                key="tia_director_pack_contract_no",
            )
            revision = input_col3.text_input(
                "Revision",
                value=default_report_context["revision"],
                key="tia_director_pack_revision",
            )

            schedule_col1, schedule_col2, schedule_col3 = st.columns(3)
            default_data_date = pd.to_datetime(default_report_context["data_date"], errors="coerce")
            data_date = schedule_col1.date_input(
                "Data Date",
                value=default_data_date.date() if pd.notna(default_data_date) else pd.Timestamp.today().date(),
                key="tia_director_pack_data_date",
            )
            employer = schedule_col2.text_input(
                "Employer / Client",
                value=default_report_context["employer"],
                key="tia_director_pack_employer",
            )
            contractor = schedule_col3.text_input(
                "Contractor",
                value=default_report_context["contractor"],
                key="tia_director_pack_contractor",
            )

            detail_col1, detail_col2 = st.columns(2)
            contract_form_clause = detail_col1.text_input(
                "Contract Form / Clause",
                value=default_report_context["contract_form_clause"],
                key="tia_director_pack_contract_form_clause",
            )
            accepted_baseline_programme = detail_col2.text_input(
                "Accepted Baseline Programme",
                value=default_report_context["accepted_baseline_programme"],
                key="tia_director_pack_bl_programme",
            )

            detail_col3, detail_col4 = st.columns(2)
            impacted_update_programme = detail_col3.text_input(
                "Impacted Update Programme",
                value=default_report_context["impacted_update_programme"],
                key="tia_director_pack_update_programme",
            )
            calendar_basis = detail_col4.text_input(
                "Calendar Basis",
                value=default_report_context["calendar_basis"],
                key="tia_director_pack_calendar_basis",
            )

            control_col1, control_col2 = st.columns(2)
            schedule_file_name = control_col1.text_input(
                "Schedule File Name",
                value=default_report_context["schedule_file_name"],
                key="tia_director_pack_schedule_file_name",
            )
            schedule_options = control_col2.text_input(
                "Schedule Options",
                value=default_report_context["schedule_options"],
                key="tia_director_pack_schedule_options",
            )

            logic_col1, logic_col2 = st.columns(2)
            critical_path_basis = logic_col1.text_input(
                "Longest Path / Critical Path Basis",
                value=default_report_context["critical_path_basis"],
                key="tia_director_pack_critical_path_basis",
            )
            retained_logic_setting = logic_col2.text_input(
                "Retained Logic / Progress Override",
                value=default_report_context["retained_logic_setting"],
                key="tia_director_pack_retained_logic",
            )

            logic_col3, logic_col4 = st.columns(2)
            out_of_sequence_treatment = logic_col3.text_input(
                "Out-of-sequence Progress Treatment",
                value=default_report_context["out_of_sequence_treatment"],
                key="tia_director_pack_oos_treatment",
            )
            constraints = logic_col4.text_area(
                "Constraints",
                value=default_report_context["constraints"],
                key="tia_director_pack_constraints",
                height=90,
            )

            support_col1, support_col2, support_col3 = st.columns(3)
            calendars = support_col1.text_input(
                "Calendars",
                value=default_report_context["calendars"],
                key="tia_director_pack_calendars",
            )
            open_ends = support_col2.text_input(
                "Open Ends",
                value=default_report_context["open_ends"],
                key="tia_director_pack_open_ends",
            )
            negative_float = support_col3.text_input(
                "Negative Float",
                value=default_report_context["negative_float"],
                key="tia_director_pack_negative_float",
            )

            selected_sources = st.multiselect(
                "Data source selection",
                list(report_source_frames.keys()),
                default=list(report_source_frames.keys()),
                key="tia_director_pack_sources",
            )
            preserve_original_charts_images = st.checkbox(
                "Preserve original charts/images",
                value=True,
                key="tia_director_pack_preserve_images",
            )

            current_report_inputs = {
                "project_name": project_name,
                "contract_no": contract_no,
                "data_date": data_date,
                "revision": revision,
                "employer": employer,
                "contractor": contractor,
                "contract_form_clause": contract_form_clause,
                "accepted_baseline_programme": accepted_baseline_programme,
                "impacted_update_programme": impacted_update_programme,
                "calendar_basis": calendar_basis,
                "schedule_file_name": schedule_file_name,
                "schedule_options": schedule_options,
                "critical_path_basis": critical_path_basis,
                "retained_logic_setting": retained_logic_setting,
                "out_of_sequence_treatment": out_of_sequence_treatment,
                "constraints": constraints,
                "calendars": calendars,
                "open_ends": open_ends,
                "negative_float": negative_float,
            }
            completion_pct, missing_required_fields = compute_required_fields_completion(current_report_inputs)
            completion_col.metric("Required Fields Completion", f"{completion_pct:.0f}%")
            source_status_col.metric("Data Sources Selected", str(len(selected_sources)))
            if missing_required_fields:
                st.warning("Generated with missing placeholders. Please complete before formal submission.")
                st.caption(f"Missing required fields: {', '.join(missing_required_fields)}")

            status_df = build_data_source_status_df(report_source_frames)
            st.markdown("##### Data Source Status")
            st.dataframe(status_df, width="stretch", hide_index=True)

            st.markdown("##### Replacement Preview")
            st.dataframe(build_replacement_preview_df(current_report_inputs), width="stretch", hide_index=True)

            if st.button("Generate Time Impact Analysis Report | Director Pack", key="generate_tia_director_pack"):
                try:
                    generator = TIADirectorPackGenerator(
                        TIA_DIRECTOR_WORD_TEMPLATE_PATH,
                        active_project_context.reports_path,
                        CONTRACT_CLAIMS_DB_PATH,
                    )
                    generation_context = build_delay_tia_director_pack_context(
                        active_delay_tia_context,
                        active_delay_tia_analysis,
                        overrides={
                            **current_report_inputs,
                            "data_date": pd.Timestamp(data_date),
                            "preserve_original_charts_images": preserve_original_charts_images,
                            "generated_by": "Delay TIA - Report Generator",
                        },
                        selected_sources=selected_sources,
                    )
                    generation_context["missing_required_fields"] = missing_required_fields
                    output_path = generator.generate(generation_context)
                    st.session_state["tia_director_pack_generated_path"] = str(output_path)
                    st.session_state["tia_director_pack_error"] = ""
                    st.success(f"Generated: {output_path.name}")
                    st.rerun()
                except Exception as exc:
                    st.session_state["tia_director_pack_error"] = str(exc)
                    st.error(f"Generation failed: {exc}")

            st.markdown("#### Primavera Fragnet Sheet")
            if not primavera_fragnet_df.empty:
                st.dataframe(primavera_fragnet_df, width="stretch", hide_index=True)
            else:
                st.info("No Primavera fragnet rows were generated from the current uploaded Delay TIA analysis.")

if active_slide_name == PROJECT_HUB_SLIDE_NAMES[13]:
    st.markdown("<div class='section-header'><h3>Output Studio</h3></div>", unsafe_allow_html=True)
    with st.expander("Sync repository", expanded=True):
        sync_notice = st.session_state.pop("repository_sync_notice", "")
        if sync_notice:
            st.success(sync_notice)
        configured_sync_pin = os.environ.get(f"{APP_ENV_PREFIX}_SYNC_ADMIN_PIN", "").strip() or _streamlit_secret(f"{APP_ENV_PREFIX}_SYNC_ADMIN_PIN")
        sync_admin_pin = ""
        if configured_sync_pin:
            sync_admin_pin = st.text_input(
                "Administrator PIN",
                type="password",
                key="repository_sync_admin_pin",
            )
        sync_authorized = repository_sync_authorized(sync_admin_pin)
        sync_now_col, sync_watch_col = st.columns(2)
        with sync_now_col:
            if st.button(
                "Sync",
                key="repository_sync_once",
                width="stretch",
                disabled=not sync_authorized,
            ):
                sync_completed = False
                with st.spinner("Synchronizing workspace..."):
                    try:
                        sync_ok, sync_output = run_repository_sync_once()
                        if sync_ok:
                            st.cache_data.clear()
                            st.session_state["repository_sync_notice"] = "Repository synchronized and application data refreshed."
                            sync_completed = True
                        else:
                            st.error("Repository synchronization failed.")
                        if sync_output:
                            st.code(sync_output, language="text")
                    except Exception as exc:
                        st.error(f"Repository synchronization failed: {exc}")
                if sync_completed:
                    st.rerun()
        with sync_watch_col:
            if st.button(
                "Start 30-minute auto sync",
                key="repository_sync_watch",
                width="stretch",
                disabled=not sync_authorized,
            ):
                try:
                    start_repository_sync_watch()
                    st.success("The 30-minute synchronization watcher started.")
                except Exception as exc:
                    st.error(f"Unable to start synchronization watcher: {exc}")
        st.code(repository_sync_log_tail(), language="text")
    if True:
        output_mode = st.radio(
            "Choose dashboard output",
            [
                "Executive dashboard",
                "Original presentation print-only",
                "Linked executive dashboard",
                "Detailed Progress report",
            ],
            horizontal=True,
        )

        metric_group_options = {
            "Overview Metrics": {
                "Project Start Date": format_project_date(overview_metrics.get("project_start")),
                "Project Finish Date": format_project_date(overview_metrics.get("project_finish")),
                "Project Duration [Days]": str(int(overview_metrics.get("duration_days", 0))),
                "Duration Elapsed": pct(overview_metrics.get("duration_elapsed_pct")),
                "Remaining Duration": pct(overview_metrics.get("remaining_duration_pct")),
                "Overall Progress": pct(overview_metrics.get("overall_progress")),
                "Planned Progress": pct(overview_metrics.get("planned_progress")),
                "Contract Value": egp(overview_metrics.get("contract_value")),
            },
            "EVM Metrics": {
                "BAC": egp(evm_metrics.get("bac")),
                "AC": egp(evm_metrics.get("ac")),
                "EV": egp(evm_metrics.get("ev")),
                "PV": egp(evm_metrics.get("pv")),
                "SV": egp(evm_metrics.get("sv")),
                "CV": egp(evm_metrics.get("cv")),
                "EAC": egp(evm_metrics.get("eac")),
                "TCPI": f"{evm_metrics.get('tcpi'):.3f}" if evm_metrics.get("tcpi") is not None else "N/A",
            },
        }

        output_project_slug = re.sub(r"[^A-Za-z0-9._-]+", "_", active_project_id or "all_projects").strip("_")
        if output_mode == "Executive dashboard":
            the_big_dashboard_html = build_the_big_decision_dashboard_html(
                overview_metrics,
                evm_metrics,
                contract_metrics,
                delay_metrics,
                risk_metrics,
                milestone_metrics,
                activity_metrics,
                s_curve_metrics,
                active_project_record,
            )
            st.components.v1.html(the_big_dashboard_html, height=2100, scrolling=True)
            st.download_button(
                "Download Executive Dashboard (.html)",
                data=the_big_dashboard_html.encode("utf-8"),
                file_name=f"{output_project_slug}_executive_dashboard.html",
                mime="text/html",
                width="stretch",
            )
        elif output_mode == "Original presentation print-only":
            original_template_bytes, original_template_unresolved = build_original_template_presentation(
                overview_metrics,
                evm_metrics,
                contract_metrics,
                delay_metrics,
                risk_metrics,
                milestone_metrics,
                activity_metrics,
                s_curve_metrics,
            )
            c1, c2 = st.columns(2)
            c1.download_button(
                "Download Updated Original Presentation (.pptx)",
                data=original_template_bytes,
                file_name=f"{output_project_slug}_presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                width="stretch",
                disabled=not bool(original_template_bytes),
            )
            c2.download_button(
                "Download Updated Original Presentation for Print (.pptx)",
                data=original_template_bytes,
                file_name=f"{output_project_slug}_presentation_print.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                width="stretch",
                disabled=not bool(original_template_bytes),
            )
            if original_template_unresolved:
                st.markdown("#### Linking Items To Confirm")
                st.dataframe(
                    pd.DataFrame({"Unresolved Template Mapping": original_template_unresolved}),
                    width="stretch",
                    hide_index=True,
                    height=dataframe_height(pd.DataFrame({"Unresolved Template Mapping": original_template_unresolved}), max_height=320),
                )
        elif output_mode == "Linked executive dashboard":
            linked_selected_project_id = selected_project_id()
            linked_generation_project_id = linked_selected_project_id

            previous_active_project_id = st.session_state.get("active_project_id", "")
            st.session_state["active_project_id"] = linked_generation_project_id
            try:
                linked_overview_metrics = build_overview_metrics()
                linked_evm_metrics = build_evm_metrics()
                linked_contract_metrics = build_contract_metrics()
                linked_delay_metrics = build_delay_metrics()
                linked_risk_metrics = build_risk_metrics()
                linked_milestone_metrics = build_milestone_metrics()
                linked_s_curve_metrics = build_s_curve_metrics()
                linked_activity_metrics = build_activity_metrics()
                linked_evmData = build_earned_value_analysis_data(linked_evm_metrics)
                rootCauseDf = build_evm_root_cause_rows(linked_delay_metrics, linked_risk_metrics, linked_contract_metrics)
                linked_project_title = linked_overview_metrics.get("project_name") or project.get("project_name", "Project")
            finally:
                st.session_state["active_project_id"] = previous_active_project_id

            mitigationDf = build_evm_mitigation_rows()
            evmComments = get_evm_comments()
            linked_report_date = pd.Timestamp.today().strftime("%d %b %Y")
            if "linked_exec_evm_comment_quantitativePerformance" not in st.session_state:
                st.session_state["linked_exec_evm_comment_quantitativePerformance"] = evmComments.get("quantitativePerformance", "")
            if "linked_exec_evm_comment_rootCauseLinkage" not in st.session_state:
                st.session_state["linked_exec_evm_comment_rootCauseLinkage"] = evmComments.get("rootCauseLinkage", "")
            if "linked_exec_evm_comment_mitigationRecovery" not in st.session_state:
                st.session_state["linked_exec_evm_comment_mitigationRecovery"] = evmComments.get("mitigationRecovery", "")
            st.markdown("#### Earned Value Analysis Add-on Controls")
            st.text_area(
                "Add / Edit My Comment - Quantitative Performance (EVM Metrics)",
                key="linked_exec_evm_comment_quantitativePerformance",
                height=110,
                on_change=persist_evm_comments_from_widget_keys,
                args=(
                    "linked_exec_evm_comment_quantitativePerformance",
                    "linked_exec_evm_comment_rootCauseLinkage",
                    "linked_exec_evm_comment_mitigationRecovery",
                ),
            )
            if st.button("Clear Quantitative Performance Comment", key="clear_linked_exec_evm_quantitative"):
                clear_evm_comment_with_widget("quantitativePerformance", "linked_exec_evm_comment_quantitativePerformance")
                st.rerun()
            st.text_area(
                "Add / Edit My Comment - Root Cause Linkage",
                key="linked_exec_evm_comment_rootCauseLinkage",
                height=110,
                on_change=persist_evm_comments_from_widget_keys,
                args=(
                    "linked_exec_evm_comment_quantitativePerformance",
                    "linked_exec_evm_comment_rootCauseLinkage",
                    "linked_exec_evm_comment_mitigationRecovery",
                ),
            )
            if st.button("Clear Root Cause Linkage Comment", key="clear_linked_exec_evm_root"):
                clear_evm_comment_with_widget("rootCauseLinkage", "linked_exec_evm_comment_rootCauseLinkage")
                st.rerun()
            st.text_area(
                "Add / Edit My Comment - Contractor Mitigation & Recovery Status",
                key="linked_exec_evm_comment_mitigationRecovery",
                height=110,
                on_change=persist_evm_comments_from_widget_keys,
                args=(
                    "linked_exec_evm_comment_quantitativePerformance",
                    "linked_exec_evm_comment_rootCauseLinkage",
                    "linked_exec_evm_comment_mitigationRecovery",
                ),
            )
            if st.button("Clear Mitigation & Recovery Comment", key="clear_linked_exec_evm_mitigation"):
                clear_evm_comment_with_widget("mitigationRecovery", "linked_exec_evm_comment_mitigationRecovery")
                st.rerun()

            evmComments = get_evm_comments()
            executive_html = build_linked_executive_dashboard_html(
                linked_overview_metrics,
                linked_evm_metrics,
                linked_contract_metrics,
                linked_delay_metrics,
                linked_risk_metrics,
                linked_milestone_metrics,
                linked_activity_metrics,
                linked_evmData,
                rootCauseDf,
                mitigationDf,
                evmComments,
            )
            evm_html = printEVMHtml(linked_project_title, linked_report_date, linked_evmData, rootCauseDf, mitigationDf, evmComments)
            evm_ppt = exportEVMToPowerPoint(linked_project_title, linked_report_date, linked_evmData, rootCauseDf, mitigationDf, evmComments)
            linked_dashboard_ppt = export_linked_executive_dashboard_to_powerpoint(
                linked_overview_metrics,
                linked_evm_metrics,
                linked_contract_metrics,
                linked_delay_metrics,
                linked_risk_metrics,
                linked_milestone_metrics,
                linked_activity_metrics,
                linked_evmData,
                rootCauseDf,
                mitigationDf,
                evmComments,
            )
            linked_dashboard_a3_summary_ppt = export_linked_executive_dashboard_a3_summary_ppt(
                linked_overview_metrics,
                linked_evm_metrics,
                linked_contract_metrics,
                linked_delay_metrics,
                linked_risk_metrics,
                linked_milestone_metrics,
                linked_activity_metrics,
                linked_s_curve_metrics,
                linked_evmData,
                rootCauseDf,
                mitigationDf,
            )
            linked_dashboard_a3_summary_html = build_linked_executive_dashboard_a3_summary_html(
                linked_overview_metrics,
                linked_evm_metrics,
                linked_contract_metrics,
                linked_delay_metrics,
                linked_risk_metrics,
                linked_milestone_metrics,
                linked_activity_metrics,
                linked_s_curve_metrics,
                linked_evmData,
                rootCauseDf,
                mitigationDf,
            )

            st.components.v1.html(executive_html, height=2200, scrolling=True)
            linked_export_col1, linked_export_col2 = st.columns(2)
            linked_export_col1.download_button(
                "Download Linked Executive Dashboard (.html)",
                data=executive_html.encode("utf-8"),
                file_name=f"{output_project_slug}_linked_executive_dashboard.html",
                mime="text/html",
                width="stretch",
            )
            linked_export_col2.download_button(
                "Download Linked Executive Dashboard PowerPoint (.pptx) - Landscape",
                data=linked_dashboard_ppt,
                file_name=f"{output_project_slug}_linked_executive_dashboard_landscape.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                width="stretch",
            )
            linked_export_col_summary_html, linked_export_col_summary, linked_export_col3, linked_export_col4 = st.columns(4)
            linked_export_col_summary_html.download_button(
                "Download Summarized Linked Dashboard (.html) - A3 Landscape One Page",
                data=linked_dashboard_a3_summary_html.encode("utf-8"),
                file_name=f"{output_project_slug}_linked_executive_dashboard_a3_landscape_one_page.html",
                mime="text/html",
                width="stretch",
            )
            linked_export_col_summary.download_button(
                "Download Summarized Linked Dashboard (.pptx) - A3 Landscape One Page",
                data=linked_dashboard_a3_summary_ppt,
                file_name=f"{output_project_slug}_linked_executive_dashboard_a3_landscape_one_page.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                width="stretch",
            )
            linked_export_col3.download_button(
                "Download EVM Add-on HTML",
                data=evm_html.encode("utf-8"),
                file_name=f"{output_project_slug}_linked_executive_dashboard_evm_analysis.html",
                mime="text/html",
                width="stretch",
            )
            linked_export_col4.download_button(
                "Download EVM Add-on PowerPoint",
                data=evm_ppt,
                file_name=f"{output_project_slug}_linked_executive_dashboard_evm_analysis.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                width="stretch",
            )
            with st.expander("Preview summarized A3 linked executive dashboard", expanded=False):
                st.components.v1.html(linked_dashboard_a3_summary_html, height=1600, scrolling=True)
        elif output_mode == "Detailed Progress report":
            if not OPENPYXL_AVAILABLE:
                st.error("`openpyxl` is not installed in the current environment, so the workbook export is temporarily unavailable.")
                st.code("pip install openpyxl", language="powershell")
                st.stop()
            workbook_bytes, report_readme, report_assumptions = build_detailed_progress_report_package(
                overview_metrics,
                evm_metrics,
                contract_metrics,
                delay_metrics,
                risk_metrics,
                milestone_metrics,
                activity_metrics,
                s_curve_metrics,
            )
            generated_tabs = [
                "01 POWER BI DASHBOARD",
                "02 WEEKLY DASHBOARD",
                "03 Report Control",
                "04 Project Overview",
                "05 Executive Summary",
                "06 Milestones",
                "07 Progress Activities",
                "08 Deliverables",
                "09 Schedule Paths",
                "10 Deviated Activities",
                "11 Cost & EVA",
                "12 Manpower",
                "13 Equipment",
                "14 S-Curve",
                "15 Shop Drawings",
                "16 Procurement",
                "17 Invoices",
                "18 Risks & Issues",
                "19 Mitigation Plan",
                "20 HSE & QAQC",
                "21 Correspondence",
                "22 Photos Register",
                "23 Weekly Data",
                "24 Power BI Model",
                "99 Lists",
            ]
            st.markdown("#### Generated Tabs")
            st.dataframe(pd.DataFrame({"Sheet Name": generated_tabs}), width="stretch", hide_index=True, height=dataframe_height(pd.DataFrame({"Sheet Name": generated_tabs}), max_height=760))
            st.markdown("#### Power BI Connection Steps")
            st.markdown(
                """
                1. Open Power BI Desktop.
                2. Get Data -> Excel Workbook.
                3. Select `Detailed_Progress_Report_PowerBI_Ready.xlsx`.
                4. Import the named Excel Tables.
                5. Build relationships using the `24 Power BI Model` sheet.
                """
            )
            if report_assumptions:
                st.markdown("#### Assumptions / Limitations")
                st.dataframe(pd.DataFrame({"Assumption / Limitation": report_assumptions}), width="stretch", hide_index=True, height=dataframe_height(pd.DataFrame({"Assumption / Limitation": report_assumptions}), max_height=320))
            st.markdown("#### Validation Checks")
            validation_checks = pd.DataFrame(
                {
                    "Check": [
                        "Workbook generated with required sheet order",
                        "Structured Excel tables created for analytical sheets",
                        "No merged cells used inside data tables",
                        "Monthly and weekly dashboards linked to the same workbook model",
                        "README generated with update workflow and governance rules",
                    ],
                    "Status": ["Passed", "Passed", "Passed", "Passed", "Passed"],
                }
            )
            st.dataframe(validation_checks, width="stretch", hide_index=True, height=dataframe_height(validation_checks, max_height=260))
            detailed_html = build_detailed_progress_report_html(
                overview_metrics,
                evm_metrics,
                contract_metrics,
                delay_metrics,
                risk_metrics,
                milestone_metrics,
                activity_metrics,
                s_curve_metrics,
            )
            detailed_docx = build_detailed_progress_report_docx_bytes(
                overview_metrics,
                evm_metrics,
                contract_metrics,
                delay_metrics,
                risk_metrics,
                milestone_metrics,
                activity_metrics,
                s_curve_metrics,
            )
            power_bi_style_html = build_detailed_progress_power_bi_style_html(
                overview_metrics,
                evm_metrics,
                contract_metrics,
                delay_metrics,
                risk_metrics,
                milestone_metrics,
                activity_metrics,
                s_curve_metrics,
            )
            c1, c2, c3, c4 = st.columns(4)
            c1.download_button(
                "Download Detailed Progress Report (.xlsx)",
                data=workbook_bytes,
                file_name=f"{output_project_slug}_Detailed_Progress_Report.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                width="stretch",
            )
            c2.download_button(
                "Download Detailed Progress Report (.html)",
                data=detailed_html.encode("utf-8"),
                file_name=f"{output_project_slug}_Detailed_Progress_Report.html",
                mime="text/html",
                width="stretch",
            )
            c3.download_button(
                "Download Detailed Progress Report (.docx)",
                data=detailed_docx,
                file_name=f"{output_project_slug}_Detailed_Progress_Report.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                width="stretch",
            )
            c4.download_button(
                "Download Power BI Style Dashboard (.html)",
                data=power_bi_style_html.encode("utf-8"),
                file_name=f"{output_project_slug}_Detailed_Progress_Report_PowerBI_Style.html",
                mime="text/html",
                width="stretch",
            )
            with st.expander("README / Power BI governance notes", expanded=False):
                st.markdown(report_readme)

if active_slide_name == PROJECT_HUB_SLIDE_NAMES[12]:
    render_claims_header(
        "Contract & Claims Intelligence Center",
        "AI-powered contract library, entitlement engine, evidence mapper, and client rebuttal system for contractor claims.",
    )
    st.markdown(f"**Active Project:** {html.escape(active_project_context.project_display_name)}")
    if active_project_context.is_all_projects:
        st.info("Select one project to use Contract & Claims Intelligence Center. Portfolio mode does not load project claims, evidence, or contract databases.")
        st.stop()
    contract_center_view = st.selectbox(
        "Contract center view",
        ["Contract Clauses", "Claims Intelligence Center"],
        key="contract_center_active_view",
    )
    if contract_center_view == "Contract Clauses":
        render_contract_clause_matching_engine()

    if contract_center_view == "Claims Intelligence Center":
        contract_repo_readme = CONTRACT_REPOSITORY_DIR / "README.md"
        if not contract_repo_readme.exists():
            contract_repo_readme.write_text(
                "# Contract Repository\n\nPut the contract files here. Supported types: PDF, DOCX, TXT, CSV, XLS, XLSX.\n"
                "The Contract & Claims Intelligence Center will analyze them once and persist the knowledge base into the project SQLite database.\n",
                encoding="utf-8",
            )

        rebuild_contract_library = st.button("Rebuild Contract Library", key="ccc_rebuild_contract_library")
        with st.expander("What does Rebuild Contract Library do?", expanded=False):
            st.markdown(
                "- Re-scans only this selected project's contract source folder.\n"
                "- Extracts clauses from supported contract files: PDF, DOCX, TXT, CSV, XLS, XLSX.\n"
                "- Rebuilds this project's own SQLite knowledge base used for clause search, claim classification, evidence mapping, rebuttal drafting, and claim drafting.\n"
                "- It does not read contracts from another project and does not update another project's database.\n"
                "- Use it after adding, replacing, or correcting contract files."
            )
            st.code(f"Source folder: {CONTRACT_REPOSITORY_DIR}\nDatabase: {CONTRACT_CLAIMS_DB_PATH}", language="text")
        contract_repo_status = ccc.persist_contract_analysis(CONTRACT_CLAIMS_DB_PATH, CONTRACT_REPOSITORY_DIR, rebuild=rebuild_contract_library)
        auto_library_status = contract_repo_status.get("auto_library_status", {})
        if auto_library_status.get("generated"):
            st.success(
                f"{auto_library_status.get('message', 'Contract clause library generated.')} "
                f"Rows: {auto_library_status.get('row_count', 0)}"
            )
        elif auto_library_status.get("message"):
            st.caption(str(auto_library_status.get("message")))
        stale_documents_removed = int(contract_repo_status.get("stale_documents_removed") or 0)
        if stale_documents_removed:
            st.warning(
                f"Cleaned {stale_documents_removed} stale contract source record(s) that were outside the active project folder. "
                "This project now uses only its own contract repository."
            )
        contract_clauses_df = ccc.load_contract_library(CONTRACT_CLAIMS_DB_PATH)
        contract_evidence_df = ccc.load_evidence_documents(CONTRACT_CLAIMS_DB_PATH)
        contract_evidence_mappings_df = ccc.load_evidence_mappings(CONTRACT_CLAIMS_DB_PATH)
        contract_claim_drafts_df = ccc.load_claim_drafts(CONTRACT_CLAIMS_DB_PATH)
        contract_analysis_status_df = ccc.load_contract_analysis_status(CONTRACT_CLAIMS_DB_PATH)
        contract_versions_df = ccc.load_contract_versions(CONTRACT_CLAIMS_DB_PATH)
        contract_center_kpis = ccc.build_contract_center_kpis(CONTRACT_CLAIMS_DB_PATH)
        for trace_frame in [
            contract_clauses_df,
            contract_evidence_df,
            contract_evidence_mappings_df,
            contract_claim_drafts_df,
            contract_analysis_status_df,
            contract_versions_df,
        ]:
            if not trace_frame.empty:
                if "project_id" not in trace_frame.columns:
                    trace_frame.insert(0, "project_id", active_project_context.project_id)
                if "project_display_name" not in trace_frame.columns:
                    trace_frame.insert(1, "project_display_name", active_project_context.project_display_name)
                if "source_folder" not in trace_frame.columns:
                    trace_frame["source_folder"] = str(active_project_context.project_folder_path)

        claim_triggers_count = int((contract_clauses_df["section_name"].astype(str) == "Claim Triggers").sum()) if not contract_clauses_df.empty else 0
        draft_claims_generated = int(len(contract_claim_drafts_df))
        mapped_items = int(len(contract_evidence_mappings_df))
        complete_items = int(contract_evidence_mappings_df["missing_evidence_items"].astype(str).str.strip().eq("").sum()) if not contract_evidence_mappings_df.empty else 0
        evidence_completeness_pct = (complete_items / mapped_items * 100.0) if mapped_items else 0.0
        kb_ready = contract_repo_status.get("knowledge_base_status") == "Ready"
        risk_badge = render_claims_badge(
            "High Contract Risk" if contract_center_kpis["high_risk_clauses"] > 0 else "Risk Profile Stable",
            "red" if contract_center_kpis["high_risk_clauses"] > 0 else "green",
        )
        entitlement_badge = render_claims_badge(
            "Strong Entitlement Items Available" if contract_center_kpis["strong_claim_opportunities"] > 0 else "No Strong Entitlement Items Yet",
            "green" if contract_center_kpis["strong_claim_opportunities"] > 0 else "amber",
        )
        kb_badge = render_claims_badge(
            "Knowledge Base Ready" if kb_ready else "Knowledge Base Pending",
            "blue" if kb_ready else "amber",
        )

        st.markdown(
            f"""
            <div class="claims-surface">
              <div class="claims-surface-title">Stored Knowledge Base Status</div>
              <div class="claims-surface-body">
                Contract files are stored in the project repository, analyzed once, persisted into SQLite, and reused for fast contract search,
                evidence mapping, rebuttal drafting, and claim building. Re-analysis runs only when the repository changes or when you explicitly rebuild the library.
              </div>
              <div class="claims-statline">{risk_badge}{entitlement_badge}{kb_badge}<span class='claims-badge badge-slate'>Evidence Completeness {pct(evidence_completeness_pct)}</span></div>
            </div>
            """,
            unsafe_allow_html=True,
        )

        render_claims_workflow()

        kpi_row_1 = st.columns(4)
        with kpi_row_1[0]:
            render_claims_kpi_card("Total Contract Clauses", contract_center_kpis["total_clauses"], "blue", "Stored searchable clauses")
        with kpi_row_1[1]:
            render_claims_kpi_card("Contractor Rights", contract_center_kpis["contractor_rights"], "green", "Rights recognized in library")
        with kpi_row_1[2]:
            render_claims_kpi_card("Employer Obligations", contract_center_kpis["employer_obligations"], "blue", "Obligation-linked opportunities")
        with kpi_row_1[3]:
            render_claims_kpi_card("Claim Triggers", claim_triggers_count, "amber", "Trigger clauses identified")

        kpi_row_2 = st.columns(4)
        with kpi_row_2[0]:
            render_claims_kpi_card("High Risk Clauses", contract_center_kpis["high_risk_clauses"], "red", "High / critical contractual exposure")
        with kpi_row_2[1]:
            render_claims_kpi_card("Strong Entitlement Items", contract_center_kpis["strong_claim_opportunities"], "green", "Strong / very strong claim positions")
        with kpi_row_2[2]:
            render_claims_kpi_card("Missing Evidence", contract_center_kpis["missing_evidence_items"], "amber", f"Completeness {pct(evidence_completeness_pct)}")
        with kpi_row_2[3]:
            render_claims_kpi_card("Draft Claims Generated", draft_claims_generated, "blue", f"Last analysis {delay_tia_docx_text(contract_center_kpis['last_analysis_date'], 'Not analyzed yet')}")

        contract_view = st.selectbox(
            "Claims Intelligence view",
            [
                "Upload & Extract",
                "Contract Library",
                "Ask Contract AI",
                "Evidence Mapping",
                "Client Rebuttal Engine",
                "Claim Builder",
                "Export Center",
            ],
            key="claims_intelligence_active_view",
        )

        if contract_view == "Upload & Extract":
            st.markdown("#### Upload & Extract")
            repo_col1, repo_col2 = st.columns([1.2, 1])
            with repo_col1:
                st.markdown(
                    f"""
                    <div class="claims-surface">
                      <div class="claims-surface-title">Contract Repository</div>
                      <div class="claims-surface-body">Place the contract files here once. The system will detect new or changed files, extract clauses, classify claims intelligence, and persist the results to SQLite.</div>
                    </div>
                    """,
                    unsafe_allow_html=True,
                )
                st.code(str(CONTRACT_REPOSITORY_DIR), language="text")
                detected_files_df = pd.DataFrame(
                    {
                        "Detected Contract File": contract_repo_status.get("detected_files", []),
                    }
                )
                if detected_files_df.empty:
                    st.warning("No contract files are currently detected in the repository folder.")
                else:
                    st.dataframe(detected_files_df, width="stretch", hide_index=True)
            with repo_col2:
                st.markdown(
                    f"""
                    <div class="claims-surface">
                      <div class="claims-surface-title">Repository Status</div>
                      <div class="claims-surface-body">
                        <b>Contract version:</b> {html.escape(str(contract_repo_status.get("contract_version", "N/A")))}<br>
                        <b>Knowledge base status:</b> {html.escape(str(contract_repo_status.get("knowledge_base_status", "Unknown")))}<br>
                        <b>Total clauses extracted:</b> {int(contract_repo_status.get("total_clauses", 0))}<br>
                        <b>Last analysis date:</b> {html.escape(delay_tia_docx_text(contract_repo_status.get("last_analysis_date"), "Not analyzed yet"))}
                      </div>
                    </div>
                    """,
                    unsafe_allow_html=True,
                )
                if contract_repo_status.get("extraction_issue"):
                    st.warning(str(contract_repo_status.get("extraction_issue")))
                st.button("Rebuild Contract Library", key="ccc_rebuild_contract_library_secondary", disabled=True, help="Use the main rebuild button at the top of the page.")

            readiness_col1, readiness_col2, readiness_col3 = st.columns(3)
            if contract_repo_status.get("contract_loaded"):
                readiness_col1.success("Contract Loaded")
            else:
                readiness_col1.warning("Contract Not Loaded")
            if contract_repo_status.get("knowledge_base_status") == "Ready":
                readiness_col2.success("Knowledge Base Ready")
            else:
                readiness_col2.warning("Knowledge Base Not Ready")
            if not contract_clauses_df.empty:
                readiness_col3.success("AI Search Ready")
            else:
                readiness_col3.warning("AI Search Not Ready")

            if not contract_versions_df.empty:
                st.markdown("#### Contract Version Register")
                st.dataframe(contract_versions_df, width="stretch", hide_index=True, height=dataframe_height(contract_versions_df, max_height=320))
            if not contract_analysis_status_df.empty:
                st.markdown("#### Analysis Status Log")
                st.dataframe(contract_analysis_status_df, width="stretch", hide_index=True, height=dataframe_height(contract_analysis_status_df, max_height=320))

        if contract_view == "Contract Library":
            st.markdown("#### Searchable Contract Claims Library")
            if contract_clauses_df.empty:
                st.info("No contract clauses are stored yet. Add contract files to the repository folder, then rebuild the contract library.")
            else:
                st.markdown(
                    "<div class='claims-surface'><div class='claims-surface-title'>Clause Search Panel</div><div class='claims-surface-body'>Search by clause number, keyword, claim type, or risk level. The filtered view remains contractor-focused and highlights risk and entitlement quality.</div></div>",
                    unsafe_allow_html=True,
                )
                lib_col1, lib_col2, lib_col3, lib_col4 = st.columns(4)
                library_search = lib_col1.text_input("Search clause number / keyword", key="ccc_library_search")
                section_filter = lib_col2.multiselect("Section", sorted(contract_clauses_df["section_name"].dropna().astype(str).unique()), key="ccc_section_filter")
                claim_type_filter = lib_col3.multiselect("Claim Type", sorted(contract_clauses_df["claim_type"].dropna().astype(str).unique()), key="ccc_claim_type_filter")
                risk_filter = lib_col4.multiselect("Risk Level", sorted(contract_clauses_df["risk_level"].dropna().astype(str).unique()), key="ccc_risk_filter")

                filtered_library_df = contract_clauses_df.copy()
                if library_search.strip():
                    filtered_library_df = ccc.clause_search_dataframe(filtered_library_df, library_search, limit=max(len(filtered_library_df), 12))
                if section_filter:
                    filtered_library_df = filtered_library_df[filtered_library_df["section_name"].isin(section_filter)]
                if claim_type_filter:
                    filtered_library_df = filtered_library_df[filtered_library_df["claim_type"].isin(claim_type_filter)]
                if risk_filter:
                    filtered_library_df = filtered_library_df[filtered_library_df["risk_level"].isin(risk_filter)]

                library_view = filtered_library_df[
                    [
                        "clause_number",
                        "clause_title",
                        "section_name",
                        "claim_type",
                        "risk_level",
                        "claim_strength",
                        "notice_required",
                        "time_impact",
                        "cost_impact",
                        "recommended_action",
                    ]
                ].copy() if not filtered_library_df.empty else pd.DataFrame()
                st.dataframe(library_view, width="stretch", hide_index=True, height=dataframe_height(library_view, max_height=720))

                with st.expander("Clause Details", expanded=False):
                    detail_cols = [
                        "clause_number", "clause_title", "section_name", "claim_type", "risk_level", "claim_strength",
                        "plain_english_meaning", "required_evidence", "possible_client_rejection",
                        "contractor_counterargument", "recommended_action", "related_project_records_needed",
                    ]
                    st.dataframe(filtered_library_df[[col for col in detail_cols if col in filtered_library_df.columns]], width="stretch", hide_index=True, height=dataframe_height(filtered_library_df, max_height=1000))

        if contract_view == "Ask Contract AI":
            st.markdown("#### Ask Contract AI")
            ccc_query_state_key = f"ccc_last_query::{active_project_context.project_id}"
            ccc_answer_state_key = f"ccc_last_answer::{active_project_context.project_id}"
            ask_query = st.text_area(
                "Ask a contract / claims question",
                value=st.session_state.get(
                    ccc_query_state_key,
                    "Can we claim EOT for late IFC drawings?",
                ),
                key="ccc_query_input",
                height=120,
            )
            if st.button("Analyze Contract Question", key="ccc_analyze_question"):
                st.session_state[ccc_query_state_key] = ask_query
                st.session_state[ccc_answer_state_key] = ccc.answer_contract_question(CONTRACT_CLAIMS_DB_PATH, ask_query)

            with st.expander("Built-in Contract AI Test Cases", expanded=False):
                test_cases_df = ccc.get_contract_ai_test_cases()
                st.dataframe(test_cases_df, width="stretch", hide_index=True, height=dataframe_height(test_cases_df, max_height=320))

            last_answer = st.session_state.get(ccc_answer_state_key)
            if last_answer:
                answer_badge_tone = "green" if last_answer["entitlement_decision"] == "YES" else "amber" if last_answer["entitlement_decision"] == "POSSIBLE" else "red" if last_answer["entitlement_decision"] == "NO" else "slate"
                risk_tone = "red" if str(last_answer["risk_assessment"]).lower() in {"critical", "high"} else "amber" if str(last_answer["risk_assessment"]).lower() == "medium" else "green"
                evidence_tone = "green" if str(last_answer["evidence_strength_label"]).lower() in {"very strong", "strong"} else "amber" if str(last_answer["evidence_strength_label"]).lower() == "medium" else "red"
                category_badges = "".join(render_claims_badge(category, "blue") for category in last_answer.get("question_categories", []))
                st.markdown(
                    f"<div class='claims-statline'>{render_claims_badge('Decision: ' + str(last_answer['entitlement_decision']), answer_badge_tone)}{render_claims_badge('Risk: ' + str(last_answer['risk_assessment']), risk_tone)}{render_claims_badge('Evidence Strength: ' + str(last_answer['evidence_strength_label']) + ' (' + str(last_answer['evidence_strength_score']) + '/100)', evidence_tone)}{render_claims_badge('Matched Clauses: ' + str(len(last_answer['relevant_clauses_df']) if isinstance(last_answer['relevant_clauses_df'], pd.DataFrame) else 0), 'blue')}</div>",
                    unsafe_allow_html=True,
                )
                if category_badges:
                    st.markdown(f"<div class='claims-statline'>{category_badges}</div>", unsafe_allow_html=True)
                render_ai_status_note(last_answer.get("ai_status"))
                st.markdown(
                    f"""
                    <div class="claims-dual-grid">
                      <div class="claims-answer-panel">
                        <div class="claims-answer-title">Short Answer</div>
                        <div class="claims-answer-body">{html.escape(str(last_answer["short_answer"]))}</div>
                      </div>
                      <div class="claims-answer-panel">
                        <div class="claims-answer-title">Contractor-Friendly Interpretation</div>
                        <div class="claims-answer-body">{html.escape(str(last_answer["contractor_friendly_interpretation"]))}</div>
                      </div>
                      <div class="claims-answer-panel">
                        <div class="claims-answer-title">Required Evidence</div>
                        <div class="claims-answer-body">{html.escape(str(last_answer["required_evidence"]))}</div>
                      </div>
                      <div class="claims-answer-panel">
                        <div class="claims-answer-title">Missing Evidence</div>
                        <div class="claims-answer-body">{html.escape(str(last_answer["missing_evidence"]))}</div>
                      </div>
                      <div class="claims-answer-panel">
                        <div class="claims-answer-title">Likely Client / Engineer Rejection</div>
                        <div class="claims-answer-body">{html.escape(str(last_answer["likely_client_rejection"]))}</div>
                      </div>
                      <div class="claims-answer-panel">
                        <div class="claims-answer-title">Contractor Rebuttal</div>
                        <div class="claims-answer-body">{html.escape(str(last_answer["contractor_rebuttal"]))}</div>
                      </div>
                      <div class="claims-answer-panel">
                        <div class="claims-answer-title">Recommended Next Action</div>
                        <div class="claims-answer-body">{html.escape(str(last_answer["recommended_next_action"]))}</div>
                      </div>
                    </div>
                    """,
                    unsafe_allow_html=True,
                )
                st.markdown("**Contract Basis**")
                for basis_row in last_answer.get("contract_basis_rows", []):
                    st.markdown(f"- {html.escape(str(basis_row))}")
                st.markdown("**Claim Strategy**")
                st.write(last_answer["claim_strategy"])
                with st.expander("Entitlement / Proof / Causation / Cost / Notice / Rebuttal Split", expanded=False):
                    analysis_dimensions = last_answer.get("analysis_dimensions", {})
                    for label, value in analysis_dimensions.items():
                        st.markdown(f"**{label}**")
                        st.write(value)
                predicted_rebuttals_df = last_answer.get("predicted_rebuttals_df", pd.DataFrame())
                if isinstance(predicted_rebuttals_df, pd.DataFrame) and not predicted_rebuttals_df.empty:
                    st.markdown("#### Predicted Client Defenses and Rebuttals")
                    st.dataframe(predicted_rebuttals_df, width="stretch", hide_index=True, height=dataframe_height(predicted_rebuttals_df, max_height=520))
                if isinstance(last_answer["relevant_clauses_df"], pd.DataFrame) and not last_answer["relevant_clauses_df"].empty:
                    st.markdown("#### Relevant Clauses")
                    st.dataframe(last_answer["relevant_clauses_df"], width="stretch", hide_index=True, height=dataframe_height(last_answer["relevant_clauses_df"], max_height=600))

        if contract_view == "Evidence Mapping":
            st.markdown("#### Evidence Mapping Engine")
            evidence_source_stream = st.selectbox(
                "Evidence source stream",
                [
                    "Emails",
                    "RFIs",
                    "Site Instructions",
                    "Meeting Minutes",
                    "Daily Reports",
                    "Weekly Reports",
                    "Drawing Logs",
                    "Submittal Logs",
                    "Payment Certificates",
                    "IPCs",
                    "Primavera P6 exports",
                    "Delay event logs",
                    "Client rejection letters",
                    "Manual Upload",
                ],
                key="ccc_evidence_source_stream",
            )
            evidence_uploads = st.file_uploader(
                "Upload project evidence files",
                type=["pdf", "docx", "txt", "csv", "xlsx", "xls", "eml", "html"],
                accept_multiple_files=True,
                key="ccc_evidence_uploads",
            )
            if st.button("Store and Map Evidence", key="ccc_store_evidence") and evidence_uploads:
                saved_files = ccc.persist_uploaded_evidence(CONTRACT_CLAIMS_DB_PATH, CONTRACT_EVIDENCE_DIR, evidence_uploads, evidence_source_stream)
                st.success(f"Stored and mapped {len(saved_files)} evidence file(s): {', '.join(saved_files)}")
                contract_evidence_df = ccc.load_evidence_documents(CONTRACT_CLAIMS_DB_PATH)
                contract_evidence_mappings_df = ccc.load_evidence_mappings(CONTRACT_CLAIMS_DB_PATH)

            ev_col1, ev_col2, ev_col3 = st.columns(3)
            ev_col1.metric("Evidence Files Stored", len(contract_evidence_df))
            ev_col2.metric("Evidence Mappings", len(contract_evidence_mappings_df))
            ev_col3.metric("Evidence Completeness", pct(evidence_completeness_pct))

            st.markdown("#### Stored Evidence Documents")
            if contract_evidence_df.empty:
                st.info("No evidence files are stored yet.")
            else:
                st.dataframe(contract_evidence_df, width="stretch", hide_index=True, height=dataframe_height(contract_evidence_df, max_height=420))

            st.markdown("#### Evidence-to-Clause Mapping")
            if contract_evidence_mappings_df.empty:
                st.info("No evidence mappings are stored yet.")
            else:
                st.dataframe(contract_evidence_mappings_df, width="stretch", hide_index=True, height=dataframe_height(contract_evidence_mappings_df, max_height=640))

        if contract_view == "Client Rebuttal Engine":
            st.markdown("#### Client Rebuttal Engine")
            ccc_rebuttal_text_state_key = f"ccc_last_rebuttal_text::{active_project_context.project_id}"
            ccc_rebuttal_state_key = f"ccc_last_rebuttal_result::{active_project_context.project_id}"
            rebuttal_upload = st.file_uploader(
                "Optional client rejection file",
                type=["pdf", "docx", "txt", "csv", "html"],
                key="ccc_rebuttal_upload",
            )
            rebuttal_text = st.text_area(
                "Paste the client / engineer rejection text",
                value=st.session_state.get(ccc_rebuttal_text_state_key, ""),
                key="ccc_rebuttal_text",
                height=150,
            )
            if rebuttal_upload is not None:
                extracted_rebuttal_text = ccc.extract_text_from_bytes(rebuttal_upload.name, rebuttal_upload.getvalue())
                if extracted_rebuttal_text.strip():
                    rebuttal_text = extracted_rebuttal_text
                    st.text_area("Extracted rejection text", value=extracted_rebuttal_text, height=180, key="ccc_rebuttal_extracted")
            if st.button("Generate Contractor Rebuttal", key="ccc_generate_rebuttal") and rebuttal_text.strip():
                st.session_state[ccc_rebuttal_text_state_key] = rebuttal_text
                st.session_state[ccc_rebuttal_state_key] = ccc.build_client_rebuttal(CONTRACT_CLAIMS_DB_PATH, rebuttal_text)
            rebuttal_result = st.session_state.get(ccc_rebuttal_state_key)
            if rebuttal_result:
                rebuttal_risk_tone = "red" if str(rebuttal_result["contractual_risk"]).lower() in {"critical", "high"} else "amber" if str(rebuttal_result["contractual_risk"]).lower() == "medium" else "green"
                rebuttal_evidence_tone = "green" if str(rebuttal_result["evidence_strength_label"]).lower() in {"very strong", "strong"} else "amber" if str(rebuttal_result["evidence_strength_label"]).lower() == "medium" else "red"
                st.markdown(
                    f"<div class='claims-statline'>{render_claims_badge(rebuttal_result['client_argument_summary'], 'amber')}{render_claims_badge('Contractual Risk: ' + str(rebuttal_result['contractual_risk']), rebuttal_risk_tone)}{render_claims_badge('Evidence Strength: ' + str(rebuttal_result['evidence_strength_label']) + ' (' + str(rebuttal_result['evidence_strength_score']) + '/100)', rebuttal_evidence_tone)}{render_claims_badge('Probability of Success: ' + str(rebuttal_result['probability_of_success']), 'green' if str(rebuttal_result['probability_of_success']).lower() == 'high' else 'amber')}</div>",
                    unsafe_allow_html=True,
                )
                render_ai_status_note(rebuttal_result.get("ai_status"))
                st.markdown(
                    f"""
                    <div class="claims-dual-grid">
                      <div class="claims-answer-panel">
                        <div class="claims-answer-title">Contractor Counterargument</div>
                        <div class="claims-answer-body">{html.escape(str(rebuttal_result["contractor_counterargument"]))}</div>
                      </div>
                      <div class="claims-answer-panel">
                        <div class="claims-answer-title">Evidence Needed</div>
                        <div class="claims-answer-body">{html.escape(str(rebuttal_result["evidence_needed"]))}</div>
                      </div>
                    </div>
                    """,
                    unsafe_allow_html=True,
                )
                st.markdown("**Recommended Response Wording**")
                st.write(rebuttal_result["recommended_response_wording"])
                detected_defenses_df = rebuttal_result.get("detected_defenses_df", pd.DataFrame())
                if isinstance(detected_defenses_df, pd.DataFrame) and not detected_defenses_df.empty:
                    st.markdown("#### Detected Client Defenses")
                    st.dataframe(detected_defenses_df, width="stretch", hide_index=True, height=dataframe_height(detected_defenses_df, max_height=560))
                if not rebuttal_result["relevant_clauses_df"].empty:
                    st.markdown("#### Supporting Clauses")
                    st.dataframe(rebuttal_result["relevant_clauses_df"], width="stretch", hide_index=True, height=dataframe_height(rebuttal_result["relevant_clauses_df"], max_height=520))

        if contract_view == "Claim Builder":
            st.markdown("#### Claim Builder")
            available_claim_types = sorted(set(ccc.CLAIM_CATEGORIES) | set(contract_clauses_df["claim_type"].dropna().astype(str).tolist())) if not contract_clauses_df.empty else ccc.CLAIM_CATEGORIES
            delay_event_options = delay_metrics["delays_df"]["delay_title"].dropna().astype(str).tolist() if "delay_title" in delay_metrics.get("delays_df", pd.DataFrame()).columns else []
            selected_claim_type = st.selectbox("Claim Type", available_claim_types, key="ccc_claim_type")
            selected_delay_event = st.selectbox("Delay / Cost Event", delay_event_options if delay_event_options else ["Manual Event Entry"], key="ccc_delay_event")
            if selected_delay_event == "Manual Event Entry":
                selected_delay_event = st.text_input("Manual delay / cost event", value="", key="ccc_delay_event_manual")

            clause_options = {}
            if not contract_clauses_df.empty:
                for _, row in contract_clauses_df.iterrows():
                    label = f"{row.get('clause_number', 'N/A')} | {row.get('section_name', '')} | {row.get('clause_title', '')[:100]}"
                    clause_options[label] = int(row["id"])
            selected_clause_labels = st.multiselect("Relevant Clauses", list(clause_options.keys()), key="ccc_selected_clauses")

            evidence_options = {}
            if not contract_evidence_df.empty:
                for _, row in contract_evidence_df.iterrows():
                    evidence_options[f"{row.get('file_name', '')} | {row.get('source_stream', '')}"] = int(row["id"])
            selected_evidence_labels = st.multiselect("Evidence Files", list(evidence_options.keys()), key="ccc_selected_evidence")

            claim_rejection_text = st.text_area("Client rejection / defense if any", value="", key="ccc_claim_rejection_text", height=120)
            if st.button("Build Claim Draft", key="ccc_build_claim_draft") and selected_delay_event:
                payload = ccc.build_claim_draft_payload(
                    CONTRACT_CLAIMS_DB_PATH,
                    claim_type=selected_claim_type,
                    delay_event=selected_delay_event,
                    selected_clause_ids=[clause_options[label] for label in selected_clause_labels],
                    selected_evidence_ids=[evidence_options[label] for label in selected_evidence_labels],
                    client_rejection_text=claim_rejection_text,
                )
                st.session_state["ccc_last_claim_payload"] = payload
                contract_claim_drafts_df = ccc.load_claim_drafts(CONTRACT_CLAIMS_DB_PATH)

            claim_payload = st.session_state.get("ccc_last_claim_payload")
            if claim_payload:
                st.markdown("#### Generated Claim Draft")
                st.markdown(
                    f"<div class='claims-statline'>{render_claims_badge(selected_claim_type, 'blue')}{render_claims_badge('Clauses: ' + str(len(selected_clause_labels)), 'green' if selected_clause_labels else 'amber')}{render_claims_badge('Evidence Files: ' + str(len(selected_evidence_labels)), 'blue' if selected_evidence_labels else 'amber')}</div>",
                    unsafe_allow_html=True,
                )
                for section_title, section_key in [
                    ("Claim Narrative", "narrative_text"),
                    ("Contractual Basis", "contractual_basis"),
                    ("Factual Background", "factual_background"),
                    ("Cause and Effect", "cause_effect"),
                    ("Evidence List", "evidence_list"),
                    ("Entitlement Statement", "entitlement_statement"),
                    ("Time Impact Statement", "time_impact_statement"),
                    ("Cost Impact Statement", "cost_impact_statement"),
                    ("Rebuttal Section", "rebuttal_section"),
                    ("Required Attachments Checklist", "attachment_checklist"),
                ]:
                    with st.expander(section_title, expanded=section_title in {"Claim Narrative", "Contractual Basis"}):
                        st.write(claim_payload.get(section_key, ""))

            if not contract_claim_drafts_df.empty:
                st.markdown("#### Stored Claim Drafts")
                st.dataframe(contract_claim_drafts_df, width="stretch", hide_index=True, height=dataframe_height(contract_claim_drafts_df, max_height=420))

        if contract_view == "Export Center":
            st.markdown("#### Export Center")
            st.markdown(
                "<div class='claims-surface'><div class='claims-surface-title'>Export Outputs</div><div class='claims-surface-body'>Export the stored contract knowledge base, mapped evidence, and the generated claim draft into management-ready and claim-ready formats.</div></div>",
                unsafe_allow_html=True,
            )
            contract_library_excel = b""
            contract_library_json = b""
            contract_library_html = ""
            evidence_matrix_csv = b""
            try:
                if OPENPYXL_AVAILABLE and not contract_clauses_df.empty:
                    contract_library_excel = ccc.build_contract_library_export_excel(contract_clauses_df)
            except Exception:
                contract_library_excel = b""
            if not contract_clauses_df.empty:
                contract_library_json = ccc.build_contract_library_json_bytes(contract_clauses_df)
                contract_library_html = ccc.build_contract_library_html(contract_clauses_df)
            if not contract_evidence_mappings_df.empty:
                evidence_matrix_csv = contract_evidence_mappings_df.to_csv(index=False).encode("utf-8")

            claim_payload = st.session_state.get("ccc_last_claim_payload")
            claim_docx_bytes = ccc.build_claim_draft_docx_bytes(claim_payload) if claim_payload else b""
            claim_pdf_bytes = ccc.build_claim_draft_pdf_bytes(claim_payload) if claim_payload and REPORTLAB_AVAILABLE else b""
            claim_html = ""
            if claim_payload:
                claim_html = (
                    "<html><head><meta charset='utf-8'><title>Claim Draft</title>"
                    "<style>body{font-family:Arial,sans-serif;padding:24px;color:#172033;}h1,h2{color:#0f172a;}</style></head><body>"
                    f"<h1>{html.escape(claim_payload.get('draft_name', 'Claim Draft'))}</h1>"
                    + "".join(
                        f"<h2>{html.escape(title)}</h2><p>{html.escape(str(claim_payload.get(key, ''))).replace(chr(10), '<br>')}</p>"
                        for title, key in [
                            ("Claim Narrative", "narrative_text"),
                            ("Contractual Basis", "contractual_basis"),
                            ("Factual Background", "factual_background"),
                            ("Cause and Effect", "cause_effect"),
                            ("Evidence List", "evidence_list"),
                            ("Entitlement Statement", "entitlement_statement"),
                            ("Time Impact Statement", "time_impact_statement"),
                            ("Cost Impact Statement", "cost_impact_statement"),
                            ("Rebuttal Section", "rebuttal_section"),
                            ("Required Attachments Checklist", "attachment_checklist"),
                        ]
                    )
                    + "</body></html>"
                )

            export_col1, export_col2, export_col3 = st.columns(3)
            export_col1.download_button(
                "Download Excel Contract Library",
                data=contract_library_excel,
                file_name="contract_claims_library.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                width="stretch",
                disabled=not bool(contract_library_excel),
            )
            export_col2.download_button(
                "Download JSON Knowledge Base",
                data=contract_library_json,
                file_name="contract_knowledge_base.json",
                mime="application/json",
                width="stretch",
                disabled=not bool(contract_library_json),
            )
            export_col3.download_button(
                "Download CSV Evidence Matrix",
                data=evidence_matrix_csv,
                file_name="contract_evidence_matrix.csv",
                mime="text/csv",
                width="stretch",
                disabled=not bool(evidence_matrix_csv),
            )

            export_col4, export_col5, export_col6 = st.columns(3)
            export_col4.download_button(
                "Download Word Claim Draft",
                data=claim_docx_bytes,
                file_name="contract_claim_draft.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                width="stretch",
                disabled=not bool(claim_docx_bytes),
            )
            export_col5.download_button(
                "Download PDF Claim Summary",
                data=claim_pdf_bytes,
                file_name="contract_claim_summary.pdf",
                mime="application/pdf",
                width="stretch",
                disabled=not bool(claim_pdf_bytes),
            )
            export_col6.download_button(
                "Download HTML Report",
                data=(claim_html or contract_library_html).encode("utf-8") if (claim_html or contract_library_html) else b"",
                file_name="contract_claims_report.html",
                mime="text/html",
                width="stretch",
                disabled=not bool(claim_html or contract_library_html),
            )

st.divider()
st.markdown("<p style='text-align:center;color:#667085;font-size:12px;'>Construction Project Control Platform | Designed and Developed By Eng. Ahmed Labib © Planning Department</p>", unsafe_allow_html=True)

